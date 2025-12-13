#!/usr/bin/env python3
"""
Create G7_C6_W5 Synthesis & Assessment presentation.
Topic: Plate Tectonics & Earth's Interior.

Assessment week structure (11 slides) - different from instructional weeks.
See PPTX_DESIGN_GUIDE.md for best practices documentation.
"""

import os
from pptx.dml.color import RGBColor
from pptx_common import (
    COLORS as BASE_COLORS,
    create_base_presentation,
    add_colored_shape,
    add_text_box,
    Inches,
    Pt,
    PP_ALIGN,
    MSO_ANCHOR,
)

# Earth/tectonic theme - extend base colors
COLORS = {**BASE_COLORS}
COLORS.update({
    'earth_brown': RGBColor(0x8B, 0x45, 0x13),
    'earth_dark': RGBColor(0x5D, 0x3A, 0x1A),
    'magma_orange': RGBColor(0xEA, 0x58, 0x0C),
    'core_red': RGBColor(0xB9, 0x1C, 0x1C),
    'seismic_blue': RGBColor(0x1E, 0x40, 0xAF),
    'light_tan_bg': RGBColor(0xFE, 0xF3, 0xE2),
    'assessment_brown': RGBColor(0x78, 0x35, 0x0F),
})


def create_presentation():
    """Create the G7_C6_W5 Assessment presentation"""
    prs = create_base_presentation()

    add_title_slide(prs)
    add_overview_slide(prs)
    add_assessed_on_slide(prs)
    add_part1_intro_slide(prs)
    add_part1_support_slide(prs)
    add_part2_intro_slide(prs)
    add_part2_support_slide(prs)
    add_part3_intro_slide(prs)
    add_part3_support_slide(prs)
    add_tips_slide(prs)
    add_summary_slide(prs)

    return prs


def add_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['assessment_brown'])

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(1.2),
                "Plate Tectonics & Earth's Interior",
                font_size=42, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")

    add_text_box(slide, Inches(0.5), Inches(2.6), Inches(9), Inches(0.5),
                "Week 5: Synthesis & Assessment",
                font_size=28, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(9), Inches(0.5),
                "Grade 7 Science | Cycle 6 | 100 Points",
                font_size=18, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(2), Inches(4.1), Inches(6), Inches(1.0), COLORS['white'])
    add_text_box(slide, Inches(2.2), Inches(4.15), Inches(5.6), Inches(0.45),
                "Part 1: Synthesis (20 pts) | Part 2: Assessment (60 pts)",
                font_size=14, bold=True, color=COLORS['assessment_brown'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(2.2), Inches(4.55), Inches(5.6), Inches(0.45),
                "Part 3: Misconception Check (20 pts)",
                font_size=14, bold=True, color=COLORS['assessment_brown'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_overview_slide(prs):
    """Slide 2: Assessment Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['assessment_brown'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Assessment Overview",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    parts = [
        ("Part 1: Synthesis Review", "20 pts", "15 min", "Connect all 4 weeks: boundaries, spreading, volcanoes, interior"),
        ("Part 2: Cumulative Assessment", "60 pts", "40 min", "Plate boundaries, seafloor evidence, volcanoes, Earth layers"),
        ("Part 3: Misconception Check", "20 pts", "20 min", "Target common errors for feedback"),
    ]

    y_pos = 0.9
    colors = [COLORS['green_start'], COLORS['seismic_blue'], COLORS['magma_orange']]
    for i, (title, pts, time, desc) in enumerate(parts):
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(1.1), COLORS['light_blue_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(1.1), colors[i])

        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.1), Inches(5.5), Inches(0.35),
                    title, font_size=16, bold=True, color=COLORS['dark_text'])
        add_text_box(slide, Inches(6.0), Inches(y_pos + 0.1), Inches(1.8), Inches(0.35),
                    pts, font_size=14, bold=True, color=colors[i], align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(7.9), Inches(y_pos + 0.1), Inches(1.8), Inches(0.35),
                    time, font_size=14, color=COLORS['gray_text'], align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.5), Inches(9.2), Inches(0.5),
                    desc, font_size=12, color=COLORS['gray_text'])
        y_pos += 1.25

    add_colored_shape(slide, Inches(0.2), Inches(4.7), Inches(9.6), Inches(0.55), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.4), Inches(4.73), Inches(9.2), Inches(0.5),
                "TOTAL: 100 Points | ~75 Minutes | Use your notecard!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_assessed_on_slide(prs):
    """Slide 3: What You'll Be Assessed On"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "What You'll Be Assessed On",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Week boxes - 2x2 grid
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(4.7), Inches(1.1), COLORS['light_tan_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(0.08), Inches(1.1), COLORS['seismic_blue'])
    add_text_box(slide, Inches(0.4), Inches(1.0), Inches(4.4), Inches(0.25),
                "W1: Plate Boundaries", font_size=12, bold=True, color=COLORS['seismic_blue'])
    add_text_box(slide, Inches(0.4), Inches(1.3), Inches(4.4), Inches(0.65),
                "• 3 boundary types\n• Convection drives motion\n• Earthquake patterns",
                font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.1), Inches(0.9), Inches(4.7), Inches(1.1), COLORS['light_blue_bg'])
    add_colored_shape(slide, Inches(5.1), Inches(0.9), Inches(0.08), Inches(1.1), COLORS['green_start'])
    add_text_box(slide, Inches(5.3), Inches(1.0), Inches(4.4), Inches(0.25),
                "W2: Seafloor Spreading", font_size=12, bold=True, color=COLORS['green_end'])
    add_text_box(slide, Inches(5.3), Inches(1.3), Inches(4.4), Inches(0.65),
                "• Magnetic stripe evidence\n• Continental drift proof\n• Pangaea evidence",
                font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(2.15), Inches(4.7), Inches(1.1), COLORS['light_orange_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(2.15), Inches(0.08), Inches(1.1), COLORS['magma_orange'])
    add_text_box(slide, Inches(0.4), Inches(2.25), Inches(4.4), Inches(0.25),
                "W3: Volcanic Eruptions", font_size=12, bold=True, color=COLORS['magma_orange'])
    add_text_box(slide, Inches(0.4), Inches(2.55), Inches(4.4), Inches(0.65),
                "• Viscosity and silica\n• Explosive vs effusive\n• Hazard prediction",
                font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.1), Inches(2.15), Inches(4.7), Inches(1.1), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(5.1), Inches(2.15), Inches(0.08), Inches(1.1), COLORS['core_red'])
    add_text_box(slide, Inches(5.3), Inches(2.25), Inches(4.4), Inches(0.25),
                "W4: Earth's Interior", font_size=12, bold=True, color=COLORS['core_red'])
    add_text_box(slide, Inches(5.3), Inches(2.55), Inches(4.4), Inches(0.65),
                "• Seismic wave evidence\n• Shadow zones\n• 4 layers and properties",
                font_size=10, color=COLORS['dark_text'])

    # Integration
    add_colored_shape(slide, Inches(0.2), Inches(3.4), Inches(9.6), Inches(0.7), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.5), Inches(9.2), Inches(0.55),
                "INTEGRATION: Connect convection → boundaries → spreading → volcanoes → interior!",
                font_size=12, bold=True, color=COLORS['green_end'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(4.25), Inches(9.6), Inches(0.5), COLORS['earth_brown'])
    add_text_box(slide, Inches(0.4), Inches(4.28), Inches(9.2), Inches(0.45),
                "KEY: Earth is a SYSTEM - all these processes connect!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_part1_intro_slide(prs):
    """Slide 4: Part 1 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['green_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Part 1 – Synthesis Review",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | ~15 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(2.3), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(2.3), COLORS['green_end'])

    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "SYNTHESIS TASK:", font_size=14, bold=True, color=COLORS['green_end'])

    task_text = """Connect ALL four weeks into one coherent understanding of Earth as a dynamic system.

You'll answer questions that require you to:
• Trace the CHAIN: Heat → Convection → Plate motion → Boundaries → Features
• Explain HOW different evidence types support our understanding
• Connect volcanic behavior to plate boundary type AND magma source"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(1.8),
                task_text, font_size=13, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(3.6), Inches(4.7), Inches(0.95), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.7), Inches(4.4), Inches(0.25),
                "FORMAT:", font_size=12, bold=True, color=COLORS['blue_accent'])
    add_text_box(slide, Inches(0.4), Inches(4.0), Inches(4.4), Inches(0.5),
                "• 4 short-answer questions\n• Trace cause-effect chains\n• Connect multiple weeks",
                font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.1), Inches(3.6), Inches(4.6), Inches(0.95), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(5.3), Inches(3.7), Inches(4.2), Inches(0.25),
                "TIPS:", font_size=12, bold=True, color=COLORS['orange_end'])
    add_text_box(slide, Inches(5.3), Inches(4.0), Inches(4.2), Inches(0.5),
                "• Think in SYSTEMS\n• Use vocabulary from all weeks\n• Draw diagrams if helpful",
                font_size=11, color=COLORS['dark_text'])


def add_part1_support_slide(prs):
    """Slide 5: Part 1 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Part 1 – The Big Picture Connections",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(9.6), Inches(2.8), COLORS['light_tan_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(9.2), Inches(0.3),
                "THE CHAIN: Earth's Heat Engine", font_size=14, bold=True, color=COLORS['earth_brown'])

    chain = """HEAT from core → CONVECTION in mantle → PLATES move

            ↓ At plate BOUNDARIES:

DIVERGENT → Seafloor spreading → Low-silica magma → SHIELD volcanoes
CONVERGENT → Subduction → High-silica magma → EXPLOSIVE volcanoes
TRANSFORM → Lateral stress → EARTHQUAKES → No volcanoes

            ↓ SEISMIC WAVES from earthquakes reveal:

LAYERS of Earth → Liquid outer core (S-wave shadow) → Solid inner core"""
    add_text_box(slide, Inches(0.4), Inches(1.3), Inches(9.2), Inches(2.2),
                chain, font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(3.8), Inches(9.6), Inches(0.7), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(3.83), Inches(9.2), Inches(0.65),
                "Everything connects! Heat drives it all, and waves reveal the structure!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_part2_intro_slide(prs):
    """Slide 6: Part 2 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['seismic_blue'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Part 2 – Cumulative Assessment",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "60 Points | ~40 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    sections = [
        ("A: Plate Boundaries", "15 pts", "Boundary types, convection, earthquake patterns"),
        ("B: Seafloor Evidence", "15 pts", "Magnetic stripes, drift evidence, spreading rates"),
        ("C: Volcanic Systems", "15 pts", "Viscosity, eruption types, hazards"),
        ("D: Earth's Interior", "15 pts", "Seismic waves, shadow zones, layer properties"),
    ]

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(2.4), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "ASSESSMENT SECTIONS:", font_size=14, bold=True, color=COLORS['seismic_blue'])

    y_pos = 1.6
    for section, pts, desc in sections:
        add_text_box(slide, Inches(0.5), Inches(y_pos), Inches(3.5), Inches(0.4),
                    section, font_size=12, bold=True, color=COLORS['dark_text'])
        add_text_box(slide, Inches(4.2), Inches(y_pos), Inches(1.2), Inches(0.4),
                    pts, font_size=12, color=COLORS['seismic_blue'])
        add_text_box(slide, Inches(5.5), Inches(y_pos), Inches(4.0), Inches(0.4),
                    desc, font_size=11, color=COLORS['gray_text'])
        y_pos += 0.5

    add_colored_shape(slide, Inches(0.2), Inches(3.7), Inches(9.6), Inches(0.85), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.8), Inches(9.2), Inches(0.25),
                "QUESTION TYPES:", font_size=12, bold=True, color=COLORS['orange_end'])
    add_text_box(slide, Inches(0.4), Inches(4.1), Inches(9.2), Inches(0.4),
                "Multiple choice • Data analysis • Short answer • Extended response (use rubric criteria)",
                font_size=11, color=COLORS['dark_text'])


def add_part2_support_slide(prs):
    """Slide 7: Part 2 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['seismic_blue'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Part 2 – Key Concepts by Section",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # 2x2 grid
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(4.7), Inches(1.35), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(4.4), Inches(0.2),
                "A: Plate Boundaries", font_size=11, bold=True, color=COLORS['blue_accent'])
    add_text_box(slide, Inches(0.4), Inches(1.2), Inches(4.4), Inches(0.95),
                "• Divergent: apart, new crust\n• Convergent: together, subduction\n• Transform: slide, earthquakes\n• Convection = driving force",
                font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.1), Inches(0.85), Inches(4.6), Inches(1.35), COLORS['light_green_bg'])
    add_text_box(slide, Inches(5.3), Inches(0.95), Inches(4.2), Inches(0.2),
                "B: Seafloor Evidence", font_size=11, bold=True, color=COLORS['green_end'])
    add_text_box(slide, Inches(5.3), Inches(1.2), Inches(4.2), Inches(0.95),
                "• Magnetic stripes = symmetrical\n• Youngest at ridge, oldest far\n• Fossils prove Pangaea\n• Spreading rate: ~2-10 cm/year",
                font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(2.35), Inches(4.7), Inches(1.35), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.45), Inches(4.4), Inches(0.2),
                "C: Volcanic Systems", font_size=11, bold=True, color=COLORS['orange_end'])
    add_text_box(slide, Inches(0.4), Inches(2.7), Inches(4.4), Inches(0.95),
                "• Silica → viscosity → style\n• Basaltic: thin, effusive\n• Rhyolitic: thick, explosive\n• Hazards: pyroclastic, lahar, ash",
                font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.1), Inches(2.35), Inches(4.6), Inches(1.35), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(5.3), Inches(2.45), Inches(4.2), Inches(0.2),
                "D: Earth's Interior", font_size=11, bold=True, color=COLORS['purple_accent'])
    add_text_box(slide, Inches(5.3), Inches(2.7), Inches(4.2), Inches(0.95),
                "• P-waves: all materials\n• S-waves: solids ONLY\n• S-shadow = liquid core proof\n• 4 layers: crust, mantle, cores",
                font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(3.85), Inches(9.6), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(3.88), Inches(9.2), Inches(0.5),
                "REMEMBER: Use vocabulary precisely! Boundary type ≠ magma type ≠ wave type",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_part3_intro_slide(prs):
    """Slide 8: Part 3 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['magma_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Part 3 – Misconception Check",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | ~20 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(1.2), COLORS['light_orange_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(1.2), COLORS['core_red'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "PURPOSE:", font_size=14, bold=True, color=COLORS['core_red'])
    purpose = """This section helps us identify common misunderstandings so we can provide targeted feedback.
Your answers help us improve instruction for everyone. Answer honestly based on your current understanding—
this is about learning, not just getting points!"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(0.7),
                purpose, font_size=12, color=COLORS['dark_text'])

    add_text_box(slide, Inches(0.3), Inches(2.5), Inches(9.4), Inches(0.3),
                "COMMON MISCONCEPTIONS WE'RE CHECKING:", font_size=13, bold=True, color=COLORS['core_red'])

    misconceptions = [
        ("plates-float-magma", "Plates float on liquid magma", "Asthenosphere is PLASTIC SOLID, not liquid"),
        ("continents-move-fast", "Continents move fast enough to notice", "~1-10 cm/year (fingernail growth rate)"),
        ("earthquakes-random", "Earthquakes happen randomly", "They cluster at plate BOUNDARIES"),
    ]

    y_pos = 2.9
    for id, wrong, correct in misconceptions:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.5), COLORS['light_pink_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.02), Inches(4.0), Inches(0.45),
                    f"❌ {wrong}", font_size=10, color=COLORS['red_accent'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(4.6), Inches(y_pos + 0.02), Inches(5.0), Inches(0.45),
                    f"✓ {correct}", font_size=10, color=COLORS['green_accent'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.55


def add_part3_support_slide(prs):
    """Slide 9: Part 3 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['core_red'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Part 3 – Think Through These Carefully",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(4.7), Inches(1.5), COLORS['light_tan_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(4.4), Inches(0.25),
                "PLASTIC SOLID, NOT LIQUID", font_size=12, bold=True, color=COLORS['earth_brown'])
    wc = """The asthenosphere is like hot taffy:
• Solid enough to carry S-waves
• Soft enough to flow SLOWLY
• Plates don't "float" - they REST on it

Think: ice cream left out - still solid, but softens!"""
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(4.4), Inches(1.0),
                wc, font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.1), Inches(0.85), Inches(4.6), Inches(1.5), COLORS['light_green_bg'])
    add_text_box(slide, Inches(5.3), Inches(0.95), Inches(4.2), Inches(0.25),
                "INCREDIBLY SLOW MOTION", font_size=12, bold=True, color=COLORS['green_end'])
    density = """Plate movement rates:
• Pacific: ~7-10 cm/year (fastest!)
• Atlantic: ~2.5 cm/year
• That's fingernail growth rate!

You'll never feel plates move - but over MILLIONS of years, continents cross oceans!"""
    add_text_box(slide, Inches(5.3), Inches(1.25), Inches(4.2), Inches(1.0),
                density, font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(2.5), Inches(9.6), Inches(1.2), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.6), Inches(9.2), Inches(0.25),
                "PATTERNS, NOT RANDOM", font_size=12, bold=True, color=COLORS['seismic_blue'])
    seasons = """Look at ANY earthquake map:
• Earthquakes cluster along LINES
• Those lines = plate BOUNDARIES
• "Ring of Fire" = Pacific Plate edges
• Random would show dots EVERYWHERE equally

Earthquake pattern = boundary pattern = NOT random!"""
    add_text_box(slide, Inches(0.4), Inches(2.9), Inches(9.2), Inches(0.75),
                seasons, font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(3.85), Inches(9.6), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(3.88), Inches(9.2), Inches(0.5),
                "Take your time—read each question carefully and think about the CORRECT explanation!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_tips_slide(prs):
    """Slide 10: Tips for Success"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Tips for Success",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    tips = [
        ("📝 USE YOUR NOTECARD", "You prepared it for this! Check all 4 weeks of content."),
        ("⏰ MANAGE YOUR TIME", "Part 1: 15 min | Part 2: 40 min | Part 3: 20 min"),
        ("📖 READ CAREFULLY", "Look for boundary types, wave types, magma types - don't mix!"),
        ("🔗 MAKE CONNECTIONS", "Everything links: heat → motion → features → evidence"),
        ("✓ CHECK YOUR WORK", "Revisit uncertain answers if time permits"),
    ]

    y_pos = 0.9
    colors = [COLORS['magma_orange'], COLORS['seismic_blue'], COLORS['core_red'],
              COLORS['green_start'], COLORS['earth_brown']]
    for i, (tip, detail) in enumerate(tips):
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.65), COLORS['light_gray_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.65), colors[i])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(3.5), Inches(0.55),
                    tip, font_size=12, bold=True, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(4.0), Inches(y_pos + 0.05), Inches(5.6), Inches(0.55),
                    detail, font_size=11, color=COLORS['gray_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.7

    add_colored_shape(slide, Inches(0.2), Inches(4.45), Inches(9.6), Inches(0.55), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.4), Inches(4.48), Inches(9.2), Inches(0.5),
                "You've got this! Show what you've learned about Earth's amazing dynamic systems!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_summary_slide(prs):
    """Slide 11: Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Cycle 6 Summary: What You Learned",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    takeaways = [
        ("Plate Dynamics", "Convection drives plates; boundaries create features"),
        ("Seafloor Evidence", "Magnetic stripes and fossils prove drift"),
        ("Volcanic Systems", "Silica → viscosity → eruption style"),
        ("Earth's Interior", "Seismic waves reveal layers we can't see"),
    ]

    y_pos = 0.9
    for title, detail in takeaways:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.65), COLORS['light_tan_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.65), COLORS['earth_brown'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.5), Inches(0.55),
                    title, font_size=13, bold=True, color=COLORS['earth_dark'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.0), Inches(y_pos + 0.05), Inches(6.6), Inches(0.55),
                    detail, font_size=12, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.72

    add_colored_shape(slide, Inches(0.2), Inches(3.85), Inches(9.6), Inches(0.65), COLORS['seismic_blue'])
    add_text_box(slide, Inches(0.4), Inches(3.88), Inches(9.2), Inches(0.6),
                "COMING UP in Cycle 7: Human Body Systems!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(4.65), Inches(9.6), Inches(0.55), COLORS['magma_orange'])
    add_text_box(slide, Inches(0.4), Inches(4.68), Inches(9.2), Inches(0.5),
                "KEY INSIGHT: Earth is a connected SYSTEM - heat drives everything we observed!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


if __name__ == "__main__":
    prs = create_presentation()
    output_path = "/home/user/Kairos.Sci.Repo/content/grade7/cycle06/week5/G7_C6_W5_Assessment_Slides_Final.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Created: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
