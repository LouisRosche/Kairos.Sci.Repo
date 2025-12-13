#!/usr/bin/env python3
"""
Create G8_C5_W1 Wave Properties & Behavior presentation.
Topic: Waves & Information Transfer.

See PPTX_DESIGN_GUIDE.md for best practices documentation.
"""

import os
from pptx.dml.color import RGBColor
from pptx_common import (
    COLORS as BASE_COLORS,
    create_base_presentation,
    add_colored_shape,
    add_text_box,
    Inches,
    Pt,
    PP_ALIGN,
    MSO_ANCHOR,
)

# Waves theme (purples and cyans) - extend base colors
COLORS = {**BASE_COLORS}
COLORS.update({
    'wave_purple': RGBColor(0x7C, 0x3A, 0xED),
    'wave_purple_dark': RGBColor(0x5B, 0x21, 0xB6),
    'wave_cyan': RGBColor(0x06, 0xB6, 0xD4),
    'wave_cyan_dark': RGBColor(0x08, 0x91, 0xB2),
    'light_wave_bg': RGBColor(0xF5, 0xF3, 0xFF),
    'em_spectrum_red': RGBColor(0xEF, 0x44, 0x44),
    'em_spectrum_blue': RGBColor(0x3B, 0x82, 0xF6),
})


def create_presentation():
    """Create the G8_C5_W1 presentation"""
    prs = create_base_presentation()

    add_title_slide(prs)
    add_phenomenon_slide(prs)
    add_driving_question_slide(prs)
    add_prior_knowledge_slide(prs)
    add_learning_targets_slide(prs)
    add_vocabulary_slide(prs)
    add_hook_intro_slide(prs)
    add_hook_support_slide(prs)
    add_station1_intro_slide(prs)
    add_station1_support_slide(prs)
    add_station2_intro_slide(prs)
    add_station2_support_slide(prs)
    add_station3_intro_slide(prs)
    add_station3_support_slide(prs)
    add_exit_ticket_slide(prs)
    add_summary_slide(prs)

    return prs


def add_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['wave_purple_dark'])

    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(1.2),
                "Waves & Information Transfer",
                font_size=44, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")

    add_text_box(slide, Inches(0.5), Inches(2.8), Inches(9), Inches(0.5),
                "Week 1: Wave Properties & Behavior",
                font_size=24, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(3.4), Inches(9), Inches(0.5),
                "Grade 8 Science | Cycle 5 | 100 Points",
                font_size=18, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Teaser box
    add_colored_shape(slide, Inches(2), Inches(4.3), Inches(6), Inches(0.8), COLORS['white'])
    add_text_box(slide, Inches(2.2), Inches(4.35), Inches(5.6), Inches(0.7),
                "How can we see through walls with WiFi but not flashlights?",
                font_size=16, bold=True, color=COLORS['wave_purple_dark'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_phenomenon_slide(prs):
    """Slide 2: Phenomenon"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['wave_purple_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "The WiFi vs Flashlight Mystery",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Main phenomenon box
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(9.6), Inches(2.4), COLORS['light_wave_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(0.08), Inches(2.4), COLORS['wave_purple'])

    phenomenon_text = """Your phone gets WiFi signal even when the router is in another room, through walls,
doors, and furniture. But shine a flashlight at that same wall—total darkness on the other side.

Both WiFi and light are WAVES. Both are electromagnetic radiation.
So why does one pass through walls while the other is blocked?

And why can X-rays see through your skin to your bones,
while radio waves pass through buildings but not mountains?"""
    add_text_box(slide, Inches(0.5), Inches(1.0), Inches(9.1), Inches(2.2),
                phenomenon_text, font_size=14, color=COLORS['dark_text'])

    # Comparison boxes
    add_colored_shape(slide, Inches(0.2), Inches(3.5), Inches(3.0), Inches(1.5), COLORS['wave_cyan'])
    add_text_box(slide, Inches(0.4), Inches(3.55), Inches(2.6), Inches(0.35),
                "WIFI", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(0.4), Inches(3.95), Inches(2.6), Inches(1.0),
                "Through walls ✓\nThrough metal ✗\nWavelength: ~12 cm",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(3.5), Inches(3.5), Inches(3.0), Inches(1.5), COLORS['wave_purple'])
    add_text_box(slide, Inches(3.7), Inches(3.55), Inches(2.6), Inches(0.35),
                "VISIBLE LIGHT", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(3.7), Inches(3.95), Inches(2.6), Inches(1.0),
                "Through walls ✗\nThrough glass ✓\nWavelength: ~500 nm",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(6.8), Inches(3.5), Inches(3.0), Inches(1.5), COLORS['em_spectrum_blue'])
    add_text_box(slide, Inches(7.0), Inches(3.55), Inches(2.6), Inches(0.35),
                "X-RAYS", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(7.0), Inches(3.95), Inches(2.6), Inches(1.0),
                "Through skin ✓\nThrough bone ✗\nWavelength: ~0.1 nm",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_driving_question_slide(prs):
    """Slide 3: Driving Question"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Driving Question",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Main question box
    add_colored_shape(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(1.2), COLORS['wave_purple_dark'])
    add_text_box(slide, Inches(0.7), Inches(1.25), Inches(8.6), Inches(1.1),
                "Why do different waves interact differently with materials?",
                font_size=28, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Sub-questions
    add_colored_shape(slide, Inches(0.2), Inches(2.7), Inches(9.6), Inches(2.0), COLORS['light_wave_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.8), Inches(9.2), Inches(0.3),
                "Questions we'll investigate:", font_size=14, bold=True, color=COLORS['wave_purple_dark'])

    questions = """• What properties do all waves share (wavelength, frequency, amplitude)?
• How are waves reflected, absorbed, and transmitted by materials?
• What makes electromagnetic waves different from mechanical waves?
• Why does wavelength determine how waves interact with matter?"""
    add_text_box(slide, Inches(0.4), Inches(3.15), Inches(9.2), Inches(1.4),
                questions, font_size=14, color=COLORS['dark_text'])

    # Connection to prior cycles
    add_colored_shape(slide, Inches(0.2), Inches(4.85), Inches(9.6), Inches(0.55), COLORS['green_start'])
    add_text_box(slide, Inches(0.4), Inches(4.88), Inches(9.2), Inches(0.5),
                "Cycle 4 Connection: How does energy transfer through ecosystems relate to wave energy?",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_prior_knowledge_slide(prs):
    """Slide 4: Prior Knowledge"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "What You Already Know",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Cycle 3 review
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(4.7), Inches(1.8), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(0.08), Inches(1.8), COLORS['purple_end'])
    add_text_box(slide, Inches(0.4), Inches(1.0), Inches(4.4), Inches(0.3),
                "From Cycle 3 (Evolution):", font_size=13, bold=True, color=COLORS['purple_end'])
    c3_review = """• Natural selection acts on variation
• Populations change over generations
• Traits can be advantageous in environment
• Evidence supports evolutionary theory"""
    add_text_box(slide, Inches(0.4), Inches(1.35), Inches(4.4), Inches(1.3),
                c3_review, font_size=12, color=COLORS['dark_text'])

    # Cycle 4 review
    add_colored_shape(slide, Inches(5.1), Inches(0.9), Inches(4.7), Inches(1.8), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(5.1), Inches(0.9), Inches(0.08), Inches(1.8), COLORS['green_end'])
    add_text_box(slide, Inches(5.3), Inches(1.0), Inches(4.4), Inches(0.3),
                "From Cycle 4 (Ecosystems):", font_size=13, bold=True, color=COLORS['green_end'])
    c4_review = """• Energy flows through systems
• Energy can be transferred/transformed
• 10% rule—energy loss at each level
• Matter cycles, energy flows"""
    add_text_box(slide, Inches(5.3), Inches(1.35), Inches(4.4), Inches(1.3),
                c4_review, font_size=12, color=COLORS['dark_text'])

    # New learning
    add_colored_shape(slide, Inches(0.2), Inches(2.9), Inches(9.6), Inches(1.7), COLORS['light_wave_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.0), Inches(9.2), Inches(0.3),
                "NEW in Cycle 5:", font_size=14, bold=True, color=COLORS['wave_purple'])
    new_learning = """Waves are how energy travels through space and matter WITHOUT moving matter itself.
Light, sound, WiFi, X-rays—all are waves carrying energy from one place to another.
Today we'll explore WHY different waves behave differently with different materials!"""
    add_text_box(slide, Inches(0.4), Inches(3.35), Inches(9.2), Inches(1.2),
                new_learning, font_size=13, color=COLORS['dark_text'])

    # Key insight
    add_colored_shape(slide, Inches(0.2), Inches(4.75), Inches(9.6), Inches(0.6), COLORS['wave_cyan'])
    add_text_box(slide, Inches(0.4), Inches(4.78), Inches(9.2), Inches(0.55),
                "KEY: Waves transfer ENERGY without transferring MATTER!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_learning_targets_slide(prs):
    """Slide 5: Learning Targets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Learning Targets",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Targets
    targets = [
        ("1", "Identify wave properties: wavelength, frequency, amplitude, speed"),
        ("2", "Explain how waves are reflected, absorbed, or transmitted through materials"),
        ("3", "Compare mechanical waves (need medium) to EM waves (no medium needed)"),
        ("4", "Design a barrier that blocks specific wave types while allowing others"),
    ]

    y_pos = 0.9
    for num, target in targets:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.7), COLORS['light_green_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.5), Inches(0.7), COLORS['green_end'])
        add_text_box(slide, Inches(0.25), Inches(y_pos + 0.05), Inches(0.4), Inches(0.6),
                    num, font_size=20, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(0.85), Inches(y_pos + 0.08), Inches(8.8), Inches(0.55),
                    target, font_size=14, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.8

    # Summary bar
    add_colored_shape(slide, Inches(0.3), Inches(4.2), Inches(9.4), Inches(0.7), COLORS['wave_purple_dark'])
    add_text_box(slide, Inches(1.5), Inches(4.25), Inches(1.5), Inches(0.35),
                "100", font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(1.5), Inches(4.55), Inches(1.5), Inches(0.3),
                "Total Points", font_size=9, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(4.25), Inches(4.25), Inches(1.5), Inches(0.35),
                "~75", font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(4.25), Inches(4.55), Inches(1.5), Inches(0.3),
                "Minutes", font_size=9, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(7.0), Inches(4.25), Inches(1.5), Inches(0.35),
                "5", font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(7.0), Inches(4.55), Inches(1.5), Inches(0.3),
                "Sections", font_size=9, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_vocabulary_slide(prs):
    """Slide 6: Key Vocabulary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.2), Inches(0.15), Inches(9.6), Inches(0.5), COLORS['wave_purple'])
    add_text_box(slide, Inches(0.4), Inches(0.18), Inches(9.2), Inches(0.45),
                "Key Vocabulary",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    vocab = [
        ("Wavelength", "Distance between wave peaks (determines type of wave)"),
        ("Frequency", "How many waves pass per second (Hz)"),
        ("Amplitude", "Height of wave—relates to energy carried"),
        ("Reflection", "Wave bounces off surface (mirror, echo)"),
        ("Absorption", "Wave energy transfers to material (heating)"),
        ("Transmission", "Wave passes through material (window)"),
    ]

    y_pos = 0.75
    for i, (term, definition) in enumerate(vocab):
        bg_color = COLORS['light_wave_bg'] if i % 2 == 0 else COLORS['white']
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.6), bg_color)
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.2), Inches(0.5),
                    term, font_size=13, bold=True, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.7), Inches(y_pos + 0.05), Inches(6.9), Inches(0.5),
                    definition, font_size=12, color=COLORS['gray_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    # Notecard prompt
    add_colored_shape(slide, Inches(0.2), Inches(4.4), Inches(9.6), Inches(0.55), COLORS['purple_start'])
    add_text_box(slide, Inches(0.4), Inches(4.43), Inches(9.2), Inches(0.5),
                "Notecard: Waves can be REFLECTED, ABSORBED, or TRANSMITTED by materials!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Quick reference
    add_colored_shape(slide, Inches(0.2), Inches(5.0), Inches(9.6), Inches(0.45), COLORS['light_wave_bg'])
    add_text_box(slide, Inches(0.4), Inches(5.03), Inches(9.2), Inches(0.4),
                "Key equation: Speed = Wavelength × Frequency (v = λf)",
                font_size=11, color=COLORS['wave_purple_dark'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_hook_intro_slide(prs):
    """Slide 7: Hook Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header with two lines
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['purple_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Hook – The WiFi vs Flashlight Mystery",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "12 Points | ~10 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Observation section
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "OBSERVE & COMPARE", font_size=14, bold=True, color=COLORS['wave_purple'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.3), COLORS['light_wave_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.3), COLORS['wave_purple'])

    observe_text = """Watch the demonstration:

1. WiFi signal strength through wall → WORKS
2. Flashlight through same wall → BLOCKED
3. Radio through building → WORKS
4. X-ray through hand → Shows BONES

All of these are electromagnetic waves!
Why do they behave so differently?"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(2.1),
                observe_text, font_size=12, color=COLORS['dark_text'])

    # Prediction section
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(1.5), COLORS['wave_cyan'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "YOUR PREDICTION", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.2), Inches(1.55), Inches(4.3), Inches(1.0),
                "Before we explain: WHY do you think some waves pass through walls while others don't?\n\nWrite your prediction in the Hook Form.",
                font_size=13, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # MTSS diagnostic
    add_colored_shape(slide, Inches(5.0), Inches(2.8), Inches(4.7), Inches(1.0), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(5.2), Inches(2.9), Inches(4.3), Inches(0.25),
                "Quick Check (0 pts):", font_size=11, bold=True, color=COLORS['orange_end'])
    add_text_box(slide, Inches(5.2), Inches(3.2), Inches(4.3), Inches(0.55),
                "Do waves carry matter or just energy?\n(This helps us know where to start!)",
                font_size=11, color=COLORS['dark_text'])

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.0), Inches(9.7), Inches(0.55), COLORS['wave_purple_dark'])
    add_text_box(slide, Inches(0.35), Inches(4.03), Inches(9.3), Inches(0.5),
                "Notecard: Write your prediction BEFORE we explain!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Cycle connection
    add_colored_shape(slide, Inches(0.15), Inches(4.65), Inches(9.7), Inches(0.55), COLORS['green_start'])
    add_text_box(slide, Inches(0.35), Inches(4.68), Inches(9.3), Inches(0.5),
                "Cycle 4 Link: How is wave energy transfer similar to energy flow in ecosystems?",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_hook_support_slide(prs):
    """Slide 8: Hook Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['wave_purple_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Hook – It's All About WAVELENGTH",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # The answer
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(9.6), Inches(1.3), COLORS['light_wave_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(0.08), Inches(1.3), COLORS['wave_purple'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(9.2), Inches(0.3),
                "THE ANSWER: Wavelength determines how waves interact with materials!", font_size=14, bold=True, color=COLORS['wave_purple'])
    answer_text = """• Long wavelengths (radio/WiFi ~cm) can "go around" small obstacles like walls
• Short wavelengths (light ~nm) are blocked by gaps smaller than their wavelength
• X-rays are so short they pass BETWEEN atoms in soft tissue, but dense bone stops them
• Material structure + wavelength = determines if wave is reflected, absorbed, or transmitted"""
    add_text_box(slide, Inches(0.4), Inches(1.3), Inches(9.2), Inches(0.8),
                answer_text, font_size=12, color=COLORS['dark_text'])

    # EM Spectrum preview
    add_colored_shape(slide, Inches(0.2), Inches(2.3), Inches(4.7), Inches(1.4), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.4), Inches(4.4), Inches(0.25),
                "THE EM SPECTRUM", font_size=12, bold=True, color=COLORS['wave_purple'])
    spectrum_text = """Long wavelength → Short wavelength
Radio → Microwave → Infrared →
Visible → UV → X-ray → Gamma

LONGER wavelength = LOWER energy
SHORTER wavelength = HIGHER energy"""
    add_text_box(slide, Inches(0.4), Inches(2.7), Inches(4.4), Inches(0.95),
                spectrum_text, font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.1), Inches(2.3), Inches(4.6), Inches(1.4), COLORS['light_cyan_bg'])
    add_text_box(slide, Inches(5.3), Inches(2.4), Inches(4.2), Inches(0.25),
                "WHY WALLS MATTER", font_size=12, bold=True, color=COLORS['wave_cyan_dark'])
    walls_text = """Wall molecules are spaced ~nm apart.

• Light (~500 nm) can't fit between → BLOCKED
• WiFi (~12 cm) is HUGE compared → passes around

Think: Basketball through fence vs
ping pong ball through same fence!"""
    add_text_box(slide, Inches(5.3), Inches(2.7), Inches(4.2), Inches(0.95),
                walls_text, font_size=11, color=COLORS['dark_text'])

    # Complete the hook
    add_colored_shape(slide, Inches(0.15), Inches(3.85), Inches(9.7), Inches(0.55), COLORS['purple_start'])
    add_text_box(slide, Inches(0.35), Inches(3.88), Inches(9.3), Inches(0.5),
                "Complete the Hook Form (12 pts) - Compare your prediction to the wavelength explanation!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station1_intro_slide(prs):
    """Slide 9: Station 1 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['orange_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 1 – Wave Tank Investigation",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | ~18 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Left - Investigation
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "YOUR INVESTIGATION", font_size=14, bold=True, color=COLORS['orange_end'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.3), COLORS['light_orange_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.3), COLORS['orange_end'])

    steps = """1. Generate waves in ripple tank simulation (2 min)
2. Observe wave reflection off barriers (3 min)
3. Observe refraction through different media (3 min)
4. Observe diffraction around obstacles (5 min)
5. Measure wavelength changes (3 min)
6. Answer analysis questions (2 min)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(2.1),
                steps, font_size=12, color=COLORS['dark_text'])

    # Right - Key behaviors
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(1.4), COLORS['orange_end'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "3 KEY WAVE BEHAVIORS", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    behaviors = """REFLECTION – Wave bounces back
REFRACTION – Wave bends at boundary
DIFFRACTION – Wave spreads around obstacle"""
    add_text_box(slide, Inches(5.2), Inches(1.55), Inches(4.3), Inches(0.9),
                behaviors, font_size=12, color=COLORS['white'])

    # Simulation link
    add_colored_shape(slide, Inches(5.0), Inches(2.65), Inches(4.7), Inches(1.15), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.2), Inches(2.75), Inches(4.3), Inches(0.25),
                "USE PhET SIMULATION:", font_size=11, bold=True, color=COLORS['blue_accent'])
    add_text_box(slide, Inches(5.2), Inches(3.0), Inches(4.3), Inches(0.75),
                "Wave Interference simulation\nCreate waves, add barriers, observe patterns!\nphet.colorado.edu/simulations/wave-interference",
                font_size=10, color=COLORS['dark_text'])

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.0), Inches(9.7), Inches(0.55), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(4.03), Inches(9.3), Inches(0.5),
                "Complete Station 1 in Form - Observe all 3 wave behaviors!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station1_support_slide(prs):
    """Slide 10: Station 1 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 1 – Understanding Wave Behaviors",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Three behaviors
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(3.1), Inches(1.5), COLORS['light_wave_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(2.7), Inches(0.3),
                "REFLECTION", font_size=12, bold=True, color=COLORS['wave_purple'])
    add_text_box(slide, Inches(0.4), Inches(1.3), Inches(2.7), Inches(1.0),
                "Wave bounces off surface\n\nExample:\n• Mirror (light)\n• Echo (sound)\n• Sonar (ships)",
                font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(3.5), Inches(0.85), Inches(3.1), Inches(1.5), COLORS['light_cyan_bg'])
    add_text_box(slide, Inches(3.7), Inches(0.95), Inches(2.7), Inches(0.3),
                "REFRACTION", font_size=12, bold=True, color=COLORS['wave_cyan_dark'])
    add_text_box(slide, Inches(3.7), Inches(1.3), Inches(2.7), Inches(1.0),
                "Wave bends at boundary\n\nExample:\n• Straw in water looks bent\n• Lens focuses light\n• Rainbow colors",
                font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(6.8), Inches(0.85), Inches(3.0), Inches(1.5), COLORS['light_green_bg'])
    add_text_box(slide, Inches(7.0), Inches(0.95), Inches(2.6), Inches(0.3),
                "DIFFRACTION", font_size=12, bold=True, color=COLORS['green_end'])
    add_text_box(slide, Inches(7.0), Inches(1.3), Inches(2.6), Inches(1.0),
                "Wave spreads around gap\n\nExample:\n• Sound around corners\n• WiFi through doorways\n• CD rainbow",
                font_size=10, color=COLORS['dark_text'])

    # Key insight
    add_colored_shape(slide, Inches(0.2), Inches(2.5), Inches(9.6), Inches(0.95), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.6), Inches(9.2), Inches(0.25),
                "KEY INSIGHT:", font_size=11, bold=True, color=COLORS['orange_end'])
    add_text_box(slide, Inches(0.4), Inches(2.9), Inches(9.2), Inches(0.5),
                "Diffraction is GREATEST when wavelength ≈ gap size. Long waves diffract more than short waves!\nThis explains why WiFi (long waves) spreads around corners better than light (short waves).",
                font_size=11, color=COLORS['dark_text'])

    # Sentence starters
    add_colored_shape(slide, Inches(0.2), Inches(3.6), Inches(9.6), Inches(0.85), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.7), Inches(9.2), Inches(0.25),
                "SENTENCE STARTERS:", font_size=11, bold=True, color=COLORS['purple_accent'])
    starters = """• "When a wave hits a barrier at an angle, it reflects because..."
• "Waves bend (refract) when entering a new medium because..."
• "Longer wavelengths diffract more because..."""
    add_text_box(slide, Inches(0.4), Inches(4.0), Inches(9.2), Inches(0.4),
                starters, font_size=10, color=COLORS['dark_text'])


def add_station2_intro_slide(prs):
    """Slide 11: Station 2 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['wave_cyan'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 2 – Electromagnetic Spectrum Exploration",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | ~15 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Left - Investigation
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "YOUR INVESTIGATION", font_size=14, bold=True, color=COLORS['wave_cyan_dark'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.3), COLORS['light_cyan_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.3), COLORS['wave_cyan_dark'])

    steps = """1. Explore the EM spectrum cards (3 min)
2. Order waves by wavelength AND energy (3 min)
3. Match wave types to their uses (3 min)
4. Compare mechanical vs EM waves (3 min)
5. Answer analysis questions (3 min)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(2.1),
                steps, font_size=12, color=COLORS['dark_text'])

    # Right - Spectrum overview
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(2.7), COLORS['wave_purple_dark'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "THE EM SPECTRUM", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    spectrum = """LONG λ → LOW energy
📻 Radio (km-m)
📶 Microwave (cm)
🔥 Infrared (μm)
🌈 Visible (nm)
☀️ Ultraviolet (nm)
🩻 X-ray (pm)
☢️ Gamma (pm)
SHORT λ → HIGH energy

ALL travel at speed of light!"""
    add_text_box(slide, Inches(5.2), Inches(1.55), Inches(4.3), Inches(2.2),
                spectrum, font_size=10, color=COLORS['white'])

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.0), Inches(9.7), Inches(0.55), COLORS['wave_cyan_dark'])
    add_text_box(slide, Inches(0.35), Inches(4.03), Inches(9.3), Inches(0.5),
                "Complete Station 2 in Form - Learn the EM spectrum order!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station2_support_slide(prs):
    """Slide 12: Station 2 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['wave_cyan_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 2 – EM vs Mechanical Waves",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Comparison
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(4.7), Inches(1.65), COLORS['light_wave_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(4.4), Inches(0.3),
                "ELECTROMAGNETIC WAVES", font_size=12, bold=True, color=COLORS['wave_purple'])
    em_text = """• DON'T need a medium
• Can travel through vacuum (space!)
• All travel at speed of light
• Examples: light, radio, X-rays

That's why sunlight reaches Earth
through empty space!"""
    add_text_box(slide, Inches(0.4), Inches(1.3), Inches(4.4), Inches(1.15),
                em_text, font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.1), Inches(0.85), Inches(4.6), Inches(1.65), COLORS['light_green_bg'])
    add_text_box(slide, Inches(5.3), Inches(0.95), Inches(4.2), Inches(0.3),
                "MECHANICAL WAVES", font_size=12, bold=True, color=COLORS['green_end'])
    mech_text = """• NEED a medium (solid/liquid/gas)
• Can't travel through vacuum
• Speed depends on medium
• Examples: sound, water waves

That's why there's no sound in space
(no air to vibrate)!"""
    add_text_box(slide, Inches(5.3), Inches(1.3), Inches(4.2), Inches(1.15),
                mech_text, font_size=10, color=COLORS['dark_text'])

    # Uses
    add_colored_shape(slide, Inches(0.2), Inches(2.65), Inches(9.6), Inches(0.9), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.75), Inches(9.2), Inches(0.25),
                "EM WAVES IN YOUR DAILY LIFE:", font_size=11, bold=True, color=COLORS['orange_end'])
    uses = """Radio (music) | Microwave (food) | Infrared (remote control) | Visible (seeing) |
UV (tanning) | X-ray (medical imaging) | Gamma (cancer treatment)"""
    add_text_box(slide, Inches(0.4), Inches(3.05), Inches(9.2), Inches(0.45),
                uses, font_size=10, color=COLORS['dark_text'])

    # Common misconception
    add_colored_shape(slide, Inches(0.2), Inches(3.7), Inches(9.6), Inches(0.7), COLORS['light_pink_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.8), Inches(9.2), Inches(0.25),
                "⚠️ COMMON MISCONCEPTION:", font_size=11, bold=True, color=COLORS['red_accent'])
    add_text_box(slide, Inches(0.4), Inches(4.05), Inches(9.2), Inches(0.3),
                "\"Waves carry matter\" → WRONG! Waves transfer ENERGY, not matter. Particles oscillate in place.",
                font_size=10, color=COLORS['dark_text'])


def add_station3_intro_slide(prs):
    """Slide 13: Station 3 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['green_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 3 – Design a Wave Barrier",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "25 Points | ~20 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Design challenge
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "ENGINEERING CHALLENGE", font_size=14, bold=True, color=COLORS['green_end'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.3), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.3), COLORS['green_end'])

    challenge = """Design a barrier for a recording studio!

Requirements:
• BLOCK sound waves (mechanical)
• ALLOW WiFi to pass (EM, ~12 cm)
• ALLOW visible light (windows!)
• Stay within $5,000 budget

Choose your materials and justify!"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(2.1),
                challenge, font_size=12, color=COLORS['dark_text'])

    # Material options
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(2.7), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "MATERIAL OPTIONS", font_size=13, bold=True, color=COLORS['blue_accent'])

    options = """Glass panel | $500
  Light ✓ | Sound ~ | WiFi ✓

Dense foam | $300
  Light ✗ | Sound ✓ | WiFi ✓

Metal mesh | $400
  Light ~ | Sound ✗ | WiFi ✗

Acoustic glass | $1,200
  Light ✓ | Sound ✓ | WiFi ✓

Concrete | $800
  Light ✗ | Sound ✓ | WiFi ~"""
    add_text_box(slide, Inches(5.2), Inches(1.55), Inches(4.3), Inches(2.2),
                options, font_size=10, color=COLORS['dark_text'])

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.0), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(4.03), Inches(9.3), Inches(0.5),
                "Complete Station 3 in Form - Design your multi-wave barrier!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station3_support_slide(prs):
    """Slide 14: Station 3 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 3 – Design Considerations",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Example calculation
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(4.7), Inches(1.5), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(4.4), Inches(0.25),
                "EXAMPLE DESIGN:", font_size=11, bold=True, color=COLORS['green_accent'])
    calc = """Recording Studio Wall:
• Acoustic glass window: $1,200
• Dense foam panels: $300 × 4 = $1,200
• Regular glass sections: $500 × 2 = $1,000
TOTAL: $3,400 ✓

Blocks sound ✓ | Allows light ✓ | Allows WiFi ✓"""
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(4.4), Inches(1.05),
                calc, font_size=10, color=COLORS['dark_text'])

    # Design questions
    add_colored_shape(slide, Inches(5.1), Inches(0.85), Inches(4.6), Inches(1.5), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(5.3), Inches(0.95), Inches(4.2), Inches(0.25),
                "THINK ABOUT:", font_size=11, bold=True, color=COLORS['orange_end'])
    questions = """• WHY does foam block sound but not WiFi?
  (Hint: wavelength differences!)

• WHY does metal mesh block WiFi?
  (Hint: it acts like a reflector)

• HOW can glass block sound but
  allow light to pass?"""
    add_text_box(slide, Inches(5.3), Inches(1.25), Inches(4.2), Inches(1.05),
                questions, font_size=10, color=COLORS['dark_text'])

    # SEP alignment
    add_colored_shape(slide, Inches(0.2), Inches(2.5), Inches(9.6), Inches(1.0), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.6), Inches(9.2), Inches(0.25),
                "SEP-2: DEVELOPING AND USING MODELS", font_size=11, bold=True, color=COLORS['purple_end'])
    sep_text = """Your design is a MODEL of wave-material interactions. Explain WHY each material choice
blocks or transmits specific wave types based on wavelength and material properties."""
    add_text_box(slide, Inches(0.4), Inches(2.9), Inches(9.2), Inches(0.5),
                sep_text, font_size=11, color=COLORS['dark_text'])

    # Final prompt
    add_colored_shape(slide, Inches(0.15), Inches(3.65), Inches(9.7), Inches(0.55), COLORS['wave_cyan'])
    add_text_box(slide, Inches(0.35), Inches(3.68), Inches(9.3), Inches(0.5),
                "KEY: Match material properties to wave properties for selective filtering!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_exit_ticket_slide(prs):
    """Slide 15: Exit Ticket"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.9), COLORS['exit_purple_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Exit Ticket – Wave Behavior Integration",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.68), Inches(9.3), Inches(0.35),
                "23 Points | ~15 min", font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Question types
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(1.5), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(0.08), Inches(1.5), COLORS['exit_purple_start'])

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.3),
                "QUESTION TYPES:", font_size=14, bold=True, color=COLORS['exit_purple_start'])

    q_types = """• 2 NEW – Wave properties, reflection/transmission, EM spectrum
• 2 SPIRAL – Cycles 3-4: Energy transfer, natural selection
• 1 INTEGRATION – Connect wave behavior to material properties
• 1 SEP-2 – Design justification (explain barrier choices)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(9.0), Inches(1.0),
                q_types, font_size=13, color=COLORS['dark_text'])

    # Tips boxes
    add_colored_shape(slide, Inches(0.3), Inches(2.85), Inches(4.6), Inches(1.1), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(2.9), Inches(4.3), Inches(0.3),
                "SUCCESS TIPS:", font_size=12, bold=True, color=COLORS['green_accent'],
                anchor=MSO_ANCHOR.MIDDLE)
    tips = """• Use vocabulary: wavelength, frequency, amplitude
• Connect wavelength to material interaction
• Remember: reflection, absorption, transmission"""
    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(4.3), Inches(0.7),
                tips, font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.1), Inches(2.85), Inches(4.6), Inches(1.1), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.3), Inches(2.9), Inches(4.2), Inches(0.3),
                "SPIRAL REMINDER:", font_size=12, bold=True, color=COLORS['blue_accent'],
                anchor=MSO_ANCHOR.MIDDLE)
    spiral = """• C3: Selection favors advantageous traits
• C4: Energy transfers through systems
• C5: Waves transfer energy without matter!"""
    add_text_box(slide, Inches(5.3), Inches(3.2), Inches(4.2), Inches(0.7),
                spiral, font_size=10, color=COLORS['dark_text'])

    # Final notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.1), Inches(9.7), Inches(0.65), COLORS['exit_purple_end'])
    add_text_box(slide, Inches(0.35), Inches(4.15), Inches(9.3), Inches(0.55),
                "FINAL Notecard: Wavelength determines wave-material interaction!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_summary_slide(prs):
    """Slide 16: Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['wave_purple_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Week 1 Summary: What You Learned",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Key takeaways
    takeaways = [
        ("Wave Properties", "Wavelength, frequency, amplitude—define wave behavior"),
        ("3 Behaviors", "Reflection (bounce), Refraction (bend), Diffraction (spread)"),
        ("EM vs Mechanical", "EM waves don't need medium; mechanical waves do"),
        ("Wavelength Matters", "Long λ diffracts more, passes through smaller gaps"),
    ]

    y_pos = 0.9
    for title, detail in takeaways:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.65), COLORS['light_wave_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.65), COLORS['wave_purple'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.4), Inches(0.55),
                    title, font_size=13, bold=True, color=COLORS['wave_purple_dark'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.9), Inches(y_pos + 0.05), Inches(6.7), Inches(0.55),
                    detail, font_size=12, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.72

    # Next week preview
    add_colored_shape(slide, Inches(0.2), Inches(3.85), Inches(9.6), Inches(0.65), COLORS['wave_cyan'])
    add_text_box(slide, Inches(0.4), Inches(3.88), Inches(9.2), Inches(0.6),
                "NEXT WEEK: How do materials block or transmit waves? Information encoding!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Key rule
    add_colored_shape(slide, Inches(0.2), Inches(4.65), Inches(9.6), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(4.68), Inches(9.2), Inches(0.5),
                "REMEMBER: Waves transfer ENERGY, not matter—wavelength is the key!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


if __name__ == "__main__":
    prs = create_presentation()
    output_path = "/home/user/Kairos.Sci.Repo/content/grade8/cycle05/week1/G8_C5_W1_Wave_Properties_Slides_Final.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Created: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
