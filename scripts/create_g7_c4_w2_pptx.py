#!/usr/bin/env python3
"""
Create G7_C4_W2 Eutrophication & Dead Zones presentation
Following established patterns - green theme
See create_g7_c3_w2_pptx.py for full PPTX DESIGN BEST PRACTICES
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

COLORS = {
    'green_primary': RGBColor(0x05, 0x96, 0x69),
    'green_dark': RGBColor(0x04, 0x78, 0x57),
    'green_light': RGBColor(0xEC, 0xFD, 0xF5),
    'blue_secondary': RGBColor(0x02, 0x84, 0xC7),
    'hook_purple': RGBColor(0x7C, 0x3A, 0xED),
    'station1_blue': RGBColor(0x3B, 0x82, 0xF6),
    'station2_purple': RGBColor(0x8B, 0x5C, 0xF6),
    'station3_pink': RGBColor(0xEC, 0x48, 0x99),
    'exit_orange': RGBColor(0xF5, 0x9E, 0x0B),
    'dark_text': RGBColor(0x2D, 0x37, 0x48),
    'gray_text': RGBColor(0x4A, 0x55, 0x68),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'light_bg': RGBColor(0xF7, 0xFA, 0xFC),
    'light_green_bg': RGBColor(0xEC, 0xFD, 0xF5),
    'light_blue_bg': RGBColor(0xDB, 0xEA, 0xFE),
    'light_purple_bg': RGBColor(0xF3, 0xE8, 0xFF),
    'light_pink_bg': RGBColor(0xFD, 0xF2, 0xF8),
    'light_orange_bg': RGBColor(0xFF, 0xFB, 0xEB),
    'red_alert': RGBColor(0xC5, 0x30, 0x30),
    'warning_bg': RGBColor(0xFE, 0xF3, 0xC7),
}

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    add_title_slide(prs)
    add_phenomenon_slide(prs)
    add_driving_question_slide(prs)
    add_prior_knowledge_slide(prs)
    add_learning_targets_slide(prs)
    add_vocabulary_slide(prs)
    add_hook_intro_slide(prs)
    add_hook_support_slide(prs)
    add_station1_intro_slide(prs)
    add_station1_support_slide(prs)
    add_station2_intro_slide(prs)
    add_station2_support_slide(prs)
    add_station3_intro_slide(prs)
    add_station3_support_slide(prs)
    add_exit_ticket_slide(prs)
    add_summary_slide(prs)

    return prs

def add_colored_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=None, align=PP_ALIGN.LEFT, font_name="Arial", anchor=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    if anchor:
        tf.anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    if color:
        p.font.color.rgb = color
    p.alignment = align
    return txBox

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['green_primary'])
    add_text_box(slide, Inches(0), Inches(1.0), Inches(10), Inches(0.8),
                "🌊 💀", font_size=48, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(1.8), Inches(9), Inches(1.0),
                "Week 2: Eutrophication & Dead Zones",
                font_size=36, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.5), Inches(2.9), Inches(9), Inches(0.5),
                "Grade 7 Science | Cycle 4 | Rosche | Kairos Academies",
                font_size=16, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(3.5), Inches(9), Inches(0.4),
                "MS-ESS3-3 Monitoring Human Impact",
                font_size=14, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(4.2), Inches(9), Inches(0.4),
                "100 Points Total | ~75 Minutes",
                font_size=14, color=COLORS['white'], align=PP_ALIGN.CENTER)

def add_phenomenon_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['green_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🌊 The Phenomenon: The Green Lake Mystery",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    add_colored_shape(slide, Inches(0.15), Inches(1.0), Inches(9.7), Inches(2.8), COLORS['warning_bg'])
    add_text_box(slide, Inches(0.4), Inches(1.1), Inches(9.2), Inches(0.5),
                "In 2014, Lake Erie turned bright green:",
                font_size=16, color=COLORS['dark_text'])
    add_text_box(slide, Inches(0.6), Inches(1.6), Inches(9.0), Inches(1.8),
                "• 500,000 people couldn't drink their tap water\n" +
                "• Fish died by the thousands\n" +
                "• Massive algae bloom covered the lake surface\n" +
                "• The culprit? Fertilizer meant to help crops\n\n" +
                "How did nutrients meant to grow food end up poisoning a lake?",
                font_size=15, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.15), Inches(3.95), Inches(9.7), Inches(1.4), COLORS['green_dark'])
    add_text_box(slide, Inches(0.35), Inches(4.15), Inches(9.3), Inches(1.0),
                "🔬 KEY CONCEPT: Excess nutrients trigger a cascade:\n" +
                "Nutrients → Algae Bloom → Death → Decomposition → Oxygen Depletion",
                font_size=16, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

def add_driving_question_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['green_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🤔 Driving Question",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    add_colored_shape(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(2.0), COLORS['green_dark'])
    add_text_box(slide, Inches(0.7), Inches(1.7), Inches(8.6), Inches(1.6),
                "What happens when nutrients run off into waterways?\n\n" +
                "How can too much of a 'good thing' (nutrients) become deadly?",
                font_size=22, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(1.4), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.7), Inches(3.9), Inches(8.6), Inches(1.2),
                "🔗 Connection to Week 1:\n" +
                "You learned how CO₂ dissolves in the ocean to cause acidification.\n" +
                "This week: A DIFFERENT chemical (N/P) causes a DIFFERENT problem (eutrophication)!",
                font_size=14, color=COLORS['green_dark'])

def add_prior_knowledge_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['green_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "💡 What You Already Know",
                font_size=22, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    concepts = [
        ("Week 1", "CO₂ causes ocean acidification\nthrough chemical reactions"),
        ("Cycle 3", "Feedback loops can make\nproblems worse (positive feedback)"),
        ("Carbon Cycle", "Matter cycles through\nEarth's systems")
    ]
    x_positions = [0.3, 3.5, 6.7]
    for i, (title, content) in enumerate(concepts):
        add_colored_shape(slide, Inches(x_positions[i]), Inches(1.1), Inches(3.0), Inches(1.8), COLORS['light_green_bg'])
        add_text_box(slide, Inches(x_positions[i] + 0.15), Inches(1.2), Inches(2.7), Inches(0.4),
                    title, font_size=14, bold=True, color=COLORS['green_dark'], align=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x_positions[i] + 0.15), Inches(1.6), Inches(2.7), Inches(1.2),
                    content, font_size=12, color=COLORS['gray_text'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.3), Inches(3.1), Inches(9.4), Inches(2.2), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(9.0), Inches(0.4),
                "🔄 NEW Connection: Eutrophication is POSITIVE FEEDBACK!",
                font_size=14, bold=True, color=COLORS['green_dark'])
    add_text_box(slide, Inches(0.5), Inches(3.6), Inches(9.0), Inches(1.6),
                "More nutrients → More algae → More dead algae → More decomposition →\n" +
                "More oxygen used → More fish die → More dead matter → MORE decomposition!\n\n" +
                "The problem makes itself WORSE - just like the feedback loops in Cycle 3!",
                font_size=13, color=COLORS['gray_text'])

def add_learning_targets_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['green_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🎯 Learning Targets",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    targets = [
        "Trace nitrogen and phosphorus through ecosystems",
        "Explain the eutrophication cascade",
        "Analyze real-world data from the Gulf of Mexico dead zone",
        "Design a solution to minimize nutrient runoff"
    ]
    y_pos = 1.0
    for i, target in enumerate(targets, 1):
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.95), COLORS['light_green_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.15), Inches(9.0), Inches(0.65),
                    f"Target {i}: {target}",
                    font_size=15, bold=True, color=COLORS['green_dark'])
        y_pos += 1.1

def add_vocabulary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['green_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "📚 Key Vocabulary This Week",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    vocab = [
        ("Eutrophication", "Excess nutrients cause algae growth → oxygen depletion"),
        ("Dead Zone", "Area where oxygen is too low for marine life (hypoxia)"),
        ("Biogeochemical Cycle", "Movement of elements (N, P, C) through living & non-living"),
        ("Nutrient Runoff", "Fertilizers/waste washing into waterways from rain"),
        ("Algae Bloom", "Rapid, explosive algae growth when nutrients abundant"),
        ("Decomposition", "Bacteria breaking down dead matter, consuming O₂")
    ]
    y_pos = 0.95
    for term, definition in vocab:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.7), COLORS['light_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.1), Inches(2.8), Inches(0.5),
                    term, font_size=12, bold=True, color=COLORS['green_dark'])
        add_text_box(slide, Inches(3.4), Inches(y_pos + 0.1), Inches(6.1), Inches(0.5),
                    definition, font_size=12, color=COLORS['gray_text'])
        y_pos += 0.75

def add_hook_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['hook_purple'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🌊 Hook – The Green Lake Mystery",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.4),
                "12 Points | ~10 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(2.2), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "🌊 The Challenge: Why did Lake Erie turn toxic green?",
                font_size=16, bold=True, color=COLORS['hook_purple'])
    add_text_box(slide, Inches(0.5), Inches(1.75), Inches(9.0), Inches(1.5),
                "1. Observe the phenomenon: Lake Erie algae bloom images (2 min)\n" +
                "2. Predict what caused the bloom (3 min)\n" +
                "3. Connect to Cycle 3: positive feedback (3 min)\n" +
                "4. Answer diagnostic questions (2 min)",
                font_size=14, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(3.55), Inches(9.4), Inches(1.8), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.65), Inches(9.0), Inches(0.4),
                "💭 Think About This:",
                font_size=14, bold=True, color=COLORS['green_dark'])
    add_text_box(slide, Inches(0.5), Inches(4.05), Inches(9.0), Inches(1.2),
                "• What do plants need to grow? (Nutrients: nitrogen, phosphorus)\n" +
                "• Where do fertilizers go when it rains?\n" +
                "• What happens when algae dies?",
                font_size=13, color=COLORS['gray_text'])

def add_hook_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['hook_purple'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🌊 Hook – The Eutrophication Cascade",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    add_colored_shape(slide, Inches(0.3), Inches(1.0), Inches(9.4), Inches(2.5), COLORS['warning_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(9.0), Inches(0.4),
                "🔄 The Cascade Process:", font_size=14, bold=True, color=COLORS['dark_text'])
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(1.9),
                "1️⃣ Nutrients (N, P) wash into water from farms, lawns, sewage\n" +
                "2️⃣ Algae EXPLODE in growth (they love nutrients!)\n" +
                "3️⃣ Algae block sunlight → underwater plants die\n" +
                "4️⃣ Algae eventually die too (overpopulation)\n" +
                "5️⃣ Bacteria decompose dead algae → USE UP OXYGEN\n" +
                "6️⃣ Fish and other organisms SUFFOCATE → Dead Zone!",
                font_size=12, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(3.7), Inches(9.4), Inches(1.65), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.78), Inches(9.0), Inches(0.5),
                "📝 Complete the Hook form on the student page",
                font_size=14, bold=True, color=COLORS['red_alert'], anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(0.5), Inches(4.2), Inches(9.0), Inches(1.0),
                "Predict what happens at each stage of the cascade.\n" +
                "How is this similar to the positive feedback loops from Cycle 3?",
                font_size=12, color=COLORS['gray_text'])

def add_station1_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['station1_blue'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🔄 Station 1 – Nutrient Cycle Modeling",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.4),
                "20 Points | ~18 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(1.3), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "🔄 Your Mission: Trace N and P Through Ecosystems",
                font_size=16, bold=True, color=COLORS['station1_blue'])
    add_text_box(slide, Inches(0.5), Inches(1.75), Inches(9.0), Inches(0.6),
                "Model how nitrogen and phosphorus move from farms to waterways\nand what happens when they accumulate.",
                font_size=14, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(2.65), Inches(4.5), Inches(2.7), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(2.75), Inches(4.1), Inches(0.4),
                "🌿 Nitrogen (N) Sources:", font_size=13, bold=True, color=COLORS['green_dark'])
    add_text_box(slide, Inches(0.5), Inches(3.15), Inches(4.1), Inches(2.1),
                "• Fertilizers (farms, lawns)\n" +
                "• Animal waste\n" +
                "• Sewage treatment plants\n" +
                "• Air pollution → rain",
                font_size=12, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(5.0), Inches(2.65), Inches(4.7), Inches(2.7), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.2), Inches(2.75), Inches(4.3), Inches(0.4),
                "💎 Phosphorus (P) Sources:", font_size=13, bold=True, color=COLORS['station1_blue'])
    add_text_box(slide, Inches(5.2), Inches(3.15), Inches(4.3), Inches(2.1),
                "• Fertilizers\n" +
                "• Detergents (laundry soap)\n" +
                "• Eroding soil\n" +
                "• Animal/human waste",
                font_size=12, color=COLORS['gray_text'])

def add_station1_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['station1_blue'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🔄 Station 1 – Biogeochemical Cycle Diagram",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    add_colored_shape(slide, Inches(0.3), Inches(1.0), Inches(9.4), Inches(3.2), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(9.0), Inches(0.4),
                "🔄 The Nutrient Journey:", font_size=14, bold=True, color=COLORS['station1_blue'])
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(2.6),
                "FARM/LAWN → RAIN RUNOFF → STREAMS → RIVERS → LAKES/OCEAN\n" +
                "     ↓              ↓              ↓           ↓           ↓\n" +
                "Fertilizer    Carries       Nutrients     Algae      DEAD\n" +
                "Applied       Nutrients     Accumulate    Bloom      ZONE\n\n" +
                "KEY CONCEPT: Nutrients don't disappear - they just MOVE!\n" +
                "What's good for crops becomes DEADLY for aquatic ecosystems.",
                font_size=12, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(4.35), Inches(9.4), Inches(1.0), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(4.38), Inches(9.0), Inches(0.95),
                "📝 Complete Station 1 form - Model the nutrient cycle and identify intervention points\n" +
                "Where could we STOP nutrients from reaching waterways?",
                font_size=12, bold=True, color=COLORS['red_alert'], anchor=MSO_ANCHOR.MIDDLE)

def add_station2_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['station2_purple'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "📊 Station 2 – Dead Zone Data Analysis",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.4),
                "20 Points | ~15 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(1.2), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "📊 Your Mission: Analyze Gulf of Mexico Dead Zone Data",
                font_size=16, bold=True, color=COLORS['station2_purple'])
    add_text_box(slide, Inches(0.5), Inches(1.75), Inches(9.0), Inches(0.5),
                "Examine real data showing the relationship between fertilizer use and dead zone size.",
                font_size=14, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(2.55), Inches(9.4), Inches(2.8), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.5), Inches(2.65), Inches(9.0), Inches(0.4),
                "📈 Gulf of Mexico Dead Zone Data:", font_size=14, bold=True, color=COLORS['station2_purple'])
    add_text_box(slide, Inches(0.5), Inches(3.1), Inches(9.0), Inches(2.1),
                "Year    |  Dead Zone Size  |  Fertilizer Use\n" +
                "1985    |    3,800 km²     |   10 million tons\n" +
                "1995    |    6,200 km²     |   12 million tons\n" +
                "2005    |    8,000 km²     |   14 million tons\n" +
                "2015    |   16,800 km²     |   15 million tons\n" +
                "2024    |   18,200 km²     |   16 million tons\n\n" +
                "Dead zone is now the size of NEW JERSEY!",
                font_size=11, color=COLORS['gray_text'])

def add_station2_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['station2_purple'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "📊 Station 2 – Correlation Analysis",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    add_colored_shape(slide, Inches(0.3), Inches(1.0), Inches(4.5), Inches(2.3), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(4.1), Inches(0.4),
                "📈 What the Data Shows:", font_size=13, bold=True, color=COLORS['green_dark'])
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4.1), Inches(1.7),
                "• Strong POSITIVE correlation\n" +
                "• More fertilizer → Bigger dead zone\n" +
                "• Dead zone grew 5× in 40 years\n" +
                "• Mississippi River carries nutrients\n  from 31 states to the Gulf",
                font_size=11, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(5.0), Inches(1.0), Inches(4.7), Inches(2.3), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.4),
                "🐟 Economic Impact:", font_size=13, bold=True, color=COLORS['station2_purple'])
    add_text_box(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(1.7),
                "• Shrimp fishery: $2.8 billion/year\n" +
                "• Fish kills hurt local economy\n" +
                "• Tourism affected\n" +
                "• Cleanup costs billions",
                font_size=11, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(3.5), Inches(9.4), Inches(1.85), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.58), Inches(9.0), Inches(0.5),
                "📝 Complete Station 2 form - Analyze the data correlation",
                font_size=14, bold=True, color=COLORS['red_alert'], anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(0.5), Inches(4.0), Inches(9.0), Inches(1.2),
                "• Calculate the rate of dead zone growth\n" +
                "• Predict dead zone size in 2030 if trends continue\n" +
                "• Is correlation = causation here? Explain your evidence.",
                font_size=12, color=COLORS['gray_text'])

def add_station3_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['station3_pink'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🔧 Station 3 – Design a Remediation Plan",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.4),
                "25 Points | ~20 min (Highest Value!)", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(1.2), COLORS['light_pink_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "🔧 Engineering Challenge: Reduce Nutrient Runoff",
                font_size=15, bold=True, color=COLORS['station3_pink'])
    add_text_box(slide, Inches(0.5), Inches(1.75), Inches(9.0), Inches(0.5),
                "Design an intervention plan to reduce dead zone size by 50% within 10 years.",
                font_size=14, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(2.55), Inches(9.4), Inches(2.8), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(2.65), Inches(9.0), Inches(0.4),
                "📋 Available Interventions:", font_size=14, bold=True, color=COLORS['green_dark'])
    add_text_box(slide, Inches(0.5), Inches(3.1), Inches(4.3), Inches(2.1),
                "Prevention (stop nutrients at source):\n" +
                "• Precision fertilizer application\n" +
                "• Buffer strips along waterways\n" +
                "• Cover crops in off-season\n" +
                "• Improved sewage treatment",
                font_size=11, color=COLORS['gray_text'])
    add_text_box(slide, Inches(5.0), Inches(3.1), Inches(4.5), Inches(2.1),
                "Remediation (treat the problem):\n" +
                "• Constructed wetlands\n" +
                "• Bioreactors in streams\n" +
                "• Algae harvesting\n" +
                "• Aeration systems",
                font_size=11, color=COLORS['gray_text'])

def add_station3_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['station3_pink'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🔧 Station 3 – Cost-Benefit Analysis",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    options = [
        ("Buffer strips", "$50/acre", "30% reduction", "Low cost, land needed"),
        ("Precision fertilizer", "$200/farm", "25% reduction", "Tech training needed"),
        ("Wetland construction", "$10,000/acre", "60% reduction", "High effectiveness"),
        ("Cover crops", "$30/acre", "20% reduction", "Extra labor required")
    ]
    y_pos = 1.0
    for option, cost, effect, notes in options:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.8), COLORS['light_pink_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.1), Inches(2.4), Inches(0.6),
                    option, font_size=12, bold=True, color=COLORS['station3_pink'])
        add_text_box(slide, Inches(3.0), Inches(y_pos + 0.1), Inches(1.5), Inches(0.6),
                    cost, font_size=11, color=COLORS['gray_text'], align=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(4.6), Inches(y_pos + 0.1), Inches(2.0), Inches(0.6),
                    effect, font_size=11, color=COLORS['gray_text'], align=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(6.7), Inches(y_pos + 0.1), Inches(2.8), Inches(0.6),
                    notes, font_size=11, color=COLORS['gray_text'])
        y_pos += 0.9

    add_colored_shape(slide, Inches(0.3), Inches(4.65), Inches(9.4), Inches(0.75), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(4.68), Inches(9.0), Inches(0.7),
                "📝 Complete Station 3 - Design a multi-intervention plan within budget!\n" +
                "Consider: What combination gives the best reduction for the lowest cost?",
                font_size=11, bold=True, color=COLORS['red_alert'], anchor=MSO_ANCHOR.MIDDLE)

def add_exit_ticket_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['exit_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🎓 Exit Ticket – Biogeochemical Systems",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.4),
                "23 Points | ~15 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(2.5), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "🎓 Show What You Learned - Question Types:",
                font_size=16, bold=True, color=COLORS['exit_orange'])
    add_text_box(slide, Inches(0.5), Inches(1.8), Inches(9.0), Inches(1.8),
                "• 2 NEW – Eutrophication content (this week)\n" +
                "• 2 SPIRAL – Week 1 review (ocean acidification, carbon budget)\n" +
                "• 1 INTEGRATION – Compare ocean acidification vs eutrophication\n" +
                "• 1 SEP-6 – Construct an explanation for dead zone formation",
                font_size=14, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(3.85), Inches(9.4), Inches(1.5), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.93), Inches(9.0), Inches(0.5),
                "📝 Complete the Exit Ticket on the student page",
                font_size=14, bold=True, color=COLORS['red_alert'], anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(0.5), Inches(4.35), Inches(9.0), Inches(0.9),
                "This is your final assessment for Week 2. Take your time!\n" +
                "Use evidence from both Week 1 and Week 2 activities.",
                font_size=13, color=COLORS['gray_text'])

def add_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['green_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🎉 Week 2 Summary: What You Learned",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    takeaways = [
        ("Eutrophication", "Excess nutrients → algae bloom → decomposition → dead zone", COLORS['green_primary']),
        ("Nutrient Sources", "Fertilizers, sewage, animal waste → runoff into waterways", COLORS['station1_blue']),
        ("Positive Feedback", "The cascade makes itself worse (like Cycle 3!)", COLORS['station2_purple']),
        ("Solutions", "Prevention (buffer strips) + Remediation (wetlands)", COLORS['station3_pink'])
    ]
    y_pos = 1.0
    for title, content, color in takeaways:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.9), COLORS['light_bg'])
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(0.1), Inches(0.9), color)
        add_text_box(slide, Inches(0.55), Inches(y_pos + 0.1), Inches(9.0), Inches(0.35),
                    title, font_size=13, bold=True, color=color)
        add_text_box(slide, Inches(0.55), Inches(y_pos + 0.45), Inches(9.0), Inches(0.4),
                    content, font_size=12, color=COLORS['gray_text'])
        y_pos += 1.0

    add_colored_shape(slide, Inches(0.3), Inches(5.05), Inches(9.4), Inches(0.45), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(5.1), Inches(9.0), Inches(0.35),
                "Next Week: Assessment – Show what you've learned about human impacts on water systems!",
                font_size=12, color=COLORS['green_dark'], align=PP_ALIGN.CENTER)

def main():
    prs = create_presentation()
    output_dir = "/home/user/Kairos.Sci.Repo/content/grade7/cycle04/week2"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "G7_C4_W2_Eutrophication_Slides_Final.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
