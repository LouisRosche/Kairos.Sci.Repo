#!/usr/bin/env python3
"""
Create G7_C4_W3 Assessment - Human Impacts on Water Systems
Following established patterns - blue/teal theme
Assessment week structure: Parts 1-3 (20/60/20 pts)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

COLORS = {
    'primary_blue': RGBColor(0x25, 0x63, 0xEB),
    'teal': RGBColor(0x08, 0x91, 0xB2),
    'green': RGBColor(0x05, 0x96, 0x69),
    'warning_orange': RGBColor(0xD9, 0x77, 0x06),
    'pink': RGBColor(0xEC, 0x48, 0x99),
    'dark_text': RGBColor(0x2D, 0x37, 0x48),
    'gray_text': RGBColor(0x4A, 0x55, 0x68),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'light_bg': RGBColor(0xF7, 0xFA, 0xFC),
    'light_blue_bg': RGBColor(0xEF, 0xF6, 0xFF),
    'light_teal_bg': RGBColor(0xEC, 0xFE, 0xFF),
    'light_green_bg': RGBColor(0xEC, 0xFD, 0xF5),
    'light_orange_bg': RGBColor(0xFE, 0xF3, 0xC7),
    'light_pink_bg': RGBColor(0xFD, 0xF2, 0xF8),
    'red_alert': RGBColor(0xC5, 0x30, 0x30),
}

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    add_title_slide(prs)
    add_overview_slide(prs)
    add_what_assessed_slide(prs)
    add_review_week1_slide(prs)
    add_review_week2_slide(prs)
    add_part1_intro_slide(prs)
    add_part1_support_slide(prs)
    add_part2_intro_slide(prs)
    add_part2_support_slide(prs)
    add_part3_intro_slide(prs)
    add_tips_slide(prs)
    add_summary_slide(prs)

    return prs

def add_colored_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=None, align=PP_ALIGN.LEFT, font_name="Arial", anchor=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    if anchor:
        tf.anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    if color:
        p.font.color.rgb = color
    p.alignment = align
    return txBox

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['primary_blue'])
    add_text_box(slide, Inches(0), Inches(1.0), Inches(10), Inches(0.8),
                "✅ 📋", font_size=48, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(1.8), Inches(9), Inches(1.0),
                "Week 3: Synthesis & Assessment",
                font_size=36, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.5), Inches(2.9), Inches(9), Inches(0.5),
                "Grade 7 Science | Cycle 4 | Rosche | Kairos Academies",
                font_size=16, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(3.5), Inches(9), Inches(0.4),
                "MS-ESS3-3 Monitoring Human Impact",
                font_size=14, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(4.2), Inches(9), Inches(0.4),
                "100 Points Total | ~75 Minutes",
                font_size=14, color=COLORS['white'], align=PP_ALIGN.CENTER)

def add_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['primary_blue'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "📋 Assessment Week Overview",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    add_colored_shape(slide, Inches(0.3), Inches(1.0), Inches(9.4), Inches(1.5), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(9.0), Inches(1.3),
                "📋 Assessment Week: This is your chance to show what you learned about\n" +
                "human impacts on water systems - both OCEAN ACIDIFICATION and EUTROPHICATION.\n" +
                "Take your time on each section!",
                font_size=14, bold=True, color=COLORS['warning_orange'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.3), Inches(2.65), Inches(9.4), Inches(2.7), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(2.75), Inches(9.0), Inches(0.4),
                "📊 Points Breakdown:", font_size=14, bold=True, color=COLORS['primary_blue'])
    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(9.0), Inches(2.0),
                "Part            |  Points  |  Time    |  What It Tests\n" +
                "────────────────────────────────────────────────────────\n" +
                "Part 1: Synthesis   |   20    |  ~15 min |  Connect Week 1 + Week 2\n" +
                "Part 2: Assessment  |   60    |  ~40 min |  Full cycle content\n" +
                "Part 3: Misconception|  20    |  ~20 min |  Common mistakes\n" +
                "────────────────────────────────────────────────────────\n" +
                "TOTAL               |  100    |  ~75 min |  Complete Assessment",
                font_size=11, color=COLORS['gray_text'])

def add_what_assessed_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "📚 What You'll Be Assessed On",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    add_colored_shape(slide, Inches(0.3), Inches(1.0), Inches(9.4), Inches(4.3), COLORS['light_teal_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(9.0), Inches(4.1),
                "Week 1: Ocean Acidification\n" +
                "• CO₂ dissolving in ocean → carbonic acid → lower pH\n" +
                "• pH changes and effects on marine life (pteropods, coral, oysters)\n" +
                "• Carbon budget calculations (sources vs sinks)\n\n" +
                "Week 2: Eutrophication\n" +
                "• Nutrient runoff → algae bloom → decomposition → dead zone\n" +
                "• Gulf of Mexico data analysis\n" +
                "• Remediation design (buffer strips, wetlands)\n\n" +
                "Connections:\n" +
                "• Both are human impacts on water systems\n" +
                "• Both involve biogeochemical cycles (C, N, P)\n" +
                "• Both show positive feedback loops (Cycle 3 spiral)\n\n" +
                "Skills: Data analysis, engineering design, constructing explanations",
                font_size=12, color=COLORS['gray_text'])

def add_review_week1_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🐚 Review: Week 1 - Ocean Acidification",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    add_colored_shape(slide, Inches(0.3), Inches(1.0), Inches(4.5), Inches(2.3), COLORS['light_teal_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(4.1), Inches(0.4),
                "🔬 Key Concepts:", font_size=13, bold=True, color=COLORS['teal'])
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4.1), Inches(1.7),
                "• CO₂ + H₂O → H₂CO₃ (carbonic acid)\n" +
                "• pH scale: lower = more acidic\n" +
                "• 0.1 pH change = 26% acidity change\n" +
                "• Ocean pH dropped from 8.25 → 8.10\n" +
                "• Shells dissolve, coral dies",
                font_size=11, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(5.0), Inches(1.0), Inches(4.7), Inches(2.3), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.4),
                "📊 Carbon Budget:", font_size=13, bold=True, color=COLORS['primary_blue'])
    add_text_box(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(1.7),
                "• Sources: 40 Gt/year (fossil fuels)\n" +
                "• Land sinks: 12 Gt (30%)\n" +
                "• Ocean sinks: 10 Gt (25%)\n" +
                "• Atmosphere: 18 Gt (45%)\n" +
                "• Trade-off: Slows climate, hurts ocean",
                font_size=11, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(3.45), Inches(9.4), Inches(1.9), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.55), Inches(9.0), Inches(0.4),
                "📝 Key Vocabulary:", font_size=13, bold=True, color=COLORS['teal'])
    add_text_box(slide, Inches(0.5), Inches(3.95), Inches(9.0), Inches(1.3),
                "Ocean Acidification | pH Scale | Carbonic Acid | Carbon Sink |\n" +
                "Logarithmic Scale | Calcifying Organisms | Mass Balance",
                font_size=12, color=COLORS['gray_text'], align=PP_ALIGN.CENTER)

def add_review_week2_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['green'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🌊 Review: Week 2 - Eutrophication",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    add_colored_shape(slide, Inches(0.3), Inches(1.0), Inches(4.5), Inches(2.3), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(4.1), Inches(0.4),
                "🔄 The Cascade:", font_size=13, bold=True, color=COLORS['green'])
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4.1), Inches(1.7),
                "1. Nutrient runoff (N, P)\n" +
                "2. Algae bloom explosion\n" +
                "3. Algae die, sink to bottom\n" +
                "4. Bacteria decompose → use O₂\n" +
                "5. Dead zone (hypoxia) forms",
                font_size=11, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(5.0), Inches(1.0), Inches(4.7), Inches(2.3), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.4),
                "📈 Gulf Dead Zone Data:", font_size=13, bold=True, color=COLORS['warning_orange'])
    add_text_box(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(1.7),
                "• 1985: 3,800 km²\n" +
                "• 2024: 18,200 km²\n" +
                "• Now = size of New Jersey!\n" +
                "• Correlation with fertilizer use\n" +
                "• 31 states drain to Gulf",
                font_size=11, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(3.45), Inches(9.4), Inches(1.9), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.55), Inches(9.0), Inches(0.4),
                "📝 Key Vocabulary:", font_size=13, bold=True, color=COLORS['green'])
    add_text_box(slide, Inches(0.5), Inches(3.95), Inches(9.0), Inches(1.3),
                "Eutrophication | Dead Zone | Biogeochemical Cycle | Nutrient Runoff |\n" +
                "Algae Bloom | Decomposition | Positive Feedback | Buffer Strips",
                font_size=12, color=COLORS['gray_text'], align=PP_ALIGN.CENTER)

def add_part1_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🔗 Part 1: Synthesis Review",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "20 Points | ~15 min | Start Here!", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(2.0), COLORS['light_teal_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "🔗 The Big Question:", font_size=16, bold=True, color=COLORS['teal'])
    add_text_box(slide, Inches(0.5), Inches(1.8), Inches(9.0), Inches(1.3),
                "How are ocean acidification and eutrophication SIMILAR human impacts?\n\n" +
                "Hint: Both involve biogeochemical cycles, both show positive feedback,\n" +
                "both require monitoring and engineering solutions.",
                font_size=14, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(3.35), Inches(9.4), Inches(2.0), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.45), Inches(9.0), Inches(0.4),
                "📝 What You'll Do (5 questions):", font_size=14, bold=True, color=COLORS['teal'])
    add_text_box(slide, Inches(0.5), Inches(3.9), Inches(9.0), Inches(1.4),
                "• Compare the cascade process in both problems\n" +
                "• Identify the positive feedback in each\n" +
                "• Connect to Cycle 3 (feedback loops)\n" +
                "• Apply concepts to a new scenario",
                font_size=12, color=COLORS['gray_text'])

def add_part1_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🔗 Part 1: Comparison Framework",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    add_colored_shape(slide, Inches(0.3), Inches(1.0), Inches(9.4), Inches(4.3), COLORS['light_teal_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(9.0), Inches(4.1),
                "Aspect           |  Ocean Acidification      |  Eutrophication\n" +
                "──────────────────────────────────────────────────────────────────\n" +
                "Pollutant        |  CO₂                      |  N and P (nutrients)\n" +
                "Source           |  Fossil fuels, industry   |  Fertilizers, sewage\n" +
                "Chemical change  |  Forms carbonic acid      |  Feeds algae growth\n" +
                "Primary effect   |  Lowers pH                |  Depletes oxygen\n" +
                "Affected life    |  Shellfish, coral         |  Fish, bottom-dwellers\n" +
                "Location         |  Global oceans            |  Coastal areas, lakes\n" +
                "Feedback type    |  Positive (warming→more)  |  Positive (more→worse)\n" +
                "Monitoring       |  pH sensors, buoys        |  Nutrient testing\n" +
                "Solutions        |  Reduce CO₂ emissions     |  Buffer strips, wetlands",
                font_size=11, color=COLORS['gray_text'])

def add_part2_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['green'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "📝 Part 2: Cumulative Assessment",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "60 Points | ~40 min | Main Assessment", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(3.0), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "📊 Question Breakdown (12 questions):", font_size=14, bold=True, color=COLORS['green'])
    add_text_box(slide, Inches(0.5), Inches(1.8), Inches(9.0), Inches(2.3),
                "Type                    |  Count  |  Topics\n" +
                "─────────────────────────────────────────────────\n" +
                "Ocean Acidification     |    3    |  pH, carbon budget\n" +
                "Eutrophication          |    3    |  Dead zones, nutrients\n" +
                "Data Analysis           |    2    |  Interpret graphs/tables\n" +
                "Engineering Design      |    2    |  Monitoring, remediation\n" +
                "Spiral (Cycle 3)        |    2    |  Feedback loops, carbon cycle",
                font_size=11, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(4.35), Inches(9.4), Inches(1.0), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(4.4), Inches(9.0), Inches(0.85),
                "⏰ Budget your time: About 3 minutes per question. Skip and return if stuck.",
                font_size=13, bold=True, color=COLORS['red_alert'], align=PP_ALIGN.CENTER)

def add_part2_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['green'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "📝 Part 2: Key Formulas & Data",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    add_colored_shape(slide, Inches(0.3), Inches(1.0), Inches(4.5), Inches(2.2), COLORS['light_teal_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(4.1), Inches(0.4),
                "🧮 Ocean Acidification:", font_size=13, bold=True, color=COLORS['teal'])
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4.1), Inches(1.6),
                "CO₂ + H₂O → H₂CO₃\n\n" +
                "pH change × 26% = acidity change\n" +
                "Example: 0.15 drop ≈ 40% more acidic\n\n" +
                "Ocean absorbs 25% of human CO₂",
                font_size=11, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(5.0), Inches(1.0), Inches(4.7), Inches(2.2), COLORS['light_green_bg'])
    add_text_box(slide, Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.4),
                "📈 Eutrophication:", font_size=13, bold=True, color=COLORS['green'])
    add_text_box(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(1.6),
                "Nutrients → Algae → Death → O₂ loss\n\n" +
                "Rate = (final - initial) / years\n" +
                "Example: (18,000-4,000)/40 = 350 km²/yr\n\n" +
                "Correlation ≠ causation (need evidence)",
                font_size=11, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(3.35), Inches(9.4), Inches(2.0), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.45), Inches(9.0), Inches(0.4),
                "📝 Complete Part 2 on the student page",
                font_size=14, bold=True, color=COLORS['red_alert'])
    add_text_box(slide, Inches(0.5), Inches(3.9), Inches(9.0), Inches(1.3),
                "This is your main assessment - 60 of your 100 points!\n" +
                "Use evidence from both Week 1 and Week 2 activities.",
                font_size=13, color=COLORS['gray_text'])

def add_part3_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['pink'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🔍 Part 3: Misconception Check",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "20 Points | ~20 min | Final Section", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(4.1), COLORS['light_pink_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "🔍 Common Misconceptions to Avoid:", font_size=14, bold=True, color=COLORS['pink'])
    add_text_box(slide, Inches(0.5), Inches(1.8), Inches(9.0), Inches(3.4),
                "❌ 'Ocean acidification means the ocean becomes acidic (below pH 7)'\n" +
                "✅ Ocean is still basic (pH 8.1), just becoming LESS basic (more acidic direction)\n\n" +
                "❌ 'Fertilizer is always bad for the environment'\n" +
                "✅ Fertilizer is helpful for crops - problems arise when it RUNS OFF into water\n\n" +
                "❌ 'Dead zones are permanent'\n" +
                "✅ Dead zones can shrink if nutrient inputs decrease (they're seasonal in Gulf)\n\n" +
                "❌ 'Ocean acidification and eutrophication are the same thing'\n" +
                "✅ Different causes (CO₂ vs nutrients), different effects (pH vs oxygen)",
                font_size=12, color=COLORS['gray_text'])

def add_tips_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['green'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "💡 Tips for Success",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    add_colored_shape(slide, Inches(0.3), Inches(1.0), Inches(9.4), Inches(4.3), COLORS['light_green_bg'])
    tips = [
        "Complete parts in order - they build on each other",
        "Take short breaks (5 min) between parts",
        "For calculations, show your work step by step",
        "For explanations, use evidence from both weeks",
        "Think about cause and effect (cascade chains)",
        "Connect to Cycle 3 feedback loops when relevant",
        "Read Part 3 feedback carefully - learn from mistakes",
        "Budget your time: ~15 + ~40 + ~20 = 75 minutes"
    ]
    y_pos = 1.1
    for tip in tips:
        add_text_box(slide, Inches(0.5), Inches(y_pos), Inches(9.0), Inches(0.45),
                    f"💡 {tip}", font_size=12, color=COLORS['green'])
        y_pos += 0.5

def add_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['primary_blue'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🎉 Cycle 4 Summary: Human Impacts on Water",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    takeaways = [
        ("Ocean Acidification", "CO₂ → carbonic acid → lower pH → shells dissolve", COLORS['teal']),
        ("Eutrophication", "Nutrients → algae → decomposition → dead zones", COLORS['green']),
        ("Biogeochemical Cycles", "C, N, P cycle through living and non-living systems", COLORS['primary_blue']),
        ("Engineering Solutions", "Monitoring + prevention + remediation", COLORS['warning_orange'])
    ]
    y_pos = 1.0
    for title, content, color in takeaways:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.9), COLORS['light_bg'])
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(0.1), Inches(0.9), color)
        add_text_box(slide, Inches(0.55), Inches(y_pos + 0.1), Inches(9.0), Inches(0.35),
                    title, font_size=13, bold=True, color=color)
        add_text_box(slide, Inches(0.55), Inches(y_pos + 0.45), Inches(9.0), Inches(0.4),
                    content, font_size=12, color=COLORS['gray_text'])
        y_pos += 1.0

    add_colored_shape(slide, Inches(0.3), Inches(5.05), Inches(9.4), Inches(0.45), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.5), Inches(5.1), Inches(9.0), Inches(0.35),
                "Good luck on your assessment! Take your time and show what you know!",
                font_size=12, color=COLORS['primary_blue'], align=PP_ALIGN.CENTER)

def main():
    prs = create_presentation()
    output_dir = "/home/user/Kairos.Sci.Repo/content/grade7/cycle04/week3"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "G7_C4_W3_Assessment_Slides_Final.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
