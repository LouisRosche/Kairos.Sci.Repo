#!/usr/bin/env python3
"""
Create G7_C3_W2 Feedback Loops & Tipping Points presentation
Following exemplar patterns from G7_C3_W1 and G8_C3_W1
Colors matched to student-page.html

=============================================================================
PPTX DESIGN BEST PRACTICES - SINGLE SOURCE OF TRUTH
=============================================================================

CYCLE STRUCTURE (All cycles follow this pattern):
-------------------------------------------------
Week 1 (Intro): Hook (12pts) + Station 1 (20pts) + Station 2 (20pts) + Station 3 (25pts) + Exit Ticket (23pts) = 100 pts
Week 2 (Deep Dive): Same structure as W1 - builds on W1 concepts
Week 3 (Assessment): Part 1 Synthesis (20pts) + Part 2 Assessment (60pts) + Part 3 Misconceptions (20pts) = 100 pts

EXCEPTION: Cycle 2 has Week 4 (extended cycle)

SLIDE STRUCTURE PATTERNS:
-------------------------
W1/W2 Standard (16 slides):
  1. Title slide
  2. Phenomenon slide
  3. Driving Question slide
  4. Prior Knowledge slide
  5. Learning Targets slide
  6. Vocabulary slide
  7. Hook intro slide
  8. Hook support slide
  9. Station 1 intro slide
  10. Station 1 support slide
  11. Station 2 intro slide
  12. Station 2 support slide
  13. Station 3 intro slide
  14. Station 3 support slide
  15. Exit Ticket intro slide
  16. Summary/What You Learned slide

W3 Assessment (11-12 slides):
  1. Title slide
  2. Assessment overview slide
  3. What you'll be assessed on slide
  4. Part 1 intro slide
  5. Part 1 support slide
  6. Part 2 intro slide
  7. Part 2 support slide
  8. Part 3 intro slide
  9. Part 3 support slide
  10. Tips for success slide
  11. Summary slide

=============================================================================

1. TEXT BOX POSITIONING:
   - Text boxes do NOT visually show actual text boundaries
   - Always add PADDING between text box edges and shape edges
   - Use vertical anchor (MSO_ANCHOR.MIDDLE) for centering in constrained heights
   - Horizontal centering via PP_ALIGN.CENTER affects text within box, not box position

2. TITLE + METADATA PATTERN:
   - Title text should NEVER share horizontal space with points/time indicators
   - Pattern: Title on one line, metadata (points, mins) on SEPARATE line below
   - OR: Place metadata as small text at far right of header bar, vertically centered

3. TEXT WRAPPING PREVENTION:
   - Calculate approximate text width: chars * font_size * 0.6 (rough estimate)
   - For long titles, either reduce font size OR split into multiple lines manually
   - Never let critical numbered items (1), 2), etc.) wrap to hidden positions

4. VERTICAL CENTERING IN SHAPES:
   - Set text_frame.paragraphs[0].alignment for horizontal
   - Set text_frame.anchor = MSO_ANCHOR.MIDDLE for vertical centering
   - Critical for small height boxes where text must appear centered

5. COLOR CONTRAST:
   - Dark text on light backgrounds, white text on dark backgrounds
   - Always test: purple/blue text on purple/blue backgrounds = BAD

6. STANDARD LAYOUT MEASUREMENTS:
   - Slide: 10" x 5.625" (16:9)
   - Header bar height: 0.6" - 0.75"
   - Standard margin: 0.15" - 0.3"
   - Standard padding inside shapes: 0.1" - 0.2"
   - Footer/notecard bar height: 0.55" - 0.7"
=============================================================================
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# HTML Color Palette (matching student-page.html)
COLORS = {
    # Main gradients
    'header_blue_start': RGBColor(0x42, 0x99, 0xE1),  # #4299E1
    'header_blue_end': RGBColor(0x2B, 0x6C, 0xB0),    # #2B6CB0
    'green_start': RGBColor(0x48, 0xBB, 0x78),        # #48BB78
    'green_end': RGBColor(0x27, 0x67, 0x49),          # #276749
    'purple_start': RGBColor(0x66, 0x7E, 0xEA),       # #667EEA (Hook)
    'purple_end': RGBColor(0x76, 0x4B, 0xA2),         # #764BA2
    'orange_start': RGBColor(0xF6, 0xAD, 0x55),       # #F6AD55 (Station 1)
    'orange_end': RGBColor(0xDD, 0x6B, 0x20),         # #DD6B20
    'exit_purple_start': RGBColor(0x9F, 0x7A, 0xEA),  # #9F7AEA
    'exit_purple_end': RGBColor(0x6B, 0x46, 0xC1),    # #6B46C1
    # Accents
    'teal': RGBColor(0x38, 0xB2, 0xAC),               # #38B2AC
    'teal_dark': RGBColor(0x23, 0x4E, 0x52),          # #234E52
    'teal_light': RGBColor(0xE6, 0xFF, 0xFA),         # #E6FFFA
    # Text colors
    'dark_text': RGBColor(0x2D, 0x37, 0x48),          # #2D3748
    'gray_text': RGBColor(0x4A, 0x55, 0x68),          # #4A5568
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    # Background colors
    'light_blue_bg': RGBColor(0xEB, 0xF8, 0xFF),      # #EBF8FF
    'light_green_bg': RGBColor(0xF0, 0xFF, 0xF4),     # #F0FFF4
    'light_orange_bg': RGBColor(0xFF, 0xFA, 0xF0),    # #FFFAF0
    'light_purple_bg': RGBColor(0xFA, 0xF5, 0xFF),    # #FAF5FF
    'light_pink_bg': RGBColor(0xFF, 0xF5, 0xF7),      # #FFF5F7
    'light_teal_bg': RGBColor(0xE6, 0xFF, 0xFA),      # #E6FFFA
    # Specific colors
    'red_accent': RGBColor(0xC5, 0x30, 0x30),         # #C53030
    'blue_accent': RGBColor(0x31, 0x82, 0xCE),        # #3182CE
    'green_accent': RGBColor(0x38, 0xA1, 0x69),       # #38A169
    'purple_accent': RGBColor(0x80, 0x5A, 0xD5),      # #805AD5
    'orange_accent': RGBColor(0xC0, 0x56, 0x21),      # #C05621
}

def create_presentation():
    """Create the G7_C3_W2 presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

    # Add all slides
    add_title_slide(prs)
    add_phenomenon_slide(prs)
    add_driving_question_slide(prs)
    add_prior_knowledge_slide(prs)
    add_learning_targets_slide(prs)
    add_vocabulary_slide(prs)
    add_hook_intro_slide(prs)
    add_hook_support_slide(prs)
    add_station1_intro_slide(prs)
    add_station1_support_slide(prs)
    add_station2_intro_slide(prs)
    add_station2_support_slide(prs)
    add_station3_intro_slide(prs)
    add_station3_support_slide(prs)
    add_exit_ticket_slide(prs)
    add_summary_slide(prs)

    return prs


def add_colored_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """Add a colored shape to the slide"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=None, align=PP_ALIGN.LEFT, font_name="Arial",
                 anchor=None):
    """
    Add a text box with specified formatting.

    IMPORTANT: Text boxes don't visually show text boundaries.
    Always add padding when positioning text inside shapes.

    Args:
        anchor: MSO_ANCHOR.MIDDLE for vertical centering, MSO_ANCHOR.TOP (default),
                MSO_ANCHOR.BOTTOM
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    # Set vertical anchor if specified (critical for centering in small boxes)
    if anchor:
        tf.anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    if color:
        p.font.color.rgb = color
    return txBox


def add_title_slide(prs):
    """Slide 1: Title slide with visual design"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background gradient shape (full slide)
    bg = add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625),
                          COLORS['header_blue_end'])

    # Main title
    title = add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(1.2),
                        "🧊 Week 2: Feedback Loops & Tipping Points 🔄",
                        font_size=42, bold=True, color=COLORS['white'],
                        align=PP_ALIGN.CENTER, font_name="Georgia")

    # Subtitle line 1
    sub1 = add_text_box(slide, Inches(0.5), Inches(2.9), Inches(9), Inches(0.5),
                       "Grade 7 Science | Rosche | Kairos Academies",
                       font_size=20, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Subtitle line 2
    sub2 = add_text_box(slide, Inches(0.5), Inches(3.4), Inches(9), Inches(0.5),
                       "MS-ESS3-5 Climate Change Investigation | 100 Points Total | ~75 Minutes",
                       font_size=16, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Driving question teaser box
    teaser_bg = add_colored_shape(slide, Inches(2), Inches(4.2), Inches(6), Inches(0.8),
                                  COLORS['white'])
    teaser_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # Make semi-transparent effect with border
    teaser_bg.line.color.rgb = COLORS['teal']
    teaser_bg.line.width = Pt(2)

    teaser = add_text_box(slide, Inches(2.1), Inches(4.35), Inches(5.8), Inches(0.5),
                         "Why does less ice = faster melting?",
                         font_size=16, color=COLORS['teal_dark'], align=PP_ALIGN.CENTER)


def add_phenomenon_slide(prs):
    """Slide 2: The Phenomenon - Melting Ice Mystery"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    add_colored_shape(slide, Inches(0.2), Inches(0.2), Inches(9.6), Inches(0.6), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.4), Inches(0.28), Inches(9.2), Inches(0.5),
                "🌡️ The Phenomenon: The Melting Ice Mystery",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia")

    # Left content area - teal background
    left_bg = add_colored_shape(slide, Inches(0.2), Inches(1.0), Inches(4.8), Inches(2.8),
                                COLORS['teal_light'])

    # Phenomenon text
    phenom_text = """Scientists discovered something strange:

• When ice STARTS melting, it doesn't melt at a steady rate
• Instead, melting SPEEDS UP over time
• The more ice that melts, the FASTER remaining ice melts

This seems backwards!"""

    add_text_box(slide, Inches(0.4), Inches(1.1), Inches(4.4), Inches(2.5),
                phenom_text, font_size=14, color=COLORS['teal_dark'])

    # Mystery question box
    mystery_bg = add_colored_shape(slide, Inches(0.2), Inches(3.9), Inches(4.8), Inches(0.7),
                                   COLORS['teal_dark'])
    add_text_box(slide, Inches(0.4), Inches(4.0), Inches(4.4), Inches(0.5),
                "Why would melting cause MORE melting?",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Right side - Quick Write box (purple like Hook)
    qw_bg = add_colored_shape(slide, Inches(5.2), Inches(1.0), Inches(4.6), Inches(3.0),
                              COLORS['purple_start'])

    add_text_box(slide, Inches(5.4), Inches(1.2), Inches(4.2), Inches(0.4),
                "📝 QUICK WRITE", font_size=18, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(5.4), Inches(1.8), Inches(4.2), Inches(1.0),
                "30 sec – silent writing", font_size=28, bold=True,
                color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(5.4), Inches(3.0), Inches(4.2), Inches(0.8),
                "What's your prediction?\nWhy might less ice = faster melting?",
                font_size=16, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Turn & Talk footer
    tt_bg = add_colored_shape(slide, Inches(5.2), Inches(4.1), Inches(4.6), Inches(0.5),
                              COLORS['light_purple_bg'])
    add_text_box(slide, Inches(5.4), Inches(4.15), Inches(4.2), Inches(0.4),
                "👥 TURN & TALK: Share with partner",
                font_size=11, bold=True, color=COLORS['purple_start'])

    # Notecard prompt at bottom
    add_colored_shape(slide, Inches(0.2), Inches(4.8), Inches(9.6), Inches(0.6), COLORS['purple_start'])
    add_text_box(slide, Inches(0.4), Inches(4.88), Inches(9.2), Inches(0.4),
                "📝 Notecard: Write your prediction - Why does melting speed up?",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_driving_question_slide(prs):
    """Slide 3: Driving Question"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Teal background
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['teal'])

    # Title
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                "DRIVING QUESTION", font_size=40, bold=True,
                color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Main question
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(1.5),
                "Why might melting ice cause MORE melting?\nWhat creates this accelerating cycle?",
                font_size=32, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")

    # Decorative divider
    add_colored_shape(slide, Inches(4), Inches(3.2), Inches(2), Inches(0.08), COLORS['white'])

    # Sub-question
    add_text_box(slide, Inches(1), Inches(3.5), Inches(8), Inches(0.6),
                "And how does this connect to what you learned about CO₂ in Week 1?",
                font_size=20, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Mission box
    mission_bg = add_colored_shape(slide, Inches(2), Inches(4.3), Inches(6), Inches(0.8), COLORS['teal_dark'])
    mission_bg.line.color.rgb = COLORS['white']
    mission_bg.line.width = Pt(2)

    add_text_box(slide, Inches(2.2), Inches(4.4), Inches(5.6), Inches(0.6),
                "🎯 Your Mission: Connect ice → albedo → feedback loops → climate",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_prior_knowledge_slide(prs):
    """Slide 4: What You Already Know (Cycle 2 + Week 1)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title - FIXED: Reduced font size to prevent wrapping
    # Original was 28pt which caused text to wrap, hiding the "1)" below boxes
    add_text_box(slide, Inches(0.3), Inches(0.1), Inches(9.4), Inches(0.5),
                "🔗 What You Already Know (Cycle 2 + Week 1)",
                font_size=24, bold=True, color=COLORS['blue_accent'],
                align=PP_ALIGN.CENTER, font_name="Georgia")

    # Top left box - Blue (Breaking bonds) - moved down slightly to ensure clearance
    add_colored_shape(slide, Inches(0.3), Inches(0.7), Inches(4.6), Inches(1.4), COLORS['light_blue_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(0.7), Inches(0.08), Inches(1.4), COLORS['blue_accent'])
    add_text_box(slide, Inches(0.5), Inches(0.8), Inches(4.3), Inches(0.3),
                "Breaking bonds...", font_size=14, bold=True, color=COLORS['blue_accent'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(4.3), Inches(0.4),
                "ABSORBS energy", font_size=18, bold=True, color=COLORS['dark_text'])
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4.3), Inches(0.4),
                "→ CO₂ vibrates, doesn't break", font_size=13, color=COLORS['gray_text'])

    # Top right box - Red (Forming bonds)
    add_colored_shape(slide, Inches(5.1), Inches(0.7), Inches(4.6), Inches(1.4), COLORS['light_pink_bg'])
    add_colored_shape(slide, Inches(5.1), Inches(0.7), Inches(0.08), Inches(1.4), COLORS['red_accent'])
    add_text_box(slide, Inches(5.3), Inches(0.8), Inches(4.3), Inches(0.3),
                "Forming bonds...", font_size=14, bold=True, color=COLORS['red_accent'])
    add_text_box(slide, Inches(5.3), Inches(1.1), Inches(4.3), Inches(0.4),
                "RELEASES energy", font_size=18, bold=True, color=COLORS['dark_text'])
    add_text_box(slide, Inches(5.3), Inches(1.5), Inches(4.3), Inches(0.4),
                "→ Combustion releases heat", font_size=13, color=COLORS['gray_text'])

    # Bottom left box - Green (Atoms conserved)
    add_colored_shape(slide, Inches(0.3), Inches(2.2), Inches(4.6), Inches(1.4), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(2.2), Inches(0.08), Inches(1.4), COLORS['green_accent'])
    add_text_box(slide, Inches(0.5), Inches(2.3), Inches(4.3), Inches(0.3),
                "Atoms are...", font_size=14, bold=True, color=COLORS['green_accent'])
    add_text_box(slide, Inches(0.5), Inches(2.6), Inches(4.3), Inches(0.4),
                "CONSERVED", font_size=18, bold=True, color=COLORS['dark_text'])
    add_text_box(slide, Inches(0.5), Inches(3.0), Inches(4.3), Inches(0.4),
                "→ Carbon cycles, never destroyed", font_size=13, color=COLORS['gray_text'])

    # Bottom right box - Purple (Week 1 greenhouse)
    add_colored_shape(slide, Inches(5.1), Inches(2.2), Inches(4.6), Inches(1.4), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(5.1), Inches(2.2), Inches(0.08), Inches(1.4), COLORS['purple_accent'])
    add_text_box(slide, Inches(5.3), Inches(2.3), Inches(4.3), Inches(0.3),
                "Week 1: Greenhouse Effect...", font_size=14, bold=True, color=COLORS['purple_accent'])
    add_text_box(slide, Inches(5.3), Inches(2.6), Inches(4.3), Inches(0.4),
                "CO₂ TRAPS heat", font_size=18, bold=True, color=COLORS['dark_text'])
    add_text_box(slide, Inches(5.3), Inches(3.0), Inches(4.3), Inches(0.4),
                "→ IR absorbed → molecules vibrate", font_size=13, color=COLORS['gray_text'])

    # Notecard prompt
    add_colored_shape(slide, Inches(0.3), Inches(3.8), Inches(9.4), Inches(0.7), COLORS['purple_start'])
    add_text_box(slide, Inches(0.5), Inches(3.9), Inches(9.0), Inches(0.5),
                "📝 Notecard: Which concept do you remember BEST? Write it + one example.",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # NEW THIS WEEK preview
    add_colored_shape(slide, Inches(0.3), Inches(4.6), Inches(9.4), Inches(0.8), COLORS['teal_light'])
    add_text_box(slide, Inches(0.5), Inches(4.7), Inches(9.0), Inches(0.6),
                "🆕 NEW THIS WEEK: How do these concepts connect to FEEDBACK LOOPS that accelerate climate change?",
                font_size=14, bold=True, color=COLORS['teal_dark'], align=PP_ALIGN.CENTER)


def add_learning_targets_slide(prs):
    """Slide 5: Learning Targets - 2x2 grid"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title with background
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['teal_light'])
    add_text_box(slide, Inches(0.3), Inches(0.1), Inches(9.4), Inches(0.5),
                "🎯 Learning Targets", font_size=32, bold=True,
                color=COLORS['teal_dark'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Target 1 - Top Left
    add_colored_shape(slide, Inches(0.3), Inches(0.7), Inches(4.6), Inches(1.3), COLORS['white'])
    add_text_box(slide, Inches(0.5), Inches(0.8), Inches(4.2), Inches(0.3),
                "TARGET 1", font_size=12, bold=True, color=COLORS['teal'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(4.2), Inches(0.7),
                "Explain how albedo affects Earth's energy balance",
                font_size=14, bold=True, color=COLORS['dark_text'])

    # Target 2 - Top Right
    add_colored_shape(slide, Inches(5.1), Inches(0.7), Inches(4.6), Inches(1.3), COLORS['white'])
    add_text_box(slide, Inches(5.3), Inches(0.8), Inches(4.2), Inches(0.3),
                "TARGET 2", font_size=12, bold=True, color=COLORS['teal'])
    add_text_box(slide, Inches(5.3), Inches(1.1), Inches(4.2), Inches(0.7),
                "Model a positive feedback loop using ice-albedo as example",
                font_size=14, bold=True, color=COLORS['dark_text'])

    # Target 3 - Bottom Left
    add_colored_shape(slide, Inches(0.3), Inches(2.1), Inches(4.6), Inches(1.3), COLORS['white'])
    add_text_box(slide, Inches(0.5), Inches(2.2), Inches(4.2), Inches(0.3),
                "TARGET 3", font_size=12, bold=True, color=COLORS['teal_dark'])
    add_text_box(slide, Inches(0.5), Inches(2.5), Inches(4.2), Inches(0.7),
                "Analyze carbon sink data and predict consequences of saturation",
                font_size=14, bold=True, color=COLORS['dark_text'])

    # Target 4 - Bottom Right
    add_colored_shape(slide, Inches(5.1), Inches(2.1), Inches(4.6), Inches(1.3), COLORS['white'])
    add_text_box(slide, Inches(5.3), Inches(2.2), Inches(4.2), Inches(0.3),
                "TARGET 4", font_size=12, bold=True, color=COLORS['teal_dark'])
    add_text_box(slide, Inches(5.3), Inches(2.5), Inches(4.2), Inches(0.7),
                "Design a carbon capture system using scientific principles",
                font_size=14, bold=True, color=COLORS['dark_text'])

    # Summary bar at bottom
    add_colored_shape(slide, Inches(0.3), Inches(3.6), Inches(9.4), Inches(0.7), COLORS['teal_dark'])

    # Stats in summary bar
    add_text_box(slide, Inches(1.5), Inches(3.7), Inches(1.5), Inches(0.5),
                "100", font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(4.0), Inches(1.5), Inches(0.25),
                "Total Points", font_size=9, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(4.25), Inches(3.7), Inches(1.5), Inches(0.5),
                "~75", font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(4.25), Inches(4.0), Inches(1.5), Inches(0.25),
                "Minutes", font_size=9, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(7.0), Inches(3.7), Inches(1.5), Inches(0.5),
                "5", font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.0), Inches(4.0), Inches(1.5), Inches(0.25),
                "Sections", font_size=9, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Points breakdown
    add_colored_shape(slide, Inches(0.3), Inches(4.5), Inches(9.4), Inches(0.9), COLORS['white'])
    add_text_box(slide, Inches(0.5), Inches(4.6), Inches(9.0), Inches(0.7),
                "🎯 Hook (12 pts)  →  ☀️ Station 1 (20 pts)  →  🌊 Station 2 (20 pts)  →  🔧 Station 3 (25 pts)  →  🎓 Exit (23 pts)",
                font_size=13, color=COLORS['dark_text'], align=PP_ALIGN.CENTER)


def add_vocabulary_slide(prs):
    """Slide 6: Key Vocabulary"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    add_colored_shape(slide, Inches(0.2), Inches(0.15), Inches(9.6), Inches(0.5), COLORS['blue_accent'])
    add_text_box(slide, Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.4),
                "📚 Key Vocabulary This Week", font_size=22, bold=True,
                color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Vocabulary table
    vocab = [
        ("Albedo", "The fraction of light a surface reflects (0 = absorbs all, 1 = reflects all)"),
        ("Positive Feedback", "When change causes MORE of the same change (amplifying loop)"),
        ("Negative Feedback", "When change causes the OPPOSITE change (stabilizing loop)"),
        ("Carbon Sink", "A reservoir that absorbs more carbon than it releases (oceans, forests)"),
        ("Tipping Point", "The threshold where a system shifts to a new, hard-to-reverse state"),
    ]

    y_pos = 0.8
    for i, (term, definition) in enumerate(vocab):
        bg_color = COLORS['light_blue_bg'] if i % 2 == 0 else COLORS['white']
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.65), bg_color)
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.1), Inches(2.2), Inches(0.4),
                    term, font_size=14, bold=True, color=COLORS['dark_text'])
        add_text_box(slide, Inches(2.7), Inches(y_pos + 0.1), Inches(6.9), Inches(0.45),
                    definition, font_size=12, color=COLORS['gray_text'])
        y_pos += 0.65

    # Notecard prompt
    add_colored_shape(slide, Inches(0.2), Inches(4.2), Inches(9.6), Inches(0.6), COLORS['purple_start'])
    add_text_box(slide, Inches(0.4), Inches(4.3), Inches(9.2), Inches(0.4),
                "📝 Notecard: Write down the term you find most confusing - we'll clarify today!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Quick reference
    add_colored_shape(slide, Inches(0.2), Inches(4.9), Inches(9.6), Inches(0.5), COLORS['teal_light'])
    add_text_box(slide, Inches(0.4), Inches(4.95), Inches(9.2), Inches(0.4),
                "💡 Remember: Positive feedback = AMPLIFIES change | Negative feedback = STABILIZES system",
                font_size=12, color=COLORS['teal_dark'], align=PP_ALIGN.CENTER)


def add_hook_intro_slide(prs):
    """Slide 7: Hook Introduction (12 pts | ~10 min)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header gradient (purple) - FIXED: Two-line header to prevent overlap
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['purple_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🎯 Hook – The Melting Ice Mystery",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "12 Points | ~10 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Left side - What You'll Do
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "WHAT YOU'LL DO", font_size=14, bold=True, color=COLORS['purple_start'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.0), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.0), COLORS['purple_start'])

    steps = """1. Observe the phenomenon: accelerating ice melt (2 min)
2. Make predictions about WHY melting speeds up (3 min)
3. Learn about feedback loops (3 min)
4. Connect to Week 1 greenhouse effect (2 min)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(1.8),
                steps, font_size=13, color=COLORS['dark_text'])

    # Key question box
    add_colored_shape(slide, Inches(0.3), Inches(3.6), Inches(4.5), Inches(0.8), COLORS['light_orange_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(3.6), Inches(4.5), Inches(0.05), COLORS['orange_accent'])
    add_text_box(slide, Inches(0.5), Inches(3.7), Inches(4.1), Inches(0.6),
                "💡 Key: What color is ice? What color is ocean?\nWhich absorbs more sunlight?",
                font_size=12, color=COLORS['orange_accent'])

    # Right side - Quick Write prompt
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(2.5), COLORS['purple_start'])
    add_text_box(slide, Inches(5.2), Inches(1.25), Inches(4.3), Inches(0.4),
                "📝 QUICK WRITE FIRST", font_size=16, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.2), Inches(1.8), Inches(4.3), Inches(1.5),
                "\"I think melting ice causes more melting because...\"",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Cold call warning
    add_colored_shape(slide, Inches(5.0), Inches(3.7), Inches(4.7), Inches(0.4), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(5.2), Inches(3.75), Inches(4.3), Inches(0.3),
                "🎲 Cold Call Coming! Be ready to share.",
                font_size=11, bold=True, color=COLORS['purple_start'])

    # Notecard prompt
    add_colored_shape(slide, Inches(0.15), Inches(4.55), Inches(9.7), Inches(0.55), COLORS['purple_start'])
    add_text_box(slide, Inches(0.35), Inches(4.6), Inches(9.3), Inches(0.4),
                "📝 Notecard: Record your prediction AND your reasoning",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_hook_support_slide(prs):
    """Slide 8: Hook Support - Think About This"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['purple_end'])
    add_text_box(slide, Inches(0.35), Inches(0.22), Inches(9.3), Inches(0.45),
                "🎯 Hook – Think About This",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia")

    # Three thinking prompts
    prompts = [
        ("What color is ice?", "White/light - reflects most sunlight", COLORS['light_blue_bg'], COLORS['blue_accent']),
        ("What color is ocean water?", "Dark blue - absorbs most sunlight", COLORS['light_teal_bg'], COLORS['teal_dark']),
        ("Which surface heats up more?", "Dark surfaces absorb more → get hotter!", COLORS['light_orange_bg'], COLORS['orange_accent']),
    ]

    x_pos = 0.25
    for question, answer, bg_color, text_color in prompts:
        add_colored_shape(slide, Inches(x_pos), Inches(0.95), Inches(3.1), Inches(1.6), bg_color)
        add_text_box(slide, Inches(x_pos + 0.15), Inches(1.05), Inches(2.8), Inches(0.5),
                    question, font_size=14, bold=True, color=text_color)
        add_text_box(slide, Inches(x_pos + 0.15), Inches(1.5), Inches(2.8), Inches(0.9),
                    answer, font_size=13, color=COLORS['dark_text'])
        x_pos += 3.2

    # Connection to Week 1 box
    add_colored_shape(slide, Inches(0.25), Inches(2.7), Inches(9.5), Inches(1.0), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.25), Inches(2.7), Inches(0.08), Inches(1.0), COLORS['green_accent'])
    add_text_box(slide, Inches(0.5), Inches(2.8), Inches(9.0), Inches(0.3),
                "🔗 Week 1 Connection", font_size=14, bold=True, color=COLORS['green_accent'])
    add_text_box(slide, Inches(0.5), Inches(3.1), Inches(9.0), Inches(0.5),
                "Last week: CO₂ absorbs infrared heat → warms atmosphere → ice melts. This week: What happens AFTER ice melts?",
                font_size=13, color=COLORS['dark_text'])

    # The Loop visualization
    add_colored_shape(slide, Inches(0.25), Inches(3.85), Inches(9.5), Inches(0.85), COLORS['teal_light'])
    add_text_box(slide, Inches(0.5), Inches(3.95), Inches(9.0), Inches(0.65),
                "🔄 THE LOOP: More CO₂ → More heat → Ice melts → Dark ocean exposed → Absorbs MORE heat → MORE ice melts → ...",
                font_size=14, bold=True, color=COLORS['teal_dark'], align=PP_ALIGN.CENTER)

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.85), Inches(9.7), Inches(0.55), COLORS['purple_start'])
    add_text_box(slide, Inches(0.35), Inches(4.9), Inches(9.3), Inches(0.4),
                "📝 Notecard: Draw a simple diagram showing this loop with arrows",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_station1_intro_slide(prs):
    """Slide 9: Station 1 Introduction - Albedo Effect (20 pts | ~18 min)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header (Orange) - FIXED: Two-line header to prevent overlap
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['orange_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "☀️ Station 1 – Albedo Effect Investigation",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "20 Points | ~18 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Left side - Investigation steps
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "YOUR INVESTIGATION", font_size=14, bold=True, color=COLORS['orange_end'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.3), COLORS['light_orange_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.3), COLORS['orange_end'])

    steps = """1. Set up heat lamp with 4 surfaces (3 min)
   • Black paper, white paper, aluminum foil, water
2. Record starting temperatures (2 min)
3. Turn on heat lamp for 3 minutes (3 min)
4. Record final temperatures (2 min)
5. Calculate temperature changes (3 min)
6. Answer analysis questions (5 min)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(2.1),
                steps, font_size=12, color=COLORS['dark_text'])

    # Right side - Discovery Goal
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(1.5), COLORS['orange_end'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "🎯 DISCOVERY GOAL", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.2), Inches(1.6), Inches(4.3), Inches(0.9),
                "Which surfaces heat up fastest? Why does surface color matter for climate?",
                font_size=15, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Prediction box
    add_colored_shape(slide, Inches(5.0), Inches(2.7), Inches(4.7), Inches(0.9), COLORS['light_orange_bg'])
    add_colored_shape(slide, Inches(5.0), Inches(2.7), Inches(4.7), Inches(0.05), COLORS['orange_end'])
    add_text_box(slide, Inches(5.2), Inches(2.8), Inches(4.3), Inches(0.7),
                "📝 Predict: \"I think _____ will heat up most because...\"",
                font_size=13, color=COLORS['orange_end'])

    # Safety reminder
    add_colored_shape(slide, Inches(5.0), Inches(3.7), Inches(4.7), Inches(0.7), COLORS['light_pink_bg'])
    add_text_box(slide, Inches(5.2), Inches(3.8), Inches(4.3), Inches(0.5),
                "⚠️ Safety: Do NOT touch heat lamp bulb!\nKeep water away from electrical equipment",
                font_size=11, color=COLORS['red_accent'])

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.55), Inches(9.7), Inches(0.55), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(4.6), Inches(9.3), Inches(0.4),
                "📝 Notecard: Record your prediction BEFORE the experiment",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_station1_support_slide(prs):
    """Slide 10: Station 1 Support - Albedo Data & Scaffolds"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.4),
                "☀️ Station 1 – Albedo Reference & Data",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia")

    # Left - Albedo definition
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(4.7), Inches(1.5), COLORS['light_orange_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(0.08), Inches(1.5), COLORS['orange_end'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(4.4), Inches(0.3),
                "ALBEDO = fraction of light REFLECTED", font_size=13, bold=True, color=COLORS['dark_text'])
    albedo_info = """• Albedo of 0 = absorbs ALL light (black)
• Albedo of 1 = reflects ALL light (mirror)
• Snow/Ice ≈ 0.8-0.9 (reflects most)
• Ocean ≈ 0.06 (absorbs most)"""
    add_text_box(slide, Inches(0.4), Inches(1.3), Inches(4.4), Inches(1.0),
                albedo_info, font_size=11, color=COLORS['dark_text'])

    # Right - Home alternative data
    add_colored_shape(slide, Inches(5.1), Inches(0.85), Inches(4.6), Inches(1.5), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.3), Inches(0.95), Inches(4.2), Inches(0.3),
                "🏠 HOME ALTERNATIVE DATA:", font_size=12, bold=True, color=COLORS['blue_accent'])
    data_text = """Surface        | Start | After 3 min | Change
Black paper   | 22°C  | 38°C        | +16°C
Water          | 22°C  | 31°C        | +9°C
White paper   | 22°C  | 27°C        | +5°C
Aluminum foil | 22°C  | 24°C        | +2°C"""
    add_text_box(slide, Inches(5.3), Inches(1.3), Inches(4.2), Inches(0.95),
                data_text, font_size=10, color=COLORS['dark_text'])

    # Sentence starters
    add_colored_shape(slide, Inches(0.2), Inches(2.5), Inches(4.7), Inches(1.3), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.6), Inches(4.4), Inches(0.25),
                "📗 SENTENCE STARTERS:", font_size=12, bold=True, color=COLORS['green_accent'])
    starters = """• \"Black paper heated most because its albedo is low, which means...\"
• \"Surfaces with high albedo, like aluminum, stayed cooler because...\"
• \"The data shows darker surfaces absorb more _____ energy...\" """
    add_text_box(slide, Inches(0.4), Inches(2.9), Inches(4.4), Inches(0.85),
                starters, font_size=10, color=COLORS['dark_text'])

    # Word bank
    add_colored_shape(slide, Inches(5.1), Inches(2.5), Inches(4.6), Inches(1.3), COLORS['teal_light'])
    add_text_box(slide, Inches(5.3), Inches(2.6), Inches(4.2), Inches(0.25),
                "🔤 WORD BANK:", font_size=12, bold=True, color=COLORS['teal_dark'])
    words = "reflects • absorbs • albedo • light energy • temperature • feedback • amplify • positive feedback"
    add_text_box(slide, Inches(5.3), Inches(2.9), Inches(4.2), Inches(0.85),
                words, font_size=11, color=COLORS['teal_dark'])

    # The key insight
    add_colored_shape(slide, Inches(0.2), Inches(3.95), Inches(9.5), Inches(0.7), COLORS['orange_end'])
    add_text_box(slide, Inches(0.4), Inches(4.05), Inches(9.1), Inches(0.5),
                "🔑 KEY INSIGHT: When ice melts → dark ocean exposed → absorbs MORE heat → MORE melting = POSITIVE FEEDBACK!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.8), Inches(9.7), Inches(0.55), COLORS['orange_start'])
    add_text_box(slide, Inches(0.35), Inches(4.85), Inches(9.3), Inches(0.4),
                "📝 Notecard: Which surface heated most? How does this explain accelerating ice melt?",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_station2_intro_slide(prs):
    """Slide 11: Station 2 Introduction - Carbon Sink Analysis (20 pts | ~15 min)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header (Blue) - FIXED: Two-line header to prevent overlap
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['header_blue_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🌊 Station 2 – Carbon Sink Analysis",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "20 Points | ~15 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Left side - Analysis steps
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "YOUR ANALYSIS", font_size=14, bold=True, color=COLORS['header_blue_end'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.0), COLORS['light_blue_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.0), COLORS['header_blue_end'])

    steps = """1. Study the carbon budget data (3 min)
2. Calculate: How much CO₂ stays in atmosphere? (3 min)
3. Analyze what happens if sinks decrease (4 min)
4. Predict consequences of sink saturation (5 min)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(1.8),
                steps, font_size=13, color=COLORS['dark_text'])

    # Right side - Discovery Goal
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(1.5), COLORS['header_blue_end'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "🎯 DISCOVERY GOAL", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.2), Inches(1.6), Inches(4.3), Inches(0.9),
                "Why is the atmosphere gaining CO₂ even though Earth has carbon sinks?",
                font_size=15, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Math preview
    add_colored_shape(slide, Inches(5.0), Inches(2.7), Inches(4.7), Inches(0.9), COLORS['light_blue_bg'])
    add_colored_shape(slide, Inches(5.0), Inches(2.7), Inches(4.7), Inches(0.05), COLORS['blue_accent'])
    add_text_box(slide, Inches(5.2), Inches(2.8), Inches(4.3), Inches(0.7),
                "🔢 Calculator needed!\nSources - Sinks = What stays in atmosphere",
                font_size=13, color=COLORS['blue_accent'], align=PP_ALIGN.CENTER)

    # Key insight preview
    add_colored_shape(slide, Inches(0.3), Inches(3.65), Inches(9.4), Inches(0.7), COLORS['teal_light'])
    add_text_box(slide, Inches(0.5), Inches(3.75), Inches(9.0), Inches(0.5),
                "💡 Hint: Earth's sinks can only absorb SO MUCH carbon. What happens when they're full?",
                font_size=13, color=COLORS['teal_dark'], align=PP_ALIGN.CENTER)

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.55), Inches(9.7), Inches(0.55), COLORS['header_blue_end'])
    add_text_box(slide, Inches(0.35), Inches(4.6), Inches(9.3), Inches(0.4),
                "📝 Notecard: Write your prediction - What happens when carbon sinks are full?",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_station2_support_slide(prs):
    """Slide 12: Station 2 Support - Carbon Budget Data"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['header_blue_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.4),
                "🌊 Station 2 – Global Carbon Budget Data",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia")

    # Left - Sources
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(4.7), Inches(1.5), COLORS['light_pink_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(4.4), Inches(0.3),
                "📤 SOURCES (Gt CO₂/year)", font_size=13, bold=True, color=COLORS['red_accent'])
    sources = """Fossil fuels:        +36 Gt/year
Deforestation:      +5 Gt/year
───────────────
TOTAL EMITTED:    41 Gt/year"""
    add_text_box(slide, Inches(0.4), Inches(1.3), Inches(4.4), Inches(1.0),
                sources, font_size=12, color=COLORS['dark_text'])

    # Right - Sinks
    add_colored_shape(slide, Inches(5.1), Inches(0.85), Inches(4.6), Inches(1.5), COLORS['light_green_bg'])
    add_text_box(slide, Inches(5.3), Inches(0.95), Inches(4.2), Inches(0.3),
                "📥 SINKS (Gt CO₂/year)", font_size=13, bold=True, color=COLORS['green_accent'])
    sinks = """Ocean absorption:    -10 Gt/year
Land/forest:            -12 Gt/year
───────────────
TOTAL ABSORBED:  22 Gt/year"""
    add_text_box(slide, Inches(5.3), Inches(1.3), Inches(4.2), Inches(1.0),
                sinks, font_size=12, color=COLORS['dark_text'])

    # The calculation
    add_colored_shape(slide, Inches(0.2), Inches(2.5), Inches(9.5), Inches(0.7), COLORS['red_accent'])
    add_text_box(slide, Inches(0.4), Inches(2.6), Inches(9.1), Inches(0.5),
                "REMAINING IN ATMOSPHERE: 41 - 22 = 19 Gt/year (and growing!)",
                font_size=18, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # What-if scenarios
    add_colored_shape(slide, Inches(0.2), Inches(3.35), Inches(4.7), Inches(1.2), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.45), Inches(4.4), Inches(0.25),
                "🤔 WHAT IF scenarios:", font_size=12, bold=True, color=COLORS['orange_accent'])
    scenarios = """• Ocean gets too warm to absorb as much?
• Forests burn or are cut down?
• Permafrost melts and RELEASES carbon?"""
    add_text_box(slide, Inches(0.4), Inches(3.75), Inches(4.4), Inches(0.7),
                scenarios, font_size=11, color=COLORS['dark_text'])

    # Another feedback loop
    add_colored_shape(slide, Inches(5.1), Inches(3.35), Inches(4.6), Inches(1.2), COLORS['teal_light'])
    add_text_box(slide, Inches(5.3), Inches(3.45), Inches(4.2), Inches(0.25),
                "🔄 ANOTHER POSITIVE FEEDBACK:", font_size=12, bold=True, color=COLORS['teal_dark'])
    feedback = """More CO₂ → Warmer oceans →
Oceans absorb LESS CO₂ →
Even MORE CO₂ stays in atmosphere!"""
    add_text_box(slide, Inches(5.3), Inches(3.75), Inches(4.2), Inches(0.7),
                feedback, font_size=11, color=COLORS['teal_dark'])

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.7), Inches(9.7), Inches(0.55), COLORS['header_blue_end'])
    add_text_box(slide, Inches(0.35), Inches(4.75), Inches(9.3), Inches(0.4),
                "📝 Notecard: Calculate - If sinks decrease by 50%, how much CO₂ stays in atmosphere?",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_station3_intro_slide(prs):
    """Slide 13: Station 3 Introduction - Engineering Carbon Capture (25 pts | ~20 min)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header (Green) - FIXED: Two-line header to prevent overlap
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🔧 Station 3 – Engineering Carbon Capture",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "25 Points | ~20 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Left side - Design challenge
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "YOUR ENGINEERING CHALLENGE", font_size=14, bold=True, color=COLORS['green_end'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.0), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.0), COLORS['green_end'])

    steps = """1. Review carbon capture options (3 min)
2. Choose your approach (2 min)
3. Calculate how much CO₂ your design captures (5 min)
4. Design your system with constraints (5 min)
5. Explain the science behind your choice (5 min)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(1.8),
                steps, font_size=12, color=COLORS['dark_text'])

    # Right side - Constraints
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(1.8), COLORS['green_end'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "📋 DESIGN CONSTRAINTS", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    constraints = """✓ Must work at SCHOOL scale (not industrial)
✓ Must be SUSTAINABLE (low energy)
✓ Must store carbon for 10+ YEARS
✓ Budget: $500 MAXIMUM"""
    add_text_box(slide, Inches(5.2), Inches(1.6), Inches(4.3), Inches(1.2),
                constraints, font_size=13, color=COLORS['white'])

    # Highest value reminder
    add_colored_shape(slide, Inches(5.0), Inches(3.0), Inches(4.7), Inches(0.6), COLORS['light_green_bg'])
    add_text_box(slide, Inches(5.2), Inches(3.1), Inches(4.3), Inches(0.4),
                "⭐ HIGHEST VALUE STATION - Show your best work!",
                font_size=12, bold=True, color=COLORS['green_end'], align=PP_ALIGN.CENTER)

    # Key question
    add_colored_shape(slide, Inches(0.3), Inches(3.65), Inches(9.4), Inches(0.7), COLORS['teal_light'])
    add_text_box(slide, Inches(0.5), Inches(3.75), Inches(9.0), Inches(0.5),
                "🔬 How can we CREATE a carbon sink? What captures CO₂ naturally?",
                font_size=14, color=COLORS['teal_dark'], align=PP_ALIGN.CENTER)

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.55), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(4.6), Inches(9.3), Inches(0.4),
                "📝 Notecard: Which carbon capture method interests you most? Why?",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_station3_support_slide(prs):
    """Slide 14: Station 3 Support - Carbon Capture Options"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.4),
                "🔧 Station 3 – Carbon Capture Options & Data",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia")

    # Options table
    options = [
        ("🌳 Tree planting", "~48 lbs CO₂/tree/year", "$10-50/tree"),
        ("🌿 Algae cultivation", "10x faster than trees!", "$100+ setup"),
        ("♻️ Composting", "Stores carbon in soil", "$50 for bins"),
        ("☀️ Solar panels", "Prevents NEW emissions", "$$$"),
    ]

    y_pos = 0.85
    for i, (approach, method, cost) in enumerate(options):
        bg_color = COLORS['light_green_bg'] if i % 2 == 0 else COLORS['white']
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.55), bg_color)
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.1), Inches(2.5), Inches(0.35),
                    approach, font_size=13, bold=True, color=COLORS['dark_text'])
        add_text_box(slide, Inches(3.0), Inches(y_pos + 0.1), Inches(3.5), Inches(0.35),
                    method, font_size=12, color=COLORS['gray_text'])
        add_text_box(slide, Inches(7.0), Inches(y_pos + 0.1), Inches(2.5), Inches(0.35),
                    cost, font_size=12, color=COLORS['green_accent'])
        y_pos += 0.55

    # Calculation help
    add_colored_shape(slide, Inches(0.2), Inches(3.15), Inches(4.7), Inches(1.2), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.25), Inches(4.4), Inches(0.25),
                "🔢 CALCULATION HELP:", font_size=12, bold=True, color=COLORS['blue_accent'])
    calc_help = """Example: 10 trees × 48 lbs/tree = 480 lbs CO₂/year
Cost: 10 trees × $25 = $250 (under budget!)
That's about 1/10 of one person's annual footprint"""
    add_text_box(slide, Inches(0.4), Inches(3.55), Inches(4.4), Inches(0.75),
                calc_help, font_size=10, color=COLORS['dark_text'])

    # Design tips
    add_colored_shape(slide, Inches(5.1), Inches(3.15), Inches(4.6), Inches(1.2), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(5.3), Inches(3.25), Inches(4.2), Inches(0.25),
                "💡 DESIGN TIPS:", font_size=12, bold=True, color=COLORS['orange_accent'])
    tips = """• Combine approaches for best results
• Consider maintenance (who takes care of it?)
• Think about WHERE to put your solution
• Explain the SCIENCE - why does it work?"""
    add_text_box(slide, Inches(5.3), Inches(3.55), Inches(4.2), Inches(0.75),
                tips, font_size=10, color=COLORS['dark_text'])

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.5), Inches(9.7), Inches(0.55), COLORS['green_start'])
    add_text_box(slide, Inches(0.35), Inches(4.55), Inches(9.3), Inches(0.4),
                "📝 Notecard: Write your design choice + calculate total CO₂ captured per year",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_exit_ticket_slide(prs):
    """Slide 15: Exit Ticket (23 pts | ~15 min)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header (Purple) - FIXED: Title and metadata on separate lines to prevent overlap
    # Using taller header bar (0.9") to accommodate two lines
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.9), COLORS['exit_purple_start'])
    # Title centered on first line
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "🎓 Exit Ticket – Feedback Loop Integration",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    # Points/time on separate line below title
    add_text_box(slide, Inches(0.35), Inches(0.7), Inches(9.3), Inches(0.3),
                "23 Points | ~15 min", font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Question types - adjusted positions to account for taller header
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(1.6), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(0.08), Inches(1.6), COLORS['exit_purple_start'])

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.3),
                "QUESTION TYPES:", font_size=14, bold=True, color=COLORS['exit_purple_start'])

    q_types = """• 2 NEW – Week 2 feedback loop content
• 2 SPIRAL – Week 1 + Cycle 2 review
• 1 INTEGRATION – Connect Week 1 & Week 2 concepts
• 1 SEP-2 – Create a feedback loop MODEL (diagram with labels)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(9.0), Inches(1.1),
                q_types, font_size=13, color=COLORS['dark_text'])

    # Tips box - adjusted positions
    add_colored_shape(slide, Inches(0.3), Inches(2.95), Inches(4.6), Inches(1.1), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.05), Inches(4.3), Inches(0.25),
                "✅ SUCCESS TIPS:", font_size=12, bold=True, color=COLORS['green_accent'])
    tips = """• Use vocabulary from today
• Reference your notecard notes
• Draw arrows in your feedback loop model
• Explain WHY each step leads to the next"""
    add_text_box(slide, Inches(0.5), Inches(3.3), Inches(4.3), Inches(0.7),
                tips, font_size=10, color=COLORS['dark_text'])

    # Model reminder - adjusted positions
    add_colored_shape(slide, Inches(5.1), Inches(2.95), Inches(4.6), Inches(1.1), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.3), Inches(3.05), Inches(4.2), Inches(0.25),
                "📊 SEP-2 MODEL MUST INCLUDE:", font_size=12, bold=True, color=COLORS['blue_accent'])
    model_req = """• Clear boxes or circles with labels
• Arrows showing cause → effect
• At least 4 steps in the loop
• \"Positive feedback\" or \"Negative feedback\" label"""
    add_text_box(slide, Inches(5.3), Inches(3.3), Inches(4.2), Inches(0.7),
                model_req, font_size=10, color=COLORS['dark_text'])

    # Final notecard - adjusted position
    add_colored_shape(slide, Inches(0.15), Inches(4.2), Inches(9.7), Inches(0.65), COLORS['exit_purple_end'])
    add_text_box(slide, Inches(0.35), Inches(4.3), Inches(9.3), Inches(0.45),
                "📝 FINAL Notecard: Draw the ice-albedo positive feedback loop with 4+ labeled steps",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_summary_slide(prs):
    """Slide 16: Week 2 Summary"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['teal_light'])

    # Title
    add_text_box(slide, Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5),
                "Week 2 Summary: What You Learned", font_size=26, bold=True,
                color=COLORS['teal_dark'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Four key takeaways
    takeaways = [
        ("Albedo Effect", "White surfaces reflect light (stay cool), dark surfaces absorb light (heat up)", COLORS['header_blue_end']),
        ("Positive Feedback", "Change that causes MORE of the same change (amplifying loop)", COLORS['green_accent']),
        ("Carbon Sinks", "Reservoirs that absorb CO₂ - can become saturated", COLORS['orange_end']),
        ("Multiple Feedbacks", "Climate change involves many interacting feedback loops", COLORS['exit_purple_start']),
    ]

    y_pos = 0.85
    for title, description, accent_color in takeaways:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.8), COLORS['white'])
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(0.1), Inches(0.8), accent_color)
        add_text_box(slide, Inches(0.55), Inches(y_pos + 0.1), Inches(9.0), Inches(0.3),
                    title + ":", font_size=14, bold=True, color=COLORS['dark_text'])
        add_text_box(slide, Inches(0.55), Inches(y_pos + 0.4), Inches(9.0), Inches(0.35),
                    description, font_size=12, color=COLORS['gray_text'])
        y_pos += 0.9

    # Completion celebration
    add_colored_shape(slide, Inches(0.3), Inches(4.5), Inches(9.4), Inches(0.8), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.5), Inches(4.6), Inches(9.0), Inches(0.3),
                "🎉 Week 2 Complete!", font_size=18, bold=True,
                color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(4.9), Inches(9.0), Inches(0.3),
                "Next Week: Synthesis & Assessment - Bringing it all together!",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)


if __name__ == "__main__":
    # Create the presentation
    prs = create_presentation()

    # Save to the correct location
    output_path = "/home/user/Kairos.Sci.Repo/content/grade7/cycle03/week2/G7_C3_W2_Feedback_Loops_Slides_Final.pptx"

    # Ensure directory exists
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    # Save
    prs.save(output_path)
    print(f"Created: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
