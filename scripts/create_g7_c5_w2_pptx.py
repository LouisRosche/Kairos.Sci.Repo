#!/usr/bin/env python3
"""
Create G7_C5_W2 Weather Prediction & Climate Patterns presentation.
Topic: Weather & Climate Systems.

See PPTX_DESIGN_GUIDE.md for best practices documentation.
"""

import os
from pptx.dml.color import RGBColor
from pptx_common import (
    COLORS as BASE_COLORS,
    create_base_presentation,
    add_colored_shape,
    add_text_box,
    Inches,
    Pt,
    PP_ALIGN,
    MSO_ANCHOR,
)

# Weather theme (blues and grays) - extend base colors
COLORS = {**BASE_COLORS}
COLORS.update({
    'storm_gray': RGBColor(0x4A, 0x55, 0x68),
    'storm_dark': RGBColor(0x2D, 0x37, 0x48),
    'cyan_start': RGBColor(0x00, 0xB5, 0xD8),
    'cyan_end': RGBColor(0x00, 0x86, 0x9B),
    'light_cyan_bg': RGBColor(0xE0, 0xF7, 0xFA),
    'light_gray_bg': RGBColor(0xF7, 0xFA, 0xFC),
})


def create_presentation():
    """Create the G7_C5_W2 presentation"""
    prs = create_base_presentation()

    add_title_slide(prs)
    add_phenomenon_slide(prs)
    add_driving_question_slide(prs)
    add_prior_knowledge_slide(prs)
    add_learning_targets_slide(prs)
    add_vocabulary_slide(prs)
    add_hook_intro_slide(prs)
    add_hook_support_slide(prs)
    add_station1_intro_slide(prs)
    add_station1_support_slide(prs)
    add_station2_intro_slide(prs)
    add_station2_support_slide(prs)
    add_station3_intro_slide(prs)
    add_station3_support_slide(prs)
    add_exit_ticket_slide(prs)
    add_summary_slide(prs)

    return prs


def add_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['header_blue_end'])

    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(1.2),
                "Weather & Climate Systems",
                font_size=44, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")

    add_text_box(slide, Inches(0.5), Inches(2.8), Inches(9), Inches(0.5),
                "Week 2: Weather Prediction & Climate Patterns",
                font_size=24, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(3.4), Inches(9), Inches(0.5),
                "Grade 7 Science | Cycle 5 | 100 Points",
                font_size=18, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Teaser box
    add_colored_shape(slide, Inches(2), Inches(4.3), Inches(6), Inches(0.8), COLORS['white'])
    add_text_box(slide, Inches(2.2), Inches(4.35), Inches(5.6), Inches(0.7),
                "Why is weather so hard to predict more than a week out?",
                font_size=16, bold=True, color=COLORS['header_blue_end'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_phenomenon_slide(prs):
    """Slide 2: Phenomenon"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['header_blue_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "The Failed Forecast Mystery",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Main phenomenon box
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(9.6), Inches(2.4), COLORS['light_blue_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(0.08), Inches(2.4), COLORS['header_blue_end'])

    phenomenon_text = """On Monday, the weatherman said "sunny skies all week!" By Thursday,
you're stuck in a thunderstorm without an umbrella.

The 10-day forecast predicted clear weather, but 3 days in, everything changed.

Meanwhile, climate scientists can tell you with confidence that summers will be
hotter in 50 years—but can't tell you if it will rain next Tuesday.

How can we predict decades ahead but struggle with next week?"""
    add_text_box(slide, Inches(0.5), Inches(1.0), Inches(9.1), Inches(2.2),
                phenomenon_text, font_size=14, color=COLORS['dark_text'])

    # Comparison boxes
    add_colored_shape(slide, Inches(0.2), Inches(3.5), Inches(4.7), Inches(1.5), COLORS['cyan_start'])
    add_text_box(slide, Inches(0.4), Inches(3.55), Inches(4.3), Inches(0.35),
                "WEATHER FORECASTS", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(0.4), Inches(3.95), Inches(4.3), Inches(1.0),
                "1-3 days: ~90% accurate\n4-7 days: ~50% accurate\n10+ days: Not much better than guessing",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(5.1), Inches(3.5), Inches(4.7), Inches(1.5), COLORS['green_start'])
    add_text_box(slide, Inches(5.3), Inches(3.55), Inches(4.3), Inches(0.35),
                "CLIMATE PREDICTIONS", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(5.3), Inches(3.95), Inches(4.3), Inches(1.0),
                "50-year trends: Very confident\nSeasonal patterns: High accuracy\nAverage conditions: Predictable",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_driving_question_slide(prs):
    """Slide 3: Driving Question"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Driving Question",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Main question box
    add_colored_shape(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(1.2), COLORS['header_blue_end'])
    add_text_box(slide, Inches(0.7), Inches(1.25), Inches(8.6), Inches(1.1),
                "Why is weather harder to predict than climate?",
                font_size=28, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Sub-questions
    add_colored_shape(slide, Inches(0.2), Inches(2.7), Inches(9.6), Inches(2.0), COLORS['light_cyan_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.8), Inches(9.2), Inches(0.3),
                "Questions we'll investigate:", font_size=14, bold=True, color=COLORS['teal_dark'])

    questions = """• What data do meteorologists collect to predict weather?
• How do small changes in initial conditions affect forecasts?
• What's the difference between weather patterns and climate trends?
• Why are long-term climate predictions more reliable than next week's forecast?"""
    add_text_box(slide, Inches(0.4), Inches(3.15), Inches(9.2), Inches(1.4),
                questions, font_size=14, color=COLORS['dark_text'])

    # Connection to Week 1
    add_colored_shape(slide, Inches(0.2), Inches(4.85), Inches(9.6), Inches(0.55), COLORS['green_start'])
    add_text_box(slide, Inches(0.4), Inches(4.88), Inches(9.2), Inches(0.5),
                "Week 1 Connection: How do air mass movements create the weather we're trying to predict?",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_prior_knowledge_slide(prs):
    """Slide 4: Prior Knowledge"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "What You Already Know",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Week 1 review
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(4.7), Inches(1.8), COLORS['light_blue_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(0.08), Inches(1.8), COLORS['cyan_start'])
    add_text_box(slide, Inches(0.4), Inches(1.0), Inches(4.4), Inches(0.3),
                "From Week 1 (Air Masses):", font_size=13, bold=True, color=COLORS['cyan_end'])
    w1_review = """• Air masses have distinct temperature/humidity
• 4 types: mT, cT, mP, cP
• Fronts form where air masses meet
• Pressure systems drive weather patterns"""
    add_text_box(slide, Inches(0.4), Inches(1.35), Inches(4.4), Inches(1.3),
                w1_review, font_size=12, color=COLORS['dark_text'])

    # Cycle 3-4 review
    add_colored_shape(slide, Inches(5.1), Inches(0.9), Inches(4.7), Inches(1.8), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(5.1), Inches(0.9), Inches(0.08), Inches(1.8), COLORS['purple_end'])
    add_text_box(slide, Inches(5.3), Inches(1.0), Inches(4.4), Inches(0.3),
                "From Cycles 3-4 (Climate):", font_size=13, bold=True, color=COLORS['purple_end'])
    c34_review = """• Greenhouse effect warms Earth
• Human activities affect climate
• CO₂ levels influence temperature
• Climate = long-term patterns"""
    add_text_box(slide, Inches(5.3), Inches(1.35), Inches(4.4), Inches(1.3),
                c34_review, font_size=12, color=COLORS['dark_text'])

    # New learning
    add_colored_shape(slide, Inches(0.2), Inches(2.9), Inches(9.6), Inches(1.7), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.0), Inches(9.2), Inches(0.3),
                "NEW in Week 2:", font_size=14, bold=True, color=COLORS['green_end'])
    new_learning = """Today we explore WHY weather is so chaotic while climate is predictable.
It's like trying to predict exactly when a single leaf will fall vs. knowing autumn comes every year.
We'll analyze real weather data and understand the chaos that makes prediction so hard!"""
    add_text_box(slide, Inches(0.4), Inches(3.35), Inches(9.2), Inches(1.2),
                new_learning, font_size=13, color=COLORS['dark_text'])

    # Key distinction
    add_colored_shape(slide, Inches(0.2), Inches(4.75), Inches(9.6), Inches(0.6), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.78), Inches(9.2), Inches(0.55),
                "KEY: Weather = the coin flip | Climate = the loaded coin (we know the odds)",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_learning_targets_slide(prs):
    """Slide 5: Learning Targets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Learning Targets",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Targets
    targets = [
        ("1", "Analyze weather data to identify patterns and variability"),
        ("2", "Explain why weather predictions become less accurate over time"),
        ("3", "Distinguish between weather variability and climate trends"),
        ("4", "Design a data collection system for climate monitoring"),
    ]

    y_pos = 0.9
    for num, target in targets:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.7), COLORS['light_green_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.5), Inches(0.7), COLORS['green_end'])
        add_text_box(slide, Inches(0.25), Inches(y_pos + 0.05), Inches(0.4), Inches(0.6),
                    num, font_size=20, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(0.85), Inches(y_pos + 0.08), Inches(8.8), Inches(0.55),
                    target, font_size=14, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.8

    # Summary bar
    add_colored_shape(slide, Inches(0.3), Inches(4.2), Inches(9.4), Inches(0.7), COLORS['teal_dark'])
    add_text_box(slide, Inches(1.5), Inches(4.25), Inches(1.5), Inches(0.35),
                "100", font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(1.5), Inches(4.55), Inches(1.5), Inches(0.3),
                "Total Points", font_size=9, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(4.25), Inches(4.25), Inches(1.5), Inches(0.35),
                "~75", font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(4.25), Inches(4.55), Inches(1.5), Inches(0.3),
                "Minutes", font_size=9, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(7.0), Inches(4.25), Inches(1.5), Inches(0.35),
                "5", font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(7.0), Inches(4.55), Inches(1.5), Inches(0.3),
                "Sections", font_size=9, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_vocabulary_slide(prs):
    """Slide 6: Key Vocabulary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.2), Inches(0.15), Inches(9.6), Inches(0.5), COLORS['blue_accent'])
    add_text_box(slide, Inches(0.4), Inches(0.18), Inches(9.2), Inches(0.45),
                "Key Vocabulary",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    vocab = [
        ("Forecast", "Prediction of future weather based on current data"),
        ("Chaos", "Small changes can cause large, unpredictable effects"),
        ("Variability", "How much something changes over time (day-to-day)"),
        ("Trend", "Long-term direction of change over years/decades"),
        ("Model", "Computer simulation using data to predict outcomes"),
        ("Climate Normal", "30-year average of weather conditions for an area"),
    ]

    y_pos = 0.75
    for i, (term, definition) in enumerate(vocab):
        bg_color = COLORS['light_blue_bg'] if i % 2 == 0 else COLORS['white']
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.6), bg_color)
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.2), Inches(0.5),
                    term, font_size=13, bold=True, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.7), Inches(y_pos + 0.05), Inches(6.9), Inches(0.5),
                    definition, font_size=12, color=COLORS['gray_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    # Notecard prompt
    add_colored_shape(slide, Inches(0.2), Inches(4.4), Inches(9.6), Inches(0.55), COLORS['purple_start'])
    add_text_box(slide, Inches(0.4), Inches(4.43), Inches(9.2), Inches(0.5),
                "Notecard: Weather = short-term VARIABILITY | Climate = long-term TREND",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Quick reference
    add_colored_shape(slide, Inches(0.2), Inches(5.0), Inches(9.6), Inches(0.45), COLORS['light_cyan_bg'])
    add_text_box(slide, Inches(0.4), Inches(5.03), Inches(9.2), Inches(0.4),
                "Remember: Chaos makes individual events hard to predict, but patterns emerge over time!",
                font_size=11, color=COLORS['teal_dark'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_hook_intro_slide(prs):
    """Slide 7: Hook Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header with two lines
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['purple_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Hook – The Failed Forecast Mystery",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "12 Points | ~10 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Data analysis section
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "ANALYZE REAL DATA", font_size=14, bold=True, color=COLORS['purple_end'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.3), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.3), COLORS['purple_end'])

    analyze_text = """Examine the forecast vs. actual data:

Day 1 forecast: 72°F, sunny
Day 1 actual:   74°F, sunny ✓

Day 5 forecast: 68°F, cloudy
Day 5 actual:   61°F, heavy rain ✗

Why did the 5-day forecast miss so badly?"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(2.1),
                analyze_text, font_size=12, color=COLORS['dark_text'])

    # Prediction section
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(1.5), COLORS['cyan_start'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "YOUR PREDICTION", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.2), Inches(1.55), Inches(4.3), Inches(1.0),
                "Before we explain: WHY do you think longer forecasts are less accurate?\n\nWrite your prediction in the Hook Form.",
                font_size=13, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # MTSS diagnostic
    add_colored_shape(slide, Inches(5.0), Inches(2.8), Inches(4.7), Inches(1.0), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(5.2), Inches(2.9), Inches(4.3), Inches(0.25),
                "Quick Check (0 pts):", font_size=11, bold=True, color=COLORS['orange_end'])
    add_text_box(slide, Inches(5.2), Inches(3.2), Inches(4.3), Inches(0.55),
                "What's the difference between weather and climate?\n(This helps us know where to start!)",
                font_size=11, color=COLORS['dark_text'])

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.0), Inches(9.7), Inches(0.55), COLORS['purple_end'])
    add_text_box(slide, Inches(0.35), Inches(4.03), Inches(9.3), Inches(0.5),
                "Notecard: Write your prediction BEFORE we explain!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Week 1 connection
    add_colored_shape(slide, Inches(0.15), Inches(4.65), Inches(9.7), Inches(0.55), COLORS['green_start'])
    add_text_box(slide, Inches(0.35), Inches(4.68), Inches(9.3), Inches(0.5),
                "Week 1 Link: How do the air masses we learned about make prediction difficult?",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_hook_support_slide(prs):
    """Slide 8: Hook Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['purple_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Hook – The Butterfly Effect",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # The answer
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(9.6), Inches(1.3), COLORS['light_blue_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(0.08), Inches(1.3), COLORS['cyan_start'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(9.2), Inches(0.3),
                "THE ANSWER: Weather is CHAOTIC", font_size=14, bold=True, color=COLORS['cyan_end'])
    answer_text = """• Small errors in initial data grow larger over time
• "Butterfly Effect" — tiny changes can lead to completely different outcomes
• After ~7-10 days, small measurement errors have amplified too much
• We CAN'T overcome this—it's a fundamental limit of the atmosphere"""
    add_text_box(slide, Inches(0.4), Inches(1.3), Inches(9.2), Inches(0.8),
                answer_text, font_size=12, color=COLORS['dark_text'])

    # Key concept boxes
    add_colored_shape(slide, Inches(0.2), Inches(2.3), Inches(4.7), Inches(1.4), COLORS['light_cyan_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.4), Inches(4.4), Inches(0.25),
                "WHY CLIMATE IS DIFFERENT", font_size=12, bold=True, color=COLORS['cyan_end'])
    climate_text = """Climate = long-term AVERAGES
• We don't predict specific events
• We predict patterns and trends
• Like coin flips: can't predict one,
  but know 50% heads over 1000 flips
• The "noise" averages out over time"""
    add_text_box(slide, Inches(0.4), Inches(2.7), Inches(4.4), Inches(0.95),
                climate_text, font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.1), Inches(2.3), Inches(4.6), Inches(1.4), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(5.3), Inches(2.4), Inches(4.2), Inches(0.25),
                "THE CHAOS MATH", font_size=12, bold=True, color=COLORS['purple_end'])
    chaos_text = """Day 1: Temp is 72.1°F (measured 72°F)
→ 0.1° error seems tiny

By Day 5: That 0.1° error has grown
→ Maybe predicts sunny instead of rain

By Day 10: Error has snowballed
→ Forecast is nearly worthless"""
    add_text_box(slide, Inches(5.3), Inches(2.7), Inches(4.2), Inches(0.95),
                chaos_text, font_size=11, color=COLORS['dark_text'])

    # Complete the hook
    add_colored_shape(slide, Inches(0.15), Inches(3.85), Inches(9.7), Inches(0.55), COLORS['purple_start'])
    add_text_box(slide, Inches(0.35), Inches(3.88), Inches(9.3), Inches(0.5),
                "Complete the Hook Form (12 pts) - Compare your prediction to the chaos explanation!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station1_intro_slide(prs):
    """Slide 9: Station 1 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['orange_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 1 – Data Collection & Analysis",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | ~18 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Left - Investigation
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "YOUR INVESTIGATION", font_size=14, bold=True, color=COLORS['orange_end'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.3), COLORS['light_orange_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.3), COLORS['orange_end'])

    steps = """1. Analyze 7-day weather data (3 min)
2. Calculate daily temperature variability (4 min)
3. Compare forecast accuracy at 1, 3, 5 days (5 min)
4. Graph error growth over forecast time (4 min)
5. Answer analysis questions (2 min)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(2.1),
                steps, font_size=12, color=COLORS['dark_text'])

    # Right - Data preview
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(1.4), COLORS['orange_end'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "DATA YOU'LL ANALYZE", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    data_text = """• Temperature readings (actual vs forecast)
• Precipitation records (mm)
• Pressure measurements (mb)
• Wind speed/direction data"""
    add_text_box(slide, Inches(5.2), Inches(1.55), Inches(4.3), Inches(0.9),
                data_text, font_size=12, color=COLORS['white'])

    # Key vocabulary
    add_colored_shape(slide, Inches(5.0), Inches(2.65), Inches(4.7), Inches(1.15), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.2), Inches(2.75), Inches(4.3), Inches(0.25),
                "KEY CALCULATION:", font_size=11, bold=True, color=COLORS['blue_accent'])
    calc_text = """Error = |Forecast - Actual|
Average error tells us forecast quality
Watch how error GROWS with time!"""
    add_text_box(slide, Inches(5.2), Inches(3.0), Inches(4.3), Inches(0.75),
                calc_text, font_size=10, color=COLORS['dark_text'])

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.0), Inches(9.7), Inches(0.55), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(4.03), Inches(9.3), Inches(0.5),
                "Complete Station 1 in Form - Analyze the weather data!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station1_support_slide(prs):
    """Slide 10: Station 1 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 1 – Understanding Error Growth",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Error growth table
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(9.6), Inches(0.45), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.88), Inches(9.2), Inches(0.4),
                "Watch how forecast error typically grows over time:",
                font_size=13, bold=True, color=COLORS['orange_end'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Data rows
    data = [
        ("1-Day Forecast", "±2°F", "~95%", "Very reliable"),
        ("3-Day Forecast", "±5°F", "~80%", "Good for planning"),
        ("5-Day Forecast", "±8°F", "~50%", "Watch for updates"),
        ("10-Day Forecast", "±12°F", "~25%", "General trends only"),
    ]

    y_pos = 1.4
    for i, (day, error, accuracy, note) in enumerate(data):
        bg_color = COLORS['light_blue_bg'] if i % 2 == 0 else COLORS['white']
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.5), bg_color)
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.2), Inches(0.4),
                    day, font_size=11, bold=True, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.7), Inches(y_pos + 0.05), Inches(1.8), Inches(0.4),
                    error, font_size=11, color=COLORS['orange_end'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(4.6), Inches(y_pos + 0.05), Inches(1.8), Inches(0.4),
                    accuracy, font_size=11, color=COLORS['gray_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(6.5), Inches(y_pos + 0.05), Inches(3.2), Inches(0.4),
                    note, font_size=10, color=COLORS['gray_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.5

    # Sentence starters
    add_colored_shape(slide, Inches(0.2), Inches(3.55), Inches(4.7), Inches(1.0), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.65), Inches(4.4), Inches(0.25),
                "SENTENCE STARTERS:", font_size=11, bold=True, color=COLORS['green_accent'])
    starters = """• "The forecast error grew from ___ to ___ because..."
• "Weather is harder to predict over time because..."
• "Small initial errors lead to...""""
    add_text_box(slide, Inches(0.4), Inches(3.9), Inches(4.4), Inches(0.6),
                starters, font_size=10, color=COLORS['dark_text'])

    # Spiral connection
    add_colored_shape(slide, Inches(5.1), Inches(3.55), Inches(4.6), Inches(1.0), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(5.3), Inches(3.65), Inches(4.2), Inches(0.25),
                "SPIRAL (Cycle 3):", font_size=11, bold=True, color=COLORS['purple_accent'])
    spiral = """How does climate change (long-term)
affect these short-term forecasts?
→ More extreme events = harder to predict
→ Historical data becomes less reliable"""
    add_text_box(slide, Inches(5.3), Inches(3.9), Inches(4.2), Inches(0.6),
                spiral, font_size=10, color=COLORS['dark_text'])


def add_station2_intro_slide(prs):
    """Slide 11: Station 2 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['cyan_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 2 – Weather vs Climate Investigation",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | ~15 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Left - Investigation
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "YOUR INVESTIGATION", font_size=14, bold=True, color=COLORS['cyan_end'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.3), COLORS['light_cyan_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.3), COLORS['cyan_end'])

    steps = """1. Compare daily temps to 30-year average (3 min)
2. Calculate climate normal for your area (3 min)
3. Identify variability vs trend patterns (4 min)
4. Explain why climate IS predictable (3 min)
5. Answer analysis questions (2 min)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(2.1),
                steps, font_size=12, color=COLORS['dark_text'])

    # Right - Key concepts
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(1.3), COLORS['cyan_end'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "KEY DIFFERENCE", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    difference = """WEATHER: Will it rain Tuesday?
→ Hard to predict! (chaos)

CLIMATE: Will summers get hotter?
→ Much easier to predict (trends)"""
    add_text_box(slide, Inches(5.2), Inches(1.55), Inches(4.3), Inches(0.8),
                difference, font_size=11, color=COLORS['white'])

    # Analogy
    add_colored_shape(slide, Inches(5.0), Inches(2.55), Inches(4.7), Inches(1.25), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(5.2), Inches(2.65), Inches(4.3), Inches(0.25),
                "ANALOGY:", font_size=11, bold=True, color=COLORS['purple_accent'])
    analogy = """Stock market:
• Tomorrow's price? Nearly impossible
• 10-year trend? Much more predictable

Weather/climate works the same way!"""
    add_text_box(slide, Inches(5.2), Inches(2.95), Inches(4.3), Inches(0.8),
                analogy, font_size=10, color=COLORS['dark_text'])

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.0), Inches(9.7), Inches(0.55), COLORS['cyan_end'])
    add_text_box(slide, Inches(0.35), Inches(4.03), Inches(9.3), Inches(0.5),
                "Complete Station 2 in Form - Compare weather variability to climate trends!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station2_support_slide(prs):
    """Slide 12: Station 2 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['cyan_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 2 – Variability vs. Trend",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Visual explanation
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(4.7), Inches(1.7), COLORS['light_cyan_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(4.4), Inches(0.25),
                "WEATHER VARIABILITY", font_size=12, bold=True, color=COLORS['cyan_end'])
    variability_text = """Day 1: 75°F
Day 2: 62°F  ← 13° drop!
Day 3: 68°F
Day 4: 71°F
Day 5: 58°F  ← 13° drop again!

The daily "noise" is HUGE
Individual days are chaotic"""
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(4.4), Inches(1.25),
                variability_text, font_size=10, color=COLORS['dark_text'])

    # Climate trend
    add_colored_shape(slide, Inches(5.1), Inches(0.85), Inches(4.6), Inches(1.7), COLORS['light_green_bg'])
    add_text_box(slide, Inches(5.3), Inches(0.95), Inches(4.2), Inches(0.25),
                "CLIMATE TREND", font_size=12, bold=True, color=COLORS['green_end'])
    trend_text = """1990s average July: 82°F
2000s average July: 83°F
2010s average July: 84°F
2020s average July: 85°F

The trend is CLEAR
0.3°F per decade = predictable"""
    add_text_box(slide, Inches(5.3), Inches(1.25), Inches(4.2), Inches(1.25),
                trend_text, font_size=10, color=COLORS['dark_text'])

    # Discussion questions
    add_colored_shape(slide, Inches(0.2), Inches(2.7), Inches(9.6), Inches(0.95), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.8), Inches(9.2), Inches(0.25),
                "DISCUSSION QUESTIONS:", font_size=11, bold=True, color=COLORS['orange_end'])
    questions = """1. Why does averaging over 30 years reveal trends that daily data hides?
2. How can we be confident about 2050's climate but not next Tuesday's weather?
3. What evidence would you need to claim "climate is changing"?"""
    add_text_box(slide, Inches(0.4), Inches(3.1), Inches(9.2), Inches(0.5),
                questions, font_size=10, color=COLORS['dark_text'])

    # Key concept
    add_colored_shape(slide, Inches(0.15), Inches(3.8), Inches(9.7), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(3.83), Inches(9.3), Inches(0.5),
                "KEY: Chaos cancels out over time—the signal (trend) emerges from the noise (variability)",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station3_intro_slide(prs):
    """Slide 13: Station 3 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['green_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 3 – Design a Climate Monitoring Station",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "25 Points | ~20 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Design challenge
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "ENGINEERING CHALLENGE", font_size=14, bold=True, color=COLORS['green_end'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.3), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.3), COLORS['green_end'])

    challenge = """Design a 30-year climate monitoring plan!

Budget: $50,000/year for 30 years
Requirements:
• Track temperature, precipitation, CO₂
• Data must be consistent for 3 decades
• Must detect 0.5°F/decade trends

Choose your equipment and justify!"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(2.1),
                challenge, font_size=12, color=COLORS['dark_text'])

    # Technology options
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(2.7), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "EQUIPMENT OPTIONS", font_size=13, bold=True, color=COLORS['blue_accent'])

    options = """Standard Weather Station | $5k/yr
  Temp, precip, wind, humidity

Precision Temp Sensors | $8k/yr
  ±0.01°F accuracy for trends

CO₂ Monitoring | $12k/yr
  Atmospheric carbon tracking

Soil Temp Sensors | $3k/yr
  Ground-level climate data

Data Storage/Analysis | $10k/yr
  Required for 30-year records"""
    add_text_box(slide, Inches(5.2), Inches(1.55), Inches(4.3), Inches(2.2),
                options, font_size=10, color=COLORS['dark_text'])

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.0), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(4.03), Inches(9.3), Inches(0.5),
                "Complete Station 3 in Form - Design your 30-year monitoring plan!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station3_support_slide(prs):
    """Slide 14: Station 3 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 3 – Long-Term Monitoring Considerations",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Example calculation
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(4.7), Inches(1.6), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(4.4), Inches(0.25),
                "EXAMPLE BUDGET:", font_size=11, bold=True, color=COLORS['green_accent'])
    calc = """Option A (Comprehensive):
• Standard Station: $5k
• Precision Temp: $8k
• CO₂ Monitoring: $12k
• Soil Sensors: $3k
• Data Storage: $10k
TOTAL: $38k/year ✓

Trade-off: Great data, under budget!"""
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(4.4), Inches(1.15),
                calc, font_size=10, color=COLORS['dark_text'])

    # Design questions
    add_colored_shape(slide, Inches(5.1), Inches(0.85), Inches(4.6), Inches(1.6), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(5.3), Inches(0.95), Inches(4.2), Inches(0.25),
                "THINK ABOUT:", font_size=11, bold=True, color=COLORS['orange_end'])
    questions = """• What's MOST important for detecting
  climate trends (not just weather)?

• How will you ensure data consistency
  over 30 YEARS of operation?

• What if equipment fails in year 15?

• Why is precision more important
  for climate than weather monitoring?"""
    add_text_box(slide, Inches(5.3), Inches(1.25), Inches(4.2), Inches(1.15),
                questions, font_size=10, color=COLORS['dark_text'])

    # SEP alignment
    add_colored_shape(slide, Inches(0.2), Inches(2.6), Inches(9.6), Inches(1.0), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.7), Inches(9.2), Inches(0.25),
                "SEP-3: PLANNING AND CARRYING OUT INVESTIGATIONS", font_size=11, bold=True, color=COLORS['purple_end'])
    sep_text = """Climate scientists must plan for DECADES of consistent data collection.
Your design should consider: Measurement precision | Data continuity | Equipment lifespan | Calibration"""
    add_text_box(slide, Inches(0.4), Inches(3.0), Inches(9.2), Inches(0.5),
                sep_text, font_size=11, color=COLORS['dark_text'])

    # Final prompt
    add_colored_shape(slide, Inches(0.15), Inches(3.75), Inches(9.7), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(3.78), Inches(9.3), Inches(0.5),
                "KEY: Climate monitoring needs PRECISION over DECADES, not just accuracy today",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_exit_ticket_slide(prs):
    """Slide 15: Exit Ticket"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.9), COLORS['exit_purple_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Exit Ticket – Prediction & Patterns Integration",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.68), Inches(9.3), Inches(0.35),
                "23 Points | ~15 min", font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Question types
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(1.5), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(0.08), Inches(1.5), COLORS['exit_purple_start'])

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.3),
                "QUESTION TYPES:", font_size=14, bold=True, color=COLORS['exit_purple_start'])

    q_types = """• 2 NEW – Chaos, forecast accuracy, weather vs climate
• 2 SPIRAL – Cycles 3-4: Climate factors, human impact on atmosphere
• 1 INTEGRATION – Connect short-term chaos to long-term predictability
• 1 SEP-3 – Data collection design (justify monitoring choices)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(9.0), Inches(1.0),
                q_types, font_size=13, color=COLORS['dark_text'])

    # Tips boxes
    add_colored_shape(slide, Inches(0.3), Inches(2.85), Inches(4.6), Inches(1.1), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(2.9), Inches(4.3), Inches(0.3),
                "SUCCESS TIPS:", font_size=12, bold=True, color=COLORS['green_accent'],
                anchor=MSO_ANCHOR.MIDDLE)
    tips = """• Use vocabulary: chaos, variability, trend
• Explain WHY forecasts fail over time
• Connect averaging to predictability"""
    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(4.3), Inches(0.7),
                tips, font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.1), Inches(2.85), Inches(4.6), Inches(1.1), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.3), Inches(2.9), Inches(4.2), Inches(0.3),
                "SPIRAL REMINDER:", font_size=12, bold=True, color=COLORS['blue_accent'],
                anchor=MSO_ANCHOR.MIDDLE)
    spiral = """• C3: CO₂ causes warming trend
• C4: Human activities add CO₂
• C5: Trend is predictable, events aren't"""
    add_text_box(slide, Inches(5.3), Inches(3.2), Inches(4.2), Inches(0.7),
                spiral, font_size=10, color=COLORS['dark_text'])

    # Final notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.1), Inches(9.7), Inches(0.65), COLORS['exit_purple_end'])
    add_text_box(slide, Inches(0.35), Inches(4.15), Inches(9.3), Inches(0.55),
                "FINAL Notecard: Weather = chaos (unpredictable) | Climate = trends (predictable)",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_summary_slide(prs):
    """Slide 16: Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Week 2 Summary: What You Learned",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Key takeaways
    takeaways = [
        ("Chaos", "Small errors grow → forecasts fail beyond ~7 days (butterfly effect)"),
        ("Variability", "Daily weather bounces around—impossible to predict exactly"),
        ("Trends", "Climate averages reveal patterns—highly predictable over decades"),
        ("Monitoring", "Climate needs precision + consistency over 30+ years"),
    ]

    y_pos = 0.9
    for title, detail in takeaways:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.65), COLORS['light_cyan_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.65), COLORS['teal'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.0), Inches(0.55),
                    title, font_size=13, bold=True, color=COLORS['teal_dark'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.5), Inches(y_pos + 0.05), Inches(7.1), Inches(0.55),
                    detail, font_size=12, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.72

    # Next week preview
    add_colored_shape(slide, Inches(0.2), Inches(3.85), Inches(9.6), Inches(0.65), COLORS['header_blue_end'])
    add_text_box(slide, Inches(0.4), Inches(3.88), Inches(9.2), Inches(0.6),
                "NEXT WEEK: Synthesis & Assessment – Bring it all together!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Key rule
    add_colored_shape(slide, Inches(0.2), Inches(4.65), Inches(9.6), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(4.68), Inches(9.2), Inches(0.5),
                "REMEMBER: Can't predict the flip, but we know the coin is loaded!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


if __name__ == "__main__":
    prs = create_presentation()
    output_path = "/home/user/Kairos.Sci.Repo/content/grade7/cycle05/week2/G7_C5_W2_Weather_Prediction_Slides_Final.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Created: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
