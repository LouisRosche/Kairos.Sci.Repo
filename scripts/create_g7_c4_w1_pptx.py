#!/usr/bin/env python3
"""
Create G7_C4_W1 Ocean Acidification Investigation presentation
Following established patterns from C3 presentations
Colors matched to student-page.html (teal/cyan theme)

See create_g7_c3_w2_pptx.py for full PPTX DESIGN BEST PRACTICES documentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# HTML Color Palette (matching student-page.html - teal/cyan theme)
COLORS = {
    # Main theme colors
    'teal_primary': RGBColor(0x08, 0x91, 0xB2),      # #0891B2
    'teal_dark': RGBColor(0x0E, 0x74, 0x90),         # #0E7490
    'teal_light': RGBColor(0xEC, 0xFE, 0xFF),        # #ECFEFF
    'cyan_light': RGBColor(0xCF, 0xFA, 0xFE),        # #CFFAFE
    # Hook purple
    'hook_purple': RGBColor(0x63, 0x66, 0xF1),       # #6366F1
    'hook_purple_dark': RGBColor(0x4F, 0x46, 0xE5),  # #4F46E5
    # Station colors
    'station1_green': RGBColor(0x10, 0xB9, 0x81),    # #10B981
    'station1_green_dark': RGBColor(0x05, 0x96, 0x69),  # #059669
    'station2_orange': RGBColor(0xF5, 0x9E, 0x0B),   # #F59E0B
    'station2_orange_dark': RGBColor(0xD9, 0x77, 0x06),  # #D97706
    'station3_purple': RGBColor(0x8B, 0x5C, 0xF6),   # #8B5CF6
    'station3_purple_dark': RGBColor(0x7C, 0x3A, 0xED),  # #7C3AED
    'exit_pink': RGBColor(0xEC, 0x48, 0x99),         # #EC4899
    'exit_pink_dark': RGBColor(0xDB, 0x27, 0x77),    # #DB2777
    # Text colors
    'dark_text': RGBColor(0x2D, 0x37, 0x48),         # #2D3748
    'gray_text': RGBColor(0x4A, 0x55, 0x68),         # #4A5568
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    # Background colors
    'light_bg': RGBColor(0xF7, 0xFA, 0xFC),          # #F7FAFC
    'light_teal_bg': RGBColor(0xE6, 0xFF, 0xFA),     # #E6FFFA
    'light_green_bg': RGBColor(0xEC, 0xFD, 0xF5),    # #ECFDF5
    'light_orange_bg': RGBColor(0xFF, 0xFB, 0xEB),   # #FFFBEB
    'light_purple_bg': RGBColor(0xF5, 0xF3, 0xFF),   # #F5F3FF
    'light_pink_bg': RGBColor(0xFD, 0xF2, 0xF8),     # #FDF2F8
    # Alert colors
    'red_alert': RGBColor(0xC5, 0x30, 0x30),         # #C53030
    'warning_yellow': RGBColor(0x92, 0x40, 0x0E),    # #92400E
}

def create_presentation():
    """Create the G7_C4_W1 presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

    # Add all slides
    add_title_slide(prs)
    add_phenomenon_slide(prs)
    add_driving_question_slide(prs)
    add_prior_knowledge_slide(prs)
    add_learning_targets_slide(prs)
    add_vocabulary_slide(prs)
    add_hook_intro_slide(prs)
    add_hook_support_slide(prs)
    add_station1_intro_slide(prs)
    add_station1_support_slide(prs)
    add_station2_intro_slide(prs)
    add_station2_support_slide(prs)
    add_station3_intro_slide(prs)
    add_station3_support_slide(prs)
    add_exit_ticket_slide(prs)
    add_summary_slide(prs)

    return prs


def add_colored_shape(slide, left, top, width, height, color):
    """Add a colored rectangle shape"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=None, align=PP_ALIGN.LEFT, font_name="Arial", anchor=None):
    """Add a text box with specified formatting"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    if anchor:
        tf.anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    if color:
        p.font.color.rgb = color
    p.alignment = align
    return txBox


def add_title_slide(prs):
    """Slide 1: Title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['teal_primary'])

    # Emoji decoration
    add_text_box(slide, Inches(0), Inches(1.0), Inches(10), Inches(0.8),
                "🐚 🌊", font_size=48, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Main title
    add_text_box(slide, Inches(0.5), Inches(1.8), Inches(9), Inches(1.0),
                "Week 1: Ocean Acidification Investigation",
                font_size=36, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Subtitle
    add_text_box(slide, Inches(0.5), Inches(2.9), Inches(9), Inches(0.5),
                "Grade 7 Science | Cycle 4 | Rosche | Kairos Academies",
                font_size=16, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Standards
    add_text_box(slide, Inches(0.5), Inches(3.5), Inches(9), Inches(0.4),
                "MS-ESS3-3 Monitoring Human Impact on the Environment",
                font_size=14, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Points/Time
    add_text_box(slide, Inches(0.5), Inches(4.2), Inches(9), Inches(0.4),
                "100 Points Total | ~75 Minutes",
                font_size=14, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_phenomenon_slide(prs):
    """Slide 2: The Phenomenon"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['teal_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🐚 The Phenomenon: The Dissolving Shells Mystery",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Main content box
    add_colored_shape(slide, Inches(0.15), Inches(1.0), Inches(9.7), Inches(2.8), COLORS['teal_light'])

    # Content
    add_text_box(slide, Inches(0.4), Inches(1.1), Inches(9.2), Inches(0.5),
                "Scientists studying pteropods ('sea butterflies') discovered something alarming:",
                font_size=16, color=COLORS['dark_text'])

    add_text_box(slide, Inches(0.6), Inches(1.6), Inches(9.0), Inches(1.8),
                "• Shells collected TODAY are thinner and more fragile than 50 years ago\n" +
                "• Some shells are literally dissolving while animals are still alive\n" +
                "• This is happening in oceans around the world\n\n" +
                "These tiny creatures haven't changed. The water has.",
                font_size=15, color=COLORS['gray_text'])

    # Key insight box
    add_colored_shape(slide, Inches(0.15), Inches(3.95), Inches(9.7), Inches(1.4), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.35), Inches(4.15), Inches(9.3), Inches(1.0),
                "🔬 KEY INSIGHT: The ocean's chemistry is changing due to increased CO₂\n" +
                "absorption from the atmosphere, making it more acidic.",
                font_size=16, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_driving_question_slide(prs):
    """Slide 3: Driving Question"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['teal_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🤔 Driving Question",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Question box
    add_colored_shape(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(2.0), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.7), Inches(1.7), Inches(8.6), Inches(1.6),
                "Why is the ocean dissolving shells NOW when it didn't 50 years ago?\n\n" +
                "What has changed in ocean chemistry?",
                font_size=22, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Connection to Cycle 3
    add_colored_shape(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(1.4), COLORS['light_teal_bg'])
    add_text_box(slide, Inches(0.7), Inches(3.9), Inches(8.6), Inches(1.2),
                "🔗 Connection to Cycle 3:\n" +
                "You learned that CO₂ is increasing in the atmosphere and causing climate change.\n" +
                "Now we'll explore what happens when CO₂ dissolves in the ocean!",
                font_size=14, color=COLORS['teal_dark'])


def add_prior_knowledge_slide(prs):
    """Slide 4: What You Already Know"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['teal_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "💡 What You Already Know (Cycle 3 Connections)",
                font_size=22, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Content boxes
    concepts = [
        ("CO₂ & Climate", "CO₂ is increasing in atmosphere\ndue to human activities"),
        ("Carbon Cycle", "Carbon moves between\natmosphere, land, and ocean"),
        ("Chemical Reactions", "Atoms are conserved in\nreactions (MS-PS1-5)")
    ]

    x_positions = [0.3, 3.5, 6.7]
    for i, (title, content) in enumerate(concepts):
        add_colored_shape(slide, Inches(x_positions[i]), Inches(1.1), Inches(3.0), Inches(1.8), COLORS['teal_light'])
        add_text_box(slide, Inches(x_positions[i] + 0.15), Inches(1.2), Inches(2.7), Inches(0.4),
                    title, font_size=14, bold=True, color=COLORS['teal_dark'], align=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x_positions[i] + 0.15), Inches(1.6), Inches(2.7), Inches(1.2),
                    content, font_size=12, color=COLORS['gray_text'], align=PP_ALIGN.CENTER)

    # New this cycle
    add_colored_shape(slide, Inches(0.3), Inches(3.1), Inches(9.4), Inches(2.2), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(9.0), Inches(0.4),
                "🆕 NEW This Cycle: MS-ESS3-3 - Monitoring Human Impact",
                font_size=14, bold=True, color=COLORS['teal_dark'])
    add_text_box(slide, Inches(0.5), Inches(3.6), Inches(9.0), Inches(1.6),
                "• How does ocean chemistry change when it absorbs CO₂?\n" +
                "• How do we measure and monitor these changes?\n" +
                "• What are the impacts on marine ecosystems?\n" +
                "• How can we design systems to track acidification?",
                font_size=13, color=COLORS['gray_text'])


def add_learning_targets_slide(prs):
    """Slide 5: Learning Targets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['teal_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🎯 Learning Targets",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    targets = [
        "Explain how atmospheric CO₂ affects ocean pH",
        "Interpret pH data and predict effects on marine life",
        "Apply mass balance to ocean-atmosphere carbon exchange",
        "Design a monitoring system for ocean acidification"
    ]

    y_pos = 1.0
    for i, target in enumerate(targets, 1):
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.95), COLORS['teal_light'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.15), Inches(9.0), Inches(0.65),
                    f"Target {i}: {target}",
                    font_size=15, bold=True, color=COLORS['teal_dark'])
        y_pos += 1.1


def add_vocabulary_slide(prs):
    """Slide 6: Key Vocabulary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['teal_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "📚 Key Vocabulary This Week",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    vocab = [
        ("Ocean Acidification", "Decrease in ocean pH caused by CO₂ absorption"),
        ("pH Scale", "Measures acidity: 0-6 acidic, 7 neutral, 8-14 basic"),
        ("Carbonic Acid (H₂CO₃)", "Formed when CO₂ dissolves: CO₂ + H₂O → H₂CO₃"),
        ("Carbon Sink", "Reservoir that absorbs more CO₂ than it releases"),
        ("Logarithmic Scale", "0.1 pH change = 26% change in acidity!"),
        ("Calcifying Organisms", "Marine life building shells from calcium carbonate")
    ]

    y_pos = 0.95
    for term, definition in vocab:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.7), COLORS['light_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.1), Inches(2.8), Inches(0.5),
                    term, font_size=12, bold=True, color=COLORS['teal_dark'])
        add_text_box(slide, Inches(3.4), Inches(y_pos + 0.1), Inches(6.1), Inches(0.5),
                    definition, font_size=12, color=COLORS['gray_text'])
        y_pos += 0.75


def add_hook_intro_slide(prs):
    """Slide 7: Hook Intro"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header - two-line pattern
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['hook_purple'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🐚 Hook – The Dissolving Shells Mystery",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "12 Points | ~10 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # What you'll do
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(2.2), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "🐚 The Challenge: What's happening to sea shells?",
                font_size=16, bold=True, color=COLORS['hook_purple_dark'])
    add_text_box(slide, Inches(0.5), Inches(1.75), Inches(9.0), Inches(1.5),
                "1. Observe the phenomenon: shells dissolving in modern oceans (2 min)\n" +
                "2. Connect to Cycle 3: What does CO₂ do in water? (3 min)\n" +
                "3. Make predictions about ocean chemistry changes (3 min)\n" +
                "4. Answer diagnostic questions (2 min)",
                font_size=14, color=COLORS['gray_text'])

    # Think about this
    add_colored_shape(slide, Inches(0.3), Inches(3.55), Inches(9.4), Inches(1.8), COLORS['light_teal_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.65), Inches(9.0), Inches(0.4),
                "💭 Think About This:",
                font_size=14, bold=True, color=COLORS['teal_dark'])
    add_text_box(slide, Inches(0.5), Inches(4.05), Inches(9.0), Inches(1.2),
                "• What gas is increasing in the atmosphere? (Cycle 3)\n" +
                "• What happens when gases dissolve in water?\n" +
                "• What do acids do to calcium carbonate (shells)?",
                font_size=13, color=COLORS['gray_text'])


def add_hook_support_slide(prs):
    """Slide 8: Hook Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['hook_purple'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🐚 Hook – Support Information",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Key facts
    add_colored_shape(slide, Inches(0.3), Inches(1.0), Inches(4.5), Inches(2.5), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(4.1), Inches(0.4),
                "🔬 Key Facts:", font_size=14, bold=True, color=COLORS['hook_purple_dark'])
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4.1), Inches(1.9),
                "• Pteropods are called 'sea butterflies'\n" +
                "• Their shells are made of calcium carbonate (CaCO₃)\n" +
                "• Same material as coral, oyster shells, and chalk\n" +
                "• Acid dissolves calcium carbonate",
                font_size=12, color=COLORS['gray_text'])

    # Chemical equation
    add_colored_shape(slide, Inches(5.0), Inches(1.0), Inches(4.7), Inches(2.5), COLORS['teal_light'])
    add_text_box(slide, Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.4),
                "⚗️ The Chemistry:", font_size=14, bold=True, color=COLORS['teal_dark'])
    add_text_box(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(1.9),
                "CO₂ + H₂O → H₂CO₃\n(carbon dioxide + water → carbonic acid)\n\n" +
                "More CO₂ in atmosphere →\nMore CO₂ dissolves in ocean →\nMore acidic ocean!",
                font_size=12, color=COLORS['gray_text'], align=PP_ALIGN.CENTER)

    # Notecard
    add_colored_shape(slide, Inches(0.3), Inches(3.7), Inches(9.4), Inches(1.65), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.8), Inches(9.0), Inches(0.4),
                "📝 Complete the Hook form on the student page",
                font_size=14, bold=True, color=COLORS['red_alert'])
    add_text_box(slide, Inches(0.5), Inches(4.2), Inches(9.0), Inches(1.0),
                "Record your observations and predictions about:\n" +
                "• Why shells might be dissolving now but not 50 years ago\n" +
                "• What connection this might have to CO₂ levels",
                font_size=12, color=COLORS['gray_text'])


def add_station1_intro_slide(prs):
    """Slide 9: Station 1 Intro"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header - two-line pattern
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['station1_green'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🧪 Station 1 – pH and Marine Life",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "20 Points | ~18 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Mission
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(1.5), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "🧪 Your Mission: Investigate pH Changes",
                font_size=16, bold=True, color=COLORS['station1_green_dark'])
    add_text_box(slide, Inches(0.5), Inches(1.75), Inches(9.0), Inches(0.8),
                "Analyze historical ocean pH data and determine which marine organisms\n" +
                "are most at risk from ocean acidification.",
                font_size=14, color=COLORS['gray_text'])

    # pH Scale info
    add_colored_shape(slide, Inches(0.3), Inches(2.85), Inches(9.4), Inches(2.5), COLORS['teal_light'])
    add_text_box(slide, Inches(0.5), Inches(2.95), Inches(9.0), Inches(0.4),
                "📊 The pH Scale (CRITICAL: It's Logarithmic!)",
                font_size=14, bold=True, color=COLORS['teal_dark'])
    add_text_box(slide, Inches(0.5), Inches(3.4), Inches(4.3), Inches(1.8),
                "• pH 7 = neutral (pure water)\n" +
                "• Below 7 = acidic (more H⁺ ions)\n" +
                "• Above 7 = basic (fewer H⁺ ions)\n" +
                "• Ocean is normally ~8.1 pH",
                font_size=13, color=COLORS['gray_text'])
    add_text_box(slide, Inches(5.0), Inches(3.4), Inches(4.5), Inches(1.8),
                "⚠️ KEY CONCEPT:\n0.1 pH change = 26% change in acidity!\n\n" +
                "Small pH numbers = BIG chemical changes",
                font_size=13, bold=True, color=COLORS['red_alert'], align=PP_ALIGN.CENTER)


def add_station1_support_slide(prs):
    """Slide 10: Station 1 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['station1_green'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🧪 Station 1 – Data Reference",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # pH Data table
    add_colored_shape(slide, Inches(0.3), Inches(1.0), Inches(4.5), Inches(2.8), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(4.1), Inches(0.4),
                "📈 Historical Ocean pH Data:", font_size=13, bold=True, color=COLORS['station1_green_dark'])
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4.1), Inches(2.2),
                "Year     |  Ocean pH  |  Change\n" +
                "1800     |    8.25    |  baseline\n" +
                "1950     |    8.18    |   -0.07\n" +
                "2000     |    8.12    |   -0.13\n" +
                "2020     |    8.10    |   -0.15\n" +
                "2100*    |    7.95    |   -0.30",
                font_size=11, color=COLORS['gray_text'])

    # Organism thresholds
    add_colored_shape(slide, Inches(5.0), Inches(1.0), Inches(4.7), Inches(2.8), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.4),
                "🐚 Organism pH Thresholds:", font_size=13, bold=True, color=COLORS['station2_orange_dark'])
    add_text_box(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(2.2),
                "Coral:\nCalcification slows below pH 8.0\n\n" +
                "Oysters:\nLarvae fail below pH 7.9\n\n" +
                "Pteropods:\nShells dissolve below pH 7.8",
                font_size=12, color=COLORS['gray_text'])

    # Instructions
    add_colored_shape(slide, Inches(0.3), Inches(4.0), Inches(9.4), Inches(1.3), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(4.1), Inches(9.0), Inches(0.4),
                "📝 Complete Station 1 on the student page",
                font_size=14, bold=True, color=COLORS['red_alert'])
    add_text_box(slide, Inches(0.5), Inches(4.5), Inches(9.0), Inches(0.7),
                "Compare pH data to organism thresholds. Which organisms are already at risk? Which will be affected by 2100?",
                font_size=12, color=COLORS['gray_text'])


def add_station2_intro_slide(prs):
    """Slide 11: Station 2 Intro"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header - two-line pattern
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['station2_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "⚖️ Station 2 – Carbon Sources and Sinks",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "20 Points | ~15 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Mission
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(1.2), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "⚖️ Your Mission: Track Carbon in the Ocean-Atmosphere System",
                font_size=16, bold=True, color=COLORS['station2_orange_dark'])
    add_text_box(slide, Inches(0.5), Inches(1.75), Inches(9.0), Inches(0.5),
                "Use mass balance to calculate how much CO₂ the ocean absorbs each year.",
                font_size=14, color=COLORS['gray_text'])

    # Key equation
    add_colored_shape(slide, Inches(0.3), Inches(2.55), Inches(9.4), Inches(0.8), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.5), Inches(2.65), Inches(9.0), Inches(0.6),
                "KEY EQUATION: Sources - Sinks = Atmospheric Accumulation",
                font_size=18, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Data preview
    add_colored_shape(slide, Inches(0.3), Inches(3.5), Inches(9.4), Inches(1.85), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.6), Inches(9.0), Inches(0.4),
                "📊 Global Carbon Budget (billions of tons CO₂/year):",
                font_size=13, bold=True, color=COLORS['station2_orange_dark'])
    add_text_box(slide, Inches(0.5), Inches(4.0), Inches(4.3), Inches(1.2),
                "SOURCES (Release CO₂):\n" +
                "• Fossil fuels + industry: +36 Gt\n" +
                "• Deforestation: +4 Gt\n" +
                "• Total Released: 40 Gt",
                font_size=11, color=COLORS['gray_text'])
    add_text_box(slide, Inches(5.0), Inches(4.0), Inches(4.5), Inches(1.2),
                "SINKS (Absorb CO₂):\n" +
                "• Land plants/soil: -12 Gt (30%)\n" +
                "• Ocean: -10 Gt (25%)\n" +
                "• Atmosphere: +18 Gt (45%)",
                font_size=11, color=COLORS['gray_text'])


def add_station2_support_slide(prs):
    """Slide 12: Station 2 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['station2_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "⚖️ Station 2 – Carbon Budget Analysis",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Trade-off analysis
    add_colored_shape(slide, Inches(0.3), Inches(1.0), Inches(4.5), Inches(2.3), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(4.1), Inches(0.4),
                "✅ BENEFIT of Ocean CO₂ Absorption:",
                font_size=13, bold=True, color=COLORS['station1_green_dark'])
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4.1), Inches(1.7),
                "• Oceans absorb 25% of human CO₂\n" +
                "• Slows climate change\n" +
                "• Without oceans, atmosphere\n  would warm much faster\n" +
                "• Ocean = Earth's 'carbon sink'",
                font_size=12, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(5.0), Inches(1.0), Inches(4.7), Inches(2.3), COLORS['light_pink_bg'])
    add_text_box(slide, Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.4),
                "❌ COST of Ocean CO₂ Absorption:",
                font_size=13, bold=True, color=COLORS['red_alert'])
    add_text_box(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(1.7),
                "• Ocean becomes more acidic\n" +
                "• Shells and coral dissolve\n" +
                "• Marine ecosystems disrupted\n" +
                "• Food webs collapse\n" +
                "• Fisheries threatened",
                font_size=12, color=COLORS['gray_text'])

    # Instructions
    add_colored_shape(slide, Inches(0.3), Inches(3.5), Inches(9.4), Inches(1.85), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.6), Inches(9.0), Inches(0.4),
                "📝 Complete Station 2 on the student page (Calculator needed!)",
                font_size=14, bold=True, color=COLORS['red_alert'])
    add_text_box(slide, Inches(0.5), Inches(4.0), Inches(9.0), Inches(1.2),
                "Use the carbon budget data to:\n" +
                "• Calculate how much CO₂ the ocean absorbs annually\n" +
                "• Explain the trade-off between climate benefits and ocean acidification",
                font_size=12, color=COLORS['gray_text'])


def add_station3_intro_slide(prs):
    """Slide 13: Station 3 Intro"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header - two-line pattern
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['station3_purple'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🔧 Station 3 – Design a Monitoring System",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "25 Points | ~20 min (Highest Value!)", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Mission
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(1.3), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "🔧 Engineering Challenge: Design an Ocean Acidification Monitoring System",
                font_size=15, bold=True, color=COLORS['station3_purple_dark'])
    add_text_box(slide, Inches(0.5), Inches(1.75), Inches(9.0), Inches(0.6),
                "Use your knowledge of ocean chemistry to design a network that can track\nacidification changes across multiple ocean regions.",
                font_size=13, color=COLORS['gray_text'])

    # Constraints
    add_colored_shape(slide, Inches(0.3), Inches(2.65), Inches(9.4), Inches(2.65), COLORS['teal_light'])
    add_text_box(slide, Inches(0.5), Inches(2.75), Inches(9.0), Inches(0.4),
                "📋 CONSTRAINTS:", font_size=14, bold=True, color=COLORS['teal_dark'])
    add_text_box(slide, Inches(0.5), Inches(3.15), Inches(4.3), Inches(2.0),
                "• Budget: $500,000 annual cost\n" +
                "• Coverage: At least 3 ocean regions\n" +
                "• Data: Must measure pH, temp, CO₂\n" +
                "• Frequency: Data collected weekly",
                font_size=12, color=COLORS['gray_text'])
    add_text_box(slide, Inches(5.0), Inches(2.75), Inches(4.5), Inches(0.4),
                "💡 Design Tips:", font_size=14, bold=True, color=COLORS['station3_purple_dark'])
    add_text_box(slide, Inches(5.0), Inches(3.15), Inches(4.5), Inches(2.0),
                "• Mix technologies to cover all needs\n" +
                "• Consider trade-offs (cost vs coverage)\n" +
                "• Justify your choices\n" +
                "• Stay under budget!",
                font_size=12, color=COLORS['gray_text'])


def add_station3_support_slide(prs):
    """Slide 14: Station 3 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['station3_purple'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🔧 Station 3 – Available Technologies",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Technology options
    technologies = [
        ("A. Fixed Buoys", "$50k/each", "pH, temp, CO₂", "Reliable, fixed location"),
        ("B. Argo Floats", "$30k/each", "pH, temp only", "Large area, no CO₂"),
        ("C. Ship Surveys", "$200k/each", "All + samples", "Most complete, infrequent"),
        ("D. Satellite", "$100k", "Surface temp", "Global, surface only")
    ]

    y_pos = 1.0
    for option, cost, measures, notes in technologies:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.85), COLORS['light_purple_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.1), Inches(2.5), Inches(0.65),
                    option, font_size=12, bold=True, color=COLORS['station3_purple_dark'])
        add_text_box(slide, Inches(3.0), Inches(y_pos + 0.1), Inches(1.5), Inches(0.65),
                    cost, font_size=11, color=COLORS['gray_text'], align=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(4.5), Inches(y_pos + 0.1), Inches(2.3), Inches(0.65),
                    measures, font_size=11, color=COLORS['gray_text'])
        add_text_box(slide, Inches(6.9), Inches(y_pos + 0.1), Inches(2.6), Inches(0.65),
                    notes, font_size=11, color=COLORS['gray_text'])
        y_pos += 0.95

    # Instructions
    add_colored_shape(slide, Inches(0.3), Inches(4.85), Inches(9.4), Inches(0.55), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(4.9), Inches(9.0), Inches(0.4),
                "📝 Complete Station 3 on the student page - Design your network and justify your choices!",
                font_size=13, bold=True, color=COLORS['red_alert'])


def add_exit_ticket_slide(prs):
    """Slide 15: Exit Ticket"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header - two-line pattern
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['exit_pink'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🎓 Exit Ticket – Ocean Chemistry Integration",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "23 Points | ~15 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Question types
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(2.5), COLORS['light_pink_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "🎓 Show What You Learned - Question Types:",
                font_size=16, bold=True, color=COLORS['exit_pink_dark'])
    add_text_box(slide, Inches(0.5), Inches(1.8), Inches(9.0), Inches(1.8),
                "• 2 NEW – Ocean acidification content (this week)\n" +
                "• 2 SPIRAL – Cycle 3 review (carbon cycle, greenhouse effect)\n" +
                "• 1 INTEGRATION – Connect atmospheric + ocean effects of CO₂\n" +
                "• 1 SEP-6 – Construct an explanation using evidence",
                font_size=14, color=COLORS['gray_text'])

    # Reminder
    add_colored_shape(slide, Inches(0.3), Inches(3.85), Inches(9.4), Inches(1.5), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.95), Inches(9.0), Inches(0.4),
                "📝 Complete the Exit Ticket on the student page",
                font_size=14, bold=True, color=COLORS['red_alert'])
    add_text_box(slide, Inches(0.5), Inches(4.35), Inches(9.0), Inches(0.9),
                "This is your final assessment for Week 1. Take your time!\n" +
                "Use evidence from today's activities to support your answers.",
                font_size=13, color=COLORS['gray_text'])


def add_summary_slide(prs):
    """Slide 16: Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['teal_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🎉 Week 1 Summary: What You Learned",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Key takeaways
    takeaways = [
        ("Ocean Acidification", "CO₂ + H₂O → H₂CO₃ (carbonic acid) → lowers pH", COLORS['teal_primary']),
        ("pH Changes", "0.1 pH = 26% more acidic (logarithmic scale)", COLORS['station1_green']),
        ("Carbon Budget", "Ocean absorbs 25% of human CO₂ - a double-edged sword", COLORS['station2_orange']),
        ("Monitoring", "Multiple technologies needed to track ocean chemistry", COLORS['station3_purple'])
    ]

    y_pos = 1.0
    for title, content, color in takeaways:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.9), COLORS['light_bg'])
        # Color bar on left
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(0.1), Inches(0.9), color)
        add_text_box(slide, Inches(0.55), Inches(y_pos + 0.1), Inches(9.0), Inches(0.35),
                    title, font_size=13, bold=True, color=color)
        add_text_box(slide, Inches(0.55), Inches(y_pos + 0.45), Inches(9.0), Inches(0.4),
                    content, font_size=12, color=COLORS['gray_text'])
        y_pos += 1.0

    # Next week preview
    add_colored_shape(slide, Inches(0.3), Inches(5.05), Inches(9.4), Inches(0.45), COLORS['teal_light'])
    add_text_box(slide, Inches(0.5), Inches(5.1), Inches(9.0), Inches(0.35),
                "Next Week: What happens when nutrients run off into waterways? (Eutrophication)",
                font_size=12, color=COLORS['teal_dark'], align=PP_ALIGN.CENTER)


def main():
    """Main function to create and save the presentation"""
    prs = create_presentation()

    # Ensure output directory exists
    output_dir = "/home/user/Kairos.Sci.Repo/content/grade7/cycle04/week1"
    os.makedirs(output_dir, exist_ok=True)

    # Save the presentation
    output_path = os.path.join(output_dir, "G7_C4_W1_Ocean_Acidification_Slides_Final.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
