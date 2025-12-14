#!/usr/bin/env python3
"""
Create G7_C6_W2 Seafloor Spreading & Continental Drift presentation.
Topic: Plate Tectonics & Earth's Interior.

Instructional week structure (16 slides).
See PPTX_DESIGN_GUIDE.md for best practices documentation.
"""

import os
from pptx.dml.color import RGBColor
from pptx_common import (
    COLORS as BASE_COLORS,
    create_base_presentation,
    add_colored_shape,
    add_text_box,
    Inches,
    Pt,
    PP_ALIGN,
    MSO_ANCHOR,
)

# Ocean/tectonic theme (deep blues and teals) - extend base colors
COLORS = {**BASE_COLORS}
COLORS.update({
    'ocean_blue': RGBColor(0x0C, 0x4A, 0x6E),
    'ocean_dark': RGBColor(0x08, 0x30, 0x4A),
    'ridge_orange': RGBColor(0xF9, 0x73, 0x16),
    'ridge_red': RGBColor(0xEF, 0x44, 0x44),
    'seafloor_teal': RGBColor(0x0D, 0x94, 0x88),
    'light_ocean_bg': RGBColor(0xE0, 0xF2, 0xFE),
    'magnetic_purple': RGBColor(0x7C, 0x3A, 0xED),
})


def create_presentation():
    """Create the G7_C6_W2 presentation"""
    prs = create_base_presentation()

    add_title_slide(prs)
    add_phenomenon_slide(prs)
    add_hook_intro_slide(prs)
    add_hook_activity_slide(prs)
    add_station1_intro_slide(prs)
    add_station1_activity_slide(prs)
    add_station2_intro_slide(prs)
    add_station2_activity_slide(prs)
    add_station3_intro_slide(prs)
    add_station3_activity_slide(prs)
    add_concepts_slide(prs)
    add_vocabulary_slide(prs)
    add_exit_intro_slide(prs)
    add_exit_tips_slide(prs)
    add_notecard_slide(prs)
    add_summary_slide(prs)

    return prs


def add_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['ocean_blue'])

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(1.2),
                "Plate Tectonics & Earth's Interior",
                font_size=40, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")

    add_text_box(slide, Inches(0.5), Inches(2.6), Inches(9), Inches(0.5),
                "Week 2: Seafloor Spreading & Continental Drift",
                font_size=24, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(9), Inches(0.5),
                "Grade 7 Science | Cycle 6 | 100 Points",
                font_size=18, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(2), Inches(4.1), Inches(6), Inches(1.0), COLORS['white'])
    add_text_box(slide, Inches(2.2), Inches(4.2), Inches(5.6), Inches(0.8),
                "Hook: 12 pts | Stations: 20+20+25 pts | Exit: 23 pts",
                font_size=14, bold=True, color=COLORS['ocean_blue'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_phenomenon_slide(prs):
    """Slide 2: Phenomenon Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['ridge_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "This Week's Phenomenon",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(9.6), Inches(1.4), COLORS['light_ocean_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(0.08), Inches(1.4), COLORS['ridge_red'])
    add_text_box(slide, Inches(0.5), Inches(1.0), Inches(9.0), Inches(1.2),
                "How can mountains exist on the ocean floor?",
                font_size=28, bold=True, color=COLORS['ocean_dark'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(2.45), Inches(9.6), Inches(1.6), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.55), Inches(9.2), Inches(0.3),
                "THINK ABOUT IT:", font_size=14, bold=True, color=COLORS['blue_accent'])
    context = """• The Mid-Atlantic Ridge is an underwater mountain range longer than any on land!
• It's 10,000 miles long and runs down the middle of the Atlantic Ocean
• Why would mountains form underwater? What's pushing them up?
• Today we'll discover how NEW ocean floor is constantly being created."""
    add_text_box(slide, Inches(0.4), Inches(2.9), Inches(9.2), Inches(1.0),
                context, font_size=13, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(4.2), Inches(9.6), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.23), Inches(9.2), Inches(0.5),
                "Building on W1: Divergent boundaries create new crust - but what's the evidence?",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_hook_intro_slide(prs):
    """Slide 3: Hook Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['purple_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Hook: The Underwater Mountain Mystery",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "12 Points | 5 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(1.6), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(1.6), COLORS['purple_end'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "YOUR RESOURCES:", font_size=14, bold=True, color=COLORS['purple_end'])
    resources = """• Mid-Atlantic Ridge bathymetry (underwater topography) map
• Volcanic footage from underwater eruptions at the ridge
• Age data showing how old the seafloor is at different distances from the ridge

OBSERVE: What patterns do you notice about the age of the ocean floor?"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(1.1),
                resources, font_size=12, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(2.9), Inches(9.6), Inches(0.95), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.0), Inches(9.2), Inches(0.25),
                "FOCUS:", font_size=12, bold=True, color=COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(3.3), Inches(9.2), Inches(0.5),
                "How do underwater mountains form? Why is the seafloor YOUNGEST at the ridge?",
                font_size=12, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(4.0), Inches(9.6), Inches(0.7), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(4.1), Inches(9.2), Inches(0.55),
                "🔗 CONNECTION to W1: We learned about divergent boundaries - now we see evidence!",
                font_size=11, color=COLORS['dark_text'], anchor=MSO_ANCHOR.MIDDLE)


def add_hook_activity_slide(prs):
    """Slide 4: Hook Activity Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['purple_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Hook Activity: What Do You Notice?",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(0.3), Inches(0.85), Inches(9.4), Inches(0.3),
                "QUESTIONS YOU'LL ANSWER:", font_size=13, bold=True, color=COLORS['purple_end'])

    questions = [
        "1. Describe the shape of the Mid-Atlantic Ridge. What does it look like?",
        "2. How does the AGE of the seafloor change as you move away from the ridge?",
        "3. What might explain why the youngest rock is at the center of the ridge?",
        "4. How is this connected to what we learned about divergent boundaries?",
        "5. What questions do you have about seafloor formation?",
    ]

    y_pos = 1.2
    for q in questions:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.5), COLORS['light_gray_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.05), Inches(9.0), Inches(0.4),
                    q, font_size=11, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.55

    add_colored_shape(slide, Inches(0.2), Inches(4.0), Inches(9.6), Inches(0.7), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.03), Inches(9.2), Inches(0.65),
                "TIP: The pattern in seafloor AGE is the key evidence for spreading!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station1_intro_slide(prs):
    """Slide 5: Station 1 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['green_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 1: Seafloor Spreading Evidence",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | 6 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(4.65), Inches(1.8), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(1.8), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(4.3), Inches(0.3),
                "RESOURCES:", font_size=12, bold=True, color=COLORS['green_end'])
    resources = """• Magnetic stripe data maps
• Age maps of ocean floor
• Spreading model animation
• Rate calculation data"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(4.3), Inches(1.2),
                resources, font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.05), Inches(1.15), Inches(4.7), Inches(1.8), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.25), Inches(1.25), Inches(4.3), Inches(0.3),
                "FOCUS:", font_size=12, bold=True, color=COLORS['blue_accent'])
    focus = """Interpret MAGNETIC stripe
evidence for spreading.

Earth's magnetic field has
REVERSED many times -
and the seafloor records it!"""
    add_text_box(slide, Inches(5.25), Inches(1.6), Inches(4.3), Inches(1.2),
                focus, font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(3.1), Inches(9.6), Inches(0.65), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.13), Inches(9.2), Inches(0.6),
                "🔗 SPIRAL from C4: Human impact on ocean systems - how do we study the deep ocean?",
                font_size=11, color=COLORS['dark_text'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(3.9), Inches(9.6), Inches(0.85), COLORS['magnetic_purple'])
    add_text_box(slide, Inches(0.4), Inches(4.0), Inches(9.2), Inches(0.7),
                "KEY EVIDENCE: Magnetic stripes are SYMMETRICAL on both sides of the ridge!\nThis proves new rock forms at the center and moves outward.",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station1_activity_slide(prs):
    """Slide 6: Station 1 Activity Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 1: Investigation Steps",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    steps = [
        ("OBSERVE", "Look at the magnetic stripe pattern - what do you notice?"),
        ("IDENTIFY", "Find matching stripes on opposite sides of the ridge"),
        ("CALCULATE", "Use distance and age to find spreading RATE (cm/year)"),
        ("EXPLAIN", "Why are stripes symmetrical? What does this prove?"),
        ("CONNECT", "How does this support the theory of plate tectonics?"),
    ]

    y_pos = 0.85
    colors = [COLORS['ocean_blue'], COLORS['magnetic_purple'], COLORS['green_end'],
              COLORS['ridge_orange'], COLORS['teal']]
    for i, (step, desc) in enumerate(steps):
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.58), COLORS['light_gray_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.58), colors[i])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(1.8), Inches(0.48),
                    step, font_size=12, bold=True, color=colors[i],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.3), Inches(y_pos + 0.05), Inches(7.3), Inches(0.48),
                    desc, font_size=11, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.63

    add_colored_shape(slide, Inches(0.2), Inches(4.05), Inches(9.6), Inches(0.65), COLORS['ocean_blue'])
    add_text_box(slide, Inches(0.4), Inches(4.08), Inches(9.2), Inches(0.6),
                "KEY: Symmetrical magnetic stripes = new rock forms at ridge, then splits and moves away!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station2_intro_slide(prs):
    """Slide 7: Station 2 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['seafloor_teal'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 2: Pangaea Puzzle Investigation",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | 5 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(1.0), COLORS['light_ocean_bg'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "RESOURCES:", font_size=12, bold=True, color=COLORS['seafloor_teal'])
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(0.5),
                "Continental cutouts • Fossil distribution maps • Rock type matching data • Climate evidence",
                font_size=11, color=COLORS['dark_text'])

    add_text_box(slide, Inches(0.3), Inches(2.3), Inches(9.4), Inches(0.3),
                "MULTIPLE LINES OF EVIDENCE:", font_size=13, bold=True, color=COLORS['seafloor_teal'])

    evidence = [
        ("Coastline Fit", "South America + Africa fit like puzzle pieces", COLORS['green_start']),
        ("Fossil Match", "Same fossils on continents now separated by ocean", COLORS['ridge_orange']),
        ("Rock Types", "Identical rock formations span multiple continents", COLORS['magnetic_purple']),
        ("Climate Clues", "Glacial deposits in now-tropical areas", COLORS['ocean_blue']),
    ]

    y_pos = 2.7
    for etype, description, color in evidence:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.45), COLORS['light_gray_bg'])
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(0.08), Inches(0.45), color)
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.02), Inches(2.5), Inches(0.4),
                    etype, font_size=11, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.1), Inches(y_pos + 0.02), Inches(6.4), Inches(0.4),
                    description, font_size=10, color=COLORS['dark_text'], anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.5

    add_colored_shape(slide, Inches(0.2), Inches(4.55), Inches(9.6), Inches(0.45), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.58), Inches(9.2), Inches(0.4),
                "FOCUS: Use MULTIPLE lines of evidence to support continental drift",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station2_activity_slide(prs):
    """Slide 8: Station 2 Activity Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['seafloor_teal'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 2: Pangaea Evidence Hunt",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(0.3), Inches(0.85), Inches(9.4), Inches(0.3),
                "FOSSILS THAT PROVE CONTINENTAL DRIFT:", font_size=13, bold=True, color=COLORS['seafloor_teal'])

    fossils = [
        ("Mesosaurus", "Freshwater reptile found in South America AND Africa - couldn't swim across ocean!"),
        ("Glossopteris", "Fern found on 5 continents - seeds too heavy to blow across oceans"),
        ("Lystrosaurus", "Land reptile found in Africa, India, AND Antarctica"),
        ("Cynognathus", "Land reptile found in South America AND Africa"),
    ]

    y_pos = 1.2
    for name, desc in fossils:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.55), COLORS['light_ocean_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.05), Inches(2.2), Inches(0.45),
                    name, font_size=11, bold=True, color=COLORS['ocean_blue'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.8), Inches(y_pos + 0.05), Inches(6.7), Inches(0.45),
                    desc, font_size=10, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    add_colored_shape(slide, Inches(0.2), Inches(3.65), Inches(9.6), Inches(1.1), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.75), Inches(9.2), Inches(0.25),
                "YOUR TASK:", font_size=12, bold=True, color=COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(4.05), Inches(9.2), Inches(0.65),
                "Argue LIKE A SCIENTIST: Use at least 3 different types of evidence to explain why continents must have been connected.",
                font_size=11, color=COLORS['dark_text'])


def add_station3_intro_slide(prs):
    """Slide 9: Station 3 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['ridge_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 3: Design a Plate Movement Detector",
                font_size=22, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "25 Points | 5 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(1.4), COLORS['light_orange_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(1.4), COLORS['ridge_red'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "ENGINEERING CHALLENGE:", font_size=14, bold=True, color=COLORS['ridge_red'])
    challenge = """Scientists need to measure how fast plates are moving. The movement is TINY -
only a few centimeters per year (about as fast as your fingernails grow).

Design a system to detect and measure this incredibly slow movement!"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(0.9),
                challenge, font_size=12, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(2.7), Inches(4.7), Inches(1.1), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.8), Inches(4.4), Inches(0.25),
                "AVAILABLE TOOLS:", font_size=11, bold=True, color=COLORS['blue_accent'])
    add_text_box(slide, Inches(0.4), Inches(3.1), Inches(4.4), Inches(0.65),
                "• GPS satellites (mm precision)\n• Laser rangefinders\n• Seismometers",
                font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.05), Inches(2.7), Inches(4.7), Inches(1.1), COLORS['light_green_bg'])
    add_text_box(slide, Inches(5.25), Inches(2.8), Inches(4.3), Inches(0.25),
                "SUCCESS CRITERIA:", font_size=11, bold=True, color=COLORS['green_end'])
    add_text_box(slide, Inches(5.25), Inches(3.1), Inches(4.3), Inches(0.65),
                "• Detect movement < 10 cm/year\n• Work across ocean distances\n• Provide continuous data",
                font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(3.95), Inches(9.6), Inches(0.7), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(3.98), Inches(9.2), Inches(0.65),
                "SEP-4: Analyzing Data | Real scientists use GPS to measure plate motion today!",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station3_activity_slide(prs):
    """Slide 10: Station 3 Activity Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['ridge_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 3: Measurement Design",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(0.3), Inches(0.85), Inches(9.4), Inches(0.3),
                "REAL PLATE VELOCITIES (cm/year):", font_size=13, bold=True, color=COLORS['ridge_orange'])

    velocities = [
        ("Pacific Plate", "~7-10 cm/year", "Fastest! Moving toward NW"),
        ("North American Plate", "~2.5 cm/year", "Moving toward SW"),
        ("African Plate", "~2 cm/year", "Moving toward NE"),
        ("Atlantic Seafloor", "~2.5 cm/year", "Spreading rate (each side)"),
    ]

    y_pos = 1.2
    for plate, rate, direction in velocities:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.48), COLORS['light_ocean_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.02), Inches(2.8), Inches(0.44),
                    plate, font_size=11, bold=True, color=COLORS['ocean_blue'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.4), Inches(y_pos + 0.02), Inches(2.0), Inches(0.44),
                    rate, font_size=11, color=COLORS['ridge_orange'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(5.5), Inches(y_pos + 0.02), Inches(4.0), Inches(0.44),
                    direction, font_size=10, color=COLORS['gray_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.53

    add_colored_shape(slide, Inches(0.2), Inches(3.4), Inches(9.6), Inches(0.85), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.5), Inches(9.2), Inches(0.25),
                "CALCULATION CHALLENGE:", font_size=12, bold=True, color=COLORS['magnetic_purple'])
    add_text_box(slide, Inches(0.4), Inches(3.8), Inches(9.2), Inches(0.4),
                "If the Atlantic is spreading at 2.5 cm/year per side, how wide was it 100 million years ago?",
                font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(4.4), Inches(9.6), Inches(0.55), COLORS['ocean_blue'])
    add_text_box(slide, Inches(0.4), Inches(4.43), Inches(9.2), Inches(0.5),
                "KEY: Plate motion is SLOW but adds up over millions of years!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_concepts_slide(prs):
    """Slide 11: Key Concepts"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Key Concepts This Week",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    concepts = [
        ("Seafloor Spreading", "New ocean floor forms at mid-ocean ridges, moves outward"),
        ("Magnetic Evidence", "Symmetrical magnetic stripes prove seafloor spreading"),
        ("Continental Drift", "Continents were once joined (Pangaea), then split apart"),
        ("Multiple Evidence", "Fossils, rocks, coastlines, climate all support drift theory"),
    ]

    y_pos = 0.9
    for title, detail in concepts:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.75), COLORS['light_ocean_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.75), COLORS['ocean_blue'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.08), Inches(3.5), Inches(0.6),
                    title, font_size=13, bold=True, color=COLORS['ocean_dark'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(4.0), Inches(y_pos + 0.08), Inches(5.6), Inches(0.6),
                    detail, font_size=11, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.82

    add_colored_shape(slide, Inches(0.2), Inches(4.2), Inches(9.6), Inches(0.55), COLORS['magnetic_purple'])
    add_text_box(slide, Inches(0.4), Inches(4.23), Inches(9.2), Inches(0.5),
                "CCC-1: Patterns | Matching patterns across continents reveal Earth's history!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_vocabulary_slide(prs):
    """Slide 12: Vocabulary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['purple_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Essential Vocabulary",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    vocab = [
        ("Seafloor spreading", "Process of new oceanic crust forming at mid-ocean ridges"),
        ("Mid-ocean ridge", "Underwater mountain range where plates diverge"),
        ("Magnetic reversal", "When Earth's magnetic poles switch (north↔south)"),
        ("Pangaea", "Supercontinent that existed ~200 million years ago"),
        ("Continental drift", "Theory that continents move slowly over time"),
        ("Subduction", "When one plate slides under another at convergent boundary"),
    ]

    y_pos = 0.9
    for term, definition in vocab:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.55), COLORS['light_purple_bg'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.5), Inches(0.45),
                    term, font_size=11, bold=True, color=COLORS['purple_end'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.0), Inches(y_pos + 0.05), Inches(6.6), Inches(0.45),
                    definition, font_size=10, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    add_colored_shape(slide, Inches(0.2), Inches(4.55), Inches(9.6), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(4.58), Inches(9.2), Inches(0.5),
                "📝 Add these terms to your NOTECARD for the exit ticket!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_exit_intro_slide(prs):
    """Slide 13: Exit Ticket Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['orange_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Exit Ticket: Earth's Dynamic Surface",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "23 Points | 6 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(2.4), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "QUESTION BREAKDOWN:", font_size=14, bold=True, color=COLORS['orange_end'])

    breakdown = [
        ("2 NEW", "Seafloor spreading, magnetic evidence, continental drift"),
        ("2 SPIRAL", "Connect to W1 boundary types + C4 human ocean study"),
        ("1 INTEGRATION", "Use multiple evidence types to support plate movement"),
        ("1 SEP", "Analyzing and interpreting data patterns"),
    ]

    y_pos = 1.65
    for qtype, desc in breakdown:
        add_text_box(slide, Inches(0.5), Inches(y_pos), Inches(1.8), Inches(0.45),
                    qtype, font_size=12, bold=True, color=COLORS['orange_end'])
        add_text_box(slide, Inches(2.4), Inches(y_pos), Inches(7.2), Inches(0.45),
                    desc, font_size=11, color=COLORS['dark_text'])
        y_pos += 0.5

    add_colored_shape(slide, Inches(0.2), Inches(3.7), Inches(9.6), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(3.73), Inches(9.2), Inches(0.5),
                "⏱️ TIME: ~15 minutes | Use your NOTECARD!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(4.4), Inches(9.6), Inches(0.55), COLORS['light_pink_bg'])
    add_text_box(slide, Inches(0.4), Inches(4.43), Inches(9.2), Inches(0.5),
                "⚠️ Remember: Plates move SLOWLY (~cm/year) but effects are HUGE over millions of years!",
                font_size=11, color=COLORS['dark_text'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_exit_tips_slide(prs):
    """Slide 14: Exit Ticket Tips"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Exit Ticket Tips",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    tips = [
        "Use the magnetic stripe pattern as EVIDENCE for spreading",
        "Connect fossils to continental drift - they couldn't swim!",
        "Remember: youngest rock at ridge, oldest at edges",
        "Use MULTIPLE evidence types in your arguments",
        "Check your notecard for key terms and concepts",
    ]

    y_pos = 0.85
    for i, tip in enumerate(tips, 1):
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.55), COLORS['light_gray_bg'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(0.5), Inches(0.45),
                    str(i), font_size=14, bold=True, color=COLORS['orange_end'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(1.0), Inches(y_pos + 0.05), Inches(8.6), Inches(0.45),
                    tip, font_size=12, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    add_colored_shape(slide, Inches(0.2), Inches(3.9), Inches(9.6), Inches(0.85), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(3.95), Inches(9.2), Inches(0.75),
                "You've got this! You found evidence, connected patterns, and calculated rates.\nNow show what you learned!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_notecard_slide(prs):
    """Slide 15: Notecard Reminder"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['ocean_blue'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "📝 Your Notecard Should Include:",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    items = [
        "Seafloor spreading process (new rock at ridge → moves outward)",
        "Magnetic stripe evidence (symmetrical pattern = proof)",
        "4+ fossils that prove continental drift (couldn't cross oceans)",
        "Plate velocity rates (~1-10 cm/year depending on plate)",
        "Connection to W1: spreading = divergent boundary in action",
    ]

    y_pos = 0.9
    for item in items:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.55), COLORS['light_ocean_bg'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(0.4), Inches(0.45),
                    "✓", font_size=14, bold=True, color=COLORS['green_end'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(0.9), Inches(y_pos + 0.05), Inches(8.7), Inches(0.45),
                    item, font_size=11, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    add_colored_shape(slide, Inches(0.2), Inches(3.95), Inches(9.6), Inches(0.85), COLORS['magnetic_purple'])
    add_text_box(slide, Inches(0.4), Inches(3.98), Inches(9.2), Inches(0.8),
                "Remember: Your notecard is the ONLY physical resource allowed!\nMake it count - organize it clearly!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_summary_slide(prs):
    """Slide 16: Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Week 2 Summary: What You Learned",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    takeaways = [
        ("Seafloor Spreading", "New crust forms at ridges, creating underwater mountains"),
        ("Magnetic Stripes", "Symmetrical patterns prove rock forms and splits"),
        ("Continental Drift", "Fossils + rocks + coastlines prove Pangaea existed"),
        ("Measurement", "GPS detects cm/year movement - adds up over millions of years"),
    ]

    y_pos = 0.9
    for title, detail in takeaways:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.65), COLORS['light_ocean_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.65), COLORS['ocean_blue'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.5), Inches(0.55),
                    title, font_size=13, bold=True, color=COLORS['ocean_dark'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.0), Inches(y_pos + 0.05), Inches(6.6), Inches(0.55),
                    detail, font_size=12, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.72

    add_colored_shape(slide, Inches(0.2), Inches(3.85), Inches(9.6), Inches(0.65), COLORS['seafloor_teal'])
    add_text_box(slide, Inches(0.4), Inches(3.88), Inches(9.2), Inches(0.6),
                "NEXT WEEK: Volcanic Eruption Styles - Why do some explode while others ooze?",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(4.65), Inches(9.6), Inches(0.55), COLORS['ridge_orange'])
    add_text_box(slide, Inches(0.4), Inches(4.68), Inches(9.2), Inches(0.5),
                "ANSWER: Underwater mountains form at divergent boundaries where new crust rises up!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


if __name__ == "__main__":
    prs = create_presentation()
    output_path = "/home/user/Kairos.Sci.Repo/content/grade7/cycle06/week2/G7_C6_W2_Slides_Final.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Created: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
