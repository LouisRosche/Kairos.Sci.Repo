#!/usr/bin/env python3
"""
Create G8_C4_W1 Energy Pyramids & Trophic Levels presentation.

See PPTX_DESIGN_GUIDE.md for best practices documentation.
"""

import os
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx_common import (
    COLORS as BASE_COLORS,
    create_base_presentation,
    add_text_box,
    Inches,
    Pt,
    PP_ALIGN,
    MSO_ANCHOR,
)

# Green theme for G8 C4 W1 - extend base colors
COLORS = {**BASE_COLORS}
COLORS.update({
    'green_primary': RGBColor(0x05, 0x96, 0x69),
    'green_dark': RGBColor(0x04, 0x78, 0x57),
    'green_light': RGBColor(0xEC, 0xFD, 0xF5),
    'green_lighter': RGBColor(0xD1, 0xFA, 0xE5),
    'hook_orange': RGBColor(0xF5, 0x9E, 0x0B),
    'hook_orange_dark': RGBColor(0xD9, 0x77, 0x06),
    'station1_blue': RGBColor(0x3B, 0x82, 0xF6),
    'station1_blue_dark': RGBColor(0x25, 0x63, 0xEB),
    'station2_purple': RGBColor(0x8B, 0x5C, 0xF6),
    'station2_purple_dark': RGBColor(0x7C, 0x3A, 0xED),
    'station3_green': RGBColor(0x05, 0x96, 0x69),
    'station3_green_dark': RGBColor(0x04, 0x78, 0x57),
    'exit_pink': RGBColor(0xEC, 0x48, 0x99),
    'exit_pink_dark': RGBColor(0xDB, 0x27, 0x77),
    'light_bg': RGBColor(0xF7, 0xFA, 0xFC),
    'warning_red_bg': RGBColor(0xFE, 0xD7, 0xD7),
})


def add_colored_shape(slide, left, top, width, height, color):
    """Add a colored rectangle shape (uses RECTANGLE for C4 theme)"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def create_presentation():
    """Create the G8_C4_W1 presentation"""
    prs = create_base_presentation()

    add_title_slide(prs)
    add_phenomenon_slide(prs)
    add_driving_question_slide(prs)
    add_prior_knowledge_slide(prs)
    add_learning_targets_slide(prs)
    add_vocabulary_slide(prs)
    add_hook_intro_slide(prs)
    add_hook_support_slide(prs)
    add_station1_intro_slide(prs)
    add_station1_support_slide(prs)
    add_station2_intro_slide(prs)
    add_station2_support_slide(prs)
    add_station3_intro_slide(prs)
    add_station3_support_slide(prs)
    add_exit_ticket_slide(prs)
    add_summary_slide(prs)

    return prs


def add_title_slide(prs):
    """Slide 1: Title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['green_primary'])

    # Emoji decoration
    add_text_box(slide, Inches(0), Inches(1.0), Inches(10), Inches(0.8),
                "🦁 📐", font_size=48, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Main title
    add_text_box(slide, Inches(0.5), Inches(1.8), Inches(9), Inches(1.0),
                "Week 1: Energy Pyramids & Trophic Levels",
                font_size=36, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Subtitle
    add_text_box(slide, Inches(0.5), Inches(2.9), Inches(9), Inches(0.5),
                "Grade 8 Science | Cycle 4 | Rosche | Kairos Academies",
                font_size=16, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Standards
    add_text_box(slide, Inches(0.5), Inches(3.5), Inches(9), Inches(0.4),
                "MS-LS2-3 Energy Flow in Ecosystems",
                font_size=14, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Points/Time
    add_text_box(slide, Inches(0.5), Inches(4.2), Inches(9), Inches(0.4),
                "100 Points Total | ~75 Minutes",
                font_size=14, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_phenomenon_slide(prs):
    """Slide 2: The Phenomenon"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['green_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🦁 The Phenomenon: The Missing Energy Mystery",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Main content box
    add_colored_shape(slide, Inches(0.15), Inches(1.0), Inches(9.7), Inches(2.8), COLORS['green_light'])

    # Content
    add_text_box(slide, Inches(0.4), Inches(1.1), Inches(9.2), Inches(0.5),
                "In the Serengeti ecosystem of Africa:",
                font_size=16, color=COLORS['dark_text'])

    add_text_box(slide, Inches(0.6), Inches(1.6), Inches(9.0), Inches(1.8),
                "• There are millions of pounds of grasses\n" +
                "• Hundreds of thousands of wildebeest, zebras, gazelles\n" +
                "• But only a few thousand lions and hyenas\n\n" +
                "Why are there SO MANY plants and herbivores, but SO FEW top predators?\n" +
                "Where does all the energy go?",
                font_size=15, color=COLORS['gray_text'])

    # Key insight box
    add_colored_shape(slide, Inches(0.15), Inches(3.95), Inches(9.7), Inches(1.4), COLORS['green_dark'])
    add_text_box(slide, Inches(0.35), Inches(4.15), Inches(9.3), Inches(1.0),
                "🔬 KEY CONCEPT: Only 10% of energy transfers between trophic levels!\n" +
                "The rest is 'lost' to heat, movement, and metabolism.",
                font_size=16, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_driving_question_slide(prs):
    """Slide 3: Driving Question"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['green_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🤔 Driving Question",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Question box
    add_colored_shape(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(2.0), COLORS['green_dark'])
    add_text_box(slide, Inches(0.7), Inches(1.7), Inches(8.6), Inches(1.6),
                "Why are there more plants than animals,\nand more herbivores than predators,\nin every ecosystem on Earth?",
                font_size=22, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Connection to Cycle 3
    add_colored_shape(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(1.4), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.7), Inches(3.9), Inches(8.6), Inches(1.2),
                "🔗 Connection to Cycle 3:\n" +
                "You learned about natural selection and adaptation.\n" +
                "Now we'll explore how energy availability limits populations at each level!",
                font_size=14, color=COLORS['green_dark'])


def add_prior_knowledge_slide(prs):
    """Slide 4: What You Already Know"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['green_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "💡 What You Already Know (Cycle 3 Connections)",
                font_size=22, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Content boxes
    concepts = [
        ("Natural Selection", "Organisms with advantageous\ntraits survive and reproduce"),
        ("Adaptation", "Traits that help organisms\nsurvive in their environment"),
        ("Population Dynamics", "Populations change based\non resources and competition")
    ]

    x_positions = [0.3, 3.5, 6.7]
    for i, (title, content) in enumerate(concepts):
        add_colored_shape(slide, Inches(x_positions[i]), Inches(1.1), Inches(3.0), Inches(1.8), COLORS['green_light'])
        add_text_box(slide, Inches(x_positions[i] + 0.15), Inches(1.2), Inches(2.7), Inches(0.4),
                    title, font_size=14, bold=True, color=COLORS['green_dark'], align=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x_positions[i] + 0.15), Inches(1.6), Inches(2.7), Inches(1.2),
                    content, font_size=12, color=COLORS['gray_text'], align=PP_ALIGN.CENTER)

    # Common mistake alert
    add_colored_shape(slide, Inches(0.3), Inches(3.1), Inches(9.4), Inches(2.2), COLORS['warning_red_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(9.0), Inches(0.4),
                "⚠️ COMMON MISTAKE ALERT: 'Energy is Destroyed'",
                font_size=14, bold=True, color=COLORS['red_alert'])
    add_text_box(slide, Inches(0.5), Inches(3.6), Inches(9.0), Inches(1.6),
                "WRONG: 'Energy is destroyed as it moves up the food chain.'\n\n" +
                "RIGHT: 'Energy is converted to heat through metabolism, movement, and digestion.\n" +
                "It's not destroyed—it just becomes unusable for the next level.'\n\n" +
                "KEY: Energy is conserved, but only ~10% transfers to the next trophic level.",
                font_size=12, color=COLORS['red_alert'])


def add_learning_targets_slide(prs):
    """Slide 5: Learning Targets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['green_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🎯 Learning Targets",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    targets = [
        "Trace energy flow through trophic levels",
        "Apply the 10% rule to calculate energy transfer",
        "Calculate and compare biomass at different levels",
        "Design a sustainable food production system"
    ]

    y_pos = 1.0
    for i, target in enumerate(targets, 1):
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.95), COLORS['green_light'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.15), Inches(9.0), Inches(0.65),
                    f"Target {i}: {target}",
                    font_size=15, bold=True, color=COLORS['green_dark'])
        y_pos += 1.1


def add_vocabulary_slide(prs):
    """Slide 6: Key Vocabulary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['green_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "📚 Key Vocabulary This Week",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    vocab = [
        ("Trophic Level", "Position in food chain (1=producer, 2=primary consumer...)"),
        ("10% Rule", "Only ~10% of energy transfers between trophic levels"),
        ("Biomass", "Total mass of organisms at a trophic level"),
        ("Energy Pyramid", "Diagram showing energy decrease at each level"),
        ("Producer", "Makes own food via photosynthesis (plants, algae)"),
        ("Consumer", "Gets energy by eating other organisms")
    ]

    y_pos = 0.95
    for term, definition in vocab:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.7), COLORS['light_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.1), Inches(2.8), Inches(0.5),
                    term, font_size=12, bold=True, color=COLORS['green_dark'])
        add_text_box(slide, Inches(3.4), Inches(y_pos + 0.1), Inches(6.1), Inches(0.5),
                    definition, font_size=12, color=COLORS['gray_text'])
        y_pos += 0.75


def add_hook_intro_slide(prs):
    """Slide 7: Hook Intro"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header - two-line pattern
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['hook_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🦁 Hook – The Missing Energy Mystery",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.4),
                "12 Points | ~10 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # What you'll do
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(2.2), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "🦁 The Challenge: Where does the energy go?",
                font_size=16, bold=True, color=COLORS['hook_orange_dark'])
    add_text_box(slide, Inches(0.5), Inches(1.75), Inches(9.0), Inches(1.5),
                "1. Observe the phenomenon: Serengeti population data (2 min)\n" +
                "2. Make predictions about why predators are rare (3 min)\n" +
                "3. Connect to Cycle 3: natural selection (3 min)\n" +
                "4. Answer diagnostic questions (2 min)",
                font_size=14, color=COLORS['gray_text'])

    # Think about this
    add_colored_shape(slide, Inches(0.3), Inches(3.55), Inches(9.4), Inches(1.8), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.65), Inches(9.0), Inches(0.4),
                "💭 Think About This:",
                font_size=14, bold=True, color=COLORS['green_dark'])
    add_text_box(slide, Inches(0.5), Inches(4.05), Inches(9.0), Inches(1.2),
                "• Where does ALL energy in an ecosystem come from?\n" +
                "• What happens to energy when an animal eats another animal?\n" +
                "• How does this connect to natural selection from Cycle 3?",
                font_size=13, color=COLORS['gray_text'])


def add_hook_support_slide(prs):
    """Slide 8: Hook Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['hook_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🦁 Hook – Serengeti Data",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Population data
    add_colored_shape(slide, Inches(0.3), Inches(1.0), Inches(4.5), Inches(2.5), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(4.1), Inches(0.4),
                "📊 Serengeti Populations:", font_size=14, bold=True, color=COLORS['hook_orange_dark'])
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4.1), Inches(1.9),
                "Grasses: Millions of pounds\n\n" +
                "Herbivores (wildebeest, zebras):\nHundreds of thousands\n\n" +
                "Predators (lions, hyenas):\nOnly a few thousand",
                font_size=12, color=COLORS['gray_text'])

    # Energy source
    add_colored_shape(slide, Inches(5.0), Inches(1.0), Inches(4.7), Inches(2.5), COLORS['green_light'])
    add_text_box(slide, Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.4),
                "☀️ Energy Source:", font_size=14, bold=True, color=COLORS['green_dark'])
    add_text_box(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(1.9),
                "Sun → Plants (photosynthesis)\n\n" +
                "Plants → Herbivores (eating)\n\n" +
                "Herbivores → Predators (eating)\n\n" +
                "But energy is LOST at each step!",
                font_size=12, color=COLORS['gray_text'], align=PP_ALIGN.CENTER)

    # Notecard
    add_colored_shape(slide, Inches(0.3), Inches(3.7), Inches(9.4), Inches(1.65), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.78), Inches(9.0), Inches(0.5),
                "📝 Complete the Hook form on the student page",
                font_size=14, bold=True, color=COLORS['red_alert'], anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(0.5), Inches(4.2), Inches(9.0), Inches(1.0),
                "Record your predictions about:\n" +
                "• Why there are more plants than herbivores\n" +
                "• Why there are more herbivores than predators",
                font_size=12, color=COLORS['gray_text'])


def add_station1_intro_slide(prs):
    """Slide 9: Station 1 Intro"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header - two-line pattern
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['station1_blue'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "📊 Station 1 – 10% Rule Investigation",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.4),
                "20 Points | ~18 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Mission
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(1.3), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "📊 Your Mission: Discover the 10% Rule",
                font_size=16, bold=True, color=COLORS['station1_blue_dark'])
    add_text_box(slide, Inches(0.5), Inches(1.75), Inches(9.0), Inches(0.6),
                "Calculate how much energy transfers between trophic levels\nand explain why ecosystems are shaped like pyramids.",
                font_size=14, color=COLORS['gray_text'])

    # 10% rule visualization
    add_colored_shape(slide, Inches(0.3), Inches(2.65), Inches(9.4), Inches(2.7), COLORS['green_light'])
    add_text_box(slide, Inches(0.5), Inches(2.75), Inches(9.0), Inches(0.4),
                "🔬 Where Does Energy Go When a Zebra Eats Grass?",
                font_size=14, bold=True, color=COLORS['green_dark'])
    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(4.3), Inches(2.0),
                "Energy in Grass Eaten: 1000 kcal\n\n" +
                "Lost to heat (metabolism): ~60%\n" +
                "Lost to waste (feces): ~25%\n" +
                "Parts not eaten: ~5%\n" +
                "Stored in zebra body: ~10% ✓",
                font_size=12, color=COLORS['gray_text'])
    add_text_box(slide, Inches(5.0), Inches(3.2), Inches(4.5), Inches(2.0),
                "⚠️ KEY INSIGHT:\n\nThe 10% stored in the zebra\nis ALL that's available\nto the NEXT level (lions)!\n\n" +
                "This explains the pyramid shape!",
                font_size=12, bold=True, color=COLORS['green_dark'], align=PP_ALIGN.CENTER)


def add_station1_support_slide(prs):
    """Slide 10: Station 1 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['station1_blue'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "📊 Station 1 – Energy Transfer Data",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Energy data table
    add_colored_shape(slide, Inches(0.3), Inches(1.0), Inches(9.4), Inches(2.5), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(9.0), Inches(0.4),
                "📈 Energy at Each Trophic Level (starting with 10,000 kcal):",
                font_size=13, bold=True, color=COLORS['station1_blue_dark'])
    add_text_box(slide, Inches(0.5), Inches(1.55), Inches(9.0), Inches(1.8),
                "Trophic Level              |  Energy (kcal)  |  % of Original\n" +
                "─────────────────────────────────────────────────────────\n" +
                "Producers (grass)          |    10,000       |     100%\n" +
                "Primary consumers (zebras) |     1,000       |      10%\n" +
                "Secondary consumers (lions)|       100       |       1%\n" +
                "Tertiary consumers         |        10       |     0.1%",
                font_size=11, color=COLORS['gray_text'])

    # Calculation help
    add_colored_shape(slide, Inches(0.3), Inches(3.65), Inches(9.4), Inches(1.7), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.75), Inches(9.0), Inches(0.4),
                "🧮 How to Calculate: Multiply by 0.10 for each level",
                font_size=13, bold=True, color=COLORS['green_dark'])
    add_text_box(slide, Inches(0.5), Inches(4.15), Inches(4.3), Inches(1.1),
                "Level 2: 10,000 × 0.10 = 1,000\n" +
                "Level 3: 1,000 × 0.10 = 100\n" +
                "Level 4: 100 × 0.10 = 10",
                font_size=11, color=COLORS['gray_text'])
    add_text_box(slide, Inches(5.0), Inches(4.13), Inches(4.5), Inches(1.2),
                "📝 Complete Station 1 form\n" +
                "Apply the 10% rule and explain\nwhy energy decreases at each level.",
                font_size=11, bold=True, color=COLORS['red_alert'], anchor=MSO_ANCHOR.MIDDLE)


def add_station2_intro_slide(prs):
    """Slide 11: Station 2 Intro"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header - two-line pattern
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['station2_purple'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "⚖️ Station 2 – Biomass Calculations",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.4),
                "20 Points | ~15 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Mission
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(1.2), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "⚖️ Your Mission: Calculate Biomass and Energy",
                font_size=16, bold=True, color=COLORS['station2_purple_dark'])
    add_text_box(slide, Inches(0.5), Inches(1.75), Inches(9.0), Inches(0.5),
                "Use real Serengeti data to calculate biomass ratios and predict population changes.",
                font_size=14, color=COLORS['gray_text'])

    # Serengeti data
    add_colored_shape(slide, Inches(0.3), Inches(2.55), Inches(9.4), Inches(2.8), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.5), Inches(2.65), Inches(9.0), Inches(0.4),
                "📊 Serengeti Ecosystem Data:",
                font_size=14, bold=True, color=COLORS['station2_purple_dark'])
    add_text_box(slide, Inches(0.5), Inches(3.1), Inches(9.0), Inches(1.3),
                "Trophic Level           |  Biomass (kg/km²)  |  Population/100 km²\n" +
                "────────────────────────────────────────────────────────────────\n" +
                "Producers (grasses)     |    4,000,000       |        N/A\n" +
                "Primary (wildebeest)    |      400,000       |      1,500\n" +
                "Secondary (lions)       |       40,000       |         25",
                font_size=11, color=COLORS['gray_text'])
    add_text_box(slide, Inches(0.5), Inches(4.5), Inches(9.0), Inches(0.7),
                "Conversion: 1 kg biomass ≈ 5 kcal stored energy",
                font_size=12, bold=True, color=COLORS['station2_purple_dark'])


def add_station2_support_slide(prs):
    """Slide 12: Station 2 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['station2_purple'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "⚖️ Station 2 – Prediction Challenge",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Scenario
    add_colored_shape(slide, Inches(0.3), Inches(1.0), Inches(9.4), Inches(1.6), COLORS['warning_red_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(9.0), Inches(0.4),
                "🌱 SCENARIO: What if grass decreases by 50%?",
                font_size=14, bold=True, color=COLORS['red_alert'])
    add_text_box(slide, Inches(0.5), Inches(1.55), Inches(9.0), Inches(0.95),
                "A drought reduces grass biomass from 4,000,000 kg to 2,000,000 kg.\n" +
                "Using the 10% rule, predict what happens to herbivores and predators.\n" +
                "Show your calculations!",
                font_size=12, color=COLORS['gray_text'])

    # Calculation hints
    add_colored_shape(slide, Inches(0.3), Inches(2.75), Inches(4.5), Inches(2.0), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(2.85), Inches(4.1), Inches(0.4),
                "🧮 Calculation Help:", font_size=13, bold=True, color=COLORS['green_dark'])
    add_text_box(slide, Inches(0.5), Inches(3.25), Inches(4.1), Inches(1.4),
                "• Ratio = larger ÷ smaller\n" +
                "• Energy = biomass × 5 kcal/kg\n" +
                "• Example: 40,000 kg × 5 = 200,000 kcal",
                font_size=11, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(5.0), Inches(2.75), Inches(4.7), Inches(2.0), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(5.2), Inches(2.85), Inches(4.3), Inches(0.4),
                "💭 Think About:", font_size=13, bold=True, color=COLORS['station2_purple_dark'])
    add_text_box(slide, Inches(5.2), Inches(3.25), Inches(4.3), Inches(1.4),
                "• If grass ↓ 50%, herbivores ↓ ?\n" +
                "• If herbivores ↓, predators ↓ ?\n" +
                "• Does the 10% rule still apply?",
                font_size=11, color=COLORS['gray_text'])

    # Instructions
    add_colored_shape(slide, Inches(0.3), Inches(4.9), Inches(9.4), Inches(0.5), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(4.88), Inches(9.0), Inches(0.5),
                "📝 Complete Station 2 form - Calculate biomass ratios and make predictions!",
                font_size=13, bold=True, color=COLORS['red_alert'], anchor=MSO_ANCHOR.MIDDLE)


def add_station3_intro_slide(prs):
    """Slide 13: Station 3 Intro"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header - two-line pattern
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['station3_green'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🌾 Station 3 – Design a Sustainable Farm",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.4),
                "25 Points | ~20 min (Highest Value!)", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Mission
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(1.2), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "🌾 Engineering Challenge: Design a Sustainable Food System",
                font_size=15, bold=True, color=COLORS['green_dark'])
    add_text_box(slide, Inches(0.5), Inches(1.75), Inches(9.0), Inches(0.5),
                "Use the 10% rule to design the most efficient farm to feed 1,000 people.",
                font_size=14, color=COLORS['gray_text'])

    # Constraints
    add_colored_shape(slide, Inches(0.3), Inches(2.55), Inches(9.4), Inches(2.8), COLORS['green_light'])
    add_text_box(slide, Inches(0.5), Inches(2.65), Inches(9.0), Inches(0.4),
                "📋 CONSTRAINTS:", font_size=14, bold=True, color=COLORS['green_dark'])
    add_text_box(slide, Inches(0.5), Inches(3.1), Inches(4.3), Inches(2.1),
                "• Goal: Feed 1,000 people\n  (2,000,000 kcal/day total)\n" +
                "• Land: 100 hectares available\n" +
                "• Budget: $500,000 to set up",
                font_size=12, color=COLORS['gray_text'])
    add_text_box(slide, Inches(5.0), Inches(2.65), Inches(4.5), Inches(0.4),
                "💡 10% Rule Insight:", font_size=14, bold=True, color=COLORS['green_dark'])
    add_text_box(slide, Inches(5.0), Inches(3.1), Inches(4.5), Inches(2.1),
                "Vegetables are 10× more efficient\nthan meat!\n\n" +
                "Why? Eating plants = 1 step\nEating meat = 2+ steps\n(more energy lost)",
                font_size=12, color=COLORS['gray_text'], align=PP_ALIGN.CENTER)


def add_station3_support_slide(prs):
    """Slide 14: Station 3 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['station3_green'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🌾 Station 3 – Available Options",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Options
    options = [
        ("A. Vegetables", "50,000 kcal/ha/day", "$2,000/ha"),
        ("B. Chicken", "5,000 kcal/ha/day", "$5,000/ha"),
        ("C. Cattle", "2,500 kcal/ha/day", "$8,000/ha"),
        ("D. Fish", "7,500 kcal/ha/day", "$10,000/ha")
    ]

    y_pos = 1.0
    for option, output, cost in options:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.8), COLORS['light_green_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.15), Inches(3.0), Inches(0.5),
                    option, font_size=13, bold=True, color=COLORS['green_dark'])
        add_text_box(slide, Inches(3.6), Inches(y_pos + 0.15), Inches(2.8), Inches(0.5),
                    output, font_size=12, color=COLORS['gray_text'], align=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(6.5), Inches(y_pos + 0.15), Inches(2.8), Inches(0.5),
                    cost, font_size=12, color=COLORS['gray_text'], align=PP_ALIGN.CENTER)
        y_pos += 0.9

    # Strategy hint
    add_colored_shape(slide, Inches(0.3), Inches(4.65), Inches(9.4), Inches(0.75), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.5), Inches(4.68), Inches(9.0), Inches(0.7),
                "💡 Strategy: 40 ha vegetables × 50,000 = 2,000,000 kcal ✓ | Cost: 40 × $2,000 = $80,000 (under budget!)\n" +
                "📝 Complete Station 3 - Design your farm and explain your choices using the 10% rule!",
                font_size=11, color=COLORS['gray_text'], anchor=MSO_ANCHOR.MIDDLE)


def add_exit_ticket_slide(prs):
    """Slide 15: Exit Ticket"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header - two-line pattern
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['exit_pink'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🎓 Exit Ticket – Energy Flow Integration",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.4),
                "23 Points | ~15 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Question types
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(2.5), COLORS['light_pink_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "🎓 Show What You Learned - Question Types:",
                font_size=16, bold=True, color=COLORS['exit_pink_dark'])
    add_text_box(slide, Inches(0.5), Inches(1.8), Inches(9.0), Inches(1.8),
                "• 2 NEW – Energy pyramids and 10% rule (this week)\n" +
                "• 2 SPIRAL – Cycle 3 review (natural selection, adaptation)\n" +
                "• 1 INTEGRATION – Connect energy limits to selection pressure\n" +
                "• 1 SEP-2 – Develop a model to explain the 10% rule",
                font_size=14, color=COLORS['gray_text'])

    # Reminder
    add_colored_shape(slide, Inches(0.3), Inches(3.85), Inches(9.4), Inches(1.5), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.93), Inches(9.0), Inches(0.5),
                "📝 Complete the Exit Ticket on the student page",
                font_size=14, bold=True, color=COLORS['red_alert'], anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(0.5), Inches(4.35), Inches(9.0), Inches(0.9),
                "This is your final assessment for Week 1. Take your time!\n" +
                "Use evidence from today's activities to support your answers.",
                font_size=13, color=COLORS['gray_text'])


def add_summary_slide(prs):
    """Slide 16: Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['green_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🎉 Week 1 Summary: What You Learned",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Key takeaways
    takeaways = [
        ("10% Rule", "Only ~10% of energy transfers between trophic levels", COLORS['green_primary']),
        ("Energy Loss", "90% lost to heat, waste, movement, and uneaten parts", COLORS['station1_blue']),
        ("Pyramid Shape", "Wide at bottom (producers), narrow at top (apex predators)", COLORS['station2_purple']),
        ("Efficiency", "Eating lower on the food chain uses less land and energy", COLORS['hook_orange'])
    ]

    y_pos = 1.0
    for title, content, color in takeaways:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.9), COLORS['light_bg'])
        # Color bar on left
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(0.1), Inches(0.9), color)
        add_text_box(slide, Inches(0.55), Inches(y_pos + 0.1), Inches(9.0), Inches(0.35),
                    title, font_size=13, bold=True, color=color)
        add_text_box(slide, Inches(0.55), Inches(y_pos + 0.45), Inches(9.0), Inches(0.4),
                    content, font_size=12, color=COLORS['gray_text'])
        y_pos += 1.0

    # Next week preview
    add_colored_shape(slide, Inches(0.3), Inches(5.05), Inches(9.4), Inches(0.45), COLORS['green_light'])
    add_text_box(slide, Inches(0.5), Inches(5.1), Inches(9.0), Inches(0.35),
                "Next Week: How do invasive species disrupt energy flow?",
                font_size=12, color=COLORS['green_dark'], align=PP_ALIGN.CENTER)


def main():
    """Main function to create and save the presentation"""
    prs = create_presentation()

    # Ensure output directory exists
    output_dir = "/home/user/Kairos.Sci.Repo/content/grade8/cycle04/week1"
    os.makedirs(output_dir, exist_ok=True)

    # Save the presentation
    output_path = os.path.join(output_dir, "G8_C4_W1_Energy_Pyramids_Slides_Final.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
