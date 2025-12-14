#!/usr/bin/env python3
"""
Create G7_C6_W4 Earth's Interior Structure & Evidence presentation.
Topic: Plate Tectonics & Earth's Interior.

Instructional week structure (16 slides).
See PPTX_DESIGN_GUIDE.md for best practices documentation.
"""

import os
from pptx.dml.color import RGBColor
from pptx_common import (
    COLORS as BASE_COLORS,
    create_base_presentation,
    add_colored_shape,
    add_text_box,
    Inches,
    Pt,
    PP_ALIGN,
    MSO_ANCHOR,
)

# Earth interior theme (deep reds and golds) - extend base colors
COLORS = {**BASE_COLORS}
COLORS.update({
    'core_red': RGBColor(0xB9, 0x1C, 0x1C),
    'core_dark': RGBColor(0x7F, 0x1D, 0x1D),
    'mantle_orange': RGBColor(0xEA, 0x58, 0x0C),
    'crust_brown': RGBColor(0x8B, 0x45, 0x13),
    'seismic_blue': RGBColor(0x1D, 0x4E, 0xD8),
    'light_core_bg': RGBColor(0xFE, 0xE2, 0xE2),
    'wave_purple': RGBColor(0x7C, 0x3A, 0xED),
})


def create_presentation():
    """Create the G7_C6_W4 presentation"""
    prs = create_base_presentation()

    add_title_slide(prs)
    add_phenomenon_slide(prs)
    add_hook_intro_slide(prs)
    add_hook_activity_slide(prs)
    add_station1_intro_slide(prs)
    add_station1_activity_slide(prs)
    add_station2_intro_slide(prs)
    add_station2_activity_slide(prs)
    add_station3_intro_slide(prs)
    add_station3_activity_slide(prs)
    add_concepts_slide(prs)
    add_vocabulary_slide(prs)
    add_exit_intro_slide(prs)
    add_exit_tips_slide(prs)
    add_notecard_slide(prs)
    add_summary_slide(prs)

    return prs


def add_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['core_red'])

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(1.2),
                "Plate Tectonics & Earth's Interior",
                font_size=40, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")

    add_text_box(slide, Inches(0.5), Inches(2.6), Inches(9), Inches(0.5),
                "Week 4: Earth's Interior Structure & Evidence",
                font_size=24, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(9), Inches(0.5),
                "Grade 7 Science | Cycle 6 | 100 Points",
                font_size=18, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(2), Inches(4.1), Inches(6), Inches(1.0), COLORS['white'])
    add_text_box(slide, Inches(2.2), Inches(4.2), Inches(5.6), Inches(0.8),
                "Hook: 12 pts | Stations: 20+20+25 pts | Exit: 23 pts",
                font_size=14, bold=True, color=COLORS['core_red'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_phenomenon_slide(prs):
    """Slide 2: Phenomenon Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['mantle_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "This Week's Phenomenon",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(9.6), Inches(1.4), COLORS['light_core_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(0.08), Inches(1.4), COLORS['core_dark'])
    add_text_box(slide, Inches(0.5), Inches(1.0), Inches(9.0), Inches(1.2),
                "How do we know what's inside Earth if no one has ever been there?",
                font_size=26, bold=True, color=COLORS['core_dark'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(2.45), Inches(9.6), Inches(1.6), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.55), Inches(9.2), Inches(0.3),
                "THINK ABOUT IT:", font_size=14, bold=True, color=COLORS['blue_accent'])
    context = """• Earth is 6,371 km to the center - but the deepest hole ever drilled is only 12 km!
• We've never seen Earth's interior directly - no cameras, no samples below ~12 km
• Yet we know there's a liquid outer core and solid inner core
• How do scientists figure out what's inside using INDIRECT evidence?"""
    add_text_box(slide, Inches(0.4), Inches(2.9), Inches(9.2), Inches(1.0),
                context, font_size=13, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(4.2), Inches(9.6), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.23), Inches(9.2), Inches(0.5),
                "Building on W1-W3: How plate tectonics connects to Earth's internal heat engine!",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_hook_intro_slide(prs):
    """Slide 3: Hook Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['purple_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Hook: The Unreachable Mystery",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "12 Points | 5 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(1.6), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(1.6), COLORS['purple_end'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "YOUR RESOURCES:", font_size=14, bold=True, color=COLORS['purple_end'])
    resources = """• Scale comparison: Deepest drill hole (12 km) vs Earth's radius (6,371 km)
• Core/mantle/crust diagram with scale
• Seismic wave animations showing how waves travel through Earth

OBSERVE: How might waves tell us about what's inside?"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(1.1),
                resources, font_size=12, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(2.9), Inches(9.6), Inches(0.95), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.0), Inches(9.2), Inches(0.25),
                "KEY QUESTION:", font_size=12, bold=True, color=COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(3.3), Inches(9.2), Inches(0.5),
                "If you can't see or touch something, what OTHER evidence could tell you about it?",
                font_size=12, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(4.0), Inches(9.6), Inches(0.7), COLORS['seismic_blue'])
    add_text_box(slide, Inches(0.4), Inches(4.1), Inches(9.2), Inches(0.55),
                "HINT: Think about doctors using ultrasound - same principle!",
                font_size=11, bold=True, color=COLORS['white'], anchor=MSO_ANCHOR.MIDDLE,
                align=PP_ALIGN.CENTER)


def add_hook_activity_slide(prs):
    """Slide 4: Hook Activity Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['purple_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Hook Activity: Imagine the Impossible",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(0.3), Inches(0.85), Inches(9.4), Inches(0.3),
                "QUESTIONS YOU'LL ANSWER:", font_size=13, bold=True, color=COLORS['purple_end'])

    questions = [
        "1. How deep is 12 km compared to Earth's total radius? What fraction?",
        "2. Why can't we just drill deeper to see what's inside?",
        "3. How do doctors 'see' inside your body without cutting you open?",
        "4. What happens to waves when they hit different materials?",
        "5. How might earthquakes help us learn about Earth's interior?",
    ]

    y_pos = 1.2
    for q in questions:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.5), COLORS['light_gray_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.05), Inches(9.0), Inches(0.4),
                    q, font_size=11, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.55

    add_colored_shape(slide, Inches(0.2), Inches(4.0), Inches(9.6), Inches(0.7), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.03), Inches(9.2), Inches(0.65),
                "TIP: Seismic waves from earthquakes act like an ultrasound for Earth!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station1_intro_slide(prs):
    """Slide 5: Station 1 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['green_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 1: Seismic Wave Investigation",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | 6 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(4.65), Inches(1.8), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(1.8), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(4.3), Inches(0.3),
                "RESOURCES:", font_size=12, bold=True, color=COLORS['green_end'])
    resources = """• Seismic wave simulation
• P-wave vs S-wave animations
• Wave travel time data
• Layer boundary diagrams"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(4.3), Inches(1.2),
                resources, font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.05), Inches(1.15), Inches(4.7), Inches(1.8), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.25), Inches(1.25), Inches(4.3), Inches(0.3),
                "TWO WAVE TYPES:", font_size=12, bold=True, color=COLORS['blue_accent'])
    focus = """P-waves: PRIMARY
• Fastest, travel through
  solids AND liquids

S-waves: SECONDARY
• Slower, travel through
  solids ONLY"""
    add_text_box(slide, Inches(5.25), Inches(1.6), Inches(4.3), Inches(1.2),
                focus, font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(3.1), Inches(9.6), Inches(0.65), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.13), Inches(9.2), Inches(0.6),
                "🔗 SPIRAL from W3: Magma is liquid! How would S-waves behave in liquid outer core?",
                font_size=11, color=COLORS['dark_text'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(3.9), Inches(9.6), Inches(0.85), COLORS['wave_purple'])
    add_text_box(slide, Inches(0.4), Inches(4.0), Inches(9.2), Inches(0.7),
                "KEY: S-waves CAN'T pass through liquid → They DISAPPEAR at the outer core!\nThis is how we KNOW the outer core is liquid!",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station1_activity_slide(prs):
    """Slide 6: Station 1 Activity Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 1: Wave Behavior Analysis",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(0.3), Inches(0.85), Inches(9.4), Inches(0.3),
                "WHAT SEISMIC WAVES TELL US:", font_size=13, bold=True, color=COLORS['green_end'])

    evidence = [
        ("P-waves bend", "Density increases with depth → waves speed up and curve"),
        ("S-wave shadow zone", "S-waves don't reach far side → liquid layer blocks them"),
        ("P-wave shadow zone", "P-waves refract through liquid → creates 'shadow' region"),
        ("Wave speed changes", "Sudden speed changes = boundary between layers"),
    ]

    y_pos = 1.2
    for observation, interpretation in evidence:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.55), COLORS['light_green_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.05), Inches(2.5), Inches(0.45),
                    observation, font_size=11, bold=True, color=COLORS['green_end'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.1), Inches(y_pos + 0.05), Inches(6.4), Inches(0.45),
                    interpretation, font_size=10, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    add_colored_shape(slide, Inches(0.2), Inches(3.65), Inches(9.6), Inches(0.65), COLORS['seismic_blue'])
    add_text_box(slide, Inches(0.4), Inches(3.68), Inches(9.2), Inches(0.6),
                "KEY: Waves REFRACT (bend) when they hit different materials!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(4.45), Inches(9.6), Inches(0.5), COLORS['core_red'])
    add_text_box(slide, Inches(0.4), Inches(4.48), Inches(9.2), Inches(0.45),
                "S-WAVE SHADOW = PROOF of liquid outer core!",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station2_intro_slide(prs):
    """Slide 7: Station 2 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['mantle_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 2: Earth's Layers Analysis",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | 5 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(1.0), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "RESOURCES:", font_size=12, bold=True, color=COLORS['mantle_orange'])
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(0.5),
                "Density data • Temperature gradient charts • Layer thickness diagrams • Material properties",
                font_size=11, color=COLORS['dark_text'])

    add_text_box(slide, Inches(0.3), Inches(2.3), Inches(9.4), Inches(0.3),
                "EARTH'S MAJOR LAYERS:", font_size=13, bold=True, color=COLORS['mantle_orange'])

    layers = [
        ("Crust", "5-70 km", "2.7-3.0 g/cm³", "Solid rock (continental/oceanic)", COLORS['crust_brown']),
        ("Mantle", "2,900 km", "3.3-5.5 g/cm³", "Plastic solid (flows slowly)", COLORS['mantle_orange']),
        ("Outer Core", "2,200 km", "9.9-12.2 g/cm³", "LIQUID iron-nickel", COLORS['core_red']),
        ("Inner Core", "1,220 km", "12.8-13.1 g/cm³", "SOLID iron-nickel (extreme pressure)", COLORS['core_dark']),
    ]

    y_pos = 2.7
    for name, thickness, density, state, color in layers:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.4), COLORS['light_gray_bg'])
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(0.08), Inches(0.4), color)
        add_text_box(slide, Inches(0.5), Inches(y_pos), Inches(1.5), Inches(0.4),
                    name, font_size=10, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.1), Inches(y_pos), Inches(1.4), Inches(0.4),
                    thickness, font_size=9, color=COLORS['dark_text'], anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.6), Inches(y_pos), Inches(1.8), Inches(0.4),
                    density, font_size=9, color=COLORS['dark_text'], anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(5.5), Inches(y_pos), Inches(4.0), Inches(0.4),
                    state, font_size=9, color=COLORS['gray_text'], anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.45

    add_colored_shape(slide, Inches(0.2), Inches(4.55), Inches(9.6), Inches(0.45), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.58), Inches(9.2), Inches(0.4),
                "FOCUS: Use DENSITY and SEISMIC data to construct a model of Earth's interior",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station2_activity_slide(prs):
    """Slide 8: Station 2 Activity Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['mantle_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 2: Evidence-Based Model Building",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(0.3), Inches(0.85), Inches(9.4), Inches(0.3),
                "EVIDENCE FOR EACH LAYER:", font_size=13, bold=True, color=COLORS['mantle_orange'])

    evidence = [
        ("Crust thin", "Seismic waves travel fast through → less dense rock"),
        ("Mantle plastic", "Convection patterns prove it flows slowly over time"),
        ("Outer core liquid", "S-waves BLOCKED completely → can't be solid"),
        ("Inner core solid", "P-waves speed UP again → denser than liquid"),
    ]

    y_pos = 1.2
    for claim, proof in evidence:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.55), COLORS['light_orange_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.05), Inches(2.5), Inches(0.45),
                    claim, font_size=11, bold=True, color=COLORS['mantle_orange'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.1), Inches(y_pos + 0.05), Inches(6.4), Inches(0.45),
                    proof, font_size=10, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    add_colored_shape(slide, Inches(0.2), Inches(3.65), Inches(9.6), Inches(1.1), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.75), Inches(9.2), Inches(0.25),
                "YOUR TASK:", font_size=12, bold=True, color=COLORS['blue_accent'])
    add_text_box(slide, Inches(0.4), Inches(4.05), Inches(9.2), Inches(0.65),
                "Draw an evidence-based model of Earth's interior. Label each layer with EVIDENCE that proves its properties.",
                font_size=11, color=COLORS['dark_text'])


def add_station3_intro_slide(prs):
    """Slide 9: Station 3 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['seismic_blue'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 3: Design an Earth Model",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "25 Points | 5 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(1.4), COLORS['light_blue_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(1.4), COLORS['seismic_blue'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "MODELING CHALLENGE:", font_size=14, bold=True, color=COLORS['seismic_blue'])
    challenge = """A museum wants to create an interactive Earth model that shows visitors why scientists
believe Earth has a liquid outer core and solid inner core.

Design a model that demonstrates how indirect evidence reveals hidden structure!"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(0.9),
                challenge, font_size=12, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(2.7), Inches(4.7), Inches(1.1), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.8), Inches(4.4), Inches(0.25),
                "MODEL REQUIREMENTS:", font_size=11, bold=True, color=COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(3.1), Inches(4.4), Inches(0.65),
                "• Show all 4 major layers\n• Demonstrate S-wave shadow zone\n• Explain evidence for each layer",
                font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.05), Inches(2.7), Inches(4.7), Inches(1.1), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(5.25), Inches(2.8), Inches(4.3), Inches(0.25),
                "SUCCESS CRITERIA:", font_size=11, bold=True, color=COLORS['wave_purple'])
    add_text_box(slide, Inches(5.25), Inches(3.1), Inches(4.3), Inches(0.65),
                "• Scientifically accurate\n• Clearly shows wave behavior\n• Visitors understand evidence",
                font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(3.95), Inches(9.6), Inches(0.7), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(3.98), Inches(9.2), Inches(0.65),
                "SEP-2: Developing and Using Models | Show how INDIRECT evidence reveals hidden structure!",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station3_activity_slide(prs):
    """Slide 10: Station 3 Activity Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['seismic_blue'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 3: Model Design Process",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(0.3), Inches(0.85), Inches(9.4), Inches(0.3),
                "MODEL FEATURES TO INCLUDE:", font_size=13, bold=True, color=COLORS['seismic_blue'])

    features = [
        ("Cross-section view", "Show layers at correct scale/proportions"),
        ("Wave paths", "Draw how P and S waves travel through layers"),
        ("Shadow zones", "Show where waves DON'T reach and WHY"),
        ("Evidence labels", "What PROVES each layer's properties"),
        ("Interactive element", "How will visitors engage with the model?"),
    ]

    y_pos = 1.2
    for feature, description in features:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.48), COLORS['light_blue_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.02), Inches(2.5), Inches(0.44),
                    feature, font_size=11, bold=True, color=COLORS['seismic_blue'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.1), Inches(y_pos + 0.02), Inches(6.4), Inches(0.44),
                    description, font_size=10, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.53

    add_colored_shape(slide, Inches(0.2), Inches(3.9), Inches(9.6), Inches(0.85), COLORS['wave_purple'])
    add_text_box(slide, Inches(0.4), Inches(4.0), Inches(9.2), Inches(0.7),
                "KEY MESSAGE: We use INDIRECT evidence (seismic waves) to 'see' what we can't directly observe!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_concepts_slide(prs):
    """Slide 11: Key Concepts"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Key Concepts This Week",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    concepts = [
        ("Seismic Waves", "P-waves through all; S-waves only through solids"),
        ("Shadow Zones", "S-wave shadow proves liquid outer core"),
        ("Four Layers", "Crust, mantle, outer core (liquid), inner core (solid)"),
        ("Indirect Evidence", "We 'see' inside Earth using wave behavior"),
    ]

    y_pos = 0.9
    for title, detail in concepts:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.75), COLORS['light_core_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.75), COLORS['core_red'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.08), Inches(2.5), Inches(0.6),
                    title, font_size=13, bold=True, color=COLORS['core_dark'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.0), Inches(y_pos + 0.08), Inches(6.6), Inches(0.6),
                    detail, font_size=11, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.82

    add_colored_shape(slide, Inches(0.2), Inches(4.2), Inches(9.6), Inches(0.55), COLORS['wave_purple'])
    add_text_box(slide, Inches(0.4), Inches(4.23), Inches(9.2), Inches(0.5),
                "CCC-3: Scale | Earth's interior is revealed through wave patterns at planetary scale!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_vocabulary_slide(prs):
    """Slide 12: Vocabulary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['purple_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Essential Vocabulary",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    vocab = [
        ("P-wave (Primary)", "Fastest seismic wave; travels through solids AND liquids"),
        ("S-wave (Secondary)", "Slower seismic wave; ONLY travels through solids"),
        ("Shadow zone", "Region where waves don't reach due to layer boundaries"),
        ("Refraction", "Bending of waves as they pass through different materials"),
        ("Lithosphere", "Rigid outer layer (crust + upper mantle)"),
        ("Asthenosphere", "Partially molten mantle layer where plates 'float'"),
    ]

    y_pos = 0.9
    for term, definition in vocab:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.55), COLORS['light_purple_bg'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.5), Inches(0.45),
                    term, font_size=10, bold=True, color=COLORS['purple_end'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.0), Inches(y_pos + 0.05), Inches(6.6), Inches(0.45),
                    definition, font_size=10, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    add_colored_shape(slide, Inches(0.2), Inches(4.55), Inches(9.6), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(4.58), Inches(9.2), Inches(0.5),
                "📝 Add these terms to your NOTECARD for the exit ticket!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_exit_intro_slide(prs):
    """Slide 13: Exit Ticket Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['orange_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Exit Ticket: Earth's Interior Integration",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "23 Points | 6 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(2.4), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "QUESTION BREAKDOWN:", font_size=14, bold=True, color=COLORS['orange_end'])

    breakdown = [
        ("2 NEW", "Seismic waves, shadow zones, Earth's layers"),
        ("2 SPIRAL", "Connect to W3 magma + W1 boundaries + C5 waves"),
        ("1 INTEGRATION", "Use seismic evidence to explain layer properties"),
        ("1 SEP", "Developing and using models with indirect evidence"),
    ]

    y_pos = 1.65
    for qtype, desc in breakdown:
        add_text_box(slide, Inches(0.5), Inches(y_pos), Inches(1.8), Inches(0.45),
                    qtype, font_size=12, bold=True, color=COLORS['orange_end'])
        add_text_box(slide, Inches(2.4), Inches(y_pos), Inches(7.2), Inches(0.45),
                    desc, font_size=11, color=COLORS['dark_text'])
        y_pos += 0.5

    add_colored_shape(slide, Inches(0.2), Inches(3.7), Inches(9.6), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(3.73), Inches(9.2), Inches(0.5),
                "⏱️ TIME: ~15 minutes | Use your NOTECARD!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(4.4), Inches(9.6), Inches(0.55), COLORS['core_red'])
    add_text_box(slide, Inches(0.4), Inches(4.43), Inches(9.2), Inches(0.5),
                "KEY: S-wave shadow = proof of liquid outer core!",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_exit_tips_slide(prs):
    """Slide 14: Exit Ticket Tips"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Exit Ticket Tips",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    tips = [
        "Know the difference: P-waves (Primary, pass through liquids) vs S-waves (Secondary, solids only)",
        "Shadow zones are EVIDENCE - explain WHY they form",
        "Connect layers to their properties (density, state, temperature)",
        "Remember: We use INDIRECT evidence - just like ultrasound!",
        "Check your notecard for wave behavior and layer characteristics",
    ]

    y_pos = 0.85
    for i, tip in enumerate(tips, 1):
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.55), COLORS['light_gray_bg'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(0.5), Inches(0.45),
                    str(i), font_size=14, bold=True, color=COLORS['orange_end'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(1.0), Inches(y_pos + 0.05), Inches(8.6), Inches(0.45),
                    tip, font_size=11, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    add_colored_shape(slide, Inches(0.2), Inches(3.9), Inches(9.6), Inches(0.85), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(3.95), Inches(9.2), Inches(0.75),
                "You've got this! You analyzed waves, interpreted shadows, and built models.\nNow show what you learned!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_notecard_slide(prs):
    """Slide 15: Notecard Reminder"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['core_red'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "📝 Your Notecard Should Include:",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    items = [
        "P-wave vs S-wave: Which travels through liquids? (P only!)",
        "Shadow zone explanation - why S-waves blocked at outer core",
        "4 layers with depths and states (solid/liquid/plastic)",
        "How indirect evidence reveals hidden structure",
        "Connection to all weeks: convection, boundaries, magma, layers",
    ]

    y_pos = 0.9
    for item in items:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.55), COLORS['light_core_bg'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(0.4), Inches(0.45),
                    "✓", font_size=14, bold=True, color=COLORS['green_end'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(0.9), Inches(y_pos + 0.05), Inches(8.7), Inches(0.45),
                    item, font_size=11, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    add_colored_shape(slide, Inches(0.2), Inches(3.95), Inches(9.6), Inches(0.85), COLORS['seismic_blue'])
    add_text_box(slide, Inches(0.4), Inches(3.98), Inches(9.2), Inches(0.8),
                "Remember: Your notecard is the ONLY physical resource allowed!\nMake it count - organize it clearly!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_summary_slide(prs):
    """Slide 16: Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Week 4 Summary: What You Learned",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    takeaways = [
        ("Seismic Waves", "P = all materials; S = solids only"),
        ("Shadow Zones", "S-wave shadow PROVES liquid outer core"),
        ("Earth's Layers", "Crust, mantle, outer core (liquid), inner core (solid)"),
        ("Indirect Evidence", "Like ultrasound - see hidden structure through waves"),
    ]

    y_pos = 0.9
    for title, detail in takeaways:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.65), COLORS['light_core_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.65), COLORS['core_red'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.5), Inches(0.55),
                    title, font_size=13, bold=True, color=COLORS['core_dark'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.0), Inches(y_pos + 0.05), Inches(6.6), Inches(0.55),
                    detail, font_size=12, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.72

    add_colored_shape(slide, Inches(0.2), Inches(3.85), Inches(9.6), Inches(0.65), COLORS['wave_purple'])
    add_text_box(slide, Inches(0.4), Inches(3.88), Inches(9.2), Inches(0.6),
                "NEXT WEEK: Synthesis & Assessment - Show what you've learned about Earth's dynamic systems!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(4.65), Inches(9.6), Inches(0.55), COLORS['mantle_orange'])
    add_text_box(slide, Inches(0.4), Inches(4.68), Inches(9.2), Inches(0.5),
                "ANSWER: Seismic waves reveal Earth's interior through their behavior in different materials!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


if __name__ == "__main__":
    prs = create_presentation()
    output_path = "/home/user/Kairos.Sci.Repo/content/grade7/cycle06/week4/G7_C6_W4_Slides_Final.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Created: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
