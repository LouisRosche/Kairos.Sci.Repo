#!/usr/bin/env python3
"""
Create G7_C6_W3 Volcanic Eruption Styles & Magma Properties presentation.
Topic: Plate Tectonics & Earth's Interior.

Instructional week structure (16 slides).
See PPTX_DESIGN_GUIDE.md for best practices documentation.
"""

import os
from pptx.dml.color import RGBColor
from pptx_common import (
    COLORS as BASE_COLORS,
    create_base_presentation,
    add_colored_shape,
    add_text_box,
    Inches,
    Pt,
    PP_ALIGN,
    MSO_ANCHOR,
)

# Volcanic theme (reds and oranges) - extend base colors
COLORS = {**BASE_COLORS}
COLORS.update({
    'lava_red': RGBColor(0xDC, 0x26, 0x26),
    'lava_dark': RGBColor(0xB9, 0x1C, 0x1C),
    'magma_orange': RGBColor(0xEA, 0x58, 0x0C),
    'ash_gray': RGBColor(0x4B, 0x55, 0x63),
    'basalt_dark': RGBColor(0x37, 0x41, 0x51),
    'light_orange_bg': RGBColor(0xFF, 0xED, 0xD5),
    'silica_blue': RGBColor(0x1D, 0x4E, 0xD8),
})


def create_presentation():
    """Create the G7_C6_W3 presentation"""
    prs = create_base_presentation()

    add_title_slide(prs)
    add_phenomenon_slide(prs)
    add_hook_intro_slide(prs)
    add_hook_activity_slide(prs)
    add_station1_intro_slide(prs)
    add_station1_activity_slide(prs)
    add_station2_intro_slide(prs)
    add_station2_activity_slide(prs)
    add_station3_intro_slide(prs)
    add_station3_activity_slide(prs)
    add_concepts_slide(prs)
    add_vocabulary_slide(prs)
    add_exit_intro_slide(prs)
    add_exit_tips_slide(prs)
    add_notecard_slide(prs)
    add_summary_slide(prs)

    return prs


def add_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['lava_red'])

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(1.2),
                "Plate Tectonics & Earth's Interior",
                font_size=40, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")

    add_text_box(slide, Inches(0.5), Inches(2.6), Inches(9), Inches(0.5),
                "Week 3: Volcanic Eruption Styles & Magma Properties",
                font_size=22, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(9), Inches(0.5),
                "Grade 7 Science | Cycle 6 | 100 Points",
                font_size=18, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(2), Inches(4.1), Inches(6), Inches(1.0), COLORS['white'])
    add_text_box(slide, Inches(2.2), Inches(4.2), Inches(5.6), Inches(0.8),
                "Hook: 12 pts | Stations: 20+20+25 pts | Exit: 23 pts",
                font_size=14, bold=True, color=COLORS['lava_red'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_phenomenon_slide(prs):
    """Slide 2: Phenomenon Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['magma_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "This Week's Phenomenon",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(9.6), Inches(1.4), COLORS['light_orange_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(0.08), Inches(1.4), COLORS['lava_dark'])
    add_text_box(slide, Inches(0.5), Inches(1.0), Inches(9.0), Inches(1.2),
                "Why do some volcanoes explode violently while others ooze slowly?",
                font_size=26, bold=True, color=COLORS['lava_dark'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(2.45), Inches(9.6), Inches(1.6), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.55), Inches(9.2), Inches(0.3),
                "THINK ABOUT IT:", font_size=14, bold=True, color=COLORS['blue_accent'])
    context = """• Mt. St. Helens (1980) exploded with the force of 500 atomic bombs!
• Hawaiian volcanoes produce slow-flowing lava you can walk alongside
• Both are volcanoes - so why such different behaviors?
• Today we discover how MAGMA PROPERTIES determine eruption style."""
    add_text_box(slide, Inches(0.4), Inches(2.9), Inches(9.2), Inches(1.0),
                context, font_size=13, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(4.2), Inches(9.6), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.23), Inches(9.2), Inches(0.5),
                "Building on W1-W2: Volcanoes occur at plate boundaries - now we learn WHY they differ!",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_hook_intro_slide(prs):
    """Slide 3: Hook Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['purple_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Hook: The Explosive vs Oozing Mystery",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "12 Points | 5 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(1.6), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(1.6), COLORS['purple_end'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "YOUR RESOURCES:", font_size=14, bold=True, color=COLORS['purple_end'])
    resources = """• Side-by-side video: Mt. St. Helens explosion vs Hawaiian lava flow
• Data table showing silica content of different magma types
• Viscosity comparison animations

OBSERVE: What differences do you notice in how the magma/lava behaves?"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(1.1),
                resources, font_size=12, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(2.9), Inches(9.6), Inches(0.95), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.0), Inches(9.2), Inches(0.25),
                "KEY QUESTION:", font_size=12, bold=True, color=COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(3.3), Inches(9.2), Inches(0.5),
                "What property of magma might explain these dramatically different eruption styles?",
                font_size=12, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(4.0), Inches(9.6), Inches(0.7), COLORS['lava_red'])
    add_text_box(slide, Inches(0.4), Inches(4.1), Inches(9.2), Inches(0.55),
                "HINT: Think about how honey vs water flows!",
                font_size=11, bold=True, color=COLORS['white'], anchor=MSO_ANCHOR.MIDDLE,
                align=PP_ALIGN.CENTER)


def add_hook_activity_slide(prs):
    """Slide 4: Hook Activity Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['purple_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Hook Activity: Compare & Contrast",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(0.3), Inches(0.85), Inches(9.4), Inches(0.3),
                "QUESTIONS YOU'LL ANSWER:", font_size=13, bold=True, color=COLORS['purple_end'])

    questions = [
        "1. Describe how the Mt. St. Helens eruption looked vs the Hawaiian eruption.",
        "2. Which lava flows faster? Which seems thicker?",
        "3. What do you think VISCOSITY means? How might it relate to eruptions?",
        "4. Looking at the silica data, what pattern do you notice?",
        "5. Predict: If you knew a magma's silica content, could you predict eruption style?",
    ]

    y_pos = 1.2
    for q in questions:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.5), COLORS['light_gray_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.05), Inches(9.0), Inches(0.4),
                    q, font_size=11, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.55

    add_colored_shape(slide, Inches(0.2), Inches(4.0), Inches(9.6), Inches(0.7), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.03), Inches(9.2), Inches(0.65),
                "TIP: Viscosity (thickness) is the KEY to understanding eruption styles!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station1_intro_slide(prs):
    """Slide 5: Station 1 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['green_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 1: Magma Viscosity Investigation",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | 6 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(4.65), Inches(1.8), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(1.8), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(4.3), Inches(0.3),
                "RESOURCES:", font_size=12, bold=True, color=COLORS['green_end'])
    resources = """• Viscosity simulation
• Magma composition data
• Temperature vs flow graphs
• Silica content charts"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(4.3), Inches(1.2),
                resources, font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.05), Inches(1.15), Inches(4.7), Inches(1.8), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.25), Inches(1.25), Inches(4.3), Inches(0.3),
                "KEY RELATIONSHIP:", font_size=12, bold=True, color=COLORS['blue_accent'])
    focus = """HIGH silica = HIGH viscosity
= THICK magma = EXPLOSIVE

LOW silica = LOW viscosity
= THIN magma = FLOWS EASILY"""
    add_text_box(slide, Inches(5.25), Inches(1.6), Inches(4.3), Inches(1.2),
                focus, font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(3.1), Inches(9.6), Inches(0.65), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.13), Inches(9.2), Inches(0.6),
                "🔗 SPIRAL from W1: Plate boundary TYPE affects magma composition!",
                font_size=11, color=COLORS['dark_text'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(3.9), Inches(9.6), Inches(0.85), COLORS['silica_blue'])
    add_text_box(slide, Inches(0.4), Inches(4.0), Inches(9.2), Inches(0.7),
                "THINK: Silica (SiO₂) forms chains in magma → More silica = more chains = thicker magma!",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station1_activity_slide(prs):
    """Slide 6: Station 1 Activity Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 1: Magma Types Comparison",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(0.3), Inches(0.85), Inches(9.4), Inches(0.3),
                "THREE MAIN MAGMA TYPES:", font_size=13, bold=True, color=COLORS['green_end'])

    magma_types = [
        ("Basaltic", "LOW silica (45-55%)", "THIN, flows easily", "Shield volcanoes (Hawaii)", COLORS['green_start']),
        ("Andesitic", "MEDIUM silica (55-65%)", "Moderate thickness", "Stratovolcanoes (Cascade Range)", COLORS['magma_orange']),
        ("Rhyolitic", "HIGH silica (65-75%)", "THICK, explosive", "Explosive eruptions (Mt. St. Helens)", COLORS['lava_red']),
    ]

    y_pos = 1.2
    for name, silica, viscosity, volcano, color in magma_types:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.7), COLORS['light_gray_bg'])
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(0.08), Inches(0.7), color)
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.05), Inches(1.5), Inches(0.6),
                    name, font_size=12, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.1), Inches(y_pos + 0.05), Inches(2.0), Inches(0.6),
                    silica, font_size=10, color=COLORS['dark_text'], anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(4.2), Inches(y_pos + 0.05), Inches(2.0), Inches(0.6),
                    viscosity, font_size=10, color=COLORS['dark_text'], anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(6.3), Inches(y_pos + 0.05), Inches(3.2), Inches(0.6),
                    volcano, font_size=10, color=COLORS['gray_text'], anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.75

    add_colored_shape(slide, Inches(0.2), Inches(3.5), Inches(9.6), Inches(0.65), COLORS['lava_red'])
    add_text_box(slide, Inches(0.4), Inches(3.53), Inches(9.2), Inches(0.6),
                "KEY: Silica content → Viscosity → Eruption style!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(4.3), Inches(9.6), Inches(0.55), COLORS['ash_gray'])
    add_text_box(slide, Inches(0.4), Inches(4.33), Inches(9.2), Inches(0.5),
                "Gas gets TRAPPED in thick magma → Pressure builds → EXPLOSION!",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station2_intro_slide(prs):
    """Slide 7: Station 2 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['magma_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 2: Eruption Type Analysis",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | 5 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(1.0), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "RESOURCES:", font_size=12, bold=True, color=COLORS['magma_orange'])
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(0.5),
                "Volcanic eruption case studies • Hazard maps • Eruption classification charts",
                font_size=11, color=COLORS['dark_text'])

    add_text_box(slide, Inches(0.3), Inches(2.3), Inches(9.4), Inches(0.3),
                "EXPLOSIVE vs EFFUSIVE:", font_size=13, bold=True, color=COLORS['magma_orange'])

    eruption_types = [
        ("EXPLOSIVE", "Violent, pyroclastic flows, ash clouds, lahars", "HIGH viscosity, gas-rich", COLORS['lava_red']),
        ("EFFUSIVE", "Slow lava flows, lava fountains, lava tubes", "LOW viscosity, gas escapes", COLORS['green_start']),
    ]

    y_pos = 2.7
    for etype, hazards, cause, color in eruption_types:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.6), COLORS['light_gray_bg'])
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(0.08), Inches(0.6), color)
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.05), Inches(1.8), Inches(0.5),
                    etype, font_size=12, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.4), Inches(y_pos + 0.05), Inches(4.2), Inches(0.5),
                    hazards, font_size=10, color=COLORS['dark_text'], anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(6.7), Inches(y_pos + 0.05), Inches(2.8), Inches(0.5),
                    cause, font_size=10, color=COLORS['gray_text'], anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.65

    add_colored_shape(slide, Inches(0.2), Inches(4.1), Inches(9.6), Inches(0.65), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.13), Inches(9.2), Inches(0.6),
                "FOCUS: Compare hazards from each eruption type. Which is more dangerous? Why?",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station2_activity_slide(prs):
    """Slide 8: Station 2 Activity Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['magma_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 2: Hazard Analysis",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(0.3), Inches(0.85), Inches(9.4), Inches(0.3),
                "VOLCANIC HAZARDS TO ANALYZE:", font_size=13, bold=True, color=COLORS['magma_orange'])

    hazards = [
        ("Pyroclastic flows", "Superheated gas + rock at 100+ mph", "Most deadly - no escape", COLORS['lava_red']),
        ("Lahars", "Volcanic mudflows down valleys", "Travel far from volcano", COLORS['ash_gray']),
        ("Ash fall", "Tiny glass particles in air", "Breathing hazard, roof collapse", COLORS['basalt_dark']),
        ("Lava flows", "Slow-moving molten rock", "Destructive but escapable", COLORS['magma_orange']),
    ]

    y_pos = 1.2
    for hazard, description, danger, color in hazards:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.55), COLORS['light_orange_bg'])
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(0.08), Inches(0.55), color)
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.05), Inches(2.2), Inches(0.45),
                    hazard, font_size=11, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.8), Inches(y_pos + 0.05), Inches(3.8), Inches(0.45),
                    description, font_size=10, color=COLORS['dark_text'], anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(6.7), Inches(y_pos + 0.05), Inches(2.8), Inches(0.45),
                    danger, font_size=10, color=COLORS['gray_text'], anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    add_colored_shape(slide, Inches(0.2), Inches(3.65), Inches(9.6), Inches(1.1), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.75), Inches(9.2), Inches(0.25),
                "YOUR TASK:", font_size=12, bold=True, color=COLORS['blue_accent'])
    add_text_box(slide, Inches(0.4), Inches(4.05), Inches(9.2), Inches(0.65),
                "Given a volcano's magma composition, predict: What hazards are most likely? How far should evacuations extend?",
                font_size=11, color=COLORS['dark_text'])


def add_station3_intro_slide(prs):
    """Slide 9: Station 3 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['lava_red'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 3: Design a Volcano Monitoring System",
                font_size=22, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "25 Points | 5 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(1.4), COLORS['light_orange_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(1.4), COLORS['lava_dark'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "ENGINEERING CHALLENGE:", font_size=14, bold=True, color=COLORS['lava_dark'])
    challenge = """A town of 50,000 people is located 30 km from an active stratovolcano.
Recent activity suggests it may erupt within the next few years.

Design a monitoring system to detect warning signs and protect the community!"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(0.9),
                challenge, font_size=12, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(2.7), Inches(4.7), Inches(1.1), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.8), Inches(4.4), Inches(0.25),
                "MONITORING TOOLS:", font_size=11, bold=True, color=COLORS['blue_accent'])
    add_text_box(slide, Inches(0.4), Inches(3.1), Inches(4.4), Inches(0.65),
                "• Seismometers (detect earthquakes)\n• Gas sensors (SO₂, CO₂)\n• GPS (ground deformation)",
                font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.05), Inches(2.7), Inches(4.7), Inches(1.1), COLORS['light_green_bg'])
    add_text_box(slide, Inches(5.25), Inches(2.8), Inches(4.3), Inches(0.25),
                "WARNING SIGNS:", font_size=11, bold=True, color=COLORS['green_end'])
    add_text_box(slide, Inches(5.25), Inches(3.1), Inches(4.3), Inches(0.65),
                "• Increased earthquake swarms\n• Rising gas emissions\n• Ground bulging upward",
                font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(3.95), Inches(9.6), Inches(0.7), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(3.98), Inches(9.2), Inches(0.65),
                "SEP-3: Planning Investigations | Design a system to DETECT eruption precursors!",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station3_activity_slide(prs):
    """Slide 10: Station 3 Activity Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['lava_red'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 3: Monitoring Design Process",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(0.3), Inches(0.85), Inches(9.4), Inches(0.3),
                "PRECURSOR SIGNS BEFORE ERUPTIONS:", font_size=13, bold=True, color=COLORS['lava_red'])

    precursors = [
        ("Earthquakes", "Magma movement cracks rock → swarms of small quakes"),
        ("Gas emissions", "SO₂ and CO₂ increase as magma rises"),
        ("Ground deformation", "Surface bulges as magma accumulates"),
        ("Temperature changes", "Nearby water/ground heats up"),
    ]

    y_pos = 1.2
    for sign, explanation in precursors:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.48), COLORS['light_orange_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.02), Inches(2.3), Inches(0.44),
                    sign, font_size=11, bold=True, color=COLORS['lava_dark'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.9), Inches(y_pos + 0.02), Inches(6.6), Inches(0.44),
                    explanation, font_size=10, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.53

    add_colored_shape(slide, Inches(0.2), Inches(3.4), Inches(9.6), Inches(1.35), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.5), Inches(9.2), Inches(0.25),
                "YOUR DESIGN MUST INCLUDE:", font_size=12, bold=True, color=COLORS['green_end'])
    design_req = """• Which sensors and where to place them
• Alert levels (green/yellow/orange/red)
• Evacuation triggers - what data means "evacuate now"?
• Communication plan to warn 50,000 people"""
    add_text_box(slide, Inches(0.4), Inches(3.8), Inches(9.2), Inches(0.9),
                design_req, font_size=11, color=COLORS['dark_text'])


def add_concepts_slide(prs):
    """Slide 11: Key Concepts"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Key Concepts This Week",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    concepts = [
        ("Viscosity", "How thick/sticky magma is - determines how it flows"),
        ("Silica Content", "More silica → higher viscosity → more explosive eruptions"),
        ("Eruption Types", "Explosive (thick magma) vs Effusive (thin magma)"),
        ("Monitoring", "Earthquakes, gas, deformation warn us before eruptions"),
    ]

    y_pos = 0.9
    for title, detail in concepts:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.75), COLORS['light_orange_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.75), COLORS['lava_red'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.08), Inches(2.5), Inches(0.6),
                    title, font_size=13, bold=True, color=COLORS['lava_dark'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.0), Inches(y_pos + 0.08), Inches(6.6), Inches(0.6),
                    detail, font_size=11, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.82

    add_colored_shape(slide, Inches(0.2), Inches(4.2), Inches(9.6), Inches(0.55), COLORS['silica_blue'])
    add_text_box(slide, Inches(0.4), Inches(4.23), Inches(9.2), Inches(0.5),
                "CCC-2: Cause and Effect | Silica content CAUSES differences in eruption behavior!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_vocabulary_slide(prs):
    """Slide 12: Vocabulary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['purple_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Essential Vocabulary",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    vocab = [
        ("Viscosity", "Measure of how thick/sticky a fluid is (resistance to flow)"),
        ("Silica (SiO₂)", "Silicon dioxide - increases magma viscosity"),
        ("Basaltic", "Low-silica magma; thin, fast-flowing; shield volcanoes"),
        ("Rhyolitic", "High-silica magma; thick, explosive; stratovolcanoes"),
        ("Pyroclastic flow", "Deadly mixture of hot gas, ash, rock fragments"),
        ("Effusive eruption", "Non-explosive; lava flows out calmly"),
    ]

    y_pos = 0.9
    for term, definition in vocab:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.55), COLORS['light_purple_bg'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.3), Inches(0.45),
                    term, font_size=11, bold=True, color=COLORS['purple_end'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.8), Inches(y_pos + 0.05), Inches(6.8), Inches(0.45),
                    definition, font_size=10, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    add_colored_shape(slide, Inches(0.2), Inches(4.55), Inches(9.6), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(4.58), Inches(9.2), Inches(0.5),
                "📝 Add these terms to your NOTECARD for the exit ticket!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_exit_intro_slide(prs):
    """Slide 13: Exit Ticket Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['orange_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Exit Ticket: Volcanic Systems Integration",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "23 Points | 6 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(2.4), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "QUESTION BREAKDOWN:", font_size=14, bold=True, color=COLORS['orange_end'])

    breakdown = [
        ("2 NEW", "Viscosity, silica, eruption styles, hazards"),
        ("2 SPIRAL", "Connect to W1-W2: plate boundaries create different magma types"),
        ("1 INTEGRATION", "Predict eruption style from magma composition"),
        ("1 SEP", "Planning monitoring investigations"),
    ]

    y_pos = 1.65
    for qtype, desc in breakdown:
        add_text_box(slide, Inches(0.5), Inches(y_pos), Inches(1.8), Inches(0.45),
                    qtype, font_size=12, bold=True, color=COLORS['orange_end'])
        add_text_box(slide, Inches(2.4), Inches(y_pos), Inches(7.2), Inches(0.45),
                    desc, font_size=11, color=COLORS['dark_text'])
        y_pos += 0.5

    add_colored_shape(slide, Inches(0.2), Inches(3.7), Inches(9.6), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(3.73), Inches(9.2), Inches(0.5),
                "⏱️ TIME: ~15 minutes | Use your NOTECARD!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(4.4), Inches(9.6), Inches(0.55), COLORS['lava_red'])
    add_text_box(slide, Inches(0.4), Inches(4.43), Inches(9.2), Inches(0.5),
                "KEY: High silica = High viscosity = Gas trapped = EXPLOSIVE!",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_exit_tips_slide(prs):
    """Slide 14: Exit Ticket Tips"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Exit Ticket Tips",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    tips = [
        "Remember the chain: Silica → Viscosity → Eruption style",
        "Connect magma type to volcano shape (shield vs stratovolcano)",
        "Know the hazards for each eruption type",
        "Explain WHY gas causes explosions in thick magma",
        "Check your notecard for key terms and relationships",
    ]

    y_pos = 0.85
    for i, tip in enumerate(tips, 1):
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.55), COLORS['light_gray_bg'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(0.5), Inches(0.45),
                    str(i), font_size=14, bold=True, color=COLORS['orange_end'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(1.0), Inches(y_pos + 0.05), Inches(8.6), Inches(0.45),
                    tip, font_size=12, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    add_colored_shape(slide, Inches(0.2), Inches(3.9), Inches(9.6), Inches(0.85), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(3.95), Inches(9.2), Inches(0.75),
                "You've got this! You compared eruptions, analyzed hazards, and designed monitoring.\nNow show what you learned!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_notecard_slide(prs):
    """Slide 15: Notecard Reminder"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['lava_red'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "📝 Your Notecard Should Include:",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    items = [
        "Viscosity definition and its relationship to silica content",
        "3 magma types: basaltic (low), andesitic (med), rhyolitic (high)",
        "Explosive vs effusive eruptions and their hazards",
        "Eruption precursors: earthquakes, gas, ground deformation",
        "Connection to W1: subduction zones → high-silica magma → explosive",
    ]

    y_pos = 0.9
    for item in items:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.55), COLORS['light_orange_bg'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(0.4), Inches(0.45),
                    "✓", font_size=14, bold=True, color=COLORS['green_end'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(0.9), Inches(y_pos + 0.05), Inches(8.7), Inches(0.45),
                    item, font_size=11, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    add_colored_shape(slide, Inches(0.2), Inches(3.95), Inches(9.6), Inches(0.85), COLORS['magma_orange'])
    add_text_box(slide, Inches(0.4), Inches(3.98), Inches(9.2), Inches(0.8),
                "Remember: Your notecard is the ONLY physical resource allowed!\nMake it count - organize it clearly!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_summary_slide(prs):
    """Slide 16: Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Week 3 Summary: What You Learned",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    takeaways = [
        ("Viscosity Matters", "Silica content determines magma thickness"),
        ("Eruption Prediction", "Composition tells us if eruption will be explosive"),
        ("Hazards Differ", "Explosive = pyroclastic; Effusive = lava flows"),
        ("Monitoring Saves Lives", "Detecting precursors enables evacuation"),
    ]

    y_pos = 0.9
    for title, detail in takeaways:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.65), COLORS['light_orange_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.65), COLORS['lava_red'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.8), Inches(0.55),
                    title, font_size=13, bold=True, color=COLORS['lava_dark'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.3), Inches(y_pos + 0.05), Inches(6.3), Inches(0.55),
                    detail, font_size=12, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.72

    add_colored_shape(slide, Inches(0.2), Inches(3.85), Inches(9.6), Inches(0.65), COLORS['silica_blue'])
    add_text_box(slide, Inches(0.4), Inches(3.88), Inches(9.2), Inches(0.6),
                "NEXT WEEK: Earth's Interior - How do we know what's inside without going there?",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(4.65), Inches(9.6), Inches(0.55), COLORS['magma_orange'])
    add_text_box(slide, Inches(0.4), Inches(4.68), Inches(9.2), Inches(0.5),
                "ANSWER: High-silica magma traps gas → Pressure builds → EXPLOSION!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


if __name__ == "__main__":
    prs = create_presentation()
    output_path = "/home/user/Kairos.Sci.Repo/content/grade7/cycle06/week3/G7_C6_W3_Slides_Final.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Created: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
