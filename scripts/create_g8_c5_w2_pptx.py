#!/usr/bin/env python3
"""
Create G8_C5_W2 Waves & Material Interactions presentation.
Topic: Waves & Information Transfer.

See PPTX_DESIGN_GUIDE.md for best practices documentation.
"""

import os
from pptx.dml.color import RGBColor
from pptx_common import (
    COLORS as BASE_COLORS,
    create_base_presentation,
    add_colored_shape,
    add_text_box,
    Inches,
    Pt,
    PP_ALIGN,
    MSO_ANCHOR,
)

# Waves theme - extend base colors
COLORS = {**BASE_COLORS}
COLORS.update({
    'wave_purple': RGBColor(0x7C, 0x3A, 0xED),
    'wave_purple_dark': RGBColor(0x5B, 0x21, 0xB6),
    'wave_cyan': RGBColor(0x06, 0xB6, 0xD4),
    'wave_cyan_dark': RGBColor(0x08, 0x91, 0xB2),
    'light_wave_bg': RGBColor(0xF5, 0xF3, 0xFF),
    'digital_green': RGBColor(0x10, 0xB9, 0x81),
})


def create_presentation():
    prs = create_base_presentation()
    add_title_slide(prs)
    add_phenomenon_slide(prs)
    add_driving_question_slide(prs)
    add_prior_knowledge_slide(prs)
    add_learning_targets_slide(prs)
    add_vocabulary_slide(prs)
    add_hook_intro_slide(prs)
    add_hook_support_slide(prs)
    add_station1_intro_slide(prs)
    add_station1_support_slide(prs)
    add_station2_intro_slide(prs)
    add_station2_support_slide(prs)
    add_station3_intro_slide(prs)
    add_station3_support_slide(prs)
    add_exit_ticket_slide(prs)
    add_summary_slide(prs)
    return prs


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['wave_purple_dark'])
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(1.2),
                "Waves & Information Transfer", font_size=44, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.5), Inches(2.8), Inches(9), Inches(0.5),
                "Week 2: Waves & Material Interactions", font_size=24, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(3.4), Inches(9), Inches(0.5),
                "Grade 8 Science | Cycle 5 | 100 Points", font_size=18, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_colored_shape(slide, Inches(2), Inches(4.3), Inches(6), Inches(0.8), COLORS['white'])
    add_text_box(slide, Inches(2.2), Inches(4.35), Inches(5.6), Inches(0.7),
                "Why do some materials block cell signals while others don't?",
                font_size=16, bold=True, color=COLORS['wave_purple_dark'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_phenomenon_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['wave_purple_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "The Signal Blocker Mystery", font_size=26, bold=True, color=COLORS['white'],
                font_name="Georgia", anchor=MSO_ANCHOR.MIDDLE)
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(9.6), Inches(2.4), COLORS['light_wave_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(0.08), Inches(2.4), COLORS['wave_purple'])
    phenomenon_text = """You're in an elevator and your phone loses signal. Step outside—full bars again.
Your friend's house has great WiFi everywhere except one room with thick concrete walls.

Some buildings are designed as "Faraday cages" that block ALL radio signals.
Yet visible light passes through glass windows in those same buildings just fine!

What makes certain materials block specific types of waves while letting others through?"""
    add_text_box(slide, Inches(0.5), Inches(1.0), Inches(9.1), Inches(2.2),
                phenomenon_text, font_size=14, color=COLORS['dark_text'])
    # Comparison boxes
    add_colored_shape(slide, Inches(0.2), Inches(3.5), Inches(4.7), Inches(1.5), COLORS['wave_cyan'])
    add_text_box(slide, Inches(0.4), Inches(3.55), Inches(4.3), Inches(0.35),
                "BLOCKS RADIO WAVES", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(0.4), Inches(3.95), Inches(4.3), Inches(1.0),
                "• Metal (reflects)\n• Concrete (absorbs)\n• Water (absorbs)\n• Faraday cages",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_colored_shape(slide, Inches(5.1), Inches(3.5), Inches(4.7), Inches(1.5), COLORS['digital_green'])
    add_text_box(slide, Inches(5.3), Inches(3.55), Inches(4.3), Inches(0.35),
                "ALLOWS RADIO WAVES", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(5.3), Inches(3.95), Inches(4.3), Inches(1.0),
                "• Wood (transmits)\n• Glass (transmits)\n• Drywall (transmits)\n• Plastic",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_driving_question_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Driving Question", font_size=26, bold=True, color=COLORS['white'],
                font_name="Georgia", anchor=MSO_ANCHOR.MIDDLE)
    add_colored_shape(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(1.2), COLORS['wave_purple_dark'])
    add_text_box(slide, Inches(0.7), Inches(1.25), Inches(8.6), Inches(1.1),
                "How do waves carry information, and why do materials affect transmission?",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_colored_shape(slide, Inches(0.2), Inches(2.7), Inches(9.6), Inches(2.0), COLORS['light_wave_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.8), Inches(9.2), Inches(0.3),
                "Questions we'll investigate:", font_size=14, bold=True, color=COLORS['wave_purple_dark'])
    questions = """• How do different materials affect wave transmission, absorption, and reflection?
• How are waves used to encode and transmit information (digital signals)?
• Why do some frequencies penetrate materials better than others?
• How can we design communication systems that work in challenging environments?"""
    add_text_box(slide, Inches(0.4), Inches(3.15), Inches(9.2), Inches(1.4),
                questions, font_size=14, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(0.2), Inches(4.85), Inches(9.6), Inches(0.55), COLORS['green_start'])
    add_text_box(slide, Inches(0.4), Inches(4.88), Inches(9.2), Inches(0.5),
                "Week 1 Connection: How do reflection, absorption, and transmission relate to materials?",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_prior_knowledge_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "What You Already Know", font_size=26, bold=True, color=COLORS['white'],
                font_name="Georgia", anchor=MSO_ANCHOR.MIDDLE)
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(4.7), Inches(1.8), COLORS['light_wave_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(0.08), Inches(1.8), COLORS['wave_purple'])
    add_text_box(slide, Inches(0.4), Inches(1.0), Inches(4.4), Inches(0.3),
                "From Week 1 (Wave Properties):", font_size=13, bold=True, color=COLORS['wave_purple'])
    w1_review = """• Wavelength determines wave behavior
• Waves reflect, refract, diffract
• EM waves don't need a medium
• Long λ diffracts more than short λ"""
    add_text_box(slide, Inches(0.4), Inches(1.35), Inches(4.4), Inches(1.3),
                w1_review, font_size=12, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(5.1), Inches(0.9), Inches(4.7), Inches(1.8), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(5.1), Inches(0.9), Inches(0.08), Inches(1.8), COLORS['green_end'])
    add_text_box(slide, Inches(5.3), Inches(1.0), Inches(4.4), Inches(0.3),
                "From Cycle 3 (Natural Selection):", font_size=13, bold=True, color=COLORS['green_end'])
    c3_review = """• Information passes between generations
• DNA encodes genetic information
• Selection acts on variation
• Information transfer enables evolution"""
    add_text_box(slide, Inches(5.3), Inches(1.35), Inches(4.4), Inches(1.3),
                c3_review, font_size=12, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(0.2), Inches(2.9), Inches(9.6), Inches(1.7), COLORS['light_cyan_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.0), Inches(9.2), Inches(0.3),
                "NEW in Week 2:", font_size=14, bold=True, color=COLORS['wave_cyan_dark'])
    new_learning = """Today we explore how waves CARRY INFORMATION! Your phone, TV, radio—all use waves to transmit
data. We'll learn how digital signals encode information and why material properties matter for transmission."""
    add_text_box(slide, Inches(0.4), Inches(3.35), Inches(9.2), Inches(1.2),
                new_learning, font_size=13, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(0.2), Inches(4.75), Inches(9.6), Inches(0.6), COLORS['wave_cyan'])
    add_text_box(slide, Inches(0.4), Inches(4.78), Inches(9.2), Inches(0.55),
                "KEY: Waves carry INFORMATION by encoding data in their properties!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_learning_targets_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Learning Targets", font_size=26, bold=True, color=COLORS['white'],
                font_name="Georgia", anchor=MSO_ANCHOR.MIDDLE)
    targets = [
        ("1", "Quantify how different materials affect wave transmission"),
        ("2", "Explain how digital signals encode information using waves"),
        ("3", "Compare how frequency and wavelength affect material penetration"),
        ("4", "Design a communication system for a challenging environment"),
    ]
    y_pos = 0.9
    for num, target in targets:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.7), COLORS['light_green_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.5), Inches(0.7), COLORS['green_end'])
        add_text_box(slide, Inches(0.25), Inches(y_pos + 0.05), Inches(0.4), Inches(0.6),
                    num, font_size=20, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(0.85), Inches(y_pos + 0.08), Inches(8.8), Inches(0.55),
                    target, font_size=14, color=COLORS['dark_text'], anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.8
    add_colored_shape(slide, Inches(0.3), Inches(4.2), Inches(9.4), Inches(0.7), COLORS['wave_purple_dark'])
    add_text_box(slide, Inches(1.5), Inches(4.25), Inches(1.5), Inches(0.35),
                "100", font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(1.5), Inches(4.55), Inches(1.5), Inches(0.3),
                "Total Points", font_size=9, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(4.25), Inches(4.25), Inches(1.5), Inches(0.35),
                "~75", font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(4.25), Inches(4.55), Inches(1.5), Inches(0.3),
                "Minutes", font_size=9, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(7.0), Inches(4.25), Inches(1.5), Inches(0.35),
                "5", font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(7.0), Inches(4.55), Inches(1.5), Inches(0.3),
                "Sections", font_size=9, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_vocabulary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.2), Inches(0.15), Inches(9.6), Inches(0.5), COLORS['wave_purple'])
    add_text_box(slide, Inches(0.4), Inches(0.18), Inches(9.2), Inches(0.45),
                "Key Vocabulary", font_size=22, bold=True, color=COLORS['white'],
                font_name="Georgia", anchor=MSO_ANCHOR.MIDDLE)
    vocab = [
        ("Transmission", "Wave passes through material (window, radio through walls)"),
        ("Absorption", "Wave energy converts to heat in material"),
        ("Digital Signal", "Information encoded as discrete values (0s and 1s)"),
        ("Analog Signal", "Information encoded as continuous wave variations"),
        ("Modulation", "Changing wave properties to encode information"),
        ("Bandwidth", "Range of frequencies a system can transmit"),
    ]
    y_pos = 0.75
    for i, (term, definition) in enumerate(vocab):
        bg_color = COLORS['light_wave_bg'] if i % 2 == 0 else COLORS['white']
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.6), bg_color)
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.2), Inches(0.5),
                    term, font_size=13, bold=True, color=COLORS['dark_text'], anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.7), Inches(y_pos + 0.05), Inches(6.9), Inches(0.5),
                    definition, font_size=12, color=COLORS['gray_text'], anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6
    add_colored_shape(slide, Inches(0.2), Inches(4.4), Inches(9.6), Inches(0.55), COLORS['purple_start'])
    add_text_box(slide, Inches(0.4), Inches(4.43), Inches(9.2), Inches(0.5),
                "Notecard: Digital = 0s and 1s | Analog = continuous waves",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_colored_shape(slide, Inches(0.2), Inches(5.0), Inches(9.6), Inches(0.45), COLORS['light_wave_bg'])
    add_text_box(slide, Inches(0.4), Inches(5.03), Inches(9.2), Inches(0.4),
                "Binary: 1 = wave ON | 0 = wave OFF (simplified digital encoding)",
                font_size=11, color=COLORS['wave_purple_dark'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_hook_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['purple_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Hook – The Signal Blocker Mystery", font_size=26, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "12 Points | ~10 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "TEST MATERIALS", font_size=14, bold=True, color=COLORS['wave_purple'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.3), COLORS['light_wave_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.3), COLORS['wave_purple'])
    test_text = """Try blocking your phone signal:

1. Wrap phone in paper → Signal?
2. Wrap phone in aluminum foil → Signal?
3. Put phone in metal box → Signal?
4. Put phone behind concrete → Signal?

Record what blocks and what doesn't!"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(2.1),
                test_text, font_size=12, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(1.5), COLORS['wave_cyan'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "YOUR PREDICTION", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.2), Inches(1.55), Inches(4.3), Inches(1.0),
                "Before testing: WHY do you think metal blocks signals but paper doesn't?\n\nWrite your prediction in the Hook Form.",
                font_size=13, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_colored_shape(slide, Inches(5.0), Inches(2.8), Inches(4.7), Inches(1.0), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(5.2), Inches(2.9), Inches(4.3), Inches(0.25),
                "Quick Check (0 pts):", font_size=11, bold=True, color=COLORS['orange_end'])
    add_text_box(slide, Inches(5.2), Inches(3.2), Inches(4.3), Inches(0.55),
                "What happens when a wave hits a material?\n(Hint: 3 possible outcomes from W1)",
                font_size=11, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(0.15), Inches(4.0), Inches(9.7), Inches(0.55), COLORS['wave_purple_dark'])
    add_text_box(slide, Inches(0.35), Inches(4.03), Inches(9.3), Inches(0.5),
                "Notecard: Write your prediction BEFORE testing!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_colored_shape(slide, Inches(0.15), Inches(4.65), Inches(9.7), Inches(0.55), COLORS['green_start'])
    add_text_box(slide, Inches(0.35), Inches(4.68), Inches(9.3), Inches(0.5),
                "Week 1 Link: Which wave behavior (R/A/T) explains why metal blocks signals?",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_hook_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['wave_purple_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Hook – Why Materials Matter", font_size=22, bold=True, color=COLORS['white'],
                font_name="Georgia", anchor=MSO_ANCHOR.MIDDLE)
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(9.6), Inches(1.3), COLORS['light_wave_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(0.08), Inches(1.3), COLORS['wave_purple'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(9.2), Inches(0.3),
                "THE ANSWER: It depends on the material's electrical properties!", font_size=14, bold=True, color=COLORS['wave_purple'])
    answer_text = """• CONDUCTORS (metal): Free electrons oscillate with wave → REFLECT almost all radio waves
• INSULATORS (plastic, wood): No free electrons → waves pass through (TRANSMIT)
• ABSORBERS (water, concrete): Convert wave energy to heat → ABSORB
• The key is how electrons in the material respond to the electromagnetic wave!"""
    add_text_box(slide, Inches(0.4), Inches(1.3), Inches(9.2), Inches(0.8),
                answer_text, font_size=12, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(0.2), Inches(2.3), Inches(4.7), Inches(1.4), COLORS['light_cyan_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.4), Inches(4.4), Inches(0.25),
                "FARADAY CAGE", font_size=12, bold=True, color=COLORS['wave_cyan_dark'])
    faraday_text = """Metal mesh or solid metal enclosure:
• Electrons in metal move to cancel
  the incoming EM field
• Creates "electromagnetic shield"
• Used in: MRI rooms, secure facilities,
  your microwave oven door!"""
    add_text_box(slide, Inches(0.4), Inches(2.7), Inches(4.4), Inches(0.95),
                faraday_text, font_size=11, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(5.1), Inches(2.3), Inches(4.6), Inches(1.4), COLORS['light_green_bg'])
    add_text_box(slide, Inches(5.3), Inches(2.4), Inches(4.2), Inches(0.25),
                "WHY FREQUENCY MATTERS", font_size=12, bold=True, color=COLORS['green_end'])
    freq_text = """Different frequencies penetrate differently:
• Low freq (radio): Through walls easily
• High freq (light): Blocked by walls
• X-rays: Through soft tissue

Material thickness also matters—
thicker = more absorption!"""
    add_text_box(slide, Inches(5.3), Inches(2.7), Inches(4.2), Inches(0.95),
                freq_text, font_size=11, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(0.15), Inches(3.85), Inches(9.7), Inches(0.55), COLORS['purple_start'])
    add_text_box(slide, Inches(0.35), Inches(3.88), Inches(9.3), Inches(0.5),
                "Complete the Hook Form (12 pts) - Compare your prediction to the explanation!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_station1_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['orange_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 1 – Transmission-Absorption Lab", font_size=26, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | ~18 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "YOUR INVESTIGATION", font_size=14, bold=True, color=COLORS['orange_end'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.3), COLORS['light_orange_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.3), COLORS['orange_end'])
    steps = """1. Use light sensor to measure transmission (3 min)
2. Test materials: glass, plastic, paper, foil (5 min)
3. Calculate % transmitted for each (4 min)
4. Graph transmission vs. material type (4 min)
5. Answer analysis questions (2 min)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(2.1),
                steps, font_size=12, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(1.4), COLORS['orange_end'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "CALCULATION", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    calc_text = """% Transmitted = (Light through / Light without) × 100

Example:
• No material: 1000 lux
• Through glass: 900 lux
• % = (900/1000) × 100 = 90%"""
    add_text_box(slide, Inches(5.2), Inches(1.55), Inches(4.3), Inches(0.9),
                calc_text, font_size=11, color=COLORS['white'])
    add_colored_shape(slide, Inches(5.0), Inches(2.65), Inches(4.7), Inches(1.15), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.2), Inches(2.75), Inches(4.3), Inches(0.25),
                "KEY INSIGHT:", font_size=11, bold=True, color=COLORS['blue_accent'])
    add_text_box(slide, Inches(5.2), Inches(3.0), Inches(4.3), Inches(0.75),
                "What's NOT transmitted is either:\n• Reflected (bounces back)\n• Absorbed (becomes heat)",
                font_size=10, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(0.15), Inches(4.0), Inches(9.7), Inches(0.55), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(4.03), Inches(9.3), Inches(0.5),
                "Complete Station 1 in Form - Calculate transmission percentages!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_station1_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 1 – Understanding Material Properties", font_size=22, bold=True, color=COLORS['white'],
                font_name="Georgia", anchor=MSO_ANCHOR.MIDDLE)
    data = [
        ("Clear glass", "~90%", "~8%", "~2%"),
        ("Frosted glass", "~70%", "~20%", "~10%"),
        ("Paper", "~30%", "~20%", "~50%"),
        ("Aluminum foil", "~0%", "~95%", "~5%"),
    ]
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(9.6), Inches(0.4), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.88), Inches(2.2), Inches(0.35),
                "Material", font_size=11, bold=True, color=COLORS['dark_text'], anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(2.7), Inches(0.88), Inches(2.2), Inches(0.35),
                "Transmitted", font_size=11, bold=True, color=COLORS['green_accent'], anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(5.0), Inches(0.88), Inches(2.2), Inches(0.35),
                "Reflected", font_size=11, bold=True, color=COLORS['blue_accent'], anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(7.3), Inches(0.88), Inches(2.4), Inches(0.35),
                "Absorbed", font_size=11, bold=True, color=COLORS['orange_end'], anchor=MSO_ANCHOR.MIDDLE)
    y_pos = 1.3
    for i, (mat, trans, refl, abso) in enumerate(data):
        bg_color = COLORS['white'] if i % 2 == 0 else COLORS['light_blue_bg']
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.45), bg_color)
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.02), Inches(2.2), Inches(0.4),
                    mat, font_size=11, color=COLORS['dark_text'], anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.7), Inches(y_pos + 0.02), Inches(2.2), Inches(0.4),
                    trans, font_size=11, color=COLORS['green_accent'], anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(5.0), Inches(y_pos + 0.02), Inches(2.2), Inches(0.4),
                    refl, font_size=11, color=COLORS['blue_accent'], anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(7.3), Inches(y_pos + 0.02), Inches(2.4), Inches(0.4),
                    abso, font_size=11, color=COLORS['orange_end'], anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.45
    add_colored_shape(slide, Inches(0.2), Inches(3.2), Inches(9.6), Inches(0.9), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.3), Inches(9.2), Inches(0.25),
                "SENTENCE STARTERS:", font_size=11, bold=True, color=COLORS['green_accent'])
    starters = """• "Glass transmits most light because..."
• "Aluminum foil reflects light because..."
• "The % transmitted + % reflected + % absorbed = 100% because..."""
    add_text_box(slide, Inches(0.4), Inches(3.6), Inches(9.2), Inches(0.45),
                starters, font_size=10, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(0.2), Inches(4.25), Inches(9.6), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.28), Inches(9.2), Inches(0.5),
                "ENERGY CONSERVATION: Transmitted + Reflected + Absorbed = 100% (energy is conserved!)",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_station2_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['wave_cyan'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 2 – Information Encoding Investigation", font_size=26, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | ~15 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "YOUR INVESTIGATION", font_size=14, bold=True, color=COLORS['wave_cyan_dark'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.3), COLORS['light_cyan_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.3), COLORS['wave_cyan_dark'])
    steps = """1. Explore binary encoding demo (3 min)
2. Convert message to binary code (4 min)
3. Decode a partner's binary message (3 min)
4. Compare digital vs analog signals (3 min)
5. Answer analysis questions (2 min)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(2.1),
                steps, font_size=12, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(2.7), COLORS['digital_green'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "BINARY CODE", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    binary_text = """Everything digital uses 0s and 1s!

A = 01000001
B = 01000010
C = 01000011
...

Each letter = 8 bits (1 byte)
Your name in binary = ?

Waves encode this as:
1 = signal ON
0 = signal OFF"""
    add_text_box(slide, Inches(5.2), Inches(1.55), Inches(4.3), Inches(2.2),
                binary_text, font_size=10, color=COLORS['white'])
    add_colored_shape(slide, Inches(0.15), Inches(4.0), Inches(9.7), Inches(0.55), COLORS['wave_cyan_dark'])
    add_text_box(slide, Inches(0.35), Inches(4.03), Inches(9.3), Inches(0.5),
                "Complete Station 2 in Form - Encode and decode messages!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_station2_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['wave_cyan_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 2 – Digital vs Analog", font_size=22, bold=True, color=COLORS['white'],
                font_name="Georgia", anchor=MSO_ANCHOR.MIDDLE)
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(4.7), Inches(1.6), COLORS['light_wave_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(4.4), Inches(0.3),
                "DIGITAL SIGNALS", font_size=12, bold=True, color=COLORS['wave_purple'])
    digital_text = """• Only two values: 0 or 1
• Easy to copy perfectly
• Resistant to noise/interference
• Used in: computers, phones, internet

Example: CD audio, streaming video,
text messages, WiFi"""
    add_text_box(slide, Inches(0.4), Inches(1.3), Inches(4.4), Inches(1.1),
                digital_text, font_size=10, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(5.1), Inches(0.85), Inches(4.6), Inches(1.6), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(5.3), Inches(0.95), Inches(4.2), Inches(0.3),
                "ANALOG SIGNALS", font_size=12, bold=True, color=COLORS['orange_end'])
    analog_text = """• Continuous range of values
• Copies degrade over time
• Susceptible to noise
• Used in: vinyl records, AM radio

Example: Vinyl records, old phones,
radio broadcasts"""
    add_text_box(slide, Inches(5.3), Inches(1.3), Inches(4.2), Inches(1.1),
                analog_text, font_size=10, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(0.2), Inches(2.6), Inches(9.6), Inches(0.9), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.7), Inches(9.2), Inches(0.25),
                "WHY DIGITAL WON:", font_size=11, bold=True, color=COLORS['green_accent'])
    why_digital = """• Copy a CD 1000 times → still perfect (digital)
• Copy a tape 1000 times → unrecognizable noise (analog)
• 0 or 1 is easy to detect even with interference; continuous values get corrupted"""
    add_text_box(slide, Inches(0.4), Inches(3.0), Inches(9.2), Inches(0.45),
                why_digital, font_size=10, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(0.2), Inches(3.65), Inches(9.6), Inches(0.7), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.75), Inches(9.2), Inches(0.25),
                "CONNECTION TO CYCLE 3:", font_size=11, bold=True, color=COLORS['purple_accent'])
    add_text_box(slide, Inches(0.4), Inches(4.0), Inches(9.2), Inches(0.3),
                "DNA uses 4 \"digits\" (A, T, G, C) to encode information—nature's digital code! Mutations = copying errors.",
                font_size=10, color=COLORS['dark_text'])


def add_station3_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['green_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 3 – Design a Communication System", font_size=26, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "25 Points | ~20 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "ENGINEERING CHALLENGE", font_size=14, bold=True, color=COLORS['green_end'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.3), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.3), COLORS['green_end'])
    challenge = """Design communication for a mine rescue!

Environment challenges:
• 500m underground
• Rock and metal everywhere
• Water in some tunnels
• Must work even if cables break

Choose your wave type and justify!"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(2.1),
                challenge, font_size=12, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(2.7), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "WAVE OPTIONS", font_size=13, bold=True, color=COLORS['blue_accent'])
    options = """Radio (VLF) | Very low frequency
  Can penetrate rock (~100m)
  Slow data transfer

Radio (UHF) | Ultra high frequency
  Fast data, poor penetration
  Needs repeaters

Sound waves | Through rock
  Travels far in solids
  Very slow data rate

Fiber optic | Light in cable
  Fast, immune to EM interference
  Needs physical connection"""
    add_text_box(slide, Inches(5.2), Inches(1.55), Inches(4.3), Inches(2.2),
                options, font_size=9, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(0.15), Inches(4.0), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(4.03), Inches(9.3), Inches(0.5),
                "Complete Station 3 in Form - Design your rescue communication system!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_station3_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 3 – Real Mine Communication Systems", font_size=22, bold=True, color=COLORS['white'],
                font_name="Georgia", anchor=MSO_ANCHOR.MIDDLE)
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(4.7), Inches(1.5), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(4.4), Inches(0.25),
                "REAL-WORLD SOLUTION:", font_size=11, bold=True, color=COLORS['green_accent'])
    real_text = """Mines use MULTIPLE systems:
1. Fiber optic backbone (main tunnels)
2. VLF radio (through-rock emergency)
3. Leaky feeder cable (UHF in tunnels)
4. Text pagers (backup)

No single system works everywhere—
redundancy saves lives!"""
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(4.4), Inches(1.05),
                real_text, font_size=10, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(5.1), Inches(0.85), Inches(4.6), Inches(1.5), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(5.3), Inches(0.95), Inches(4.2), Inches(0.25),
                "DESIGN CONSIDERATIONS:", font_size=11, bold=True, color=COLORS['orange_end'])
    design_text = """• What frequency penetrates rock best?
• How do you handle water absorption?
• What's the backup if main fails?
• Speed vs. reliability trade-off?

Apply W1 knowledge:
Low freq = longer λ = better diffraction
around obstacles!"""
    add_text_box(slide, Inches(5.3), Inches(1.25), Inches(4.2), Inches(1.05),
                design_text, font_size=10, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(0.2), Inches(2.5), Inches(9.6), Inches(1.0), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.6), Inches(9.2), Inches(0.25),
                "SEP-3: PLANNING INVESTIGATIONS", font_size=11, bold=True, color=COLORS['purple_end'])
    sep_text = """Your design should address: Wave type selection | Material penetration | Redundancy/backup |
Data rate requirements | Power constraints. Justify each choice with wave physics!"""
    add_text_box(slide, Inches(0.4), Inches(2.9), Inches(9.2), Inches(0.5),
                sep_text, font_size=10, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(0.15), Inches(3.65), Inches(9.7), Inches(0.55), COLORS['wave_cyan'])
    add_text_box(slide, Inches(0.35), Inches(3.68), Inches(9.3), Inches(0.5),
                "KEY: Match wave properties to environment challenges!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_exit_ticket_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.9), COLORS['exit_purple_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Exit Ticket – Material Interaction Integration", font_size=26, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.68), Inches(9.3), Inches(0.35),
                "23 Points | ~15 min", font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(1.5), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(0.08), Inches(1.5), COLORS['exit_purple_start'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.3),
                "QUESTION TYPES:", font_size=14, bold=True, color=COLORS['exit_purple_start'])
    q_types = """• 2 NEW – Material properties, digital encoding, transmission %
• 2 SPIRAL – Cycles 3-4: Information transfer (DNA), energy conservation
• 1 INTEGRATION – Connect wave properties to communication design
• 1 SEP-3 – Design justification (explain system choices)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(9.0), Inches(1.0),
                q_types, font_size=13, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(0.3), Inches(2.85), Inches(4.6), Inches(1.1), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(2.9), Inches(4.3), Inches(0.3),
                "SUCCESS TIPS:", font_size=12, bold=True, color=COLORS['green_accent'], anchor=MSO_ANCHOR.MIDDLE)
    tips = """• Use: transmission, absorption, reflection
• Remember: T + R + A = 100%
• Digital: 0s and 1s, perfect copies"""
    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(4.3), Inches(0.7),
                tips, font_size=10, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(5.1), Inches(2.85), Inches(4.6), Inches(1.1), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.3), Inches(2.9), Inches(4.2), Inches(0.3),
                "SPIRAL REMINDER:", font_size=12, bold=True, color=COLORS['blue_accent'], anchor=MSO_ANCHOR.MIDDLE)
    spiral = """• C3: DNA = nature's digital code (ATGC)
• C4: Energy is conserved in transfers
• C5: Waves encode + transfer information!"""
    add_text_box(slide, Inches(5.3), Inches(3.2), Inches(4.2), Inches(0.7),
                spiral, font_size=10, color=COLORS['dark_text'])
    add_colored_shape(slide, Inches(0.15), Inches(4.1), Inches(9.7), Inches(0.65), COLORS['exit_purple_end'])
    add_text_box(slide, Inches(0.35), Inches(4.15), Inches(9.3), Inches(0.55),
                "FINAL Notecard: Digital signals use waves to carry 0s and 1s!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['wave_purple_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Week 2 Summary: What You Learned", font_size=24, bold=True, color=COLORS['white'],
                font_name="Georgia", anchor=MSO_ANCHOR.MIDDLE)
    takeaways = [
        ("Material Effects", "Conductors reflect, insulators transmit, absorbers convert to heat"),
        ("Energy Balance", "Transmitted + Reflected + Absorbed = 100% (always!)"),
        ("Digital Signals", "0s and 1s encoded in waves—perfect copies, noise resistant"),
        ("Design Matching", "Choose wave frequency based on material/environment"),
    ]
    y_pos = 0.9
    for title, detail in takeaways:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.65), COLORS['light_wave_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.65), COLORS['wave_purple'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.4), Inches(0.55),
                    title, font_size=13, bold=True, color=COLORS['wave_purple_dark'], anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.9), Inches(y_pos + 0.05), Inches(6.7), Inches(0.55),
                    detail, font_size=12, color=COLORS['dark_text'], anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.72
    add_colored_shape(slide, Inches(0.2), Inches(3.85), Inches(9.6), Inches(0.65), COLORS['wave_cyan'])
    add_text_box(slide, Inches(0.4), Inches(3.88), Inches(9.2), Inches(0.6),
                "NEXT WEEK: Synthesis & Assessment – Show what you know!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_colored_shape(slide, Inches(0.2), Inches(4.65), Inches(9.6), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(4.68), Inches(9.2), Inches(0.5),
                "REMEMBER: Information travels on waves—digital encoding makes it reliable!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


if __name__ == "__main__":
    prs = create_presentation()
    output_path = "/home/user/Kairos.Sci.Repo/content/grade8/cycle05/week2/G8_C5_W2_Material_Interactions_Slides_Final.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Created: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
