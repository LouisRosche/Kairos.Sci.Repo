#!/usr/bin/env python3
"""
Shared utilities for PPTX generation scripts.

This module extracts common functionality from the individual create_g*_c*_w*_pptx.py
scripts to reduce code duplication (~400 KB → ~80 KB savings across 11 scripts).

Usage:
    from pptx_common import COLORS, add_colored_shape, add_text_box, create_base_presentation

See PPTX_DESIGN_GUIDE.md for best practices documentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# =============================================================================
# STANDARD COLOR PALETTE
# =============================================================================
# These colors are matched to the student-page.html design system.
# Override individual colors in your script if needed for grade-specific theming.

COLORS = {
    # Main gradients
    'header_blue_start': RGBColor(0x42, 0x99, 0xE1),  # #4299E1
    'header_blue_end': RGBColor(0x2B, 0x6C, 0xB0),    # #2B6CB0
    'green_start': RGBColor(0x48, 0xBB, 0x78),        # #48BB78
    'green_end': RGBColor(0x27, 0x67, 0x49),          # #276749
    'purple_start': RGBColor(0x66, 0x7E, 0xEA),       # #667EEA (Hook)
    'purple_end': RGBColor(0x76, 0x4B, 0xA2),         # #764BA2
    'orange_start': RGBColor(0xF6, 0xAD, 0x55),       # #F6AD55 (Station 1)
    'orange_end': RGBColor(0xDD, 0x6B, 0x20),         # #DD6B20
    'exit_purple_start': RGBColor(0x9F, 0x7A, 0xEA),  # #9F7AEA
    'exit_purple_end': RGBColor(0x6B, 0x46, 0xC1),    # #6B46C1

    # Accents
    'teal': RGBColor(0x38, 0xB2, 0xAC),               # #38B2AC
    'teal_dark': RGBColor(0x23, 0x4E, 0x52),          # #234E52
    'teal_light': RGBColor(0xE6, 0xFF, 0xFA),         # #E6FFFA

    # Text colors
    'dark_text': RGBColor(0x2D, 0x37, 0x48),          # #2D3748
    'gray_text': RGBColor(0x4A, 0x55, 0x68),          # #4A5568
    'white': RGBColor(0xFF, 0xFF, 0xFF),

    # Background colors
    'light_blue_bg': RGBColor(0xEB, 0xF8, 0xFF),      # #EBF8FF
    'light_green_bg': RGBColor(0xF0, 0xFF, 0xF4),     # #F0FFF4
    'light_orange_bg': RGBColor(0xFF, 0xFA, 0xF0),    # #FFFAF0
    'light_purple_bg': RGBColor(0xFA, 0xF5, 0xFF),    # #FAF5FF
    'light_pink_bg': RGBColor(0xFF, 0xF5, 0xF7),      # #FFF5F7
    'light_teal_bg': RGBColor(0xE6, 0xFF, 0xFA),      # #E6FFFA

    # Specific accents
    'red_accent': RGBColor(0xC5, 0x30, 0x30),         # #C53030
    'blue_accent': RGBColor(0x31, 0x82, 0xCE),        # #3182CE
    'green_accent': RGBColor(0x38, 0xA1, 0x69),       # #38A169
    'purple_accent': RGBColor(0x80, 0x5A, 0xD5),      # #805AD5
    'orange_accent': RGBColor(0xC0, 0x56, 0x21),      # #C05621
}


# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def create_base_presentation():
    """Create a new 16:9 presentation with standard dimensions."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 aspect ratio
    return prs


def add_colored_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """
    Add a colored shape to the slide.

    Args:
        slide: The slide to add the shape to
        left, top, width, height: Position and dimensions (use Inches())
        color: RGBColor from COLORS dict
        shape_type: MSO_SHAPE enum (default: ROUNDED_RECTANGLE)

    Returns:
        The created shape object
    """
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=None, align=PP_ALIGN.LEFT, font_name="Arial", anchor=None):
    """
    Add a text box with specified formatting.

    IMPORTANT: Text boxes don't visually show text boundaries.
    Always add padding when positioning text inside shapes.

    Args:
        slide: The slide to add the text box to
        left, top, width, height: Position and dimensions (use Inches())
        text: The text content
        font_size: Font size in points (default: 18)
        bold: Whether text is bold (default: False)
        color: RGBColor for text (default: None = black)
        align: PP_ALIGN enum for horizontal alignment (default: LEFT)
        font_name: Font family name (default: "Arial")
        anchor: MSO_ANCHOR enum for vertical alignment (default: None = TOP)
                Use MSO_ANCHOR.MIDDLE for vertical centering

    Returns:
        The created text box object
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Set vertical anchor if specified (critical for centering in small boxes)
    if anchor:
        tf.anchor = anchor

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align

    if color:
        p.font.color.rgb = color

    return txBox


def add_title_bar(slide, text, color=None):
    """
    Add a standard title bar at the top of a slide.

    Args:
        slide: The slide to add the title bar to
        text: The title text
        color: Background color (default: teal_dark)

    Returns:
        Tuple of (bar_shape, text_box)
    """
    bar_color = color or COLORS['teal_dark']
    bar = add_colored_shape(slide, Inches(0.2), Inches(0.2), Inches(9.6), Inches(0.6), bar_color)
    title = add_text_box(slide, Inches(0.4), Inches(0.28), Inches(9.2), Inches(0.5),
                         text, font_size=24, bold=True, color=COLORS['white'], font_name="Georgia")
    return bar, title


def add_notecard_bar(slide, text="📝 Write your learning on your notecard!", color=None):
    """
    Add a standard notecard reminder bar at the bottom of a slide.

    Args:
        slide: The slide to add the bar to
        text: The reminder text
        color: Background color (default: purple_start)

    Returns:
        Tuple of (bar_shape, text_box)
    """
    bar_color = color or COLORS['purple_start']
    bar = add_colored_shape(slide, Inches(0.2), Inches(4.8), Inches(9.6), Inches(0.6), bar_color)
    note = add_text_box(slide, Inches(0.4), Inches(4.88), Inches(9.2), Inches(0.4),
                        text, font_size=14, color=COLORS['white'], align=PP_ALIGN.CENTER,
                        anchor=MSO_ANCHOR.MIDDLE)
    return bar, note


# =============================================================================
# STANDARD SLIDE DIMENSIONS
# =============================================================================

SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)
MARGIN = Inches(0.2)
CONTENT_WIDTH = Inches(9.6)

# Re-export commonly used items for convenience
__all__ = [
    'COLORS',
    'create_base_presentation',
    'add_colored_shape',
    'add_text_box',
    'add_title_bar',
    'add_notecard_bar',
    'SLIDE_WIDTH',
    'SLIDE_HEIGHT',
    'MARGIN',
    'CONTENT_WIDTH',
    # Re-exports from pptx
    'Presentation',
    'Inches',
    'Pt',
    'RGBColor',
    'PP_ALIGN',
    'MSO_ANCHOR',
    'MSO_SHAPE',
]
