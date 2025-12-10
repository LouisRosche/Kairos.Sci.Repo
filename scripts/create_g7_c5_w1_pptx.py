#!/usr/bin/env python3
"""
Create G7_C5_W1 Air Masses & Storm Formation presentation
Topic: Weather & Climate Systems
NGSS: MS-ESS2-5 (Primary), MS-ESS3-3 & MS-ESS3-5 (Spiral)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color Palette - Weather theme (blues and grays)
COLORS = {
    # Main gradients
    'header_blue_start': RGBColor(0x42, 0x99, 0xE1),  # #4299E1
    'header_blue_end': RGBColor(0x2B, 0x6C, 0xB0),    # #2B6CB0
    'storm_gray': RGBColor(0x4A, 0x55, 0x68),         # #4A5568
    'storm_dark': RGBColor(0x2D, 0x37, 0x48),         # #2D3748
    'purple_start': RGBColor(0x66, 0x7E, 0xEA),       # #667EEA (Hook)
    'purple_end': RGBColor(0x76, 0x4B, 0xA2),         # #764BA2
    'orange_start': RGBColor(0xF6, 0xAD, 0x55),       # #F6AD55 (Station 1)
    'orange_end': RGBColor(0xDD, 0x6B, 0x20),         # #DD6B20
    'cyan_start': RGBColor(0x00, 0xB5, 0xD8),         # #00B5D8 (Station 2)
    'cyan_end': RGBColor(0x00, 0x86, 0x9B),           # #00869B
    'green_start': RGBColor(0x48, 0xBB, 0x78),        # #48BB78 (Station 3)
    'green_end': RGBColor(0x27, 0x67, 0x49),          # #276749
    'exit_purple_start': RGBColor(0x9F, 0x7A, 0xEA),  # #9F7AEA
    'exit_purple_end': RGBColor(0x6B, 0x46, 0xC1),    # #6B46C1
    # Accents
    'teal': RGBColor(0x38, 0xB2, 0xAC),
    'teal_dark': RGBColor(0x23, 0x4E, 0x52),
    # Text colors
    'dark_text': RGBColor(0x2D, 0x37, 0x48),
    'gray_text': RGBColor(0x4A, 0x55, 0x68),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    # Background colors
    'light_blue_bg': RGBColor(0xEB, 0xF8, 0xFF),
    'light_cyan_bg': RGBColor(0xE0, 0xF7, 0xFA),
    'light_green_bg': RGBColor(0xF0, 0xFF, 0xF4),
    'light_orange_bg': RGBColor(0xFF, 0xFA, 0xF0),
    'light_purple_bg': RGBColor(0xFA, 0xF5, 0xFF),
    'light_gray_bg': RGBColor(0xF7, 0xFA, 0xFC),
    # Specific colors
    'red_accent': RGBColor(0xE5, 0x3E, 0x3E),
    'blue_accent': RGBColor(0x31, 0x82, 0xCE),
    'green_accent': RGBColor(0x38, 0xA1, 0x69),
    'orange_accent': RGBColor(0xC0, 0x56, 0x21),
}


def create_presentation():
    """Create the G7_C5_W1 presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    add_title_slide(prs)
    add_phenomenon_slide(prs)
    add_driving_question_slide(prs)
    add_prior_knowledge_slide(prs)
    add_learning_targets_slide(prs)
    add_vocabulary_slide(prs)
    add_hook_intro_slide(prs)
    add_hook_support_slide(prs)
    add_station1_intro_slide(prs)
    add_station1_support_slide(prs)
    add_station2_intro_slide(prs)
    add_station2_support_slide(prs)
    add_station3_intro_slide(prs)
    add_station3_support_slide(prs)
    add_exit_ticket_slide(prs)
    add_summary_slide(prs)

    return prs


def add_colored_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=None, align=PP_ALIGN.LEFT, font_name="Arial", anchor=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    if anchor:
        tf.anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    if color:
        p.font.color.rgb = color
    return txBox


def add_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['header_blue_end'])

    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(1.2),
                "Weather & Climate Systems",
                font_size=44, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")

    add_text_box(slide, Inches(0.5), Inches(2.8), Inches(9), Inches(0.5),
                "Week 1: Air Masses & Storm Formation",
                font_size=24, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(3.4), Inches(9), Inches(0.5),
                "Grade 7 Science | Cycle 5 | 100 Points",
                font_size=18, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Teaser box
    add_colored_shape(slide, Inches(2), Inches(4.3), Inches(6), Inches(0.8), COLORS['white'])
    add_text_box(slide, Inches(2.2), Inches(4.35), Inches(5.6), Inches(0.7),
                "Why do some storms spin while others don't?",
                font_size=16, bold=True, color=COLORS['header_blue_end'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_phenomenon_slide(prs):
    """Slide 2: Phenomenon"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['header_blue_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "The Spinning Storm Mystery",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Main phenomenon box
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(9.6), Inches(2.4), COLORS['light_blue_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(0.08), Inches(2.4), COLORS['header_blue_end'])

    phenomenon_text = """Hurricane Katrina spun counterclockwise for days over the Gulf of Mexico, while tornadoes in Kansas spin much faster but only for minutes.

Meanwhile, the weather front that brought rain yesterday just moved straight across—no spinning at all.

What makes these storms so different?"""
    add_text_box(slide, Inches(0.5), Inches(1.0), Inches(9.1), Inches(2.2),
                phenomenon_text, font_size=15, color=COLORS['dark_text'])

    # Visual comparison boxes
    add_colored_shape(slide, Inches(0.2), Inches(3.5), Inches(3.0), Inches(1.5), COLORS['cyan_start'])
    add_text_box(slide, Inches(0.4), Inches(3.55), Inches(2.6), Inches(0.35),
                "HURRICANES", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(0.4), Inches(3.95), Inches(2.6), Inches(1.0),
                "Days to weeks\nCounterclockwise\nForms over warm ocean",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(3.5), Inches(3.5), Inches(3.0), Inches(1.5), COLORS['storm_gray'])
    add_text_box(slide, Inches(3.7), Inches(3.55), Inches(2.6), Inches(0.35),
                "TORNADOES", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(3.7), Inches(3.95), Inches(2.6), Inches(1.0),
                "Minutes to hours\nExtreme spin\nForms from thunderstorms",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(6.8), Inches(3.5), Inches(3.0), Inches(1.5), COLORS['green_start'])
    add_text_box(slide, Inches(7.0), Inches(3.55), Inches(2.6), Inches(0.35),
                "FRONTS", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(7.0), Inches(3.95), Inches(2.6), Inches(1.0),
                "Hours to days\nNo spin\nAir mass boundaries",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_driving_question_slide(prs):
    """Slide 3: Driving Question"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Driving Question",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Main question box
    add_colored_shape(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(1.2), COLORS['header_blue_end'])
    add_text_box(slide, Inches(0.7), Inches(1.25), Inches(8.6), Inches(1.1),
                "Why do some storms spin while others don't?",
                font_size=28, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Sub-questions
    add_colored_shape(slide, Inches(0.2), Inches(2.7), Inches(9.6), Inches(2.0), COLORS['light_cyan_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.8), Inches(9.2), Inches(0.3),
                "Questions we'll investigate:", font_size=14, bold=True, color=COLORS['teal_dark'])

    questions = """• What are air masses and how do they get their properties?
• What happens when different air masses collide?
• How do pressure systems affect weather patterns?
• Why does spinning happen in some storms but not others?"""
    add_text_box(slide, Inches(0.4), Inches(3.15), Inches(9.2), Inches(1.4),
                questions, font_size=14, color=COLORS['dark_text'])

    # Connection to prior cycles
    add_colored_shape(slide, Inches(0.2), Inches(4.85), Inches(9.6), Inches(0.55), COLORS['green_start'])
    add_text_box(slide, Inches(0.4), Inches(4.88), Inches(9.2), Inches(0.5),
                "Cycle 4 Connection: How do the ocean and atmosphere work together?",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_prior_knowledge_slide(prs):
    """Slide 4: Prior Knowledge"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "What You Already Know",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Cycle 3 review
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(4.7), Inches(1.8), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(0.08), Inches(1.8), COLORS['purple_end'])
    add_text_box(slide, Inches(0.4), Inches(1.0), Inches(4.4), Inches(0.3),
                "From Cycle 3 (Climate Change):", font_size=13, bold=True, color=COLORS['purple_end'])
    c3_review = """• Greenhouse effect traps heat
• CO₂ affects atmospheric temperature
• Feedback loops amplify changes
• Human activity impacts atmosphere"""
    add_text_box(slide, Inches(0.4), Inches(1.35), Inches(4.4), Inches(1.3),
                c3_review, font_size=12, color=COLORS['dark_text'])

    # Cycle 4 review
    add_colored_shape(slide, Inches(5.1), Inches(0.9), Inches(4.7), Inches(1.8), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(5.1), Inches(0.9), Inches(0.08), Inches(1.8), COLORS['green_end'])
    add_text_box(slide, Inches(5.3), Inches(1.0), Inches(4.4), Inches(0.3),
                "From Cycle 4 (Human Impact):", font_size=13, bold=True, color=COLORS['green_end'])
    c4_review = """• Ocean absorbs CO₂ and heat
• Water cycle moves energy
• Ecosystems respond to changes
• Humans affect Earth systems"""
    add_text_box(slide, Inches(5.3), Inches(1.35), Inches(4.4), Inches(1.3),
                c4_review, font_size=12, color=COLORS['dark_text'])

    # New learning
    add_colored_shape(slide, Inches(0.2), Inches(2.9), Inches(9.6), Inches(1.7), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.0), Inches(9.2), Inches(0.3),
                "NEW in Cycle 5:", font_size=14, bold=True, color=COLORS['header_blue_end'])
    new_learning = """The atmosphere creates weather—the short-term patterns we experience every day.
The same processes that drive climate (long-term) also create individual weather events (short-term).
Today we'll explore HOW air masses form and interact to create the weather we see."""
    add_text_box(slide, Inches(0.4), Inches(3.35), Inches(9.2), Inches(1.2),
                new_learning, font_size=13, color=COLORS['dark_text'])

    # Key distinction
    add_colored_shape(slide, Inches(0.2), Inches(4.75), Inches(9.6), Inches(0.6), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.78), Inches(9.2), Inches(0.55),
                "KEY: Weather is what's happening NOW | Climate is the long-term AVERAGE pattern",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_learning_targets_slide(prs):
    """Slide 5: Learning Targets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Learning Targets",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Targets
    targets = [
        ("1", "Identify characteristics of different air mass types (mT, cT, mP, cP)"),
        ("2", "Explain how frontal boundaries form where air masses meet"),
        ("3", "Interpret weather maps showing pressure systems and fronts"),
        ("4", "Predict weather changes based on approaching fronts"),
    ]

    y_pos = 0.9
    for num, target in targets:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.7), COLORS['light_green_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.5), Inches(0.7), COLORS['green_end'])
        add_text_box(slide, Inches(0.25), Inches(y_pos + 0.05), Inches(0.4), Inches(0.6),
                    num, font_size=20, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(0.85), Inches(y_pos + 0.08), Inches(8.8), Inches(0.55),
                    target, font_size=14, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.8

    # Summary bar
    add_colored_shape(slide, Inches(0.3), Inches(4.2), Inches(9.4), Inches(0.7), COLORS['teal_dark'])
    add_text_box(slide, Inches(1.5), Inches(4.25), Inches(1.5), Inches(0.35),
                "100", font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(1.5), Inches(4.55), Inches(1.5), Inches(0.3),
                "Total Points", font_size=9, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(4.25), Inches(4.25), Inches(1.5), Inches(0.35),
                "~75", font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(4.25), Inches(4.55), Inches(1.5), Inches(0.3),
                "Minutes", font_size=9, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(7.0), Inches(4.25), Inches(1.5), Inches(0.35),
                "5", font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(7.0), Inches(4.55), Inches(1.5), Inches(0.3),
                "Sections", font_size=9, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_vocabulary_slide(prs):
    """Slide 6: Key Vocabulary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.2), Inches(0.15), Inches(9.6), Inches(0.5), COLORS['blue_accent'])
    add_text_box(slide, Inches(0.4), Inches(0.18), Inches(9.2), Inches(0.45),
                "Key Vocabulary",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    vocab = [
        ("Air Mass", "Large body of air with uniform temperature and humidity"),
        ("Front", "Boundary where two different air masses meet"),
        ("Maritime (m)", "Air mass formed over water—humid"),
        ("Continental (c)", "Air mass formed over land—dry"),
        ("Polar (P)", "Air mass from high latitudes—cold"),
        ("Tropical (T)", "Air mass from low latitudes—warm"),
    ]

    y_pos = 0.75
    for i, (term, definition) in enumerate(vocab):
        bg_color = COLORS['light_blue_bg'] if i % 2 == 0 else COLORS['white']
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.6), bg_color)
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.2), Inches(0.5),
                    term, font_size=13, bold=True, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.7), Inches(y_pos + 0.05), Inches(6.9), Inches(0.5),
                    definition, font_size=12, color=COLORS['gray_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    # Notecard prompt
    add_colored_shape(slide, Inches(0.2), Inches(4.4), Inches(9.6), Inches(0.55), COLORS['purple_start'])
    add_text_box(slide, Inches(0.4), Inches(4.43), Inches(9.2), Inches(0.5),
                "Notecard: Write down mT, cT, mP, cP with what each letter means!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Quick reference
    add_colored_shape(slide, Inches(0.2), Inches(5.0), Inches(9.6), Inches(0.45), COLORS['light_cyan_bg'])
    add_text_box(slide, Inches(0.4), Inches(5.03), Inches(9.2), Inches(0.4),
                "Memory Tip: m = moisture (water) | c = continental (land) | P = polar (cold) | T = tropical (warm)",
                font_size=11, color=COLORS['teal_dark'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_hook_intro_slide(prs):
    """Slide 7: Hook Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header with two lines
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['purple_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Hook – The Spinning Storm Mystery",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "12 Points | ~10 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Video comparison section
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "OBSERVE & COMPARE", font_size=14, bold=True, color=COLORS['purple_end'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.3), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.3), COLORS['purple_end'])

    observe_text = """Watch the video comparison:
• Hurricane satellite view
• Tornado ground footage
• Front radar image

Note the KEY differences:
• Rotation speed
• Size
• Duration
• Where they form"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(2.1),
                observe_text, font_size=12, color=COLORS['dark_text'])

    # Prediction section
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(1.5), COLORS['cyan_start'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "YOUR PREDICTION", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.2), Inches(1.55), Inches(4.3), Inches(1.0),
                "Before we explain: WHY do you think hurricanes spin but fronts don't?\n\nWrite your prediction in the Hook Form.",
                font_size=13, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # MTSS diagnostic
    add_colored_shape(slide, Inches(5.0), Inches(2.8), Inches(4.7), Inches(1.0), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(5.2), Inches(2.9), Inches(4.3), Inches(0.25),
                "Quick Check (0 pts):", font_size=11, bold=True, color=COLORS['orange_end'])
    add_text_box(slide, Inches(5.2), Inches(3.2), Inches(4.3), Inches(0.55),
                "Does cold air rise or sink?\n(This helps us know where to start!)",
                font_size=11, color=COLORS['dark_text'])

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.0), Inches(9.7), Inches(0.55), COLORS['purple_end'])
    add_text_box(slide, Inches(0.35), Inches(4.03), Inches(9.3), Inches(0.5),
                "Notecard: Write your prediction BEFORE we explain!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Cycle connection
    add_colored_shape(slide, Inches(0.15), Inches(4.65), Inches(9.7), Inches(0.55), COLORS['green_start'])
    add_text_box(slide, Inches(0.35), Inches(4.68), Inches(9.3), Inches(0.5),
                "Cycle 4 Link: How does atmospheric CO₂ affect the energy available for storms?",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_hook_support_slide(prs):
    """Slide 8: Hook Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['purple_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Hook – Why Storms Spin (or Don't)",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # The answer
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(9.6), Inches(1.3), COLORS['light_blue_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(0.08), Inches(1.3), COLORS['cyan_start'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(9.2), Inches(0.3),
                "THE ANSWER: It's all about ROTATION and SCALE", font_size=14, bold=True, color=COLORS['cyan_end'])
    answer_text = """• Large systems (hurricanes) are affected by Earth's rotation (Coriolis effect)
• Small systems (tornadoes) spin due to wind shear—winds at different speeds/directions
• Fronts are boundaries, not rotating systems—they just push through like a wall"""
    add_text_box(slide, Inches(0.4), Inches(1.3), Inches(9.2), Inches(0.8),
                answer_text, font_size=12, color=COLORS['dark_text'])

    # Key concept boxes
    add_colored_shape(slide, Inches(0.2), Inches(2.3), Inches(4.7), Inches(1.4), COLORS['light_cyan_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.4), Inches(4.4), Inches(0.25),
                "KEY: DENSITY DRIVES WEATHER", font_size=12, bold=True, color=COLORS['cyan_end'])
    density_text = """• Warm air is LESS dense → rises
• Cold air is MORE dense → sinks
• When air masses meet, dense cold air
  pushes under light warm air
• This creates fronts and storms!"""
    add_text_box(slide, Inches(0.4), Inches(2.7), Inches(4.4), Inches(0.95),
                density_text, font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.1), Inches(2.3), Inches(4.6), Inches(1.4), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(5.3), Inches(2.4), Inches(4.2), Inches(0.25),
                "TODAY'S FOCUS: AIR MASSES", font_size=12, bold=True, color=COLORS['purple_end'])
    focus_text = """Air masses are like giant invisible
blankets of air with similar:
• Temperature (warm or cold)
• Humidity (wet or dry)

They get these properties from WHERE
they form (source region)."""
    add_text_box(slide, Inches(5.3), Inches(2.7), Inches(4.2), Inches(0.95),
                focus_text, font_size=11, color=COLORS['dark_text'])

    # Complete the hook
    add_colored_shape(slide, Inches(0.15), Inches(3.85), Inches(9.7), Inches(0.55), COLORS['purple_start'])
    add_text_box(slide, Inches(0.35), Inches(3.88), Inches(9.3), Inches(0.5),
                "Complete the Hook Form (12 pts) - Compare your prediction to the actual answer!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station1_intro_slide(prs):
    """Slide 9: Station 1 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['orange_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 1 – Air Mass Properties",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | ~18 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Left - Investigation
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "YOUR INVESTIGATION", font_size=14, bold=True, color=COLORS['orange_end'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.3), COLORS['light_orange_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.3), COLORS['orange_end'])

    steps = """1. Watch the convection demonstration (3 min)
2. Analyze temperature data from demo (5 min)
3. Match air mass types to source regions (5 min)
4. Identify which air masses affect your area (3 min)
5. Answer analysis questions (2 min)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(2.1),
                steps, font_size=12, color=COLORS['dark_text'])

    # Right - Air mass types
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(1.4), COLORS['orange_end'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "4 AIR MASS TYPES", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    types_text = """mT = maritime Tropical (warm & wet)
cT = continental Tropical (warm & dry)
mP = maritime Polar (cold & wet)
cP = continental Polar (cold & dry)"""
    add_text_box(slide, Inches(5.2), Inches(1.55), Inches(4.3), Inches(0.9),
                types_text, font_size=12, color=COLORS['white'])

    # Key vocabulary
    add_colored_shape(slide, Inches(5.0), Inches(2.65), Inches(4.7), Inches(1.15), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.2), Inches(2.75), Inches(4.3), Inches(0.25),
                "KEY VOCABULARY:", font_size=11, bold=True, color=COLORS['blue_accent'])
    vocab_text = """• Air mass: Large body of air with uniform properties
• Source region: Where air mass forms & gets properties
• Density: How tightly packed molecules are"""
    add_text_box(slide, Inches(5.2), Inches(3.0), Inches(4.3), Inches(0.75),
                vocab_text, font_size=10, color=COLORS['dark_text'])

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.0), Inches(9.7), Inches(0.55), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(4.03), Inches(9.3), Inches(0.5),
                "Complete Station 1 in Form - Identify all 4 air mass types!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station1_support_slide(prs):
    """Slide 10: Station 1 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 1 – Air Mass Source Regions",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # US map region descriptions
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(9.6), Inches(0.45), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.88), Inches(9.2), Inches(0.4),
                "Where do air masses form? The SOURCE REGION determines the properties!",
                font_size=13, bold=True, color=COLORS['orange_end'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Air mass table
    air_masses = [
        ("cP (Continental Polar)", "Canada/Arctic", "Cold & Dry", "Winter cold snaps"),
        ("mP (Maritime Polar)", "N. Pacific/Atlantic", "Cold & Wet", "Cool, rainy weather"),
        ("cT (Continental Tropical)", "SW Desert/Mexico", "Hot & Dry", "Summer heat waves"),
        ("mT (Maritime Tropical)", "Gulf of Mexico", "Warm & Wet", "Humidity, storms"),
    ]

    y_pos = 1.4
    for i, (name, source, properties, effect) in enumerate(air_masses):
        bg_color = COLORS['light_blue_bg'] if i % 2 == 0 else COLORS['white']
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.55), bg_color)
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.5), Inches(0.45),
                    name, font_size=11, bold=True, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.0), Inches(y_pos + 0.05), Inches(2.0), Inches(0.45),
                    source, font_size=11, color=COLORS['gray_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(5.1), Inches(y_pos + 0.05), Inches(1.8), Inches(0.45),
                    properties, font_size=11, color=COLORS['orange_end'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(7.0), Inches(y_pos + 0.05), Inches(2.7), Inches(0.45),
                    effect, font_size=10, color=COLORS['gray_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.55

    # Sentence starters
    add_colored_shape(slide, Inches(0.2), Inches(3.65), Inches(4.7), Inches(1.0), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.75), Inches(4.4), Inches(0.25),
                "SENTENCE STARTERS:", font_size=11, bold=True, color=COLORS['green_accent'])
    starters = """• "The ___ air mass is cold because it formed over..."
• "Maritime air masses are humid because..."
• "When a cP air mass moves south, it brings...""""
    add_text_box(slide, Inches(0.4), Inches(4.0), Inches(4.4), Inches(0.6),
                starters, font_size=10, color=COLORS['dark_text'])

    # Misconception alert
    add_colored_shape(slide, Inches(5.1), Inches(3.65), Inches(4.6), Inches(1.0), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(5.3), Inches(3.75), Inches(4.2), Inches(0.25),
                "COMMON MISCONCEPTION:", font_size=11, bold=True, color=COLORS['red_accent'])
    misconception = """Cold air RISES → WRONG!
Cold air SINKS because it's denser.
Warm air RISES because it's less dense."""
    add_text_box(slide, Inches(5.3), Inches(4.0), Inches(4.2), Inches(0.6),
                misconception, font_size=10, color=COLORS['dark_text'])


def add_station2_intro_slide(prs):
    """Slide 11: Station 2 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['cyan_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 2 – Pressure Systems Mapping",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | ~15 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Left - Investigation
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "YOUR INVESTIGATION", font_size=14, bold=True, color=COLORS['cyan_end'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.3), COLORS['light_cyan_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.3), COLORS['cyan_end'])

    steps = """1. Analyze weather maps with H and L (3 min)
2. Draw isobars connecting equal pressure (4 min)
3. Identify wind direction patterns (3 min)
4. Predict weather from pressure systems (3 min)
5. Answer analysis questions (2 min)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(2.1),
                steps, font_size=12, color=COLORS['dark_text'])

    # Right - Key concepts
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(1.3), COLORS['cyan_end'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "PRESSURE SYSTEM BASICS", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    basics = """H = High Pressure → Fair weather
     Clockwise winds (N. Hemisphere)

L = Low Pressure → Stormy weather
     Counterclockwise winds"""
    add_text_box(slide, Inches(5.2), Inches(1.55), Inches(4.3), Inches(0.8),
                basics, font_size=11, color=COLORS['white'])

    # Wind rule
    add_colored_shape(slide, Inches(5.0), Inches(2.55), Inches(4.7), Inches(1.25), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.2), Inches(2.65), Inches(4.3), Inches(0.25),
                "KEY RULE:", font_size=11, bold=True, color=COLORS['blue_accent'])
    wind_rule = """Winds flow from HIGH to LOW pressure.

Closer isobars = Stronger winds
(Steeper pressure gradient!)"""
    add_text_box(slide, Inches(5.2), Inches(2.95), Inches(4.3), Inches(0.8),
                wind_rule, font_size=11, color=COLORS['dark_text'])

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.0), Inches(9.7), Inches(0.55), COLORS['cyan_end'])
    add_text_box(slide, Inches(0.35), Inches(4.03), Inches(9.3), Inches(0.5),
                "Complete Station 2 in Form - Map pressure systems and predict weather!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station2_support_slide(prs):
    """Slide 12: Station 2 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['cyan_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 2 – Reading Weather Maps",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Pressure explanation
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(4.7), Inches(1.7), COLORS['light_cyan_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(4.4), Inches(0.25),
                "WHAT IS BAROMETRIC PRESSURE?", font_size=12, bold=True, color=COLORS['cyan_end'])
    pressure_text = """The weight of air pushing down on Earth.
• Measured in millibars (mb) or inches Hg
• Average sea level: ~1013 mb
• High pressure: >1013 mb (air sinking)
• Low pressure: <1013 mb (air rising)

Rising air → clouds → precipitation
Sinking air → clear skies → dry weather"""
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(4.4), Inches(1.2),
                pressure_text, font_size=10, color=COLORS['dark_text'])

    # Front types
    add_colored_shape(slide, Inches(5.1), Inches(0.85), Inches(4.6), Inches(1.7), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(5.3), Inches(0.95), Inches(4.2), Inches(0.25),
                "FRONT TYPES ON MAPS", font_size=12, bold=True, color=COLORS['purple_end'])
    fronts = """COLD FRONT (blue triangles →)
• Cold air pushes under warm air
• Fast-moving, steep slope
• Heavy rain, then clearing

WARM FRONT (red semicircles →)
• Warm air slides over cold air
• Slow-moving, gentle slope
• Light rain for longer time"""
    add_text_box(slide, Inches(5.3), Inches(1.25), Inches(4.2), Inches(1.2),
                fronts, font_size=10, color=COLORS['dark_text'])

    # Discussion questions
    add_colored_shape(slide, Inches(0.2), Inches(2.7), Inches(9.6), Inches(0.95), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.8), Inches(9.2), Inches(0.25),
                "DISCUSSION QUESTIONS:", font_size=11, bold=True, color=COLORS['green_accent'])
    questions = """1. Why does low pressure bring clouds and storms? (Hint: What happens to rising air?)
2. How do meteorologists use pressure readings to predict weather changes?
3. What weather would you expect when a cold front approaches your location?"""
    add_text_box(slide, Inches(0.4), Inches(3.1), Inches(9.2), Inches(0.5),
                questions, font_size=10, color=COLORS['dark_text'])

    # Cycle 3 connection
    add_colored_shape(slide, Inches(0.15), Inches(3.8), Inches(9.7), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(3.83), Inches(9.3), Inches(0.5),
                "Cycle 3 Link: How does atmospheric heating (greenhouse effect) power these pressure systems?",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station3_intro_slide(prs):
    """Slide 13: Station 3 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['green_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 3 – Design a Storm Warning System",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "25 Points | ~20 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Design challenge
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "ENGINEERING CHALLENGE", font_size=14, bold=True, color=COLORS['green_end'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.3), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.3), COLORS['green_end'])

    challenge = """Design an early warning network for your community!

Budget: $300,000/year
Requirements:
• At least 4 monitoring stations
• Must measure: pressure, temp, humidity, wind
• Data updated at least hourly

Choose your technologies and justify your choices!"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(2.1),
                challenge, font_size=12, color=COLORS['dark_text'])

    # Technology options
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(2.7), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "TECHNOLOGY OPTIONS", font_size=13, bold=True, color=COLORS['blue_accent'])

    options = """Weather Station | $25k/yr each
  Reliable, all measurements, fixed location

Doppler Radar | $150k/yr
  Storm detection, wide area, not ground data

Weather Balloons | $15k/yr
  Upper atmosphere, daily only

Satellite Feed | $50k/yr
  Regional view, no ground-level data"""
    add_text_box(slide, Inches(5.2), Inches(1.55), Inches(4.3), Inches(2.2),
                options, font_size=10, color=COLORS['dark_text'])

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.0), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(4.03), Inches(9.3), Inches(0.5),
                "Complete Station 3 in Form - Design and justify your system!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station3_support_slide(prs):
    """Slide 14: Station 3 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 3 – Design Tips & Trade-offs",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Example calculation
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(4.7), Inches(1.6), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(4.4), Inches(0.25),
                "EXAMPLE BUDGET CALCULATION:", font_size=11, bold=True, color=COLORS['green_accent'])
    calc = """Option A:
• 4 Weather Stations: 4 × $25k = $100k
• 1 Doppler Radar: $150k
• Satellite Feed: $50k
TOTAL: $300k ✓ (exactly on budget!)

Trade-off: Great coverage but no upper atmosphere data"""
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(4.4), Inches(1.15),
                calc, font_size=10, color=COLORS['dark_text'])

    # Design questions
    add_colored_shape(slide, Inches(5.1), Inches(0.85), Inches(4.6), Inches(1.6), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(5.3), Inches(0.95), Inches(4.2), Inches(0.25),
                "THINK ABOUT:", font_size=11, bold=True, color=COLORS['orange_end'])
    questions = """• What data is MOST important for your area?
• What severe weather threatens your community?
• How would emergency managers USE this data?
• What's the trade-off if you skip radar?

Remember: There's no single "right" answer—
justify your choices with evidence!"""
    add_text_box(slide, Inches(5.3), Inches(1.25), Inches(4.2), Inches(1.15),
                questions, font_size=10, color=COLORS['dark_text'])

    # SEP alignment
    add_colored_shape(slide, Inches(0.2), Inches(2.6), Inches(9.6), Inches(1.0), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.7), Inches(9.2), Inches(0.25),
                "SEP-3: PLANNING AND CARRYING OUT INVESTIGATIONS", font_size=11, bold=True, color=COLORS['purple_end'])
    sep_text = """Scientists and engineers must make decisions about data collection within constraints.
Your design should balance: Coverage area | Data types | Update frequency | Cost"""
    add_text_box(slide, Inches(0.4), Inches(3.0), Inches(9.2), Inches(0.5),
                sep_text, font_size=11, color=COLORS['dark_text'])

    # Final prompt
    add_colored_shape(slide, Inches(0.15), Inches(3.75), Inches(9.7), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(3.78), Inches(9.3), Inches(0.5),
                "KEY QUESTION: How would emergency managers use YOUR data to save lives?",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_exit_ticket_slide(prs):
    """Slide 15: Exit Ticket"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.9), COLORS['exit_purple_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Exit Ticket – Air Mass Integration",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.68), Inches(9.3), Inches(0.35),
                "23 Points | ~15 min", font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Question types
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(1.5), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(0.08), Inches(1.5), COLORS['exit_purple_start'])

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.3),
                "QUESTION TYPES:", font_size=14, bold=True, color=COLORS['exit_purple_start'])

    q_types = """• 2 NEW – Air masses and fronts (today's content)
• 2 SPIRAL – Cycles 3-4: Climate change, human impact on atmosphere
• 1 INTEGRATION – Connect atmospheric heating to weather patterns
• 1 SEP-3 – Data collection planning (design rationale)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(9.0), Inches(1.0),
                q_types, font_size=13, color=COLORS['dark_text'])

    # Tips boxes
    add_colored_shape(slide, Inches(0.3), Inches(2.85), Inches(4.6), Inches(1.1), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(2.9), Inches(4.3), Inches(0.3),
                "SUCCESS TIPS:", font_size=12, bold=True, color=COLORS['green_accent'],
                anchor=MSO_ANCHOR.MIDDLE)
    tips = """• Use vocabulary: mT, cP, front, pressure
• Connect to source regions
• Explain density → air movement"""
    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(4.3), Inches(0.7),
                tips, font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.1), Inches(2.85), Inches(4.6), Inches(1.1), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.3), Inches(2.9), Inches(4.2), Inches(0.3),
                "SPIRAL REMINDER:", font_size=12, bold=True, color=COLORS['blue_accent'],
                anchor=MSO_ANCHOR.MIDDLE)
    spiral = """• C3: Greenhouse effect heats atmosphere
• C4: Ocean stores heat, affects coast
• C5: Heat drives air mass movement!"""
    add_text_box(slide, Inches(5.3), Inches(3.2), Inches(4.2), Inches(0.7),
                spiral, font_size=10, color=COLORS['dark_text'])

    # Final notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.1), Inches(9.7), Inches(0.65), COLORS['exit_purple_end'])
    add_text_box(slide, Inches(0.35), Inches(4.15), Inches(9.3), Inches(0.55),
                "FINAL Notecard: Draw a simple diagram showing what happens when mT and cP air masses collide",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_summary_slide(prs):
    """Slide 16: Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Week 1 Summary: What You Learned",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Key takeaways
    takeaways = [
        ("Air Masses", "Large bodies of air with uniform temperature & humidity from source regions"),
        ("4 Types", "mT (warm/wet), cT (warm/dry), mP (cold/wet), cP (cold/dry)"),
        ("Fronts", "Boundaries where air masses meet—dense cold air pushes under warm air"),
        ("Pressure", "High = sinking air = fair | Low = rising air = storms"),
    ]

    y_pos = 0.9
    for title, detail in takeaways:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.65), COLORS['light_cyan_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.65), COLORS['teal'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.0), Inches(0.55),
                    title, font_size=13, bold=True, color=COLORS['teal_dark'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.5), Inches(y_pos + 0.05), Inches(7.1), Inches(0.55),
                    detail, font_size=12, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.72

    # Next week preview
    add_colored_shape(slide, Inches(0.2), Inches(3.85), Inches(9.6), Inches(0.65), COLORS['header_blue_end'])
    add_text_box(slide, Inches(0.4), Inches(3.88), Inches(9.2), Inches(0.6),
                "NEXT WEEK: How do meteorologists use data to predict severe weather?",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Key rule
    add_colored_shape(slide, Inches(0.2), Inches(4.65), Inches(9.6), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(4.68), Inches(9.2), Inches(0.5),
                "REMEMBER: Density drives weather—warm air rises, cold air sinks!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


if __name__ == "__main__":
    prs = create_presentation()
    output_path = "/home/user/Kairos.Sci.Repo/content/grade7/cycle05/week1/G7_C5_W1_Air_Masses_Slides_Final.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Created: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
