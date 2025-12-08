#!/usr/bin/env python3
"""
G8 C4 W3 - Ecosystems & Energy Transfer Assessment
MS-LS2-3 & MS-LS2-4 | 100 Points | ~75 Minutes

Part 1: Synthesis Review (20 pts) - Connect Week 1 & Week 2
Part 2: Cumulative Assessment (60 pts) - Full cycle content + data analysis
Part 3: Misconception Check (20 pts) - Common misconceptions + explanations

Color scheme: Green (#059669, #047857) with cyan, purple accents
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# Theme colors
GREEN_PRIMARY = RGBColor(5, 150, 105)    # #059669
GREEN_DARK = RGBColor(4, 120, 87)        # #047857
CYAN_PRIMARY = RGBColor(8, 145, 178)     # #0891B2
CYAN_DARK = RGBColor(14, 116, 144)       # #0E7490
PURPLE_PRIMARY = RGBColor(139, 92, 246)  # #8B5CF6
PURPLE_DARK = RGBColor(124, 58, 237)     # #7C3AED
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(45, 55, 72)
LIGHT_BG = RGBColor(240, 253, 244)       # #F0FDF4

def add_title_slide(prs, title, subtitle, bg_color=GREEN_PRIMARY):
    """Add a title slide with colored background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    background.fill.solid()
    background.fill.fore_color.rgb = bg_color
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_header(prs, title, metadata, bg_color=GREEN_PRIMARY):
    """Add a section header slide with two-line format."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    background.fill.solid()
    background.fill.fore_color.rgb = bg_color
    background.line.fill.background()

    # Title (line 1)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Metadata (line 2)
    meta_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.6))
    tf = meta_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = metadata
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, metadata, bullets, title_color=GREEN_PRIMARY):
    """Add a content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Light background
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    background.fill.solid()
    background.fill.fore_color.rgb = LIGHT_BG
    background.line.fill.background()

    # Title with metadata (two-line header)
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(1.1))
    tf = title_box.text_frame
    tf.word_wrap = True

    # Line 1: Title
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.alignment = PP_ALIGN.LEFT

    # Line 2: Metadata
    p2 = tf.add_paragraph()
    p2.text = metadata
    p2.font.size = Pt(14)
    p2.font.color.rgb = DARK_TEXT
    p2.alignment = PP_ALIGN.LEFT

    # Content area
    content_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.4), Inches(9.2), Inches(3.8))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(10)
        p.alignment = PP_ALIGN.LEFT

    return slide

def add_two_column_slide(prs, title, metadata, left_title, left_items, right_title, right_items, title_color=GREEN_PRIMARY):
    """Add a two-column content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Light background
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    background.fill.solid()
    background.fill.fore_color.rgb = LIGHT_BG
    background.line.fill.background()

    # Title with metadata
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(1.1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.alignment = PP_ALIGN.LEFT

    p2 = tf.add_paragraph()
    p2.text = metadata
    p2.font.size = Pt(14)
    p2.font.color.rgb = DARK_TEXT

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.4), Inches(4.4), Inches(3.8))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = title_color

    for item in left_items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.4), Inches(3.8))
    tf = right_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = title_color

    for item in right_items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)

    return slide

def create_presentation():
    """Create the G8 C4 W3 Assessment presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Week 3: Synthesis & Assessment",
        "Grade 8 Science | Ecosystems & Energy Transfer | 100 Points",
        GREEN_PRIMARY
    )

    # Slide 2: Assessment Overview
    add_content_slide(
        prs,
        "Assessment Week Overview",
        "MS-LS2-3 & MS-LS2-4 | ~75 Minutes Total",
        [
            "Part 1: Synthesis Review (20 pts, ~15 min)",
            "Part 2: Cumulative Assessment (60 pts, ~40 min)",
            "Part 3: Misconception Check (20 pts, ~20 min)",
            "Covers: Week 1 (10% rule) + Week 2 (invasive species)",
            "Spiral content: Cycle 3 natural selection concepts"
        ],
        GREEN_PRIMARY
    )

    # Slide 3: What You'll Be Assessed On
    add_two_column_slide(
        prs,
        "What You'll Be Assessed On",
        "Content from Weeks 1 & 2 plus Cycle 3 spiral",
        "Week 1 Topics",
        [
            "Energy pyramids & 10% rule",
            "Biomass calculations",
            "Why pyramids are pyramid-shaped",
            "Efficiency of eating at lower levels"
        ],
        "Week 2 Topics",
        [
            "Invasive species & trophic cascades",
            "Keystone species & their roles",
            "Top-down vs bottom-up effects",
            "Why native species are vulnerable"
        ],
        GREEN_PRIMARY
    )

    # Slide 4: Part 1 Header
    add_section_header(
        prs,
        "Part 1: Synthesis Review",
        "20 Points | ~15 Minutes | Connect Week 1 & Week 2",
        CYAN_PRIMARY
    )

    # Slide 5: Synthesis Big Question
    add_content_slide(
        prs,
        "The Big Synthesis Question",
        "Part 1 | 20 Points | Connect your learning",
        [
            "How does the 10% rule explain why invasive species are so damaging?",
            "Think: Energy is limited at each level",
            "Think: Invasives enter at a particular trophic level",
            "Think: Limited energy makes ecosystems vulnerable",
            "Connect: Energy flow explains cascade effects"
        ],
        CYAN_PRIMARY
    )

    # Slide 6: Part 2 Header
    add_section_header(
        prs,
        "Part 2: Cumulative Assessment",
        "60 Points | ~40 Minutes | Main Assessment",
        GREEN_PRIMARY
    )

    # Slide 7: Assessment Question Types
    add_content_slide(
        prs,
        "Question Breakdown (12 Questions)",
        "Part 2 | 60 Points | Show what you know",
        [
            "Energy Pyramids (3 questions): 10% rule calculations, biomass",
            "Ecosystem Disruption (3 questions): Invasives, trophic cascades",
            "Data Analysis (2 questions): Interpret graphs and tables",
            "Engineering Design (2 questions): Sustainable systems",
            "Spiral - Cycle 3 (2 questions): Natural selection, adaptation"
        ],
        GREEN_PRIMARY
    )

    # Slide 8: Calculation Reminder
    add_content_slide(
        prs,
        "10% Rule Calculation Reminder",
        "Use this method for energy transfer questions",
        [
            "Example: If producers have 100,000 kcal, what reaches tertiary consumers?",
            "Producers: 100,000 kcal (starting point)",
            "Primary consumers: 100,000 × 0.10 = 10,000 kcal",
            "Secondary consumers: 10,000 × 0.10 = 1,000 kcal",
            "Tertiary consumers: 1,000 × 0.10 = 100 kcal (answer!)"
        ],
        GREEN_PRIMARY
    )

    # Slide 9: Part 3 Header
    add_section_header(
        prs,
        "Part 3: Misconception Check",
        "20 Points | ~20 Minutes | Correct Common Mistakes",
        PURPLE_PRIMARY
    )

    # Slide 10: Common Misconceptions
    add_two_column_slide(
        prs,
        "Common Misconceptions to Avoid",
        "Part 3 | 20 Points | Know the difference",
        "WRONG",
        [
            "Energy is destroyed at each level",
            "Invasive = poisonous/dangerous",
            "Cascades only go top to bottom",
            "Ecosystems always recover naturally"
        ],
        "RIGHT",
        [
            "Energy converted to heat (still exists!)",
            "Invasive = non-native + ecological harm",
            "Cascades can be top-down OR bottom-up",
            "Some damage is permanent w/o intervention"
        ],
        PURPLE_PRIMARY
    )

    # Slide 11: Cycle 4 Summary
    add_content_slide(
        prs,
        "Cycle 4 Complete: Key Takeaways",
        "What you learned about Ecosystems & Energy Transfer",
        [
            "10% Rule: Only ~10% transfers; 90% lost to heat/waste",
            "Energy Pyramids: Wide at bottom, narrow at top",
            "Invasive Species: Outcompete natives due to no predators",
            "Trophic Cascades: Changes ripple through entire food webs",
            "The Connection: Limited energy makes ecosystems vulnerable"
        ],
        GREEN_PRIMARY
    )

    # Slide 12: Completion
    add_title_slide(
        prs,
        "Cycle 4 Complete!",
        "You've mastered Ecosystems & Energy Transfer. Great work!",
        GREEN_PRIMARY
    )

    return prs

def main():
    # Create presentation
    prs = create_presentation()

    # Save path
    output_dir = "/home/user/Kairos.Sci.Repo/content/grade8/cycle04/week3"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "G8_C4_W3_Assessment_Slides_Final.pptx")

    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
