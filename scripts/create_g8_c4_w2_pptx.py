#!/usr/bin/env python3
"""
Create G8_C4_W2 Ecosystem Disruption & Invasive Species presentation
Following established patterns - red theme
See create_g7_c3_w2_pptx.py for full PPTX DESIGN BEST PRACTICES
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

COLORS = {
    'red_primary': RGBColor(0xDC, 0x26, 0x26),
    'red_dark': RGBColor(0xB9, 0x1C, 0x1C),
    'red_light': RGBColor(0xFE, 0xF2, 0xF2),
    'hook_orange': RGBColor(0xF5, 0x9E, 0x0B),
    'station1_blue': RGBColor(0x3B, 0x82, 0xF6),
    'station2_purple': RGBColor(0x8B, 0x5C, 0xF6),
    'station3_green': RGBColor(0x05, 0x96, 0x69),
    'exit_pink': RGBColor(0xEC, 0x48, 0x99),
    'dark_text': RGBColor(0x2D, 0x37, 0x48),
    'gray_text': RGBColor(0x4A, 0x55, 0x68),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'light_bg': RGBColor(0xF7, 0xFA, 0xFC),
    'light_red_bg': RGBColor(0xFE, 0xF2, 0xF2),
    'light_orange_bg': RGBColor(0xFF, 0xFB, 0xEB),
    'light_blue_bg': RGBColor(0xEF, 0xF6, 0xFF),
    'light_purple_bg': RGBColor(0xF5, 0xF3, 0xFF),
    'light_green_bg': RGBColor(0xEC, 0xFD, 0xF5),
    'light_pink_bg': RGBColor(0xFD, 0xF2, 0xF8),
    'warning_bg': RGBColor(0xFE, 0xD7, 0xD7),
}

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    add_title_slide(prs)
    add_phenomenon_slide(prs)
    add_driving_question_slide(prs)
    add_prior_knowledge_slide(prs)
    add_learning_targets_slide(prs)
    add_vocabulary_slide(prs)
    add_hook_intro_slide(prs)
    add_hook_support_slide(prs)
    add_station1_intro_slide(prs)
    add_station1_support_slide(prs)
    add_station2_intro_slide(prs)
    add_station2_support_slide(prs)
    add_station3_intro_slide(prs)
    add_station3_support_slide(prs)
    add_exit_ticket_slide(prs)
    add_summary_slide(prs)

    return prs

def add_colored_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=None, align=PP_ALIGN.LEFT, font_name="Arial", anchor=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    if anchor:
        tf.anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    if color:
        p.font.color.rgb = color
    p.alignment = align
    return txBox

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['red_primary'])
    add_text_box(slide, Inches(0), Inches(1.0), Inches(10), Inches(0.8),
                "🐸 ⚠️", font_size=48, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(1.8), Inches(9), Inches(1.0),
                "Week 2: Ecosystem Disruption & Invasive Species",
                font_size=34, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.5), Inches(2.9), Inches(9), Inches(0.5),
                "Grade 8 Science | Cycle 4 | Rosche | Kairos Academies",
                font_size=16, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(3.5), Inches(9), Inches(0.4),
                "MS-LS2-3 & MS-LS2-4 Ecosystem Disruption",
                font_size=14, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(4.2), Inches(9), Inches(0.4),
                "100 Points Total | ~75 Minutes",
                font_size=14, color=COLORS['white'], align=PP_ALIGN.CENTER)

def add_phenomenon_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['red_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🐸 The Phenomenon: The Cane Toad Disaster",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    add_colored_shape(slide, Inches(0.15), Inches(1.0), Inches(9.7), Inches(2.8), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(1.1), Inches(9.2), Inches(0.5),
                "In 1935, Australia had a 'brilliant' idea:",
                font_size=16, color=COLORS['dark_text'])
    add_text_box(slide, Inches(0.6), Inches(1.6), Inches(9.0), Inches(1.8),
                "• They introduced 102 cane toads from Hawaii to eat beetles\n" +
                "• Today there are 200+ MILLION cane toads\n" +
                "• Native predators that try to eat them DIE from their poison\n" +
                "• The beetles? Still destroying crops - toads don't eat them!\n\n" +
                "How did 102 toads become an ecological catastrophe?",
                font_size=15, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.15), Inches(3.95), Inches(9.7), Inches(1.4), COLORS['red_dark'])
    add_text_box(slide, Inches(0.35), Inches(4.15), Inches(9.3), Inches(1.0),
                "🔬 KEY CONCEPT: Invasive species disrupt energy flow because they\n" +
                "don't fit into existing food webs - no predators, no competition!",
                font_size=16, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

def add_driving_question_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['red_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🤔 Driving Question",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    add_colored_shape(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(2.0), COLORS['red_dark'])
    add_text_box(slide, Inches(0.7), Inches(1.7), Inches(8.6), Inches(1.6),
                "How do invasive species disrupt energy flow?\n\n" +
                "Why is it so hard to stop them once they're established?",
                font_size=22, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(1.4), COLORS['light_red_bg'])
    add_text_box(slide, Inches(0.7), Inches(3.9), Inches(8.6), Inches(1.2),
                "🔗 Connection to Week 1:\n" +
                "You learned about the 10% rule and energy pyramids.\n" +
                "Invasive species break these patterns - they take energy without giving it back!",
                font_size=14, color=COLORS['red_dark'])

def add_prior_knowledge_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['red_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "💡 What You Already Know",
                font_size=22, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    concepts = [
        ("Week 1", "10% rule limits populations\nat each trophic level"),
        ("Cycle 3", "Natural selection shapes\nadaptations over time"),
        ("Cycle 3", "Organisms without defenses\nare vulnerable to new threats")
    ]
    x_positions = [0.3, 3.5, 6.7]
    for i, (title, content) in enumerate(concepts):
        add_colored_shape(slide, Inches(x_positions[i]), Inches(1.1), Inches(3.0), Inches(1.8), COLORS['light_red_bg'])
        add_text_box(slide, Inches(x_positions[i] + 0.15), Inches(1.2), Inches(2.7), Inches(0.4),
                    title, font_size=14, bold=True, color=COLORS['red_dark'], align=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x_positions[i] + 0.15), Inches(1.6), Inches(2.7), Inches(1.2),
                    content, font_size=12, color=COLORS['gray_text'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.3), Inches(3.1), Inches(9.4), Inches(2.2), COLORS['warning_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(9.0), Inches(0.4),
                "⚠️ COMMON MISTAKE: 'Invasive = Bad Animal'",
                font_size=14, bold=True, color=COLORS['red_dark'])
    add_text_box(slide, Inches(0.5), Inches(3.6), Inches(9.0), Inches(1.6),
                "WRONG: 'Cane toads are invasive because they're ugly and poisonous.'\n\n" +
                "RIGHT: 'Cane toads are invasive because they're NON-NATIVE and outcompete\n" +
                "natives for resources. Native predators haven't evolved defenses against their poison.'\n\n" +
                "KEY: 'Invasive' means non-native AND causing ecological harm - not just 'harmful'",
                font_size=12, color=COLORS['red_dark'])

def add_learning_targets_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['red_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🎯 Learning Targets",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    targets = [
        "Explain how invasive species disrupt energy flow",
        "Model the cascade effects of species removal/addition",
        "Analyze real-world invasive species case studies",
        "Design and evaluate intervention strategies"
    ]
    y_pos = 1.0
    for i, target in enumerate(targets, 1):
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.95), COLORS['light_red_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.15), Inches(9.0), Inches(0.65),
                    f"Target {i}: {target}",
                    font_size=15, bold=True, color=COLORS['red_dark'])
        y_pos += 1.1

def add_vocabulary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['red_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "📚 Key Vocabulary This Week",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    vocab = [
        ("Invasive Species", "Non-native species that outcompetes natives"),
        ("Trophic Cascade", "Effects rippling through multiple food web levels"),
        ("Keystone Species", "Species with outsized ecosystem effects"),
        ("Ecological Niche", "A species' role (what it eats, where it lives)"),
        ("Top-Down Cascade", "Predator changes affecting all levels below"),
        ("Biocontrol", "Using living organisms to control pests")
    ]
    y_pos = 0.95
    for term, definition in vocab:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.7), COLORS['light_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.1), Inches(2.8), Inches(0.5),
                    term, font_size=12, bold=True, color=COLORS['red_dark'])
        add_text_box(slide, Inches(3.4), Inches(y_pos + 0.1), Inches(6.1), Inches(0.5),
                    definition, font_size=12, color=COLORS['gray_text'])
        y_pos += 0.75

def add_hook_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['hook_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🐸 Hook – The Cane Toad Disaster",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "12 Points | ~10 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(2.2), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "🐸 The Challenge: Why did the population explode?",
                font_size=16, bold=True, color=COLORS['hook_orange'])
    add_text_box(slide, Inches(0.5), Inches(1.75), Inches(9.0), Inches(1.5),
                "1. Watch the cane toad population explosion (2 min)\n" +
                "2. Predict WHY using Week 1 concepts - 10% rule (3 min)\n" +
                "3. Connect to Cycle 3: natural selection and adaptation (3 min)\n" +
                "4. Answer diagnostic questions (2 min)",
                font_size=14, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(3.55), Inches(9.4), Inches(1.8), COLORS['light_red_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.65), Inches(9.0), Inches(0.4),
                "💭 Think About This:",
                font_size=14, bold=True, color=COLORS['red_dark'])
    add_text_box(slide, Inches(0.5), Inches(4.05), Inches(9.0), Inches(1.2),
                "• What normally limits a population? (predators, competition, disease)\n" +
                "• What happens when none of those limits exist?\n" +
                "• Why haven't Australian predators evolved to eat toads?",
                font_size=13, color=COLORS['gray_text'])

def add_hook_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['hook_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🐸 Hook – Cane Toad Population Data",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    add_colored_shape(slide, Inches(0.3), Inches(1.0), Inches(4.5), Inches(2.5), COLORS['light_red_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(4.1), Inches(0.4),
                "📊 Population Explosion:", font_size=14, bold=True, color=COLORS['red_dark'])
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4.1), Inches(1.9),
                "Year     |  Population\n" +
                "1935     |       102\n" +
                "1950     |    ~10,000\n" +
                "1980     |   ~1 million\n" +
                "2010     |  ~100 million\n" +
                "Today    |  200+ million",
                font_size=11, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(5.0), Inches(1.0), Inches(4.7), Inches(2.5), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.4),
                "❓ Why No Control?", font_size=14, bold=True, color=COLORS['hook_orange'])
    add_text_box(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(1.9),
                "• No natural predators in Australia\n" +
                "• Native predators die eating them\n" +
                "• Toads eat almost anything\n" +
                "• Breed rapidly (30,000 eggs/year)\n" +
                "• Evolution takes too long!",
                font_size=11, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(3.7), Inches(9.4), Inches(1.65), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.8), Inches(9.0), Inches(0.4),
                "📝 Complete the Hook form on the student page",
                font_size=14, bold=True, color=COLORS['red_dark'])
    add_text_box(slide, Inches(0.5), Inches(4.2), Inches(9.0), Inches(1.0),
                "Explain WHY the population doubled every 5-10 years using:\n" +
                "• The 10% rule (Week 1) and • Natural selection (Cycle 3)",
                font_size=12, color=COLORS['gray_text'])

def add_station1_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['station1_blue'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🔗 Station 1 – Modeling Trophic Cascades",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "20 Points | ~18 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(1.3), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "🔗 Your Mission: Model Cascade Effects",
                font_size=16, bold=True, color=COLORS['station1_blue'])
    add_text_box(slide, Inches(0.5), Inches(1.75), Inches(9.0), Inches(0.6),
                "When one species changes, effects RIPPLE through ALL trophic levels -\n" +
                "like dominoes falling!",
                font_size=14, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(2.65), Inches(4.5), Inches(2.7), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(2.75), Inches(4.1), Inches(0.4),
                "⬇️ Top-Down Cascade:", font_size=13, bold=True, color=COLORS['station3_green'])
    add_text_box(slide, Inches(0.5), Inches(3.15), Inches(4.1), Inches(2.1),
                "Predator → Prey → Plants\n\n" +
                "Example: Wolves removed →\n" +
                "Deer overpopulate →\n" +
                "Trees die →\n" +
                "Rivers change course!",
                font_size=12, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(5.0), Inches(2.65), Inches(4.7), Inches(2.7), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.2), Inches(2.75), Inches(4.3), Inches(0.4),
                "⬆️ Bottom-Up Cascade:", font_size=13, bold=True, color=COLORS['station1_blue'])
    add_text_box(slide, Inches(5.2), Inches(3.15), Inches(4.3), Inches(2.1),
                "Plants → Herbivores → Predators\n\n" +
                "Example: Drought kills grass →\n" +
                "Zebras starve →\n" +
                "Lions decline →\n" +
                "Ecosystem collapses!",
                font_size=12, color=COLORS['gray_text'])

def add_station1_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['station1_blue'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🔗 Station 1 – Yellowstone Wolf Example",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    add_colored_shape(slide, Inches(0.3), Inches(1.0), Inches(9.4), Inches(2.5), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(9.0), Inches(0.4),
                "🐺 Case Study: Yellowstone Wolves", font_size=14, bold=True, color=COLORS['hook_orange'])
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(1.9),
                "When wolves were REMOVED (1926):          When wolves RETURNED (1995):\n" +
                "• Elk population exploded                    • Elk population decreased\n" +
                "• Elk ate all the willows/aspens            • Trees grew back near rivers\n" +
                "• Riverbanks eroded                          • Riverbanks stabilized\n" +
                "• Beavers couldn't build dams               • Beavers returned\n" +
                "• Fish populations crashed                   • Fish populations recovered\n" +
                "• Songbirds left (no trees)                  • Songbirds returned",
                font_size=11, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(3.65), Inches(9.4), Inches(1.7), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.75), Inches(9.0), Inches(0.4),
                "📝 Complete Station 1 form - Model a cascade!", font_size=14, bold=True, color=COLORS['red_dark'])
    add_text_box(slide, Inches(0.5), Inches(4.15), Inches(9.0), Inches(1.1),
                "Draw a food web and predict what happens when one species is:\n" +
                "• Removed (like wolves in Yellowstone) OR • Added (like cane toads in Australia)",
                font_size=12, color=COLORS['gray_text'])

def add_station2_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['station2_purple'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "📊 Station 2 – Case Study Analysis",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "20 Points | ~15 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(1.2), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "📊 Your Mission: Analyze Real Invasive Species Data",
                font_size=16, bold=True, color=COLORS['station2_purple'])
    add_text_box(slide, Inches(0.5), Inches(1.75), Inches(9.0), Inches(0.5),
                "Compare four famous invasive species and identify patterns.",
                font_size=14, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(2.55), Inches(9.4), Inches(2.8), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.5), Inches(2.65), Inches(9.0), Inches(0.4),
                "📈 Case Study Comparison:", font_size=14, bold=True, color=COLORS['station2_purple'])
    add_text_box(slide, Inches(0.5), Inches(3.1), Inches(9.0), Inches(2.1),
                "Species         |  Location      |  Why Invasive?             |  Impact\n" +
                "───────────────────────────────────────────────────────────────────────\n" +
                "Cane Toad       |  Australia     |  Poisonous, no predators   |  Kills native predators\n" +
                "Zebra Mussel    |  Great Lakes   |  Reproduces fast           |  Clogs pipes, starves fish\n" +
                "Kudzu           |  Southern US   |  Grows 1 ft/day            |  Smothers native trees\n" +
                "Brown Tree Snake|  Guam          |  No native defenses        |  10/12 bird species EXTINCT",
                font_size=10, color=COLORS['gray_text'])

def add_station2_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['station2_purple'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "📊 Station 2 – Pattern Analysis",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    add_colored_shape(slide, Inches(0.3), Inches(1.0), Inches(4.5), Inches(2.3), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(4.1), Inches(0.4),
                "✅ Common Patterns:", font_size=13, bold=True, color=COLORS['station3_green'])
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4.1), Inches(1.7),
                "• Non-native (from elsewhere)\n" +
                "• No natural predators\n" +
                "• Reproduce rapidly\n" +
                "• Outcompete natives\n" +
                "• Natives lack defenses",
                font_size=11, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(5.0), Inches(1.0), Inches(4.7), Inches(2.3), COLORS['light_red_bg'])
    add_text_box(slide, Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.4),
                "🔗 Connection to Cycle 3:", font_size=13, bold=True, color=COLORS['red_dark'])
    add_text_box(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(1.7),
                "Why don't natives evolve?\n\n" +
                "• Evolution takes THOUSANDS\n  of generations\n" +
                "• Invasives spread in YEARS\n" +
                "• Not enough time to adapt!",
                font_size=11, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(3.5), Inches(9.4), Inches(1.85), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.6), Inches(9.0), Inches(0.4),
                "📝 Complete Station 2 form - Identify patterns across case studies",
                font_size=14, bold=True, color=COLORS['red_dark'])
    add_text_box(slide, Inches(0.5), Inches(4.0), Inches(9.0), Inches(1.2),
                "• What do ALL invasive species have in common?\n" +
                "• Why are islands especially vulnerable?\n" +
                "• How does this connect to the 10% rule and natural selection?",
                font_size=12, color=COLORS['gray_text'])

def add_station3_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['station3_green'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🔧 Station 3 – Design Intervention Strategies",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "25 Points | ~20 min (Highest Value!)", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(1.2), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "🔧 Engineering Challenge: Design a Biocontrol Solution",
                font_size=15, bold=True, color=COLORS['station3_green'])
    add_text_box(slide, Inches(0.5), Inches(1.75), Inches(9.0), Inches(0.5),
                "Propose a strategy to control cane toads while minimizing unintended consequences.",
                font_size=14, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(2.55), Inches(9.4), Inches(2.8), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(2.65), Inches(9.0), Inches(0.4),
                "📋 Available Intervention Options:", font_size=14, bold=True, color=COLORS['station3_green'])
    add_text_box(slide, Inches(0.5), Inches(3.1), Inches(4.3), Inches(2.1),
                "Physical removal:\n" +
                "• Trapping (labor-intensive)\n" +
                "• Fencing (limited areas)\n" +
                "• Community hunts ($$ prizes)",
                font_size=11, color=COLORS['gray_text'])
    add_text_box(slide, Inches(5.0), Inches(3.1), Inches(4.5), Inches(2.1),
                "Biological control:\n" +
                "• Introduce predator (RISKY!)\n" +
                "• Species-specific virus\n" +
                "• Gene drive technology\n" +
                "• Train natives to avoid toads",
                font_size=11, color=COLORS['gray_text'])

def add_station3_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['station3_green'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🔧 Station 3 – Risk Assessment",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    options = [
        ("Trapping/Removal", "Low risk", "High cost, slow", "Can't keep up with breeding"),
        ("Introduce predator", "HIGH RISK", "Could backfire!", "Might become next invasive"),
        ("Species-specific virus", "Medium risk", "High effectiveness", "What if it mutates?"),
        ("Train native predators", "Low risk", "Moderate effect", "Takes time, limited reach")
    ]
    y_pos = 1.0
    for option, risk, effect, notes in options:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.8), COLORS['light_green_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.1), Inches(2.8), Inches(0.6),
                    option, font_size=11, bold=True, color=COLORS['station3_green'])
        add_text_box(slide, Inches(3.4), Inches(y_pos + 0.1), Inches(1.5), Inches(0.6),
                    risk, font_size=10, color=COLORS['red_dark'] if "HIGH" in risk else COLORS['gray_text'], align=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(5.0), Inches(y_pos + 0.1), Inches(2.0), Inches(0.6),
                    effect, font_size=10, color=COLORS['gray_text'], align=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(7.1), Inches(y_pos + 0.1), Inches(2.5), Inches(0.6),
                    notes, font_size=10, color=COLORS['gray_text'])
        y_pos += 0.9

    add_colored_shape(slide, Inches(0.3), Inches(4.65), Inches(9.4), Inches(0.75), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(4.7), Inches(9.0), Inches(0.6),
                "📝 Complete Station 3 - Design your intervention plan!\n" +
                "Consider: What are the trade-offs? Could your solution cause NEW problems?",
                font_size=11, bold=True, color=COLORS['red_dark'])

def add_exit_ticket_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['exit_pink'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🎓 Exit Ticket – Ecosystem Integration",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "23 Points | ~15 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(2.5), COLORS['light_pink_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4),
                "🎓 Show What You Learned - Question Types:",
                font_size=16, bold=True, color=COLORS['exit_pink'])
    add_text_box(slide, Inches(0.5), Inches(1.8), Inches(9.0), Inches(1.8),
                "• 2 NEW – Invasive species and trophic cascades (this week)\n" +
                "• 2 SPIRAL – Week 1 review (10% rule, energy pyramids)\n" +
                "• 1 INTEGRATION – Connect energy limits to cascade effects\n" +
                "• 1 SEP-2 – Develop a model showing cascade effects",
                font_size=14, color=COLORS['gray_text'])

    add_colored_shape(slide, Inches(0.3), Inches(3.85), Inches(9.4), Inches(1.5), COLORS['light_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.95), Inches(9.0), Inches(0.4),
                "📝 Complete the Exit Ticket on the student page",
                font_size=14, bold=True, color=COLORS['red_dark'])
    add_text_box(slide, Inches(0.5), Inches(4.35), Inches(9.0), Inches(0.9),
                "This is your final assessment for Week 2. Take your time!\n" +
                "Use evidence from both Week 1 and Week 2 activities.",
                font_size=13, color=COLORS['gray_text'])

def add_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.7), COLORS['red_primary'])
    add_text_box(slide, Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.5),
                "🎉 Week 2 Summary: What You Learned",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    takeaways = [
        ("Invasive Species", "Non-native + no predators + fast breeding = population explosion", COLORS['red_primary']),
        ("Trophic Cascades", "Effects ripple through ALL levels - up AND down", COLORS['station1_blue']),
        ("Evolution Too Slow", "Natural selection can't keep up with invasive spread", COLORS['station2_purple']),
        ("Interventions", "All solutions have trade-offs - some could backfire!", COLORS['station3_green'])
    ]
    y_pos = 1.0
    for title, content, color in takeaways:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.9), COLORS['light_bg'])
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(0.1), Inches(0.9), color)
        add_text_box(slide, Inches(0.55), Inches(y_pos + 0.1), Inches(9.0), Inches(0.35),
                    title, font_size=13, bold=True, color=color)
        add_text_box(slide, Inches(0.55), Inches(y_pos + 0.45), Inches(9.0), Inches(0.4),
                    content, font_size=12, color=COLORS['gray_text'])
        y_pos += 1.0

    add_colored_shape(slide, Inches(0.3), Inches(5.05), Inches(9.4), Inches(0.45), COLORS['light_red_bg'])
    add_text_box(slide, Inches(0.5), Inches(5.1), Inches(9.0), Inches(0.35),
                "Next Week: Assessment – Show what you've learned about ecosystems and energy!",
                font_size=12, color=COLORS['red_dark'], align=PP_ALIGN.CENTER)

def main():
    prs = create_presentation()
    output_dir = "/home/user/Kairos.Sci.Repo/content/grade8/cycle04/week2"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "G8_C4_W2_Invasive_Species_Slides_Final.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
