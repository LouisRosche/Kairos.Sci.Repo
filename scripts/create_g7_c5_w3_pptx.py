#!/usr/bin/env python3
"""
Create G7_C5_W3 Synthesis & Assessment presentation.
Topic: Weather & Climate Systems.

Assessment week structure (11 slides) - different from instructional weeks.
See PPTX_DESIGN_GUIDE.md for best practices documentation.
"""

import os
from pptx.dml.color import RGBColor
from pptx_common import (
    COLORS as BASE_COLORS,
    create_base_presentation,
    add_colored_shape,
    add_text_box,
    Inches,
    Pt,
    PP_ALIGN,
    MSO_ANCHOR,
)

# Weather theme (blues and grays) - extend base colors
COLORS = {**BASE_COLORS}
COLORS.update({
    'storm_gray': RGBColor(0x4A, 0x55, 0x68),
    'storm_dark': RGBColor(0x2D, 0x37, 0x48),
    'cyan_start': RGBColor(0x00, 0xB5, 0xD8),
    'cyan_end': RGBColor(0x00, 0x86, 0x9B),
    'light_cyan_bg': RGBColor(0xE0, 0xF7, 0xFA),
    'assessment_blue': RGBColor(0x2C, 0x5F, 0x9E),
})


def create_presentation():
    """Create the G7_C5_W3 Assessment presentation"""
    prs = create_base_presentation()

    add_title_slide(prs)
    add_overview_slide(prs)
    add_assessed_on_slide(prs)
    add_part1_intro_slide(prs)
    add_part1_support_slide(prs)
    add_part2_intro_slide(prs)
    add_part2_support_slide(prs)
    add_part3_intro_slide(prs)
    add_part3_support_slide(prs)
    add_tips_slide(prs)
    add_summary_slide(prs)

    return prs


def add_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['assessment_blue'])

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(1.2),
                "Weather & Climate Systems",
                font_size=44, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")

    add_text_box(slide, Inches(0.5), Inches(2.6), Inches(9), Inches(0.5),
                "Week 3: Synthesis & Assessment",
                font_size=28, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(9), Inches(0.5),
                "Grade 7 Science | Cycle 5 | 100 Points",
                font_size=18, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Assessment info box
    add_colored_shape(slide, Inches(2), Inches(4.1), Inches(6), Inches(1.0), COLORS['white'])
    add_text_box(slide, Inches(2.2), Inches(4.15), Inches(5.6), Inches(0.45),
                "Part 1: Synthesis (20 pts) | Part 2: Assessment (60 pts)",
                font_size=14, bold=True, color=COLORS['assessment_blue'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(2.2), Inches(4.55), Inches(5.6), Inches(0.45),
                "Part 3: Misconception Check (20 pts)",
                font_size=14, bold=True, color=COLORS['assessment_blue'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_overview_slide(prs):
    """Slide 2: Assessment Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['assessment_blue'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Assessment Overview",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Three parts
    parts = [
        ("Part 1: Synthesis Review", "20 pts", "15 min", "Connect W1 air masses + W2 prediction/climate"),
        ("Part 2: Cumulative Assessment", "60 pts", "40 min", "Air masses, fronts, chaos, climate trends"),
        ("Part 3: Misconception Check", "20 pts", "20 min", "Target common errors for feedback"),
    ]

    y_pos = 0.9
    colors = [COLORS['green_start'], COLORS['cyan_start'], COLORS['purple_start']]
    for i, (title, pts, time, desc) in enumerate(parts):
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(1.1), COLORS['light_blue_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(1.1), colors[i])

        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.1), Inches(5.5), Inches(0.35),
                    title, font_size=16, bold=True, color=COLORS['dark_text'])
        add_text_box(slide, Inches(6.0), Inches(y_pos + 0.1), Inches(1.8), Inches(0.35),
                    pts, font_size=14, bold=True, color=colors[i], align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(7.9), Inches(y_pos + 0.1), Inches(1.8), Inches(0.35),
                    time, font_size=14, color=COLORS['gray_text'], align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.5), Inches(9.2), Inches(0.5),
                    desc, font_size=12, color=COLORS['gray_text'])
        y_pos += 1.25

    # Total bar
    add_colored_shape(slide, Inches(0.2), Inches(4.7), Inches(9.6), Inches(0.55), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.4), Inches(4.73), Inches(9.2), Inches(0.5),
                "TOTAL: 100 Points | ~75 Minutes | Use your notecard!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_assessed_on_slide(prs):
    """Slide 3: What You'll Be Assessed On"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "What You'll Be Assessed On",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Week 1 content
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(4.7), Inches(1.7), COLORS['light_cyan_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(0.08), Inches(1.7), COLORS['cyan_start'])
    add_text_box(slide, Inches(0.4), Inches(1.0), Inches(4.4), Inches(0.3),
                "WEEK 1: Air Masses & Storms", font_size=13, bold=True, color=COLORS['cyan_end'])
    w1_content = """• 4 air mass types (mT, cT, mP, cP)
• How source regions determine properties
• Front formation and weather effects
• High/low pressure systems
• Design trade-offs for monitoring"""
    add_text_box(slide, Inches(0.4), Inches(1.35), Inches(4.4), Inches(1.2),
                w1_content, font_size=11, color=COLORS['dark_text'])

    # Week 2 content
    add_colored_shape(slide, Inches(5.1), Inches(0.9), Inches(4.7), Inches(1.7), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(5.1), Inches(0.9), Inches(0.08), Inches(1.7), COLORS['purple_start'])
    add_text_box(slide, Inches(5.3), Inches(1.0), Inches(4.4), Inches(0.3),
                "WEEK 2: Prediction & Climate", font_size=13, bold=True, color=COLORS['purple_end'])
    w2_content = """• Why weather is chaotic (butterfly effect)
• Forecast accuracy vs. time
• Weather variability vs. climate trends
• Why climate IS predictable
• Long-term monitoring design"""
    add_text_box(slide, Inches(5.3), Inches(1.35), Inches(4.4), Inches(1.2),
                w2_content, font_size=11, color=COLORS['dark_text'])

    # Integration skills
    add_colored_shape(slide, Inches(0.2), Inches(2.75), Inches(9.6), Inches(1.3), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.85), Inches(9.2), Inches(0.3),
                "INTEGRATION SKILLS:", font_size=13, bold=True, color=COLORS['green_end'])
    integration = """• Connect air mass movement to weather prediction challenges
• Explain why short-term chaos leads to long-term predictability
• Apply SEP-3: Planning and Carrying Out Investigations to data collection design
• Connect to Cycles 3-4: How does climate change affect weather patterns?"""
    add_text_box(slide, Inches(0.4), Inches(3.2), Inches(9.2), Inches(0.8),
                integration, font_size=11, color=COLORS['dark_text'])

    # Key reminder
    add_colored_shape(slide, Inches(0.2), Inches(4.2), Inches(9.6), Inches(0.55), COLORS['orange_end'])
    add_text_box(slide, Inches(0.4), Inches(4.23), Inches(9.2), Inches(0.5),
                "KEY: Weather = chaotic events | Climate = predictable trends | Your notecard is your friend!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_part1_intro_slide(prs):
    """Slide 4: Part 1 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['green_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Part 1 – Synthesis Review",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | ~15 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # What to do
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(2.3), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(2.3), COLORS['green_end'])

    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "SYNTHESIS TASK:", font_size=14, bold=True, color=COLORS['green_end'])

    task_text = """Connect what you learned in Week 1 (air masses, fronts, pressure) with Week 2
(chaos, prediction, climate trends).

You'll answer questions that require you to:
• Explain HOW air mass behavior creates the chaos that makes weather hard to predict
• Describe WHY long-term climate trends are predictable despite short-term chaos
• Connect atmospheric science to real-world forecasting challenges"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(1.8),
                task_text, font_size=13, color=COLORS['dark_text'])

    # Format
    add_colored_shape(slide, Inches(0.2), Inches(3.6), Inches(4.7), Inches(0.95), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.7), Inches(4.4), Inches(0.25),
                "FORMAT:", font_size=12, bold=True, color=COLORS['blue_accent'])
    add_text_box(slide, Inches(0.4), Inches(4.0), Inches(4.4), Inches(0.5),
                "• 4 short-answer questions\n• Use specific vocabulary\n• Connect W1 to W2 concepts",
                font_size=11, color=COLORS['dark_text'])

    # Tips
    add_colored_shape(slide, Inches(5.1), Inches(3.6), Inches(4.6), Inches(0.95), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(5.3), Inches(3.7), Inches(4.2), Inches(0.25),
                "TIPS:", font_size=12, bold=True, color=COLORS['orange_end'])
    add_text_box(slide, Inches(5.3), Inches(4.0), Inches(4.2), Inches(0.5),
                "• Review your notecard from W1-W2\n• Think about CAUSE → EFFECT\n• Use vocabulary: chaos, trend, front",
                font_size=11, color=COLORS['dark_text'])


def add_part1_support_slide(prs):
    """Slide 5: Part 1 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Part 1 – Key Connections to Make",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Connection 1
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(9.6), Inches(1.2), COLORS['light_cyan_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(9.2), Inches(0.3),
                "CONNECTION 1: Air Masses → Chaos", font_size=13, bold=True, color=COLORS['cyan_end'])
    conn1 = """Air masses are HUGE (thousands of miles) and interact in complex ways.
Small differences in where fronts form, how fast they move, or exact temperatures
lead to VERY different weather outcomes → This is WHY forecasts fail beyond 7 days."""
    add_text_box(slide, Inches(0.4), Inches(1.3), Inches(9.2), Inches(0.7),
                conn1, font_size=11, color=COLORS['dark_text'])

    # Connection 2
    add_colored_shape(slide, Inches(0.2), Inches(2.2), Inches(9.6), Inches(1.2), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.3), Inches(9.2), Inches(0.3),
                "CONNECTION 2: Chaos → Climate Predictability", font_size=13, bold=True, color=COLORS['purple_end'])
    conn2 = """Even though individual storms are chaotic, the AVERAGE of many storms is predictable.
Climate = many years of weather averaged together → Chaos "cancels out" over time
→ We can predict trends even when we can't predict individual events."""
    add_text_box(slide, Inches(0.4), Inches(2.65), Inches(9.2), Inches(0.7),
                conn2, font_size=11, color=COLORS['dark_text'])

    # Sentence starters
    add_colored_shape(slide, Inches(0.2), Inches(3.55), Inches(9.6), Inches(1.0), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.65), Inches(9.2), Inches(0.25),
                "SENTENCE STARTERS FOR SYNTHESIS:", font_size=11, bold=True, color=COLORS['green_accent'])
    starters = """• "Air mass interactions cause weather chaos because..."
• "Climate is more predictable than weather because..."
• "The connection between short-term variability and long-term trends is..."""
    add_text_box(slide, Inches(0.4), Inches(3.95), Inches(9.2), Inches(0.55),
                starters, font_size=10, color=COLORS['dark_text'])


def add_part2_intro_slide(prs):
    """Slide 6: Part 2 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['cyan_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Part 2 – Cumulative Assessment",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "60 Points | ~40 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Sections
    sections = [
        ("A: Air Mass Interactions", "15 pts", "Types, properties, fronts"),
        ("B: Weather vs Climate", "15 pts", "Variability, trends, chaos"),
        ("C: Data Analysis", "15 pts", "Interpret forecast data, calculate error"),
        ("D: System Models", "15 pts", "Design monitoring systems (SEP-3)"),
    ]

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(2.4), COLORS['light_cyan_bg'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "ASSESSMENT SECTIONS:", font_size=14, bold=True, color=COLORS['cyan_end'])

    y_pos = 1.6
    for section, pts, desc in sections:
        add_text_box(slide, Inches(0.5), Inches(y_pos), Inches(3.5), Inches(0.4),
                    section, font_size=12, bold=True, color=COLORS['dark_text'])
        add_text_box(slide, Inches(4.2), Inches(y_pos), Inches(1.2), Inches(0.4),
                    pts, font_size=12, color=COLORS['cyan_end'])
        add_text_box(slide, Inches(5.5), Inches(y_pos), Inches(4.0), Inches(0.4),
                    desc, font_size=11, color=COLORS['gray_text'])
        y_pos += 0.5

    # Question types
    add_colored_shape(slide, Inches(0.2), Inches(3.7), Inches(9.6), Inches(0.85), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.8), Inches(9.2), Inches(0.25),
                "QUESTION TYPES:", font_size=12, bold=True, color=COLORS['orange_end'])
    add_text_box(slide, Inches(0.4), Inches(4.1), Inches(9.2), Inches(0.4),
                "Multiple choice • Data analysis • Short answer • Extended response (use rubric criteria)",
                font_size=11, color=COLORS['dark_text'])


def add_part2_support_slide(prs):
    """Slide 7: Part 2 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['cyan_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Part 2 – Key Concepts to Review",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Section A
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(4.7), Inches(1.45), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(4.4), Inches(0.25),
                "SECTION A: Air Masses", font_size=11, bold=True, color=COLORS['blue_accent'])
    a_content = """• mT: warm/wet (Gulf of Mexico)
• cT: warm/dry (SW desert)
• mP: cold/wet (N. Pacific/Atlantic)
• cP: cold/dry (Canada/Arctic)
• Fronts form where masses meet"""
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(4.4), Inches(1.0),
                a_content, font_size=10, color=COLORS['dark_text'])

    # Section B
    add_colored_shape(slide, Inches(5.1), Inches(0.85), Inches(4.6), Inches(1.45), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(5.3), Inches(0.95), Inches(4.2), Inches(0.25),
                "SECTION B: Weather vs Climate", font_size=11, bold=True, color=COLORS['purple_accent'])
    b_content = """• Weather: short-term, chaotic
• Climate: long-term average, predictable
• Chaos: small errors → big changes
• Variability averages out over time
• Trends emerge from many events"""
    add_text_box(slide, Inches(5.3), Inches(1.25), Inches(4.2), Inches(1.0),
                b_content, font_size=10, color=COLORS['dark_text'])

    # Section C
    add_colored_shape(slide, Inches(0.2), Inches(2.45), Inches(4.7), Inches(1.3), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.55), Inches(4.4), Inches(0.25),
                "SECTION C: Data Analysis", font_size=11, bold=True, color=COLORS['green_accent'])
    c_content = """• Error = |Forecast - Actual|
• Error grows with forecast time
• 1-day: ~2°F error | 5-day: ~8°F
• Calculate averages from data
• Identify patterns in variability"""
    add_text_box(slide, Inches(0.4), Inches(2.85), Inches(4.4), Inches(0.85),
                c_content, font_size=10, color=COLORS['dark_text'])

    # Section D
    add_colored_shape(slide, Inches(5.1), Inches(2.45), Inches(4.6), Inches(1.3), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(5.3), Inches(2.55), Inches(4.2), Inches(0.25),
                "SECTION D: System Design", font_size=11, bold=True, color=COLORS['orange_end'])
    d_content = """• Weather: frequent updates, many sensors
• Climate: precision over decades
• Trade-offs: cost vs. coverage
• Justify your choices with evidence
• Consider equipment lifespan"""
    add_text_box(slide, Inches(5.3), Inches(2.85), Inches(4.2), Inches(0.85),
                d_content, font_size=10, color=COLORS['dark_text'])

    # Final reminder
    add_colored_shape(slide, Inches(0.2), Inches(3.9), Inches(9.6), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(3.93), Inches(9.2), Inches(0.5),
                "REMEMBER: Use vocabulary precisely! mT ≠ mP | Weather ≠ Climate | Variability ≠ Trend",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_part3_intro_slide(prs):
    """Slide 8: Part 3 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['purple_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Part 3 – Misconception Check",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | ~20 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Purpose
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(1.2), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(1.2), COLORS['purple_end'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "PURPOSE:", font_size=14, bold=True, color=COLORS['purple_end'])
    purpose = """This section helps us identify common misunderstandings so we can provide targeted feedback.
Your answers help us improve instruction for everyone. Answer honestly based on your current understanding—
this is about learning, not just getting points!"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(0.7),
                purpose, font_size=12, color=COLORS['dark_text'])

    # Targeted misconceptions
    add_text_box(slide, Inches(0.3), Inches(2.5), Inches(9.4), Inches(0.3),
                "COMMON MISCONCEPTIONS WE'RE CHECKING:", font_size=13, bold=True, color=COLORS['purple_end'])

    misconceptions = [
        ("weather-climate-same", "Weather and climate are the same thing", "Weather = days | Climate = decades"),
        ("cold-air-rises", "Cold air rises, warm air sinks", "WRONG! Warm rises (less dense), cold sinks"),
        ("seasons-distance", "Seasons caused by Earth-Sun distance", "Actually caused by axial TILT"),
    ]

    y_pos = 2.9
    for id, wrong, correct in misconceptions:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.5), COLORS['light_pink_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.02), Inches(4.5), Inches(0.45),
                    f"❌ {wrong}", font_size=10, color=COLORS['red_accent'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(5.1), Inches(y_pos + 0.02), Inches(4.5), Inches(0.45),
                    f"✓ {correct}", font_size=10, color=COLORS['green_accent'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.55


def add_part3_support_slide(prs):
    """Slide 9: Part 3 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['purple_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Part 3 – Think Through These Carefully",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Weather vs Climate
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(4.7), Inches(1.5), COLORS['light_cyan_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(4.4), Inches(0.3),
                "WEATHER vs CLIMATE", font_size=12, bold=True, color=COLORS['cyan_end'])
    wc = """Weather: What's happening NOW
• Today's temperature, rain, wind
• Changes hour to hour, day to day

Climate: What USUALLY happens
• 30-year average conditions
• Changes over decades"""
    add_text_box(slide, Inches(0.4), Inches(1.3), Inches(4.4), Inches(1.0),
                wc, font_size=10, color=COLORS['dark_text'])

    # Density
    add_colored_shape(slide, Inches(5.1), Inches(0.85), Inches(4.6), Inches(1.5), COLORS['light_green_bg'])
    add_text_box(slide, Inches(5.3), Inches(0.95), Inches(4.2), Inches(0.3),
                "AIR DENSITY & MOVEMENT", font_size=12, bold=True, color=COLORS['green_end'])
    density = """WARM air:
• Molecules spread out = LESS dense
• Less dense = RISES (like a balloon)

COLD air:
• Molecules packed tight = MORE dense
• More dense = SINKS"""
    add_text_box(slide, Inches(5.3), Inches(1.3), Inches(4.2), Inches(1.0),
                density, font_size=10, color=COLORS['dark_text'])

    # Seasons
    add_colored_shape(slide, Inches(0.2), Inches(2.5), Inches(9.6), Inches(1.3), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.6), Inches(9.2), Inches(0.3),
                "SEASONS: TILT, NOT DISTANCE", font_size=12, bold=True, color=COLORS['orange_end'])
    seasons = """Earth is actually CLOSER to the sun in Northern Hemisphere winter! (January)
Seasons are caused by Earth's 23.5° TILT:
• Summer: Your hemisphere tilts TOWARD sun → direct sunlight → warm
• Winter: Your hemisphere tilts AWAY → indirect sunlight → cold
Distance only varies by ~3% and has minimal effect on temperatures."""
    add_text_box(slide, Inches(0.4), Inches(2.95), Inches(9.2), Inches(0.8),
                seasons, font_size=10, color=COLORS['dark_text'])

    # Final note
    add_colored_shape(slide, Inches(0.2), Inches(3.95), Inches(9.6), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(3.98), Inches(9.2), Inches(0.5),
                "Take your time—read each question carefully and think about the CORRECT explanation!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_tips_slide(prs):
    """Slide 10: Tips for Success"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Tips for Success",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Tips
    tips = [
        ("📝 USE YOUR NOTECARD", "You prepared it for this! Check definitions and key concepts."),
        ("⏰ MANAGE YOUR TIME", "Part 1: 15 min | Part 2: 40 min | Part 3: 20 min"),
        ("📖 READ CAREFULLY", "Look for key words: weather vs climate, cause vs effect"),
        ("🔗 MAKE CONNECTIONS", "Link air masses → chaos → prediction → climate trends"),
        ("✓ CHECK YOUR WORK", "Revisit uncertain answers if time permits"),
    ]

    y_pos = 0.9
    colors = [COLORS['purple_start'], COLORS['cyan_start'], COLORS['orange_start'],
              COLORS['green_start'], COLORS['blue_accent']]
    for i, (tip, detail) in enumerate(tips):
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.65), COLORS['light_gray_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.65), colors[i])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(3.5), Inches(0.55),
                    tip, font_size=12, bold=True, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(4.0), Inches(y_pos + 0.05), Inches(5.6), Inches(0.55),
                    detail, font_size=11, color=COLORS['gray_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.7

    # Good luck
    add_colored_shape(slide, Inches(0.2), Inches(4.45), Inches(9.6), Inches(0.55), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.4), Inches(4.48), Inches(9.2), Inches(0.5),
                "You've got this! Show what you've learned about weather, climate, and prediction!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_summary_slide(prs):
    """Slide 11: Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Cycle 5 Summary: What You Learned",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Key takeaways
    takeaways = [
        ("Air Masses", "4 types based on source region → determine local weather"),
        ("Fronts", "Where air masses meet → weather changes occur"),
        ("Chaos", "Small errors grow → forecasts fail beyond ~7 days"),
        ("Climate", "Long-term trends emerge from averaging → predictable!"),
    ]

    y_pos = 0.9
    for title, detail in takeaways:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.65), COLORS['light_cyan_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.65), COLORS['teal'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.0), Inches(0.55),
                    title, font_size=13, bold=True, color=COLORS['teal_dark'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.5), Inches(y_pos + 0.05), Inches(7.1), Inches(0.55),
                    detail, font_size=12, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.72

    # Next cycle preview
    add_colored_shape(slide, Inches(0.2), Inches(3.85), Inches(9.6), Inches(0.65), COLORS['header_blue_end'])
    add_text_box(slide, Inches(0.4), Inches(3.88), Inches(9.2), Inches(0.6),
                "COMING UP in Cycle 6: Plate Tectonics & Earth's Interior!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Final takeaway
    add_colored_shape(slide, Inches(0.2), Inches(4.65), Inches(9.6), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(4.68), Inches(9.2), Inches(0.5),
                "KEY INSIGHT: Chaos makes weather unpredictable, but patterns make climate knowable!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


if __name__ == "__main__":
    prs = create_presentation()
    output_path = "/home/user/Kairos.Sci.Repo/content/grade7/cycle05/week3/G7_C5_W3_Assessment_Slides_Final.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Created: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
