#!/usr/bin/env python3
"""
Create G8_C3_W2 Evidence of Evolution presentation.

See PPTX_DESIGN_GUIDE.md for best practices documentation.
"""

import os
from pptx.dml.color import RGBColor
from pptx_common import (
    COLORS as BASE_COLORS,
    create_base_presentation,
    add_colored_shape,
    add_text_box,
    Inches,
    Pt,
    PP_ALIGN,
    MSO_ANCHOR,
)

# G8 Evolution theme - extend base colors
COLORS = {**BASE_COLORS}
COLORS.update({
    'header_purple_start': RGBColor(0x9F, 0x7A, 0xEA),  # #9F7AEA
    'header_purple_end': RGBColor(0x6B, 0x46, 0xC1),    # #6B46C1
    'hook_purple_start': RGBColor(0x66, 0x7E, 0xEA),    # #667EEA
    'hook_purple_end': RGBColor(0x76, 0x4B, 0xA2),      # #764BA2
    'blue_start': RGBColor(0x42, 0x99, 0xE1),           # #4299E1
    'blue_end': RGBColor(0x2B, 0x6C, 0xB0),             # #2B6CB0
    'light_pink_bg': RGBColor(0xFE, 0xD7, 0xD7),        # #FED7D7
})


def create_presentation():
    """Create the G8_C3_W2 presentation"""
    prs = create_base_presentation()

    add_title_slide(prs)
    add_phenomenon_slide(prs)
    add_driving_question_slide(prs)
    add_prior_knowledge_slide(prs)
    add_learning_targets_slide(prs)
    add_vocabulary_slide(prs)
    add_hook_intro_slide(prs)
    add_hook_support_slide(prs)
    add_station1_intro_slide(prs)
    add_station1_support_slide(prs)
    add_station2_intro_slide(prs)
    add_station2_support_slide(prs)
    add_station3_intro_slide(prs)
    add_station3_support_slide(prs)
    add_exit_ticket_slide(prs)
    add_summary_slide(prs)

    return prs


def add_title_slide(prs):
    """Slide 1: Title slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Purple background
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625),
                      COLORS['header_purple_end'])

    # Main title
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(1.2),
                "🐋 Week 2: Evidence of Evolution 🧬",
                font_size=42, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")

    # Subtitle line 1
    add_text_box(slide, Inches(0.5), Inches(2.9), Inches(9), Inches(0.5),
                "Grade 8 Science | Rosche | Kairos Academies",
                font_size=20, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Subtitle line 2
    add_text_box(slide, Inches(0.5), Inches(3.4), Inches(9), Inches(0.5),
                "MS-LS4-2 & MS-LS4-4 Evidence of Evolution | 100 Points Total | ~75 Minutes",
                font_size=16, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Driving question teaser box (simplified - no border)
    teaser_bg = add_colored_shape(slide, Inches(2), Inches(4.2), Inches(6), Inches(0.8),
                                  COLORS['white'])

    add_text_box(slide, Inches(2.1), Inches(4.35), Inches(5.8), Inches(0.5),
                "Why do whale flippers have finger bones?",
                font_size=16, color=COLORS['teal_dark'], align=PP_ALIGN.CENTER)


def add_phenomenon_slide(prs):
    """Slide 2: The Phenomenon - The Whale Finger Mystery"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    add_colored_shape(slide, Inches(0.2), Inches(0.2), Inches(9.6), Inches(0.6), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.4), Inches(0.28), Inches(9.2), Inches(0.5),
                "🐋 The Phenomenon: The Whale Finger Mystery",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia")

    # Left content area
    left_bg = add_colored_shape(slide, Inches(0.2), Inches(1.0), Inches(4.8), Inches(2.8),
                                COLORS['light_teal_bg'])

    phenom_text = """Look at an X-ray of a whale's flipper. Inside that smooth flipper, there are BONES:

• Upper arm bone (humerus)
• Two forearm bones (radius and ulna)
• Wrist bones (carpals)
• FIVE FINGERS with knuckles!

But whales don't have hands. They swim."""

    add_text_box(slide, Inches(0.4), Inches(1.1), Inches(4.4), Inches(2.5),
                phenom_text, font_size=13, color=COLORS['teal_dark'])

    # Mystery question box
    add_colored_shape(slide, Inches(0.2), Inches(3.9), Inches(4.8), Inches(0.7),
                      COLORS['teal_dark'])
    add_text_box(slide, Inches(0.4), Inches(4.0), Inches(4.4), Inches(0.5),
                "So why would a whale have finger bones?",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Right side - Quick Write box
    add_colored_shape(slide, Inches(5.2), Inches(1.0), Inches(4.6), Inches(3.0),
                      COLORS['hook_purple_start'])

    add_text_box(slide, Inches(5.4), Inches(1.2), Inches(4.2), Inches(0.4),
                "📝 QUICK WRITE", font_size=18, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(5.4), Inches(1.8), Inches(4.2), Inches(1.0),
                "30 sec – silent writing", font_size=28, bold=True,
                color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(5.4), Inches(3.0), Inches(4.2), Inches(0.8),
                "What's your prediction?\nWhy might whales have finger bones?",
                font_size=15, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Turn & Talk footer
    add_colored_shape(slide, Inches(5.2), Inches(4.1), Inches(4.6), Inches(0.5),
                      COLORS['light_purple_bg'])
    add_text_box(slide, Inches(5.4), Inches(4.15), Inches(4.2), Inches(0.4),
                "👥 TURN & TALK: Share with partner",
                font_size=11, bold=True, color=COLORS['hook_purple_start'])

    # Notecard prompt
    add_colored_shape(slide, Inches(0.2), Inches(4.8), Inches(9.6), Inches(0.6), COLORS['hook_purple_start'])
    add_text_box(slide, Inches(0.4), Inches(4.88), Inches(9.2), Inches(0.4),
                "📝 Notecard: Write your prediction - Why do whales have finger bones?",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_driving_question_slide(prs):
    """Slide 3: Driving Question"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Purple background
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['header_purple_end'])

    # Title
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                "DRIVING QUESTION", font_size=40, bold=True,
                color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Main question
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(1.5),
                "Why do whale flippers have finger bones?\nWhat does this tell us about whale ancestors?",
                font_size=32, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")


    # Sub-question
    add_text_box(slide, Inches(1), Inches(3.5), Inches(8), Inches(0.6),
                "And how does this connect to natural selection from Week 1?",
                font_size=20, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Mission box (simplified - no border)
    mission_bg = add_colored_shape(slide, Inches(2), Inches(4.3), Inches(6), Inches(0.8), COLORS['teal_dark'])

    add_text_box(slide, Inches(2.2), Inches(4.4), Inches(5.6), Inches(0.6),
                "🎯 Your Mission: Use anatomical & fossil evidence to explain evolution",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_prior_knowledge_slide(prs):
    """Slide 4: What You Already Know (Week 1 Natural Selection)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title - reduced font to prevent wrapping
    add_text_box(slide, Inches(0.3), Inches(0.1), Inches(9.4), Inches(0.5),
                "🔗 What You Already Know (Week 1 Natural Selection)",
                font_size=24, bold=True, color=COLORS['purple_accent'],
                align=PP_ALIGN.CENTER, font_name="Georgia")

    # Four concept boxes in 2x2 grid
    # Top left - Variation
    add_colored_shape(slide, Inches(0.3), Inches(0.7), Inches(4.6), Inches(1.4), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(0.7), Inches(0.08), Inches(1.4), COLORS['header_purple_start'])
    add_text_box(slide, Inches(0.5), Inches(0.8), Inches(4.3), Inches(0.3),
                "1. Variation exists...", font_size=14, bold=True, color=COLORS['header_purple_start'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(4.3), Inches(0.4),
                "in every population", font_size=18, bold=True, color=COLORS['dark_text'])
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4.3), Inches(0.4),
                "→ Different traits = different genes", font_size=13, color=COLORS['gray_text'])

    # Top right - Selection
    add_colored_shape(slide, Inches(5.1), Inches(0.7), Inches(4.6), Inches(1.4), COLORS['light_orange_bg'])
    add_colored_shape(slide, Inches(5.1), Inches(0.7), Inches(0.08), Inches(1.4), COLORS['orange_end'])
    add_text_box(slide, Inches(5.3), Inches(0.8), Inches(4.3), Inches(0.3),
                "2. Selection pressure...", font_size=14, bold=True, color=COLORS['orange_end'])
    add_text_box(slide, Inches(5.3), Inches(1.1), Inches(4.3), Inches(0.4),
                "favors some traits", font_size=18, bold=True, color=COLORS['dark_text'])
    add_text_box(slide, Inches(5.3), Inches(1.5), Inches(4.3), Inches(0.4),
                "→ Environment determines fitness", font_size=13, color=COLORS['gray_text'])

    # Bottom left - Survival
    add_colored_shape(slide, Inches(0.3), Inches(2.2), Inches(4.6), Inches(1.4), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(2.2), Inches(0.08), Inches(1.4), COLORS['green_accent'])
    add_text_box(slide, Inches(0.5), Inches(2.3), Inches(4.3), Inches(0.3),
                "3. Survivors reproduce...", font_size=14, bold=True, color=COLORS['green_accent'])
    add_text_box(slide, Inches(0.5), Inches(2.6), Inches(4.3), Inches(0.4),
                "passing on their genes", font_size=18, bold=True, color=COLORS['dark_text'])
    add_text_box(slide, Inches(0.5), Inches(3.0), Inches(4.3), Inches(0.4),
                "→ Offspring inherit advantageous traits", font_size=13, color=COLORS['gray_text'])

    # Bottom right - Population change
    add_colored_shape(slide, Inches(5.1), Inches(2.2), Inches(4.6), Inches(1.4), COLORS['light_blue_bg'])
    add_colored_shape(slide, Inches(5.1), Inches(2.2), Inches(0.08), Inches(1.4), COLORS['blue_accent'])
    add_text_box(slide, Inches(5.3), Inches(2.3), Inches(4.3), Inches(0.3),
                "4. Population changes...", font_size=14, bold=True, color=COLORS['blue_accent'])
    add_text_box(slide, Inches(5.3), Inches(2.6), Inches(4.3), Inches(0.4),
                "over generations", font_size=18, bold=True, color=COLORS['dark_text'])
    add_text_box(slide, Inches(5.3), Inches(3.0), Inches(4.3), Inches(0.4),
                "→ Evolution = population change over time", font_size=13, color=COLORS['gray_text'])

    # Notecard prompt
    add_colored_shape(slide, Inches(0.3), Inches(3.8), Inches(9.4), Inches(0.7), COLORS['hook_purple_start'])
    add_text_box(slide, Inches(0.5), Inches(3.9), Inches(9.0), Inches(0.5),
                "📝 Notecard: Which step of natural selection do you remember BEST?",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # NEW THIS WEEK preview
    add_colored_shape(slide, Inches(0.3), Inches(4.6), Inches(9.4), Inches(0.8), COLORS['light_teal_bg'])
    add_text_box(slide, Inches(0.5), Inches(4.7), Inches(9.0), Inches(0.6),
                "🆕 NEW THIS WEEK: What EVIDENCE do we have that evolution actually happened?",
                font_size=14, bold=True, color=COLORS['teal_dark'], align=PP_ALIGN.CENTER)


def add_learning_targets_slide(prs):
    """Slide 5: Learning Targets"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.3), Inches(0.1), Inches(9.4), Inches(0.5),
                "🎯 Learning Targets", font_size=32, bold=True,
                color=COLORS['teal_dark'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Target 1
    add_colored_shape(slide, Inches(0.3), Inches(0.7), Inches(4.6), Inches(1.3), COLORS['white'])
    add_text_box(slide, Inches(0.5), Inches(0.8), Inches(4.2), Inches(0.3),
                "TARGET 1", font_size=12, bold=True, color=COLORS['header_purple_start'])
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(4.2), Inches(0.7),
                "Distinguish between homologous and analogous structures",
                font_size=14, bold=True, color=COLORS['dark_text'])

    # Target 2
    add_colored_shape(slide, Inches(5.1), Inches(0.7), Inches(4.6), Inches(1.3), COLORS['white'])
    add_text_box(slide, Inches(5.3), Inches(0.8), Inches(4.2), Inches(0.3),
                "TARGET 2", font_size=12, bold=True, color=COLORS['header_purple_start'])
    add_text_box(slide, Inches(5.3), Inches(1.1), Inches(4.2), Inches(0.7),
                "Use anatomical evidence to infer evolutionary relationships",
                font_size=14, bold=True, color=COLORS['dark_text'])

    # Target 3
    add_colored_shape(slide, Inches(0.3), Inches(2.1), Inches(4.6), Inches(1.3), COLORS['white'])
    add_text_box(slide, Inches(0.5), Inches(2.2), Inches(4.2), Inches(0.3),
                "TARGET 3", font_size=12, bold=True, color=COLORS['header_purple_end'])
    add_text_box(slide, Inches(0.5), Inches(2.5), Inches(4.2), Inches(0.7),
                "Interpret fossil evidence including transitional forms",
                font_size=14, bold=True, color=COLORS['dark_text'])

    # Target 4
    add_colored_shape(slide, Inches(5.1), Inches(2.1), Inches(4.6), Inches(1.3), COLORS['white'])
    add_text_box(slide, Inches(5.3), Inches(2.2), Inches(4.2), Inches(0.3),
                "TARGET 4", font_size=12, bold=True, color=COLORS['header_purple_end'])
    add_text_box(slide, Inches(5.3), Inches(2.5), Inches(4.2), Inches(0.7),
                "Predict features of transitional organisms based on evidence",
                font_size=14, bold=True, color=COLORS['dark_text'])

    # Summary bar - FIXED: Added vertical centering for proper alignment
    add_colored_shape(slide, Inches(0.3), Inches(3.6), Inches(9.4), Inches(0.7), COLORS['teal_dark'])
    add_text_box(slide, Inches(1.5), Inches(3.65), Inches(1.5), Inches(0.35),
                "100", font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(1.5), Inches(4.0), Inches(1.5), Inches(0.25),
                "Total Points", font_size=9, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(4.25), Inches(3.65), Inches(1.5), Inches(0.35),
                "~75", font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(4.25), Inches(4.0), Inches(1.5), Inches(0.25),
                "Minutes", font_size=9, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(7.0), Inches(3.65), Inches(1.5), Inches(0.35),
                "5", font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(7.0), Inches(4.0), Inches(1.5), Inches(0.25),
                "Sections", font_size=9, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Points breakdown - FIXED: Added vertical centering
    add_colored_shape(slide, Inches(0.3), Inches(4.5), Inches(9.4), Inches(0.9), COLORS['white'])
    add_text_box(slide, Inches(0.5), Inches(4.55), Inches(9.0), Inches(0.8),
                "🎯 Hook (12 pts)  →  🦴 Station 1 (20 pts)  →  🦕 Station 2 (20 pts)  →  🔧 Station 3 (25 pts)  →  🎓 Exit (23 pts)",
                font_size=13, color=COLORS['dark_text'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_vocabulary_slide(prs):
    """Slide 6: Key Vocabulary"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    add_colored_shape(slide, Inches(0.2), Inches(0.15), Inches(9.6), Inches(0.5), COLORS['blue_accent'])
    add_text_box(slide, Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.4),
                "📚 Key Vocabulary This Week", font_size=22, bold=True,
                color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Vocabulary table
    vocab = [
        ("Homologous Structures", "Same bone structure, different function (evidence of common ancestor)"),
        ("Analogous Structures", "Different structure, same function (NOT from common ancestor)"),
        ("Transitional Fossil", "Fossil with features of both ancestral and descendant groups"),
        ("Vestigial Structure", "Reduced or non-functional structure left over from ancestors"),
        ("Common Ancestor", "An ancestral species from which multiple species evolved"),
    ]

    # FIXED: Added vertical centering for proper text alignment in rows
    y_pos = 0.8
    for i, (term, definition) in enumerate(vocab):
        bg_color = COLORS['light_blue_bg'] if i % 2 == 0 else COLORS['white']
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.65), bg_color)
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.08), Inches(2.4), Inches(0.5),
                    term, font_size=13, bold=True, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.9), Inches(y_pos + 0.08), Inches(6.7), Inches(0.5),
                    definition, font_size=11, color=COLORS['gray_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.65

    # Notecard prompt
    add_colored_shape(slide, Inches(0.2), Inches(4.15), Inches(9.6), Inches(0.6), COLORS['hook_purple_start'])
    add_text_box(slide, Inches(0.4), Inches(4.25), Inches(9.2), Inches(0.4),
                "📝 Notecard: Write down the term you find most confusing - we'll clarify today!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Quick reference
    add_colored_shape(slide, Inches(0.2), Inches(4.85), Inches(9.6), Inches(0.5), COLORS['light_teal_bg'])
    add_text_box(slide, Inches(0.4), Inches(4.9), Inches(9.2), Inches(0.4),
                "💡 KEY: Same STRUCTURE = common ancestor | Same FUNCTION only = evolved independently",
                font_size=12, color=COLORS['teal_dark'], align=PP_ALIGN.CENTER)


def add_hook_intro_slide(prs):
    """Slide 7: Hook Introduction"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header - two-line pattern to prevent overlap
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['hook_purple_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🎯 Hook – The Whale Finger Mystery",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "12 Points | ~10 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Left side - What You'll Do
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "WHAT YOU'LL DO", font_size=14, bold=True, color=COLORS['hook_purple_start'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.0), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.0), COLORS['hook_purple_start'])

    steps = """1. Observe the whale flipper bone structure (2 min)
2. Compare to human arm bones (2 min)
3. Identify which other animals share this pattern (3 min)
4. Connect to natural selection from Week 1 (3 min)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(1.8),
                steps, font_size=13, color=COLORS['dark_text'])

    # Key insight box
    add_colored_shape(slide, Inches(0.3), Inches(3.6), Inches(4.5), Inches(0.8), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(3.6), Inches(4.5), Inches(0.05), COLORS['green_accent'])
    add_text_box(slide, Inches(0.5), Inches(3.7), Inches(4.1), Inches(0.6),
                "💡 Key: Same bones, different functions = evidence of common ancestor!",
                font_size=12, color=COLORS['green_accent'])

    # Right side - Quick Write prompt
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(2.5), COLORS['hook_purple_start'])
    add_text_box(slide, Inches(5.2), Inches(1.25), Inches(4.3), Inches(0.4),
                "📝 QUICK WRITE FIRST", font_size=16, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.2), Inches(1.8), Inches(4.3), Inches(1.5),
                "\"Whales have finger bones because...\"",
                font_size=22, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Cold call warning
    add_colored_shape(slide, Inches(5.0), Inches(3.7), Inches(4.7), Inches(0.4), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(5.2), Inches(3.75), Inches(4.3), Inches(0.3),
                "🎲 Cold Call Coming! Be ready to share.",
                font_size=11, bold=True, color=COLORS['hook_purple_start'])

    # Notecard prompt
    add_colored_shape(slide, Inches(0.15), Inches(4.55), Inches(9.7), Inches(0.55), COLORS['hook_purple_start'])
    add_text_box(slide, Inches(0.35), Inches(4.6), Inches(9.3), Inches(0.4),
                "📝 Notecard: Record your prediction AND your reasoning",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_hook_support_slide(prs):
    """Slide 8: Hook Support - Bone Pattern Comparison"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['hook_purple_end'])
    add_text_box(slide, Inches(0.35), Inches(0.22), Inches(9.3), Inches(0.45),
                "🎯 Hook – The \"1-2-Many-5\" Pattern",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia")

    # Bone pattern explanation
    add_colored_shape(slide, Inches(0.25), Inches(0.95), Inches(9.5), Inches(1.5), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.25), Inches(0.95), Inches(0.08), Inches(1.5), COLORS['hook_purple_start'])
    add_text_box(slide, Inches(0.5), Inches(1.05), Inches(9.0), Inches(0.3),
                "🦴 The Same Bone Pattern in ALL These Animals:", font_size=14, bold=True, color=COLORS['hook_purple_start'])

    pattern_text = """• Humerus – 1 upper arm/shoulder bone
• Radius & Ulna – 2 forearm bones
• Carpals – Multiple small wrist bones
• Metacarpals & Phalanges – Hand/finger bones (5 digits)"""
    add_text_box(slide, Inches(0.5), Inches(1.4), Inches(4.5), Inches(0.95),
                pattern_text, font_size=12, color=COLORS['dark_text'])

    # Animals with this pattern
    add_text_box(slide, Inches(5.5), Inches(1.4), Inches(4.0), Inches(0.95),
                "Human arm ✓\nWhale flipper ✓\nBat wing ✓\nCat leg ✓\nDog leg ✓",
                font_size=12, bold=True, color=COLORS['dark_text'])

    # Key insight
    add_colored_shape(slide, Inches(0.25), Inches(2.6), Inches(9.5), Inches(0.7), COLORS['green_accent'])
    add_text_box(slide, Inches(0.5), Inches(2.7), Inches(9.0), Inches(0.5),
                "🔑 KEY: Same bones = HOMOLOGOUS = inherited from common ancestor!",
                font_size=16, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Contrast - Analogous
    add_colored_shape(slide, Inches(0.25), Inches(3.45), Inches(9.5), Inches(0.9), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.55), Inches(9.0), Inches(0.25),
                "⚠️ But bird wings and butterfly wings are ANALOGOUS:", font_size=12, bold=True, color=COLORS['orange_end'])
    add_text_box(slide, Inches(0.5), Inches(3.85), Inches(9.0), Inches(0.4),
                "Same function (flying) but DIFFERENT structure = evolved independently, NOT from common ancestor",
                font_size=12, color=COLORS['dark_text'])

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.5), Inches(9.7), Inches(0.55), COLORS['hook_purple_start'])
    add_text_box(slide, Inches(0.35), Inches(4.55), Inches(9.3), Inches(0.4),
                "📝 Notecard: Why is same STRUCTURE more important than same FUNCTION for proving ancestry?",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_station1_intro_slide(prs):
    """Slide 9: Station 1 Introduction - Comparative Anatomy"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header - two-line pattern
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['orange_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🦴 Station 1 – Comparative Anatomy Analysis",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "20 Points | ~18 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Left side - Analysis steps
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "YOUR ANALYSIS", font_size=14, bold=True, color=COLORS['orange_end'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.0), COLORS['light_orange_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.0), COLORS['orange_end'])

    steps = """1. Compare bone structures in 5 animals (5 min)
2. Identify homologous vs analogous structures (5 min)
3. Determine which animals share ancestors (4 min)
4. Explain WHY same bones = common ancestry (4 min)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(1.8),
                steps, font_size=13, color=COLORS['dark_text'])

    # Right side - Discovery Goal
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(1.5), COLORS['orange_end'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "🎯 DISCOVERY GOAL", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.2), Inches(1.6), Inches(4.3), Inches(0.9),
                "What does shared bone structure tell us about animal ancestry?",
                font_size=15, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Key distinction box
    add_colored_shape(slide, Inches(5.0), Inches(2.7), Inches(4.7), Inches(1.0), COLORS['light_orange_bg'])
    add_colored_shape(slide, Inches(5.0), Inches(2.7), Inches(4.7), Inches(0.05), COLORS['orange_end'])
    add_text_box(slide, Inches(5.2), Inches(2.8), Inches(4.3), Inches(0.8),
                "🔑 Remember:\nHomologous = same structure\nAnalogous = same function",
                font_size=12, color=COLORS['orange_end'])

    # Key question
    add_colored_shape(slide, Inches(0.3), Inches(3.65), Inches(9.4), Inches(0.7), COLORS['light_teal_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.75), Inches(9.0), Inches(0.5),
                "💡 Why would a whale, bat, and human all have the SAME bones arranged the same way?",
                font_size=13, color=COLORS['teal_dark'], align=PP_ALIGN.CENTER)

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.55), Inches(9.7), Inches(0.55), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(4.6), Inches(9.3), Inches(0.4),
                "📝 Notecard: Write your prediction BEFORE analyzing the diagrams",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_station1_support_slide(prs):
    """Slide 10: Station 1 Support - Homologous vs Analogous Table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.4),
                "🦴 Station 1 – Homologous vs Analogous Reference",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia")

    # Comparison table header - FIXED: Added vertical centering
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(9.6), Inches(0.5), COLORS['orange_end'])
    add_text_box(slide, Inches(0.4), Inches(0.88), Inches(2.2), Inches(0.45),
                "Type", font_size=13, bold=True, color=COLORS['white'],
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(2.7), Inches(0.88), Inches(2.0), Inches(0.45),
                "Structure", font_size=13, bold=True, color=COLORS['white'],
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(4.8), Inches(0.88), Inches(2.0), Inches(0.45),
                "Function", font_size=13, bold=True, color=COLORS['white'],
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(6.9), Inches(0.88), Inches(2.7), Inches(0.45),
                "Indicates", font_size=13, bold=True, color=COLORS['white'],
                anchor=MSO_ANCHOR.MIDDLE)

    # Row 1 - Homologous - FIXED: Added vertical centering
    add_colored_shape(slide, Inches(0.2), Inches(1.4), Inches(9.6), Inches(0.55), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(1.43), Inches(2.2), Inches(0.48),
                "Homologous", font_size=13, bold=True, color=COLORS['dark_text'],
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(2.7), Inches(1.43), Inches(2.0), Inches(0.48),
                "SAME bones", font_size=12, color=COLORS['dark_text'],
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(4.8), Inches(1.43), Inches(2.0), Inches(0.48),
                "Different", font_size=12, color=COLORS['dark_text'],
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(6.9), Inches(1.43), Inches(2.7), Inches(0.48),
                "Common ancestor ✓", font_size=12, bold=True, color=COLORS['green_accent'],
                anchor=MSO_ANCHOR.MIDDLE)

    # Row 2 - Analogous - FIXED: Added vertical centering
    add_colored_shape(slide, Inches(0.2), Inches(2.0), Inches(9.6), Inches(0.55), COLORS['light_pink_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.03), Inches(2.2), Inches(0.48),
                "Analogous", font_size=13, bold=True, color=COLORS['dark_text'],
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(2.7), Inches(2.03), Inches(2.0), Inches(0.48),
                "Different", font_size=12, color=COLORS['dark_text'],
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(4.8), Inches(2.03), Inches(2.0), Inches(0.48),
                "SAME", font_size=12, color=COLORS['dark_text'],
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(6.9), Inches(2.03), Inches(2.7), Inches(0.48),
                "Independent evolution ✗", font_size=12, bold=True, color=COLORS['red_accent'],
                anchor=MSO_ANCHOR.MIDDLE)

    # Examples
    add_colored_shape(slide, Inches(0.2), Inches(2.7), Inches(4.7), Inches(1.0), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.8), Inches(4.4), Inches(0.25),
                "✓ HOMOLOGOUS Examples:", font_size=12, bold=True, color=COLORS['green_accent'])
    add_text_box(slide, Inches(0.4), Inches(3.1), Inches(4.4), Inches(0.55),
                "• Whale flipper + bat wing + human arm\n• Dog leg + cat leg + horse leg",
                font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.1), Inches(2.7), Inches(4.7), Inches(1.0), COLORS['light_pink_bg'])
    add_text_box(slide, Inches(5.3), Inches(2.8), Inches(4.4), Inches(0.25),
                "✗ ANALOGOUS Examples:", font_size=12, bold=True, color=COLORS['red_accent'])
    add_text_box(slide, Inches(5.3), Inches(3.1), Inches(4.4), Inches(0.55),
                "• Bird wing + butterfly wing (both fly)\n• Shark fin + dolphin fin (both swim)",
                font_size=11, color=COLORS['dark_text'])

    # Sentence starters
    add_colored_shape(slide, Inches(0.2), Inches(3.85), Inches(9.6), Inches(0.75), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.95), Inches(9.2), Inches(0.25),
                "📗 SENTENCE STARTERS:", font_size=11, bold=True, color=COLORS['blue_accent'])
    add_text_box(slide, Inches(0.4), Inches(4.2), Inches(9.2), Inches(0.35),
                "\"These structures are homologous because they have the same _____ but different _____.\"",
                font_size=10, color=COLORS['dark_text'])

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.75), Inches(9.7), Inches(0.55), COLORS['orange_start'])
    add_text_box(slide, Inches(0.35), Inches(4.8), Inches(9.3), Inches(0.4),
                "📝 Notecard: Is a dolphin fin homologous or analogous to a shark fin? Explain WHY.",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_station2_intro_slide(prs):
    """Slide 11: Station 2 Introduction - Fossil Record"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header - two-line pattern
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['blue_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🦕 Station 2 – Fossil Record Investigation",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "20 Points | ~15 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Left side - Investigation steps
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "YOUR INVESTIGATION", font_size=14, bold=True, color=COLORS['blue_end'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.0), COLORS['light_blue_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.0), COLORS['blue_end'])

    steps = """1. Study the whale evolution timeline (4 min)
2. Identify transitional features in each fossil (4 min)
3. Explain what changed from Pakicetus to modern whales (4 min)
4. Connect to natural selection (3 min)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(1.8),
                steps, font_size=13, color=COLORS['dark_text'])

    # Right side - Discovery Goal
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(1.5), COLORS['blue_end'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "🎯 DISCOVERY GOAL", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.2), Inches(1.6), Inches(4.3), Inches(0.9),
                "How do transitional fossils show evolution from land to water?",
                font_size=15, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Key term box
    add_colored_shape(slide, Inches(5.0), Inches(2.7), Inches(4.7), Inches(1.0), COLORS['light_blue_bg'])
    add_colored_shape(slide, Inches(5.0), Inches(2.7), Inches(4.7), Inches(0.05), COLORS['blue_accent'])
    add_text_box(slide, Inches(5.2), Inches(2.8), Inches(4.3), Inches(0.8),
                "🔑 Transitional fossil = has features of BOTH ancestor and descendant groups",
                font_size=12, color=COLORS['blue_end'])

    # Key question
    add_colored_shape(slide, Inches(0.3), Inches(3.65), Inches(9.4), Inches(0.7), COLORS['light_teal_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.75), Inches(9.0), Inches(0.5),
                "💡 Why do whales still have tiny hip bones if they don't have legs?",
                font_size=13, color=COLORS['teal_dark'], align=PP_ALIGN.CENTER)

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.55), Inches(9.7), Inches(0.55), COLORS['blue_end'])
    add_text_box(slide, Inches(0.35), Inches(4.6), Inches(9.3), Inches(0.4),
                "📝 Notecard: Write your prediction - Why would whales have vestigial hip bones?",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_station2_support_slide(prs):
    """Slide 12: Station 2 Support - Whale Evolution Timeline"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['blue_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.4),
                "🦕 Station 2 – Whale Evolution Timeline",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia")

    # Timeline table header - FIXED: Added vertical centering
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(9.6), Inches(0.45), COLORS['blue_end'])
    add_text_box(slide, Inches(0.4), Inches(0.88), Inches(2.2), Inches(0.4),
                "Fossil", font_size=12, bold=True, color=COLORS['white'],
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(2.7), Inches(0.88), Inches(1.2), Inches(0.4),
                "Age", font_size=12, bold=True, color=COLORS['white'],
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(4.0), Inches(0.88), Inches(5.6), Inches(0.4),
                "Key Features (What Changed)", font_size=12, bold=True, color=COLORS['white'],
                anchor=MSO_ANCHOR.MIDDLE)

    # Timeline rows
    fossils = [
        ("Pakicetus", "50 mya", "4 legs, lived on LAND near water, wolf-like"),
        ("Ambulocetus", "49 mya", "Could walk AND swim, webbed feet, \"walking whale\""),
        ("Rodhocetus", "47 mya", "Short legs, large tail, MOSTLY aquatic"),
        ("Basilosaurus", "40 mya", "Tiny back legs, FULLY aquatic, long body"),
        ("Modern whales", "Today", "Flippers, NO visible legs, vestigial hip bones inside"),
    ]

    # FIXED: Added vertical centering for proper text alignment in rows
    y_pos = 1.35
    for i, (name, age, features) in enumerate(fossils):
        bg_color = COLORS['light_blue_bg'] if i % 2 == 0 else COLORS['white']
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.55), bg_color)
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.2), Inches(0.45),
                    name, font_size=11, bold=True, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.7), Inches(y_pos + 0.05), Inches(1.2), Inches(0.45),
                    age, font_size=11, color=COLORS['gray_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(4.0), Inches(y_pos + 0.05), Inches(5.6), Inches(0.45),
                    features, font_size=10, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.55

    # Key insight
    add_colored_shape(slide, Inches(0.2), Inches(4.15), Inches(9.6), Inches(0.5), COLORS['green_accent'])
    add_text_box(slide, Inches(0.4), Inches(4.23), Inches(9.2), Inches(0.35),
                "🔑 Notice: Legs got SMALLER over 10 million years. Each fossil is a \"snapshot\" of evolution!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.8), Inches(9.7), Inches(0.55), COLORS['blue_start'])
    add_text_box(slide, Inches(0.35), Inches(4.85), Inches(9.3), Inches(0.4),
                "📝 Notecard: Why is Ambulocetus called a \"transitional fossil\"? What 2 groups does it bridge?",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_station3_intro_slide(prs):
    """Slide 13: Station 3 Introduction - Design a Transitional Form"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header - two-line pattern
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🔧 Station 3 – Design a Transitional Form",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "25 Points | ~20 min (Highest Value!)", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Left side - Design challenge
    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "YOUR DESIGN CHALLENGE", font_size=14, bold=True, color=COLORS['green_end'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.0), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.0), COLORS['green_end'])

    steps = """1. Review ancestor and descendant features (3 min)
2. Predict what the in-between looked like (5 min)
3. Draw/describe your transitional form (7 min)
4. Compare to real fossil (Tiktaalik) (5 min)"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(1.8),
                steps, font_size=13, color=COLORS['dark_text'])

    # Right side - The Challenge
    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(1.8), COLORS['green_end'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "📋 YOUR SCENARIO", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    challenge = """ANCESTOR (100 mya):
Fully land mammal – 4 legs, fur, small ears

DESCENDANT (today):
Fully aquatic – flippers, smooth skin

YOUR DESIGN (60 mya): ???"""
    add_text_box(slide, Inches(5.2), Inches(1.55), Inches(4.3), Inches(1.3),
                challenge, font_size=11, color=COLORS['white'])

    # Highest value reminder
    add_colored_shape(slide, Inches(5.0), Inches(3.0), Inches(4.7), Inches(0.6), COLORS['light_green_bg'])
    add_text_box(slide, Inches(5.2), Inches(3.1), Inches(4.3), Inches(0.4),
                "⭐ HIGHEST VALUE STATION - Show your best work!",
                font_size=12, bold=True, color=COLORS['green_end'], align=PP_ALIGN.CENTER)

    # Key question
    add_colored_shape(slide, Inches(0.3), Inches(3.65), Inches(9.4), Inches(0.7), COLORS['light_teal_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.75), Inches(9.0), Inches(0.5),
                "💡 Each body part should be HALFWAY between ancestor and descendant!",
                font_size=13, color=COLORS['teal_dark'], align=PP_ALIGN.CENTER)

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.55), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(4.6), Inches(9.3), Inches(0.4),
                "📝 Notecard: What features would YOU design for limbs, body covering, and habitat?",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_station3_support_slide(prs):
    """Slide 14: Station 3 Support - Worked Example (Tiktaalik)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.4),
                "🔧 Station 3 – Worked Example: Tiktaalik",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia")

    # Worked example
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(9.6), Inches(1.6), COLORS['light_teal_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(0.08), Inches(1.6), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(9.2), Inches(0.3),
                "📝 REAL TRANSITIONAL FOSSIL: Fish → Tetrapod (375 mya)", font_size=13, bold=True, color=COLORS['teal_dark'])

    example_text = """TIKTAALIK had:
• Limbs: Fins with WRIST BONES (in-between fins and legs) – could prop itself up
• Body: Fish scales + flat head like tetrapods
• Habitat: Shallow water – not fully aquatic, not fully land"""
    add_text_box(slide, Inches(0.4), Inches(1.3), Inches(9.2), Inches(1.0),
                example_text, font_size=12, color=COLORS['dark_text'])

    # Design tips - three columns
    add_text_box(slide, Inches(0.3), Inches(2.6), Inches(9.4), Inches(0.3),
                "🎯 HOW TO DESIGN YOUR TRANSITIONAL FORM:", font_size=13, bold=True, color=COLORS['green_end'])

    # Column 1 - Limbs
    add_colored_shape(slide, Inches(0.3), Inches(3.0), Inches(3.0), Inches(0.9), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.1), Inches(2.6), Inches(0.25),
                "LIMBS", font_size=11, bold=True, color=COLORS['green_accent'])
    add_text_box(slide, Inches(0.5), Inches(3.35), Inches(2.6), Inches(0.5),
                "Halfway: Shortened legs with webbed feet?",
                font_size=10, color=COLORS['dark_text'])

    # Column 2 - Body
    add_colored_shape(slide, Inches(3.5), Inches(3.0), Inches(3.0), Inches(0.9), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(3.7), Inches(3.1), Inches(2.6), Inches(0.25),
                "BODY COVERING", font_size=11, bold=True, color=COLORS['blue_accent'])
    add_text_box(slide, Inches(3.7), Inches(3.35), Inches(2.6), Inches(0.5),
                "Halfway: Sparse fur? Smooth but not slick?",
                font_size=10, color=COLORS['dark_text'])

    # Column 3 - Habitat
    add_colored_shape(slide, Inches(6.7), Inches(3.0), Inches(3.0), Inches(0.9), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(6.9), Inches(3.1), Inches(2.6), Inches(0.25),
                "HABITAT", font_size=11, bold=True, color=COLORS['orange_end'])
    add_text_box(slide, Inches(6.9), Inches(3.35), Inches(2.6), Inches(0.5),
                "Halfway: Lives in shallow water near shore?",
                font_size=10, color=COLORS['dark_text'])

    # Key reminder
    add_colored_shape(slide, Inches(0.2), Inches(4.05), Inches(9.6), Inches(0.5), COLORS['green_accent'])
    add_text_box(slide, Inches(0.4), Inches(4.13), Inches(9.2), Inches(0.35),
                "🔑 DON'T make fully land OR fully water – your design should have features of BOTH!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Notecard
    add_colored_shape(slide, Inches(0.15), Inches(4.7), Inches(9.7), Inches(0.55), COLORS['green_start'])
    add_text_box(slide, Inches(0.35), Inches(4.75), Inches(9.3), Inches(0.4),
                "📝 Notecard: Draw or describe your transitional form with labels for 3 features",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_exit_ticket_slide(prs):
    """Slide 15: Exit Ticket"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header - two-line pattern
    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.9), COLORS['header_purple_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "🎓 Exit Ticket – Evidence Integration",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.7), Inches(9.3), Inches(0.3),
                "23 Points | ~15 min", font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Question types
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(1.6), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(0.08), Inches(1.6), COLORS['header_purple_start'])

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.3),
                "QUESTION TYPES:", font_size=14, bold=True, color=COLORS['header_purple_start'])

    q_types = """• 2 NEW – Week 2 evolutionary evidence content
• 2 SPIRAL – Week 1 natural selection review
• 1 INTEGRATION – Connect evidence to natural selection mechanism
• 1 SEP-6 – Construct a scientific explanation using evidence"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(9.0), Inches(1.1),
                q_types, font_size=13, color=COLORS['dark_text'])

    # Tips box - FIXED: Increased heights and added vertical centering
    add_colored_shape(slide, Inches(0.3), Inches(2.95), Inches(4.6), Inches(1.2), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.5), Inches(3.0), Inches(4.3), Inches(0.3),
                "✅ SUCCESS TIPS:", font_size=12, bold=True, color=COLORS['green_accent'],
                anchor=MSO_ANCHOR.MIDDLE)
    tips = """• Use vocabulary: homologous, analogous, transitional
• Reference specific evidence from today
• Connect structure → ancestry → natural selection"""
    add_text_box(slide, Inches(0.5), Inches(3.3), Inches(4.3), Inches(0.8),
                tips, font_size=10, color=COLORS['dark_text'])

    # SEP reminder - FIXED: Increased heights and added vertical centering
    add_colored_shape(slide, Inches(5.1), Inches(2.95), Inches(4.6), Inches(1.2), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.3), Inches(3.0), Inches(4.2), Inches(0.3),
                "📊 SEP-6 EXPLANATION FORMAT:", font_size=12, bold=True, color=COLORS['blue_accent'],
                anchor=MSO_ANCHOR.MIDDLE)
    model_req = """• Claim: State your answer clearly
• Evidence: Cite specific observations
• Reasoning: Explain WHY evidence supports claim"""
    add_text_box(slide, Inches(5.3), Inches(3.3), Inches(4.2), Inches(0.8),
                model_req, font_size=10, color=COLORS['dark_text'])

    # Final notecard - FIXED: Added vertical centering
    add_colored_shape(slide, Inches(0.15), Inches(4.2), Inches(9.7), Inches(0.65), COLORS['header_purple_end'])
    add_text_box(slide, Inches(0.35), Inches(4.25), Inches(9.3), Inches(0.55),
                "📝 FINAL Notecard: How does the whale finger mystery connect to natural selection?",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_summary_slide(prs):
    """Slide 16: Week 2 Summary"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['light_teal_bg'])

    # Title
    add_text_box(slide, Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5),
                "Week 2 Summary: What You Learned", font_size=26, bold=True,
                color=COLORS['teal_dark'], align=PP_ALIGN.CENTER, font_name="Georgia")

    # Four key takeaways
    takeaways = [
        ("Homologous Structures", "Same bones, different functions = evidence of common ancestor", COLORS['header_purple_start']),
        ("Analogous Structures", "Different structures, same function = evolved independently", COLORS['orange_end']),
        ("Transitional Fossils", "Show intermediate features between ancestral and descendant groups", COLORS['blue_end']),
        ("Key Misconception", "Individuals don't evolve – POPULATIONS change over generations!", COLORS['green_end']),
    ]

    y_pos = 0.85
    for title, description, accent_color in takeaways:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.8), COLORS['white'])
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(0.1), Inches(0.8), accent_color)
        add_text_box(slide, Inches(0.55), Inches(y_pos + 0.1), Inches(9.0), Inches(0.3),
                    title + ":", font_size=14, bold=True, color=COLORS['dark_text'])
        add_text_box(slide, Inches(0.55), Inches(y_pos + 0.4), Inches(9.0), Inches(0.35),
                    description, font_size=12, color=COLORS['gray_text'])
        y_pos += 0.9

    # Completion celebration
    add_colored_shape(slide, Inches(0.3), Inches(4.5), Inches(9.4), Inches(0.8), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.5), Inches(4.6), Inches(9.0), Inches(0.3),
                "🎉 Week 2 Complete!", font_size=18, bold=True,
                color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(4.9), Inches(9.0), Inches(0.3),
                "Next Week: Synthesis & Assessment - Bringing it all together!",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)


if __name__ == "__main__":
    prs = create_presentation()
    output_path = "/home/user/Kairos.Sci.Repo/content/grade8/cycle03/week2/G8_C3_W2_Evidence_Evolution_Slides_Final.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Created: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
