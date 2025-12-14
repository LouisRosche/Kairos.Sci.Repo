#!/usr/bin/env python3
"""
Create G7_C6_W1 Plate Boundaries & Seismic Patterns presentation.
Topic: Plate Tectonics & Earth's Interior.

Instructional week structure (16 slides).
See PPTX_DESIGN_GUIDE.md for best practices documentation.
"""

import os
from pptx.dml.color import RGBColor
from pptx_common import (
    COLORS as BASE_COLORS,
    create_base_presentation,
    add_colored_shape,
    add_text_box,
    Inches,
    Pt,
    PP_ALIGN,
    MSO_ANCHOR,
)

# Earth/tectonic theme (earthy browns and oranges) - extend base colors
COLORS = {**BASE_COLORS}
COLORS.update({
    'earth_brown': RGBColor(0x8B, 0x45, 0x13),
    'earth_dark': RGBColor(0x5D, 0x3A, 0x1A),
    'magma_orange': RGBColor(0xEA, 0x58, 0x0C),
    'magma_red': RGBColor(0xDC, 0x26, 0x26),
    'crust_tan': RGBColor(0xD4, 0xA5, 0x73),
    'light_tan_bg': RGBColor(0xFE, 0xF3, 0xE2),
    'seismic_blue': RGBColor(0x1E, 0x40, 0xAF),
})


def create_presentation():
    """Create the G7_C6_W1 presentation"""
    prs = create_base_presentation()

    add_title_slide(prs)
    add_phenomenon_slide(prs)
    add_hook_intro_slide(prs)
    add_hook_activity_slide(prs)
    add_station1_intro_slide(prs)
    add_station1_activity_slide(prs)
    add_station2_intro_slide(prs)
    add_station2_activity_slide(prs)
    add_station3_intro_slide(prs)
    add_station3_activity_slide(prs)
    add_concepts_slide(prs)
    add_vocabulary_slide(prs)
    add_exit_intro_slide(prs)
    add_exit_tips_slide(prs)
    add_notecard_slide(prs)
    add_summary_slide(prs)

    return prs


def add_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['earth_brown'])

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(1.2),
                "Plate Tectonics & Earth's Interior",
                font_size=40, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")

    add_text_box(slide, Inches(0.5), Inches(2.6), Inches(9), Inches(0.5),
                "Week 1: Plate Boundaries & Seismic Patterns",
                font_size=24, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(9), Inches(0.5),
                "Grade 7 Science | Cycle 6 | 100 Points",
                font_size=18, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Point breakdown box
    add_colored_shape(slide, Inches(2), Inches(4.1), Inches(6), Inches(1.0), COLORS['white'])
    add_text_box(slide, Inches(2.2), Inches(4.2), Inches(5.6), Inches(0.8),
                "Hook: 12 pts | Stations: 20+20+25 pts | Exit: 23 pts",
                font_size=14, bold=True, color=COLORS['earth_brown'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_phenomenon_slide(prs):
    """Slide 2: Phenomenon Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['magma_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "This Week's Phenomenon",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Main phenomenon question
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(9.6), Inches(1.4), COLORS['light_tan_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(0.08), Inches(1.4), COLORS['magma_red'])
    add_text_box(slide, Inches(0.5), Inches(1.0), Inches(9.0), Inches(1.2),
                "Why do earthquakes happen in the same places over and over?",
                font_size=28, bold=True, color=COLORS['earth_dark'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Context
    add_colored_shape(slide, Inches(0.2), Inches(2.45), Inches(9.6), Inches(1.6), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.55), Inches(9.2), Inches(0.3),
                "THINK ABOUT IT:", font_size=14, bold=True, color=COLORS['blue_accent'])
    context = """• California, Japan, and Chile have earthquakes frequently - but why THOSE places?
• If earthquakes were random, every place would have equal risk
• There's a pattern - the "Ring of Fire" - but what causes it?
• Today we'll investigate why earthquakes cluster at specific locations."""
    add_text_box(slide, Inches(0.4), Inches(2.9), Inches(9.2), Inches(1.0),
                context, font_size=13, color=COLORS['dark_text'])

    # Standard connection
    add_colored_shape(slide, Inches(0.2), Inches(4.2), Inches(9.6), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.23), Inches(9.2), Inches(0.5),
                "MS-ESS2-2: Explain how geoscience processes change Earth's surface at varying scales",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_hook_intro_slide(prs):
    """Slide 3: Hook Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['purple_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Hook: The Earthquake Pattern Mystery",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "12 Points | 5 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Resource description
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(1.6), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(1.6), COLORS['purple_end'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "YOUR RESOURCES:", font_size=14, bold=True, color=COLORS['purple_end'])
    resources = """• Global earthquake map showing ALL earthquakes from the last 30 days
• Ring of Fire visual highlighting Pacific Rim earthquake zones
• Data table showing earthquake frequency by region

OBSERVE the patterns! What do you notice about WHERE earthquakes occur?"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(1.1),
                resources, font_size=12, color=COLORS['dark_text'])

    # Focus
    add_colored_shape(slide, Inches(0.2), Inches(2.9), Inches(9.6), Inches(0.95), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.0), Inches(9.2), Inches(0.25),
                "FOCUS:", font_size=12, bold=True, color=COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(3.3), Inches(9.2), Inches(0.5),
                "Observe clustering patterns on the map. Generate hypotheses about WHY earthquakes cluster.",
                font_size=12, color=COLORS['dark_text'])

    # Prior knowledge activation
    add_colored_shape(slide, Inches(0.2), Inches(4.0), Inches(9.6), Inches(0.7), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(4.1), Inches(9.2), Inches(0.55),
                "🔗 SPIRAL: How is this pattern similar to the weather patterns we saw in Cycle 5?",
                font_size=11, color=COLORS['dark_text'], anchor=MSO_ANCHOR.MIDDLE)


def add_hook_activity_slide(prs):
    """Slide 4: Hook Activity Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['purple_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Hook Activity: What Do You Notice?",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Questions preview
    add_text_box(slide, Inches(0.3), Inches(0.85), Inches(9.4), Inches(0.3),
                "QUESTIONS YOU'LL ANSWER:", font_size=13, bold=True, color=COLORS['purple_end'])

    questions = [
        "1. What pattern do you observe in where earthquakes occur?",
        "2. Which regions have the MOST earthquakes? Which have the FEWEST?",
        "3. What might cause earthquakes to cluster in certain areas?",
        "4. How does this pattern connect to what you know about Earth's surface?",
        "5. What questions do you have about earthquake patterns?",
    ]

    y_pos = 1.2
    for q in questions:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.5), COLORS['light_gray_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.05), Inches(9.0), Inches(0.4),
                    q, font_size=11, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.55

    # Tips
    add_colored_shape(slide, Inches(0.2), Inches(4.0), Inches(9.6), Inches(0.7), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.03), Inches(9.2), Inches(0.65),
                "TIP: Look for lines and clusters on the map. The pattern IS the clue!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station1_intro_slide(prs):
    """Slide 5: Station 1 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['green_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 1: Convection & Plate Movement",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | 6 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Resources
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(4.65), Inches(1.8), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(1.8), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(4.3), Inches(0.3),
                "RESOURCES:", font_size=12, bold=True, color=COLORS['green_end'])
    resources = """• Convection demo (heated water)
• PhET Plate Tectonics sim
• Plate models showing motion
• Interactive visualization"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(4.3), Inches(1.2),
                resources, font_size=11, color=COLORS['dark_text'])

    # Focus
    add_colored_shape(slide, Inches(5.05), Inches(1.15), Inches(4.7), Inches(1.8), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.25), Inches(1.25), Inches(4.3), Inches(0.3),
                "FOCUS:", font_size=12, bold=True, color=COLORS['blue_accent'])
    focus = """Connect HEAT transfer to
PLATE motion mechanism.

Just like warm air rises
(Cycle 5), hot rock slowly
rises in Earth's mantle!"""
    add_text_box(slide, Inches(5.25), Inches(1.6), Inches(4.3), Inches(1.2),
                focus, font_size=11, color=COLORS['dark_text'])

    # Spiral connection
    add_colored_shape(slide, Inches(0.2), Inches(3.1), Inches(9.6), Inches(0.65), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.13), Inches(9.2), Inches(0.6),
                "🔗 SPIRAL from C5: Air mass movement is driven by convection. So is plate movement!",
                font_size=11, color=COLORS['dark_text'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Key misconception alert
    add_colored_shape(slide, Inches(0.2), Inches(3.9), Inches(9.6), Inches(0.85), COLORS['light_pink_bg'])
    add_text_box(slide, Inches(0.4), Inches(4.0), Inches(9.2), Inches(0.25),
                "⚠️ MISCONCEPTION ALERT:", font_size=11, bold=True, color=COLORS['red_accent'])
    add_text_box(slide, Inches(0.4), Inches(4.3), Inches(9.2), Inches(0.4),
                "Plates do NOT float on liquid magma! The asthenosphere is a PLASTIC SOLID (like hot taffy).",
                font_size=11, color=COLORS['dark_text'])


def add_station1_activity_slide(prs):
    """Slide 6: Station 1 Activity Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 1: Investigation Steps",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Steps
    steps = [
        ("OBSERVE", "Watch the convection demo - where does hot water go? Cold water?"),
        ("CONNECT", "How is mantle convection similar to the water demo?"),
        ("MODEL", "Use the PhET sim to see how convection drives plate motion"),
        ("ANALYZE", "What happens where convection currents meet?"),
        ("PREDICT", "Where would earthquakes happen in your model?"),
    ]

    y_pos = 0.85
    colors = [COLORS['magma_orange'], COLORS['seismic_blue'], COLORS['green_end'],
              COLORS['purple_start'], COLORS['teal']]
    for i, (step, desc) in enumerate(steps):
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.58), COLORS['light_gray_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.58), colors[i])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(1.8), Inches(0.48),
                    step, font_size=12, bold=True, color=colors[i],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.3), Inches(y_pos + 0.05), Inches(7.3), Inches(0.48),
                    desc, font_size=11, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.63

    # Key concept
    add_colored_shape(slide, Inches(0.2), Inches(4.05), Inches(9.6), Inches(0.65), COLORS['earth_brown'])
    add_text_box(slide, Inches(0.4), Inches(4.08), Inches(9.2), Inches(0.6),
                "KEY: Heat from Earth's core drives convection → Convection moves plates → Plate edges = earthquakes!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station2_intro_slide(prs):
    """Slide 7: Station 2 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['seismic_blue'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 2: Plate Boundary Analysis",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | 5 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Resources
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(1.0), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "RESOURCES:", font_size=12, bold=True, color=COLORS['blue_accent'])
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(0.5),
                "Boundary type cards • Interactive world map • Data tables with boundary locations",
                font_size=11, color=COLORS['dark_text'])

    # Three boundary types
    add_text_box(slide, Inches(0.3), Inches(2.3), Inches(9.4), Inches(0.3),
                "THREE TYPES OF PLATE BOUNDARIES:", font_size=13, bold=True, color=COLORS['seismic_blue'])

    boundaries = [
        ("DIVERGENT", "Plates PULL APART", "Mid-ocean ridges, new crust forms", COLORS['green_start']),
        ("CONVERGENT", "Plates PUSH TOGETHER", "Mountains, trenches, subduction", COLORS['magma_orange']),
        ("TRANSFORM", "Plates SLIDE PAST", "Earthquakes along fault lines", COLORS['purple_start']),
    ]

    y_pos = 2.7
    for btype, motion, result, color in boundaries:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.5), COLORS['light_gray_bg'])
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(0.08), Inches(0.5), color)
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.02), Inches(2.2), Inches(0.45),
                    btype, font_size=11, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.8), Inches(y_pos + 0.02), Inches(2.5), Inches(0.45),
                    motion, font_size=10, color=COLORS['dark_text'], anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(5.4), Inches(y_pos + 0.02), Inches(4.2), Inches(0.45),
                    result, font_size=10, color=COLORS['gray_text'], anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.55

    # Focus
    add_colored_shape(slide, Inches(0.2), Inches(4.4), Inches(9.6), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.43), Inches(9.2), Inches(0.5),
                "FOCUS: Classify boundaries by their motion → Predict what geological features form there",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station2_activity_slide(prs):
    """Slide 8: Station 2 Activity Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['seismic_blue'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 2: Boundary Classification",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Real examples
    add_text_box(slide, Inches(0.3), Inches(0.85), Inches(9.4), Inches(0.3),
                "REAL-WORLD EXAMPLES TO ANALYZE:", font_size=13, bold=True, color=COLORS['seismic_blue'])

    examples = [
        ("San Andreas Fault (CA)", "Transform", "Plates slide past - major earthquakes"),
        ("Mid-Atlantic Ridge", "Divergent", "Plates pull apart - new ocean floor"),
        ("Himalayas", "Convergent", "Plates collide - tallest mountains"),
        ("Marianas Trench", "Convergent", "Oceanic plate subducts - deepest ocean"),
    ]

    y_pos = 1.2
    for location, btype, description in examples:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.55), COLORS['light_blue_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.05), Inches(3.0), Inches(0.45),
                    location, font_size=11, bold=True, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.6), Inches(y_pos + 0.05), Inches(1.8), Inches(0.45),
                    btype, font_size=11, color=COLORS['seismic_blue'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(5.5), Inches(y_pos + 0.05), Inches(4.0), Inches(0.45),
                    description, font_size=10, color=COLORS['gray_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    # Prediction task
    add_colored_shape(slide, Inches(0.2), Inches(3.65), Inches(9.6), Inches(1.1), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.75), Inches(9.2), Inches(0.25),
                "YOUR TASK:", font_size=12, bold=True, color=COLORS['orange_end'])
    add_text_box(slide, Inches(0.4), Inches(4.05), Inches(9.2), Inches(0.65),
                "Given a NEW boundary location, predict: What type is it? What features form there? Why do earthquakes occur there?",
                font_size=11, color=COLORS['dark_text'])


def add_station3_intro_slide(prs):
    """Slide 9: Station 3 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['magma_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 3: Design an Earthquake-Resistant Structure",
                font_size=22, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "25 Points | 5 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Engineering challenge
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(1.4), COLORS['light_orange_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(1.4), COLORS['magma_red'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "ENGINEERING CHALLENGE:", font_size=14, bold=True, color=COLORS['magma_red'])
    challenge = """Your community is located near a transform plate boundary. Recent data shows
increased seismic activity. Design a building that can withstand earthquake shaking.

You have: Virtual building materials, shake table simulator, budget constraints"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(0.9),
                challenge, font_size=12, color=COLORS['dark_text'])

    # Constraints
    add_colored_shape(slide, Inches(0.2), Inches(2.7), Inches(4.7), Inches(1.1), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.8), Inches(4.4), Inches(0.25),
                "CONSTRAINTS:", font_size=11, bold=True, color=COLORS['blue_accent'])
    add_text_box(slide, Inches(0.4), Inches(3.1), Inches(4.4), Inches(0.65),
                "• Budget: $500,000\n• Must be 3+ stories\n• Must house 100+ people",
                font_size=10, color=COLORS['dark_text'])

    # Success criteria
    add_colored_shape(slide, Inches(5.05), Inches(2.7), Inches(4.7), Inches(1.1), COLORS['light_green_bg'])
    add_text_box(slide, Inches(5.25), Inches(2.8), Inches(4.3), Inches(0.25),
                "SUCCESS CRITERIA:", font_size=11, bold=True, color=COLORS['green_end'])
    add_text_box(slide, Inches(5.25), Inches(3.1), Inches(4.3), Inches(0.65),
                "• Survives magnitude 7.0\n• Minimal structural damage\n• Occupants can evacuate safely",
                font_size=10, color=COLORS['dark_text'])

    # SEP connection
    add_colored_shape(slide, Inches(0.2), Inches(3.95), Inches(9.6), Inches(0.7), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(3.98), Inches(9.2), Inches(0.65),
                "SEP-6: Constructing Explanations | Use your boundary knowledge to justify design choices!",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station3_activity_slide(prs):
    """Slide 10: Station 3 Activity Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['magma_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 3: Design Process",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Design features to consider
    add_text_box(slide, Inches(0.3), Inches(0.85), Inches(9.4), Inches(0.3),
                "DESIGN FEATURES TO CONSIDER:", font_size=13, bold=True, color=COLORS['magma_orange'])

    features = [
        ("Foundation Type", "Deep vs shallow, isolated vs connected"),
        ("Building Shape", "Symmetric vs asymmetric, tall vs wide"),
        ("Materials", "Steel frame, reinforced concrete, wood"),
        ("Dampers", "Base isolators, tuned mass dampers"),
        ("Escape Routes", "Multiple exits, wide stairwells"),
    ]

    y_pos = 1.2
    for feature, options in features:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.48), COLORS['light_tan_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.02), Inches(2.5), Inches(0.44),
                    feature, font_size=11, bold=True, color=COLORS['earth_brown'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.1), Inches(y_pos + 0.02), Inches(6.4), Inches(0.44),
                    options, font_size=10, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.53

    # Justification reminder
    add_colored_shape(slide, Inches(0.2), Inches(3.9), Inches(9.6), Inches(0.85), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(4.0), Inches(9.2), Inches(0.25),
                "CRITICAL: JUSTIFY YOUR CHOICES!", font_size=12, bold=True, color=COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(4.3), Inches(9.2), Inches(0.4),
                "Your design must explain WHY each choice helps during an earthquake based on plate boundary science.",
                font_size=11, color=COLORS['dark_text'])


def add_concepts_slide(prs):
    """Slide 11: Key Concepts"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Key Concepts This Week",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    concepts = [
        ("Convection Drives Movement", "Heat from Earth's core creates convection currents in the mantle"),
        ("Plate Boundaries", "Where plates meet - divergent (apart), convergent (together), transform (slide)"),
        ("Earthquake Patterns", "Earthquakes cluster at plate boundaries where stress builds up"),
        ("Engineering Applications", "Understanding boundaries helps us design safer structures"),
    ]

    y_pos = 0.9
    for title, detail in concepts:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.75), COLORS['light_blue_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.75), COLORS['seismic_blue'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.08), Inches(3.5), Inches(0.6),
                    title, font_size=13, bold=True, color=COLORS['earth_dark'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(4.0), Inches(y_pos + 0.08), Inches(5.6), Inches(0.6),
                    detail, font_size=11, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.82

    # CCC connection
    add_colored_shape(slide, Inches(0.2), Inches(4.2), Inches(9.6), Inches(0.55), COLORS['earth_brown'])
    add_text_box(slide, Inches(0.4), Inches(4.23), Inches(9.2), Inches(0.5),
                "CCC-1: Patterns | Earthquake patterns reveal plate boundary locations!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_vocabulary_slide(prs):
    """Slide 12: Vocabulary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['purple_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Essential Vocabulary",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    vocab = [
        ("Tectonic plate", "Large slab of Earth's lithosphere that moves on the asthenosphere"),
        ("Convection", "Heat transfer by movement of heated material (drives plate motion)"),
        ("Divergent", "Boundary where plates move apart (spreading)"),
        ("Convergent", "Boundary where plates move together (collision/subduction)"),
        ("Transform", "Boundary where plates slide past each other"),
        ("Asthenosphere", "Partially molten layer below lithosphere (plastic solid, not liquid!)"),
    ]

    y_pos = 0.9
    for term, definition in vocab:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.55), COLORS['light_purple_bg'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.3), Inches(0.45),
                    term, font_size=11, bold=True, color=COLORS['purple_end'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.8), Inches(y_pos + 0.05), Inches(6.8), Inches(0.45),
                    definition, font_size=10, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    # Notecard reminder
    add_colored_shape(slide, Inches(0.2), Inches(4.55), Inches(9.6), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(4.58), Inches(9.2), Inches(0.5),
                "📝 Add these terms to your NOTECARD for the exit ticket!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_exit_intro_slide(prs):
    """Slide 13: Exit Ticket Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['orange_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Exit Ticket: Plate Dynamics Integration",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "23 Points | 6 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Question breakdown
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(2.4), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "QUESTION BREAKDOWN:", font_size=14, bold=True, color=COLORS['orange_end'])

    breakdown = [
        ("2 NEW", "Today's content: boundaries, convection, earthquake patterns"),
        ("2 SPIRAL", "Connect to C5: convection in air masses, weather patterns"),
        ("1 INTEGRATION", "Connect boundary type to earthquake risk assessment"),
        ("1 SEP", "Constructing explanations using evidence"),
    ]

    y_pos = 1.65
    for qtype, desc in breakdown:
        add_text_box(slide, Inches(0.5), Inches(y_pos), Inches(1.8), Inches(0.45),
                    qtype, font_size=12, bold=True, color=COLORS['orange_end'])
        add_text_box(slide, Inches(2.4), Inches(y_pos), Inches(7.2), Inches(0.45),
                    desc, font_size=11, color=COLORS['dark_text'])
        y_pos += 0.5

    # Timer
    add_colored_shape(slide, Inches(0.2), Inches(3.7), Inches(9.6), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(3.73), Inches(9.2), Inches(0.5),
                "⏱️ TIME: ~15 minutes | Use your NOTECARD!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Misconception check
    add_colored_shape(slide, Inches(0.2), Inches(4.4), Inches(9.6), Inches(0.55), COLORS['light_pink_bg'])
    add_text_box(slide, Inches(0.4), Inches(4.43), Inches(9.2), Inches(0.5),
                "⚠️ Remember: Plates rest on PLASTIC SOLID (asthenosphere), not liquid magma!",
                font_size=11, color=COLORS['dark_text'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_exit_tips_slide(prs):
    """Slide 14: Exit Ticket Tips"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Exit Ticket Tips",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    tips = [
        "Use EVIDENCE from today's activities in your answers",
        "Connect boundary TYPE to specific geological features",
        "Remember: convection in mantle = similar to convection in atmosphere",
        "Justify engineering choices with plate boundary science",
        "Check your notecard for vocabulary and key concepts",
    ]

    y_pos = 0.85
    for i, tip in enumerate(tips, 1):
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.55), COLORS['light_gray_bg'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(0.5), Inches(0.45),
                    str(i), font_size=14, bold=True, color=COLORS['orange_end'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(1.0), Inches(y_pos + 0.05), Inches(8.6), Inches(0.45),
                    tip, font_size=12, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    # Confidence reminder
    add_colored_shape(slide, Inches(0.2), Inches(3.9), Inches(9.6), Inches(0.85), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(3.95), Inches(9.2), Inches(0.75),
                "You've got this! You observed patterns, investigated causes, and designed solutions.\nNow show what you learned!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_notecard_slide(prs):
    """Slide 15: Notecard Reminder"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['seismic_blue'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "📝 Your Notecard Should Include:",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # What to include
    items = [
        "3 boundary types with motions (divergent=apart, convergent=together, transform=slide)",
        "Convection definition and connection to plate movement",
        "Why earthquakes cluster at plate boundaries",
        "Key vocabulary: asthenosphere, lithosphere, subduction",
        "Connection to C5: convection drives both weather AND plate motion",
    ]

    y_pos = 0.9
    for item in items:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.55), COLORS['light_blue_bg'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(0.4), Inches(0.45),
                    "✓", font_size=14, bold=True, color=COLORS['green_end'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(0.9), Inches(y_pos + 0.05), Inches(8.7), Inches(0.45),
                    item, font_size=11, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    # Physical notecard reminder
    add_colored_shape(slide, Inches(0.2), Inches(3.95), Inches(9.6), Inches(0.85), COLORS['earth_brown'])
    add_text_box(slide, Inches(0.4), Inches(3.98), Inches(9.2), Inches(0.8),
                "Remember: Your notecard is the ONLY physical resource allowed!\nMake it count - organize it clearly!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_summary_slide(prs):
    """Slide 16: Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Week 1 Summary: What You Learned",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Key takeaways
    takeaways = [
        ("Earthquake Patterns", "Cluster at plate boundaries - NOT random!"),
        ("Convection", "Heat drives plate movement (like air masses in C5)"),
        ("Boundary Types", "Divergent, convergent, transform - each has unique features"),
        ("Engineering", "Boundary knowledge helps design safer structures"),
    ]

    y_pos = 0.9
    for title, detail in takeaways:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.65), COLORS['light_tan_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.65), COLORS['earth_brown'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.5), Inches(0.55),
                    title, font_size=13, bold=True, color=COLORS['earth_dark'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.0), Inches(y_pos + 0.05), Inches(6.6), Inches(0.55),
                    detail, font_size=12, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.72

    # Next week preview
    add_colored_shape(slide, Inches(0.2), Inches(3.85), Inches(9.6), Inches(0.65), COLORS['seismic_blue'])
    add_text_box(slide, Inches(0.4), Inches(3.88), Inches(9.2), Inches(0.6),
                "NEXT WEEK: Seafloor Spreading & Continental Drift - How mountains form underwater!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Phenomenon answer
    add_colored_shape(slide, Inches(0.2), Inches(4.65), Inches(9.6), Inches(0.55), COLORS['magma_orange'])
    add_text_box(slide, Inches(0.4), Inches(4.68), Inches(9.2), Inches(0.5),
                "ANSWER: Earthquakes repeat in the same places because that's where plate boundaries are!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


if __name__ == "__main__":
    prs = create_presentation()
    output_path = "/home/user/Kairos.Sci.Repo/content/grade7/cycle06/week1/G7_C6_W1_Slides_Final.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Created: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
