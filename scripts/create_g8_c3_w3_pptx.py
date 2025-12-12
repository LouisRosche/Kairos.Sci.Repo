#!/usr/bin/env python3
"""
Create G8_C3_W3 Synthesis & Assessment presentation.
Assessment week - Evolution unit cumulative assessment.

See PPTX_DESIGN_GUIDE.md for best practices documentation.
"""

import os
from pptx.dml.color import RGBColor
from pptx_common import (
    COLORS as BASE_COLORS,
    create_base_presentation,
    add_colored_shape,
    add_text_box,
    Inches,
    Pt,
    PP_ALIGN,
    MSO_ANCHOR,
)

# G8 Assessment week uses purple theme - extend base colors
COLORS = {**BASE_COLORS}
COLORS.update({
    'purple_start': RGBColor(0x80, 0x5A, 0xD5),
    'purple_end': RGBColor(0x55, 0x3C, 0x9A),
    'purple_dark': RGBColor(0x44, 0x33, 0x7A),
    'hook_purple_start': RGBColor(0x66, 0x7E, 0xEA),
    'hook_purple_end': RGBColor(0x76, 0x4B, 0xA2),
    'blue_start': RGBColor(0x42, 0x99, 0xE1),
    'blue_end': RGBColor(0x2B, 0x6C, 0xB0),
    'orange': RGBColor(0xD6, 0x9E, 0x2E),
    'light_orange_bg': RGBColor(0xFF, 0xFB, 0xEB),
    'light_pink_bg': RGBColor(0xFE, 0xD7, 0xD7),
})


def create_presentation():
    prs = create_base_presentation()

    add_title_slide(prs)
    add_assessment_overview_slide(prs)
    add_review_guide_slide(prs)
    add_common_mistakes_slide(prs)
    add_whale_timeline_slide(prs)
    add_part1_intro_slide(prs)
    add_part1_support_slide(prs)
    add_part2_intro_slide(prs)
    add_part2_support_slide(prs)
    add_part3_intro_slide(prs)
    add_part3_support_slide(prs)
    add_completion_slide(prs)

    return prs


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['purple_end'])

    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(1.2),
                "📋 Week 3: Synthesis & Assessment 🎯",
                font_size=42, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")

    add_text_box(slide, Inches(0.5), Inches(2.9), Inches(9), Inches(0.5),
                "Grade 8 Science | Rosche | Kairos Academies",
                font_size=20, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(3.4), Inches(9), Inches(0.5),
                "Cycle 3 Cumulative Assessment | 100 Points Total | ~75 Minutes",
                font_size=16, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Teaser box (simplified - no border)
    teaser_bg = add_colored_shape(slide, Inches(2), Inches(4.2), Inches(6), Inches(0.8), COLORS['white'])
    add_text_box(slide, Inches(2.1), Inches(4.35), Inches(5.8), Inches(0.5),
                "⚠️ Assessment Week - Show What You Know!",
                font_size=16, color=COLORS['purple_dark'], align=PP_ALIGN.CENTER)


def add_assessment_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.2), Inches(0.2), Inches(9.6), Inches(0.6), COLORS['purple_end'])
    add_text_box(slide, Inches(0.4), Inches(0.28), Inches(9.2), Inches(0.5),
                "⚠️ Assessment Week - Important Information",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia")

    add_colored_shape(slide, Inches(0.2), Inches(1.0), Inches(9.6), Inches(2.2), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.0), Inches(0.08), Inches(2.2), COLORS['purple_start'])

    info = """• This is your cumulative assessment for Cycle 3 Evolution unit
• You may NOT go back once you submit each part
• Complete all 3 parts in order: Synthesis → Assessment → Misconception Check
• Take breaks between parts - your brain works better with rest!
• Ask Mr. Rosche if you need clarification on any question"""
    add_text_box(slide, Inches(0.5), Inches(1.15), Inches(9.0), Inches(2.0),
                info, font_size=15, color=COLORS['purple_dark'])

    # Three parts - FIXED: Added vertical centering
    add_colored_shape(slide, Inches(0.2), Inches(3.4), Inches(3.0), Inches(1.0), COLORS['hook_purple_start'])
    add_text_box(slide, Inches(0.4), Inches(3.45), Inches(2.6), Inches(0.4),
                "🔗 Part 1: Synthesis", font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(0.4), Inches(3.85), Inches(2.6), Inches(0.5),
                "20 pts | ~15 min", font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(3.5), Inches(3.4), Inches(3.0), Inches(1.0), COLORS['purple_end'])
    add_text_box(slide, Inches(3.7), Inches(3.45), Inches(2.6), Inches(0.4),
                "📝 Part 2: Assessment", font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(3.7), Inches(3.85), Inches(2.6), Inches(0.5),
                "60 pts | ~40 min", font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(6.8), Inches(3.4), Inches(3.0), Inches(1.0), COLORS['green_end'])
    add_text_box(slide, Inches(7.0), Inches(3.45), Inches(2.6), Inches(0.4),
                "🎯 Part 3: Misconceptions", font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(7.0), Inches(3.85), Inches(2.6), Inches(0.5),
                "20 pts | ~20 min", font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(4.6), Inches(9.6), Inches(0.6), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.63), Inches(9.2), Inches(0.55),
                "💡 Take 5-minute breaks between parts for best performance!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_review_guide_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.2), Inches(0.15), Inches(9.6), Inches(0.5), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.4),
                "📚 What You Need to Know (Review Guide)",
                font_size=22, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Week 1
    add_colored_shape(slide, Inches(0.2), Inches(0.8), Inches(4.7), Inches(1.5), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.8), Inches(0.08), Inches(1.5), COLORS['purple_start'])
    add_text_box(slide, Inches(0.4), Inches(0.9), Inches(4.4), Inches(0.3),
                "Week 1: Natural Selection", font_size=13, bold=True, color=COLORS['purple_start'])
    w1 = """• Variation → Selection → Survival → Reproduction
• POPULATIONS evolve, not individuals
• Variation must exist BEFORE selection
• Individuals cannot change genes through actions"""
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(4.4), Inches(0.95),
                w1, font_size=11, color=COLORS['dark_text'])

    # Week 2
    add_colored_shape(slide, Inches(5.1), Inches(0.8), Inches(4.7), Inches(1.5), COLORS['light_blue_bg'])
    add_colored_shape(slide, Inches(5.1), Inches(0.8), Inches(0.08), Inches(1.5), COLORS['blue_end'])
    add_text_box(slide, Inches(5.3), Inches(0.9), Inches(4.4), Inches(0.3),
                "Week 2: Evidence of Evolution", font_size=13, bold=True, color=COLORS['blue_end'])
    w2 = """• Homologous = same bones, different function
• Analogous = different bones, same function
• Transitional fossils = intermediate features
• Vestigial structures = reduced from ancestors"""
    add_text_box(slide, Inches(5.3), Inches(1.25), Inches(4.4), Inches(0.95),
                w2, font_size=11, color=COLORS['dark_text'])

    # Cycle 2 Spiral
    add_colored_shape(slide, Inches(0.2), Inches(2.45), Inches(9.6), Inches(1.1), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(2.45), Inches(0.08), Inches(1.1), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(2.55), Inches(9.2), Inches(0.3),
                "Cycle 2 Spiral Review (Forces):", font_size=13, bold=True, color=COLORS['green_end'])
    c2 = """• Newton's 3rd Law: Equal and opposite forces (size doesn't matter!)  •  F=ma: Same force + less mass = more acceleration  •  Physics determines what traits help survival"""
    add_text_box(slide, Inches(0.4), Inches(2.9), Inches(9.2), Inches(0.55),
                c2, font_size=12, color=COLORS['dark_text'])

    # Key connection
    add_colored_shape(slide, Inches(0.2), Inches(3.7), Inches(9.6), Inches(1.0), COLORS['purple_start'])
    add_text_box(slide, Inches(0.4), Inches(3.8), Inches(9.2), Inches(0.25),
                "🔑 KEY CONNECTION:", font_size=12, bold=True, color=COLORS['white'])
    add_text_box(slide, Inches(0.4), Inches(4.1), Inches(9.2), Inches(0.5),
                "Natural selection (W1 mechanism) PRODUCES the evidence (W2): homologous structures, transitional fossils, vestigial parts",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.2), Inches(4.85), Inches(9.6), Inches(0.55), COLORS['purple_end'])
    add_text_box(slide, Inches(0.4), Inches(4.88), Inches(9.2), Inches(0.5),
                "📝 Use your notecard from W1 & W2 during the assessment!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_common_mistakes_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.2), Inches(0.15), Inches(9.6), Inches(0.5), COLORS['orange'])
    add_text_box(slide, Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.4),
                "⚠️ Common Mistakes to Avoid (Lamarckian Thinking)",
                font_size=22, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    mistakes = [
        ("\"Giraffes stretched necks to reach leaves, so babies had longer necks\"",
         "Variation in neck length already existed; longer-necked giraffes survived better"),
        ("\"Whales evolved flippers BECAUSE they moved to the ocean\"",
         "Some individuals had limb variations better suited for swimming; they survived better"),
        ("\"The elephant pushes harder because it's bigger\"",
         "Newton's 3rd Law: Forces are EQUAL and opposite regardless of size"),
    ]

    y = 0.8
    for wrong, correct in mistakes:
        add_colored_shape(slide, Inches(0.2), Inches(y), Inches(4.7), Inches(1.0), COLORS['light_pink_bg'])
        add_text_box(slide, Inches(0.4), Inches(y + 0.1), Inches(4.4), Inches(0.2),
                    "❌ LAMARCKIAN:", font_size=9, bold=True, color=COLORS['red_accent'])
        add_text_box(slide, Inches(0.4), Inches(y + 0.3), Inches(4.4), Inches(0.6),
                    wrong, font_size=10, color=COLORS['dark_text'])

        add_colored_shape(slide, Inches(5.1), Inches(y), Inches(4.7), Inches(1.0), COLORS['light_green_bg'])
        add_text_box(slide, Inches(5.3), Inches(y + 0.1), Inches(4.4), Inches(0.2),
                    "✅ DARWINIAN:", font_size=9, bold=True, color=COLORS['green_end'])
        add_text_box(slide, Inches(5.3), Inches(y + 0.3), Inches(4.4), Inches(0.6),
                    correct, font_size=10, color=COLORS['dark_text'])
        y += 1.1

    add_colored_shape(slide, Inches(0.2), Inches(4.15), Inches(9.6), Inches(0.7), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.25), Inches(9.2), Inches(0.5),
                "💡 REMEMBER: Organisms can NOT change their genes by trying! Variation exists FIRST,\nthen selection happens. POPULATIONS change over generations, not individuals.",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_whale_timeline_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.2), Inches(0.15), Inches(9.6), Inches(0.5), COLORS['blue_end'])
    add_text_box(slide, Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.4),
                "🐋 Whale Evolution Timeline (Reference)",
                font_size=22, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Header
    add_colored_shape(slide, Inches(0.2), Inches(0.8), Inches(9.6), Inches(0.4), COLORS['blue_end'])
    add_text_box(slide, Inches(0.4), Inches(0.88), Inches(2.0), Inches(0.25),
                "Stage", font_size=11, bold=True, color=COLORS['white'])
    add_text_box(slide, Inches(2.5), Inches(0.88), Inches(1.5), Inches(0.25),
                "MYA", font_size=11, bold=True, color=COLORS['white'])
    add_text_box(slide, Inches(4.2), Inches(0.88), Inches(5.4), Inches(0.25),
                "Key Features", font_size=11, bold=True, color=COLORS['white'])

    fossils = [
        ("Pakicetus", "50", "4 legs, land animal, whale-like inner ear"),
        ("Ambulocetus", "48", "4 legs, webbed feet, swam + walked (\"walking whale\")"),
        ("Rodhocetus", "46", "Smaller legs, larger tail, more aquatic"),
        ("Basilosaurus", "37", "Tiny vestigial legs, fully aquatic"),
        ("Modern Whales", "Today", "Flippers with finger bones, vestigial hip bones inside"),
    ]

    y = 1.25
    for i, (name, age, features) in enumerate(fossils):
        bg = COLORS['light_blue_bg'] if i % 2 == 0 else COLORS['white']
        add_colored_shape(slide, Inches(0.2), Inches(y), Inches(9.6), Inches(0.5), bg)
        add_text_box(slide, Inches(0.4), Inches(y + 0.1), Inches(2.0), Inches(0.3),
                    name, font_size=11, bold=True, color=COLORS['dark_text'])
        add_text_box(slide, Inches(2.5), Inches(y + 0.1), Inches(1.5), Inches(0.3),
                    age, font_size=11, color=COLORS['gray_text'])
        add_text_box(slide, Inches(4.2), Inches(y + 0.1), Inches(5.4), Inches(0.3),
                    features, font_size=10, color=COLORS['dark_text'])
        y += 0.5

    add_colored_shape(slide, Inches(0.2), Inches(3.85), Inches(9.6), Inches(0.8), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(3.95), Inches(9.2), Inches(0.6),
                "🔑 This is TRANSITIONAL FOSSIL evidence: Each stage shows intermediate features!\nAmbulocetus = \"walking whale\" - has features of BOTH land mammals and whales.",
                font_size=13, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.4), Inches(4.85), Inches(9.2), Inches(0.4),
                "💡 Vestigial hip bones in modern whales = evidence of land-dwelling ancestors!",
                font_size=12, color=COLORS['purple_start'], align=PP_ALIGN.CENTER)


def add_part1_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['hook_purple_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🔗 Part 1: Synthesis Review",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "20 Points | ~15 min | Connects Week 1 + Week 2", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "WHAT YOU'LL DO", font_size=14, bold=True, color=COLORS['hook_purple_start'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.0), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.0), COLORS['hook_purple_start'])

    steps = """1. Explain how natural selection (W1)
   produces the evidence (W2)
2. Identify best evidence for whale evolution
3. Apply natural selection to Tiktaalik
4. Connect forces (C2) to survival advantages"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(1.8),
                steps, font_size=12, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(2.4), COLORS['hook_purple_end'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "🎯 KEY CONNECTION", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.2), Inches(1.6), Inches(4.3), Inches(1.8),
                "MECHANISM (Week 1):\nNatural Selection\n↓\nPRODUCES\n↓\nEVIDENCE (Week 2):\nHomologous structures\nTransitional fossils\nVestigial structures",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.15), Inches(3.7), Inches(9.7), Inches(0.5), COLORS['light_teal_bg'])
    add_text_box(slide, Inches(0.35), Inches(3.78), Inches(9.3), Inches(0.35),
                "💡 This is about CONNECTING concepts - explain HOW the mechanism creates the evidence!",
                font_size=13, color=COLORS['teal_dark'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.15), Inches(4.35), Inches(9.7), Inches(0.55), COLORS['hook_purple_start'])
    add_text_box(slide, Inches(0.35), Inches(4.38), Inches(9.3), Inches(0.5),
                "📝 Complete Part 1 Form, then take a 5-minute break",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_part1_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['hook_purple_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.4),
                "🔗 Part 1 – How Mechanism Creates Evidence",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia")

    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(9.6), Inches(2.3), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(9.2), Inches(0.3),
                "🔄 Whale Evolution Example:", font_size=13, bold=True, color=COLORS['hook_purple_start'])

    diagram = """1. VARIATION: Some early whale ancestors had limbs slightly better suited for swimming
                    ↓
2. SELECTION: Aquatic environment favored individuals who could swim better
                    ↓
3. SURVIVAL: Better swimmers caught more food, escaped predators
                    ↓
4. REPRODUCTION: Better swimmers had more offspring, passed on genes
                    ↓
5. POPULATION CHANGE: Over millions of years, legs → flippers
                    ↓
6. EVIDENCE: We see transitional fossils + vestigial hip bones + homologous limb structure"""
    add_text_box(slide, Inches(0.4), Inches(1.3), Inches(9.2), Inches(1.75),
                diagram, font_size=11, color=COLORS['dark_text'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.2), Inches(3.3), Inches(9.6), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(3.38), Inches(9.2), Inches(0.4),
                "🔑 The MECHANISM (natural selection) over TIME creates the EVIDENCE (fossils, structures)!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.15), Inches(4.0), Inches(9.7), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(4.03), Inches(9.3), Inches(0.5),
                "⏸️ After Part 1: Take a 5-minute break before Part 2!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_part2_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['purple_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "📝 Part 2: Cumulative Assessment",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "60 Points | ~40 min | All Learning Targets", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    sections = [
        ("Section A", "Natural Selection (15 pts)", "Apply the mechanism", COLORS['light_purple_bg'], COLORS['purple_start']),
        ("Section B", "Evidence (15 pts)", "Homologous, analogous, transitional", COLORS['light_blue_bg'], COLORS['blue_end']),
        ("Section C", "Connecting (15 pts)", "Mechanism to evidence", COLORS['light_green_bg'], COLORS['green_end']),
        ("Section D", "SEP-6 (15 pts)", "Construct explanations", COLORS['light_orange_bg'], COLORS['orange']),
    ]

    x = 0.2
    for title, subtitle, desc, bg, accent in sections:
        add_colored_shape(slide, Inches(x), Inches(1.15), Inches(2.3), Inches(1.5), bg)
        add_colored_shape(slide, Inches(x), Inches(1.15), Inches(2.3), Inches(0.05), accent)
        add_text_box(slide, Inches(x + 0.1), Inches(1.25), Inches(2.1), Inches(0.25),
                    title, font_size=11, bold=True, color=accent)
        add_text_box(slide, Inches(x + 0.1), Inches(1.5), Inches(2.1), Inches(0.3),
                    subtitle, font_size=10, bold=True, color=COLORS['dark_text'])
        add_text_box(slide, Inches(x + 0.1), Inches(1.85), Inches(2.1), Inches(0.7),
                    desc, font_size=10, color=COLORS['gray_text'])
        x += 2.4

    add_colored_shape(slide, Inches(0.2), Inches(2.8), Inches(9.6), Inches(0.8), COLORS['light_pink_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(2.8), Inches(0.08), Inches(0.8), COLORS['red_accent'])
    add_text_box(slide, Inches(0.4), Inches(2.9), Inches(9.2), Inches(0.6),
                "⚠️ WARNING: You cannot go back once you submit. Read each question carefully!\nWatch out for LAMARCKIAN trap answers!",
                font_size=12, color=COLORS['red_accent'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.15), Inches(3.75), Inches(9.7), Inches(0.55), COLORS['purple_end'])
    add_text_box(slide, Inches(0.35), Inches(3.78), Inches(9.3), Inches(0.5),
                "📝 Complete Part 2 Form - Take your time!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_part2_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['purple_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.4),
                "📝 Part 2 – Quick Reference",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia")

    # Left - Evidence types
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(4.7), Inches(2.0), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(4.4), Inches(0.25),
                "🦴 EVIDENCE TYPES:", font_size=12, bold=True, color=COLORS['blue_end'])
    evidence = """HOMOLOGOUS = Same structure, different function
→ Whale flipper + bat wing + human arm
→ Evidence of COMMON ANCESTOR

ANALOGOUS = Different structure, same function
→ Bird wing + butterfly wing (both fly)
→ Evidence of INDEPENDENT evolution

TRANSITIONAL = Intermediate features
→ Tiktaalik, Ambulocetus, Archaeopteryx"""
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(4.4), Inches(1.5),
                evidence, font_size=10, color=COLORS['dark_text'])

    # Right - SEP-6 format
    add_colored_shape(slide, Inches(5.1), Inches(0.85), Inches(4.7), Inches(2.0), COLORS['light_green_bg'])
    add_text_box(slide, Inches(5.3), Inches(0.95), Inches(4.4), Inches(0.25),
                "📊 SEP-6 EXPLANATION FORMAT:", font_size=12, bold=True, color=COLORS['green_end'])
    sep6 = """CLAIM: State your answer clearly
"Whales and bats share a common ancestor"

EVIDENCE: Cite specific observations
"Both have humerus, radius, ulna, carpals, 5 digits"

REASONING: Explain WHY evidence supports claim
"Same bone pattern despite different functions
indicates shared ancestry, not independent evolution" """
    add_text_box(slide, Inches(5.3), Inches(1.25), Inches(4.4), Inches(1.5),
                sep6, font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(3.0), Inches(9.6), Inches(0.6), COLORS['orange'])
    add_text_box(slide, Inches(0.4), Inches(3.08), Inches(9.2), Inches(0.4),
                "🔢 Forces: F=ma → Same force, less mass = MORE acceleration (helps small predators catch prey!)",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.15), Inches(3.75), Inches(9.7), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(3.78), Inches(9.3), Inches(0.5),
                "⏸️ After Part 2: Take a 5-minute break before Part 3!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_part3_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🎯 Part 3: Misconception Check",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "20 Points | ~20 min | Final Chance!", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "WHAT THIS TESTS", font_size=14, bold=True, color=COLORS['green_end'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(1.8), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(1.8), COLORS['green_end'])

    tests = """• Lamarckian Misconception: "Use it or
  lose it" / "Organisms change BECAUSE
  they need to"

• Newton's 3rd Law: Equal and opposite
  forces (bigger ≠ more force!)

• F=ma Application: Same force + less
  mass = more acceleration"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(1.6),
                tests, font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(2.2), COLORS['orange'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "⚠️ #1 TRAP:", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.2), Inches(1.6), Inches(4.3), Inches(1.6),
                "\"Organisms evolved traits BECAUSE they needed them\"\n\n❌ WRONG! (Lamarckian)\n\n✅ Variation existed FIRST\n✅ Selection happened AFTER\n✅ Populations change, not individuals",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.15), Inches(3.5), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(3.53), Inches(9.3), Inches(0.5),
                "📝 Complete Part 3 Form - Last section! You've got this!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_part3_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['green_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.4),
                "🎯 Part 3 – Misconception Traps to Avoid",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia")

    # Two columns
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(4.7), Inches(1.6), COLORS['light_pink_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(4.4), Inches(0.25),
                "❌ TRAP ANSWERS:", font_size=12, bold=True, color=COLORS['red_accent'])
    traps = """• "Giraffes stretched their necks..."
• "Whales evolved flippers because..."
• "The bigger animal pushes harder"
• "Individuals evolved over their lifetime" """
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(4.4), Inches(1.1),
                traps, font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.1), Inches(0.85), Inches(4.7), Inches(1.6), COLORS['light_green_bg'])
    add_text_box(slide, Inches(5.3), Inches(0.95), Inches(4.4), Inches(0.25),
                "✅ CORRECT ANSWERS:", font_size=12, bold=True, color=COLORS['green_end'])
    correct = """• "Variation already existed in the population"
• "Selection favored certain pre-existing traits"
• "Forces are equal regardless of size (N3L)"
• "Populations change over generations" """
    add_text_box(slide, Inches(5.3), Inches(1.25), Inches(4.4), Inches(1.1),
                correct, font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(2.6), Inches(9.6), Inches(0.7), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(2.7), Inches(9.2), Inches(0.5),
                "💡 TIP: If the answer suggests organisms TRIED to evolve or NEEDED to change, it's WRONG!\nEvolution is NOT goal-directed. Variation comes FIRST, then selection happens.",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_completion_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['green_end'])

    add_text_box(slide, Inches(0.5), Inches(0.8), Inches(9), Inches(0.8),
                "🎉 Cycle 3 Complete!",
                font_size=40, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    add_text_box(slide, Inches(0.5), Inches(1.7), Inches(9), Inches(0.5),
                "Congratulations on finishing the Evolution unit!",
                font_size=20, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(2), Inches(2.4), Inches(6), Inches(2.2), COLORS['white'])
    add_text_box(slide, Inches(2.2), Inches(2.5), Inches(5.6), Inches(0.3),
                "Key Takeaways:", font_size=16, bold=True, color=COLORS['green_end'])
    takeaways = """• Natural selection requires variation + differential survival
• POPULATIONS evolve, not individuals
• Homologous structures prove common ancestry
• Transitional fossils show evolution in progress
• Physics (F=ma, N3L) determines survival advantages"""
    add_text_box(slide, Inches(2.2), Inches(2.85), Inches(5.6), Inches(1.6),
                takeaways, font_size=13, color=COLORS['dark_text'])

    add_text_box(slide, Inches(0.5), Inches(4.8), Inches(9), Inches(0.5),
                "Great work this cycle! 🧬",
                font_size=18, color=COLORS['white'], align=PP_ALIGN.CENTER)


if __name__ == "__main__":
    prs = create_presentation()
    output_path = "/home/user/Kairos.Sci.Repo/content/grade8/cycle03/week3/G8_C3_W3_Assessment_Slides_Final.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Created: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
