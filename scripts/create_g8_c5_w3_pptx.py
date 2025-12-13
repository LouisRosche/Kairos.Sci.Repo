#!/usr/bin/env python3
"""
Create G8_C5_W3 Synthesis & Assessment presentation.
Topic: Waves & Information Transfer.

Assessment week structure (11 slides) - different from instructional weeks.
See PPTX_DESIGN_GUIDE.md for best practices documentation.
"""

import os
from pptx.dml.color import RGBColor
from pptx_common import (
    COLORS as BASE_COLORS,
    create_base_presentation,
    add_colored_shape,
    add_text_box,
    Inches,
    Pt,
    PP_ALIGN,
    MSO_ANCHOR,
)

# Wave theme (purples and cyans) - extend base colors
COLORS = {**BASE_COLORS}
COLORS.update({
    'wave_purple': RGBColor(0x7C, 0x3A, 0xED),
    'wave_purple_dark': RGBColor(0x5B, 0x21, 0xB6),
    'wave_cyan': RGBColor(0x06, 0xB6, 0xD4),
    'wave_cyan_dark': RGBColor(0x08, 0x91, 0xB2),
    'light_purple_bg': RGBColor(0xF3, 0xE8, 0xFF),
    'light_cyan_bg': RGBColor(0xE0, 0xF7, 0xFA),
    'assessment_purple': RGBColor(0x6D, 0x28, 0xD9),
})


def create_presentation():
    """Create the G8_C5_W3 Assessment presentation"""
    prs = create_base_presentation()

    add_title_slide(prs)
    add_overview_slide(prs)
    add_assessed_on_slide(prs)
    add_part1_intro_slide(prs)
    add_part1_support_slide(prs)
    add_part2_intro_slide(prs)
    add_part2_support_slide(prs)
    add_part3_intro_slide(prs)
    add_part3_support_slide(prs)
    add_tips_slide(prs)
    add_summary_slide(prs)

    return prs


def add_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['assessment_purple'])

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(1.2),
                "Waves & Information Transfer",
                font_size=44, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")

    add_text_box(slide, Inches(0.5), Inches(2.6), Inches(9), Inches(0.5),
                "Week 3: Synthesis & Assessment",
                font_size=28, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(9), Inches(0.5),
                "Grade 8 Science | Cycle 5 | 100 Points",
                font_size=18, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Assessment info box
    add_colored_shape(slide, Inches(2), Inches(4.1), Inches(6), Inches(1.0), COLORS['white'])
    add_text_box(slide, Inches(2.2), Inches(4.15), Inches(5.6), Inches(0.45),
                "Part 1: Synthesis (20 pts) | Part 2: Assessment (60 pts)",
                font_size=14, bold=True, color=COLORS['assessment_purple'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(2.2), Inches(4.55), Inches(5.6), Inches(0.45),
                "Part 3: Misconception Check (20 pts)",
                font_size=14, bold=True, color=COLORS['assessment_purple'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_overview_slide(prs):
    """Slide 2: Assessment Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['assessment_purple'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Assessment Overview",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Three parts
    parts = [
        ("Part 1: Synthesis Review", "20 pts", "15 min", "Connect W1 wave properties + W2 material interactions"),
        ("Part 2: Cumulative Assessment", "60 pts", "40 min", "Wave properties, materials, information transfer, models"),
        ("Part 3: Misconception Check", "20 pts", "20 min", "Target common errors for feedback"),
    ]

    y_pos = 0.9
    colors = [COLORS['green_start'], COLORS['wave_cyan'], COLORS['wave_purple']]
    for i, (title, pts, time, desc) in enumerate(parts):
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(1.1), COLORS['light_blue_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(1.1), colors[i])

        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.1), Inches(5.5), Inches(0.35),
                    title, font_size=16, bold=True, color=COLORS['dark_text'])
        add_text_box(slide, Inches(6.0), Inches(y_pos + 0.1), Inches(1.8), Inches(0.35),
                    pts, font_size=14, bold=True, color=colors[i], align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(7.9), Inches(y_pos + 0.1), Inches(1.8), Inches(0.35),
                    time, font_size=14, color=COLORS['gray_text'], align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.5), Inches(9.2), Inches(0.5),
                    desc, font_size=12, color=COLORS['gray_text'])
        y_pos += 1.25

    # Total bar
    add_colored_shape(slide, Inches(0.2), Inches(4.7), Inches(9.6), Inches(0.55), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.4), Inches(4.73), Inches(9.2), Inches(0.5),
                "TOTAL: 100 Points | ~75 Minutes | Use your notecard!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_assessed_on_slide(prs):
    """Slide 3: What You'll Be Assessed On"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "What You'll Be Assessed On",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Week 1 content
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(4.7), Inches(1.7), COLORS['light_cyan_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(0.08), Inches(1.7), COLORS['wave_cyan'])
    add_text_box(slide, Inches(0.4), Inches(1.0), Inches(4.4), Inches(0.3),
                "WEEK 1: Wave Properties", font_size=13, bold=True, color=COLORS['wave_cyan_dark'])
    w1_content = """• Energy transfer (not matter)
• Reflection, absorption, transmission
• Wavelength and frequency relationship
• EM spectrum and wave types
• Wave barrier design trade-offs"""
    add_text_box(slide, Inches(0.4), Inches(1.35), Inches(4.4), Inches(1.2),
                w1_content, font_size=11, color=COLORS['dark_text'])

    # Week 2 content
    add_colored_shape(slide, Inches(5.1), Inches(0.9), Inches(4.7), Inches(1.7), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(5.1), Inches(0.9), Inches(0.08), Inches(1.7), COLORS['wave_purple'])
    add_text_box(slide, Inches(5.3), Inches(1.0), Inches(4.4), Inches(0.3),
                "WEEK 2: Material Interactions", font_size=13, bold=True, color=COLORS['wave_purple_dark'])
    w2_content = """• Why materials block some waves
• Transmission vs absorption coefficients
• Digital signal encoding (binary)
• Information transfer via waves
• Communication system design"""
    add_text_box(slide, Inches(5.3), Inches(1.35), Inches(4.4), Inches(1.2),
                w2_content, font_size=11, color=COLORS['dark_text'])

    # Integration skills
    add_colored_shape(slide, Inches(0.2), Inches(2.75), Inches(9.6), Inches(1.3), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.85), Inches(9.2), Inches(0.3),
                "INTEGRATION SKILLS:", font_size=13, bold=True, color=COLORS['green_end'])
    integration = """• Connect wave properties to how materials interact with different wavelengths
• Explain why WiFi passes through walls but light doesn't (wavelength matters!)
• Apply SEP-2: Developing and Using Models to wave behavior predictions
• Connect to Cycle 4: How does energy transfer in waves compare to ecosystems?"""
    add_text_box(slide, Inches(0.4), Inches(3.2), Inches(9.2), Inches(0.8),
                integration, font_size=11, color=COLORS['dark_text'])

    # Key reminder
    add_colored_shape(slide, Inches(0.2), Inches(4.2), Inches(9.6), Inches(0.55), COLORS['wave_purple'])
    add_text_box(slide, Inches(0.4), Inches(4.23), Inches(9.2), Inches(0.5),
                "KEY: Waves transfer ENERGY, not matter | Wavelength determines material interaction!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_part1_intro_slide(prs):
    """Slide 4: Part 1 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['green_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Part 1 – Synthesis Review",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | ~15 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # What to do
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(2.3), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(2.3), COLORS['green_end'])

    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "SYNTHESIS TASK:", font_size=14, bold=True, color=COLORS['green_end'])

    task_text = """Connect what you learned in Week 1 (wave properties, reflection/absorption/transmission)
with Week 2 (material interactions, information transfer).

You'll answer questions that require you to:
• Explain HOW wave properties determine whether waves pass through materials
• Describe WHY different waves (WiFi vs light) interact differently with the same material
• Connect wave behavior to real-world communication and technology"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(1.8),
                task_text, font_size=13, color=COLORS['dark_text'])

    # Format
    add_colored_shape(slide, Inches(0.2), Inches(3.6), Inches(4.7), Inches(0.95), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.7), Inches(4.4), Inches(0.25),
                "FORMAT:", font_size=12, bold=True, color=COLORS['blue_accent'])
    add_text_box(slide, Inches(0.4), Inches(4.0), Inches(4.4), Inches(0.5),
                "• 4 short-answer questions\n• Use specific vocabulary\n• Connect W1 to W2 concepts",
                font_size=11, color=COLORS['dark_text'])

    # Tips
    add_colored_shape(slide, Inches(5.1), Inches(3.6), Inches(4.6), Inches(0.95), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(5.3), Inches(3.7), Inches(4.2), Inches(0.25),
                "TIPS:", font_size=12, bold=True, color=COLORS['orange_end'])
    add_text_box(slide, Inches(5.3), Inches(4.0), Inches(4.2), Inches(0.5),
                "• Review your notecard from W1-W2\n• Think about wavelength differences\n• Use vocab: transmit, absorb, reflect",
                font_size=11, color=COLORS['dark_text'])


def add_part1_support_slide(prs):
    """Slide 5: Part 1 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Part 1 – Key Connections to Make",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Connection 1
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(9.6), Inches(1.2), COLORS['light_cyan_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(9.2), Inches(0.3),
                "CONNECTION 1: Wave Properties → Material Interaction", font_size=13, bold=True, color=COLORS['wave_cyan_dark'])
    conn1 = """Wavelength determines how waves interact with materials. Longer wavelengths (radio/WiFi)
can pass around obstacles and through walls. Shorter wavelengths (light) are blocked by
most solid materials → This is WHY your phone works through walls but you can't see through them."""
    add_text_box(slide, Inches(0.4), Inches(1.3), Inches(9.2), Inches(0.7),
                conn1, font_size=11, color=COLORS['dark_text'])

    # Connection 2
    add_colored_shape(slide, Inches(0.2), Inches(2.2), Inches(9.6), Inches(1.2), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.3), Inches(9.2), Inches(0.3),
                "CONNECTION 2: Energy Transfer → Information Transfer", font_size=13, bold=True, color=COLORS['wave_purple_dark'])
    conn2 = """Waves transfer ENERGY, not matter. This energy can carry patterns (information).
Digital signals encode info as patterns of wave pulses (on/off = 1/0).
→ WiFi, cell signals, fiber optics all use wave energy to carry information."""
    add_text_box(slide, Inches(0.4), Inches(2.65), Inches(9.2), Inches(0.7),
                conn2, font_size=11, color=COLORS['dark_text'])

    # Sentence starters
    add_colored_shape(slide, Inches(0.2), Inches(3.55), Inches(9.6), Inches(1.0), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.65), Inches(9.2), Inches(0.25),
                "SENTENCE STARTERS FOR SYNTHESIS:", font_size=11, bold=True, color=COLORS['green_accent'])
    starters = """• "WiFi passes through walls because its wavelength is..."
• "Information is transferred by waves through..."
• "The relationship between wavelength and material transmission is..."""
    add_text_box(slide, Inches(0.4), Inches(3.95), Inches(9.2), Inches(0.55),
                starters, font_size=10, color=COLORS['dark_text'])


def add_part2_intro_slide(prs):
    """Slide 6: Part 2 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['wave_cyan'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Part 2 – Cumulative Assessment",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "60 Points | ~40 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Sections
    sections = [
        ("A: Wave Properties", "15 pts", "Energy transfer, types, wavelength/frequency"),
        ("B: Material Interactions", "15 pts", "Transmit, absorb, reflect based on wavelength"),
        ("C: Information Transfer", "15 pts", "Digital encoding, signal patterns"),
        ("D: Model Development", "15 pts", "Predict wave behavior (SEP-2)"),
    ]

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(2.4), COLORS['light_cyan_bg'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "ASSESSMENT SECTIONS:", font_size=14, bold=True, color=COLORS['wave_cyan_dark'])

    y_pos = 1.6
    for section, pts, desc in sections:
        add_text_box(slide, Inches(0.5), Inches(y_pos), Inches(3.5), Inches(0.4),
                    section, font_size=12, bold=True, color=COLORS['dark_text'])
        add_text_box(slide, Inches(4.2), Inches(y_pos), Inches(1.2), Inches(0.4),
                    pts, font_size=12, color=COLORS['wave_cyan_dark'])
        add_text_box(slide, Inches(5.5), Inches(y_pos), Inches(4.0), Inches(0.4),
                    desc, font_size=11, color=COLORS['gray_text'])
        y_pos += 0.5

    # Question types
    add_colored_shape(slide, Inches(0.2), Inches(3.7), Inches(9.6), Inches(0.85), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.8), Inches(9.2), Inches(0.25),
                "QUESTION TYPES:", font_size=12, bold=True, color=COLORS['orange_end'])
    add_text_box(slide, Inches(0.4), Inches(4.1), Inches(9.2), Inches(0.4),
                "Multiple choice • Data analysis • Short answer • Extended response (use rubric criteria)",
                font_size=11, color=COLORS['dark_text'])


def add_part2_support_slide(prs):
    """Slide 7: Part 2 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['wave_cyan_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Part 2 – Key Concepts to Review",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Section A
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(4.7), Inches(1.45), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(4.4), Inches(0.25),
                "SECTION A: Wave Properties", font_size=11, bold=True, color=COLORS['blue_accent'])
    a_content = """• Waves transfer energy, NOT matter
• Particles oscillate in place
• Wavelength × frequency = speed
• EM spectrum: radio → visible → X-ray
• All EM waves: same speed in vacuum"""
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(4.4), Inches(1.0),
                a_content, font_size=10, color=COLORS['dark_text'])

    # Section B
    add_colored_shape(slide, Inches(5.1), Inches(0.85), Inches(4.6), Inches(1.45), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(5.3), Inches(0.95), Inches(4.2), Inches(0.25),
                "SECTION B: Material Interactions", font_size=11, bold=True, color=COLORS['wave_purple'])
    b_content = """• Transmit: wave passes through
• Absorb: wave energy transferred
• Reflect: wave bounces back
• Wavelength determines behavior
• Different materials = different effects"""
    add_text_box(slide, Inches(5.3), Inches(1.25), Inches(4.2), Inches(1.0),
                b_content, font_size=10, color=COLORS['dark_text'])

    # Section C
    add_colored_shape(slide, Inches(0.2), Inches(2.45), Inches(4.7), Inches(1.3), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.55), Inches(4.4), Inches(0.25),
                "SECTION C: Information Transfer", font_size=11, bold=True, color=COLORS['green_accent'])
    c_content = """• Digital: on/off patterns (binary)
• Analog: continuous wave changes
• Encoding: info → wave patterns
• Decoding: wave patterns → info
• Bandwidth = data rate capacity"""
    add_text_box(slide, Inches(0.4), Inches(2.85), Inches(4.4), Inches(0.85),
                c_content, font_size=10, color=COLORS['dark_text'])

    # Section D
    add_colored_shape(slide, Inches(5.1), Inches(2.45), Inches(4.6), Inches(1.3), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(5.3), Inches(2.55), Inches(4.2), Inches(0.25),
                "SECTION D: Model Development", font_size=11, bold=True, color=COLORS['orange_end'])
    d_content = """• Predict wave behavior in materials
• Draw/explain reflection angles
• Model wave interference patterns
• Justify predictions with evidence
• Apply to real-world scenarios"""
    add_text_box(slide, Inches(5.3), Inches(2.85), Inches(4.2), Inches(0.85),
                d_content, font_size=10, color=COLORS['dark_text'])

    # Final reminder
    add_colored_shape(slide, Inches(0.2), Inches(3.9), Inches(9.6), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(3.93), Inches(9.2), Inches(0.5),
                "REMEMBER: Energy transfers, matter stays! | Wavelength determines material interaction!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_part3_intro_slide(prs):
    """Slide 8: Part 3 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['wave_purple'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Part 3 – Misconception Check",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | ~20 min", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Purpose
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(1.2), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(1.2), COLORS['wave_purple_dark'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "PURPOSE:", font_size=14, bold=True, color=COLORS['wave_purple_dark'])
    purpose = """This section helps us identify common misunderstandings so we can provide targeted feedback.
Your answers help us improve instruction for everyone. Answer honestly based on your current understanding—
this is about learning, not just getting points!"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(0.7),
                purpose, font_size=12, color=COLORS['dark_text'])

    # Targeted misconceptions
    add_text_box(slide, Inches(0.3), Inches(2.5), Inches(9.4), Inches(0.3),
                "COMMON MISCONCEPTIONS WE'RE CHECKING:", font_size=13, bold=True, color=COLORS['wave_purple_dark'])

    misconceptions = [
        ("waves-move-matter", "Waves carry matter from place to place", "Waves transfer ENERGY; particles oscillate in place"),
        ("all-waves-same", "All waves behave the same way", "Different wavelengths interact differently with materials"),
        ("light-instant", "Light travels instantly", "Light has finite speed: ~3×10⁸ m/s"),
    ]

    y_pos = 2.9
    for id, wrong, correct in misconceptions:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.5), COLORS['light_pink_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.02), Inches(4.5), Inches(0.45),
                    f"❌ {wrong}", font_size=10, color=COLORS['red_accent'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(5.1), Inches(y_pos + 0.02), Inches(4.5), Inches(0.45),
                    f"✓ {correct}", font_size=10, color=COLORS['green_accent'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.55


def add_part3_support_slide(prs):
    """Slide 9: Part 3 Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['wave_purple_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Part 3 – Think Through These Carefully",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Energy not matter
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(4.7), Inches(1.5), COLORS['light_cyan_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(4.4), Inches(0.3),
                "ENERGY TRANSFER, NOT MATTER", font_size=12, bold=True, color=COLORS['wave_cyan_dark'])
    wc = """Think of a stadium wave:
• People stand up, then sit down
• The WAVE moves around the stadium
• But PEOPLE stay in their seats!

Same with all waves - particles oscillate,
energy moves forward."""
    add_text_box(slide, Inches(0.4), Inches(1.3), Inches(4.4), Inches(1.0),
                wc, font_size=10, color=COLORS['dark_text'])

    # Different waves
    add_colored_shape(slide, Inches(5.1), Inches(0.85), Inches(4.6), Inches(1.5), COLORS['light_green_bg'])
    add_text_box(slide, Inches(5.3), Inches(0.95), Inches(4.2), Inches(0.3),
                "DIFFERENT WAVES, DIFFERENT BEHAVIOR", font_size=12, bold=True, color=COLORS['green_end'])
    density = """Why WiFi passes through walls but light doesn't:
• WiFi: ~12 cm wavelength (long)
• Light: ~500 nm wavelength (tiny!)

Longer wavelengths → pass around/through
smaller obstacles more easily."""
    add_text_box(slide, Inches(5.3), Inches(1.3), Inches(4.2), Inches(1.0),
                density, font_size=10, color=COLORS['dark_text'])

    # Light speed
    add_colored_shape(slide, Inches(0.2), Inches(2.5), Inches(9.6), Inches(1.3), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.6), Inches(9.2), Inches(0.3),
                "LIGHT HAS A SPEED!", font_size=12, bold=True, color=COLORS['orange_end'])
    seasons = """Light travels at ~300,000,000 m/s (186,000 miles/second) - FAST but not instant!
• Sunlight takes ~8 minutes to reach Earth
• Moon's reflected light takes ~1.3 seconds
• When you see a star 100 light-years away, you're seeing it 100 years in the past!

We can measure light's travel time with precise instruments."""
    add_text_box(slide, Inches(0.4), Inches(2.95), Inches(9.2), Inches(0.8),
                seasons, font_size=10, color=COLORS['dark_text'])

    # Final note
    add_colored_shape(slide, Inches(0.2), Inches(3.95), Inches(9.6), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(3.98), Inches(9.2), Inches(0.5),
                "Take your time—read each question carefully and think about the CORRECT explanation!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_tips_slide(prs):
    """Slide 10: Tips for Success"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Tips for Success",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Tips
    tips = [
        ("📝 USE YOUR NOTECARD", "You prepared it for this! Check definitions and key concepts."),
        ("⏰ MANAGE YOUR TIME", "Part 1: 15 min | Part 2: 40 min | Part 3: 20 min"),
        ("📖 READ CAREFULLY", "Look for key words: transmit, absorb, wavelength, energy"),
        ("🔗 MAKE CONNECTIONS", "Link wave properties → material behavior → information transfer"),
        ("✓ CHECK YOUR WORK", "Revisit uncertain answers if time permits"),
    ]

    y_pos = 0.9
    colors = [COLORS['wave_purple'], COLORS['wave_cyan'], COLORS['orange_start'],
              COLORS['green_start'], COLORS['blue_accent']]
    for i, (tip, detail) in enumerate(tips):
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.65), COLORS['light_gray_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.65), colors[i])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(3.5), Inches(0.55),
                    tip, font_size=12, bold=True, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(4.0), Inches(y_pos + 0.05), Inches(5.6), Inches(0.55),
                    detail, font_size=11, color=COLORS['gray_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.7

    # Good luck
    add_colored_shape(slide, Inches(0.2), Inches(4.45), Inches(9.6), Inches(0.55), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.4), Inches(4.48), Inches(9.2), Inches(0.5),
                "You've got this! Show what you've learned about waves, materials, and information!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_summary_slide(prs):
    """Slide 11: Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Cycle 5 Summary: What You Learned",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    # Key takeaways
    takeaways = [
        ("Wave Properties", "Energy transfer, reflection/absorption/transmission"),
        ("Material Interactions", "Wavelength determines how waves interact with matter"),
        ("Information Transfer", "Waves carry encoded information (digital/analog)"),
        ("Applications", "WiFi, cell signals, fiber optics all use wave principles!"),
    ]

    y_pos = 0.9
    for title, detail in takeaways:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.65), COLORS['light_purple_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.65), COLORS['wave_purple'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.2), Inches(0.55),
                    title, font_size=13, bold=True, color=COLORS['wave_purple_dark'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.7), Inches(y_pos + 0.05), Inches(6.9), Inches(0.55),
                    detail, font_size=12, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.72

    # Next cycle preview
    add_colored_shape(slide, Inches(0.2), Inches(3.85), Inches(9.6), Inches(0.65), COLORS['header_blue_end'])
    add_text_box(slide, Inches(0.4), Inches(3.88), Inches(9.2), Inches(0.6),
                "COMING UP in Cycle 6: Electricity & Magnetism!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # Final takeaway
    add_colored_shape(slide, Inches(0.2), Inches(4.65), Inches(9.6), Inches(0.55), COLORS['wave_purple'])
    add_text_box(slide, Inches(0.4), Inches(4.68), Inches(9.2), Inches(0.5),
                "KEY INSIGHT: Waves transfer energy to carry information through our connected world!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


if __name__ == "__main__":
    prs = create_presentation()
    output_path = "/home/user/Kairos.Sci.Repo/content/grade8/cycle05/week3/G8_C5_W3_Assessment_Slides_Final.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Created: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
