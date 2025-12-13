#!/usr/bin/env python3
"""
Create G8_C6_W1 Magnetic Forces & Fields presentation.
Topic: Electricity & Magnetism.

Instructional week structure (16 slides).
See PPTX_DESIGN_GUIDE.md for best practices documentation.
"""

import os
from pptx.dml.color import RGBColor
from pptx_common import (
    COLORS as BASE_COLORS,
    create_base_presentation,
    add_colored_shape,
    add_text_box,
    Inches,
    Pt,
    PP_ALIGN,
    MSO_ANCHOR,
)

# Electromagnetic theme (blues and silvers) - extend base colors
COLORS = {**BASE_COLORS}
COLORS.update({
    'magnet_red': RGBColor(0xDC, 0x26, 0x26),
    'magnet_blue': RGBColor(0x1D, 0x4E, 0xD8),
    'field_purple': RGBColor(0x7C, 0x3A, 0xED),
    'steel_gray': RGBColor(0x4B, 0x55, 0x63),
    'force_orange': RGBColor(0xEA, 0x58, 0x0C),
    'light_magnet_bg': RGBColor(0xEE, 0xF2, 0xFF),
    'electric_blue': RGBColor(0x00, 0x77, 0xB6),
})


def create_presentation():
    """Create the G8_C6_W1 presentation"""
    prs = create_base_presentation()

    add_title_slide(prs)
    add_phenomenon_slide(prs)
    add_hook_intro_slide(prs)
    add_hook_activity_slide(prs)
    add_station1_intro_slide(prs)
    add_station1_activity_slide(prs)
    add_station2_intro_slide(prs)
    add_station2_activity_slide(prs)
    add_station3_intro_slide(prs)
    add_station3_activity_slide(prs)
    add_concepts_slide(prs)
    add_vocabulary_slide(prs)
    add_exit_intro_slide(prs)
    add_exit_tips_slide(prs)
    add_notecard_slide(prs)
    add_summary_slide(prs)

    return prs


def add_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['magnet_blue'])

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(1.2),
                "Electricity & Magnetism",
                font_size=44, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")

    add_text_box(slide, Inches(0.5), Inches(2.6), Inches(9), Inches(0.5),
                "Week 1: Magnetic Forces & Fields",
                font_size=26, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(9), Inches(0.5),
                "Grade 8 Science | Cycle 6 | 100 Points",
                font_size=18, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(2), Inches(4.1), Inches(6), Inches(1.0), COLORS['white'])
    add_text_box(slide, Inches(2.2), Inches(4.2), Inches(5.6), Inches(0.8),
                "Hook: 12 pts | Stations: 20+20+25 pts | Exit: 23 pts",
                font_size=14, bold=True, color=COLORS['magnet_blue'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_phenomenon_slide(prs):
    """Slide 2: Phenomenon Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['field_purple'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "This Week's Phenomenon",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(9.6), Inches(1.4), COLORS['light_magnet_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.9), Inches(0.08), Inches(1.4), COLORS['magnet_red'])
    add_text_box(slide, Inches(0.5), Inches(1.0), Inches(9.0), Inches(1.2),
                "Why do some magnets attract through tables while others don't?",
                font_size=26, bold=True, color=COLORS['magnet_blue'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(2.45), Inches(9.6), Inches(1.6), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.55), Inches(9.2), Inches(0.3),
                "THINK ABOUT IT:", font_size=14, bold=True, color=COLORS['blue_accent'])
    context = """• A strong neodymium magnet can hold papers through a 1-inch wooden desk
• A fridge magnet barely sticks through a piece of paper
• Both are magnets - what makes some so much stronger?
• Today we'll investigate what determines magnetic force strength!"""
    add_text_box(slide, Inches(0.4), Inches(2.9), Inches(9.2), Inches(1.0),
                context, font_size=13, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(4.2), Inches(9.6), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.23), Inches(9.2), Inches(0.5),
                "MS-PS2-3: Ask questions about data to determine factors affecting force strength",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_hook_intro_slide(prs):
    """Slide 3: Hook Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['purple_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Hook: The Invisible Force Mystery",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "12 Points | 5 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(1.6), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(1.6), COLORS['purple_end'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "YOUR RESOURCES:", font_size=14, bold=True, color=COLORS['purple_end'])
    resources = """• Demonstration: magnets through various materials (paper, wood, glass, metal)
• Different magnet types: fridge magnet vs bar magnet vs neodymium
• Distance testing apparatus

OBSERVE: When does the magnet work through the material? When doesn't it?"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(1.1),
                resources, font_size=12, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(2.9), Inches(9.6), Inches(0.95), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.0), Inches(9.2), Inches(0.25),
                "FOCUS:", font_size=12, bold=True, color=COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(3.3), Inches(9.2), Inches(0.5),
                "What factors might affect how far a magnetic force can reach?",
                font_size=12, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(4.0), Inches(9.6), Inches(0.7), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(4.1), Inches(9.2), Inches(0.55),
                "🔗 SPIRAL: Connect to C5 wave behavior - both are invisible forces acting at a distance!",
                font_size=11, color=COLORS['dark_text'], anchor=MSO_ANCHOR.MIDDLE)


def add_hook_activity_slide(prs):
    """Slide 4: Hook Activity Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['purple_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Hook Activity: What Do You Notice?",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(0.3), Inches(0.85), Inches(9.4), Inches(0.3),
                "QUESTIONS YOU'LL ANSWER:", font_size=13, bold=True, color=COLORS['purple_end'])

    questions = [
        "1. Which magnet can affect objects from farthest away?",
        "2. What happens to the force as you increase distance?",
        "3. Does the material BETWEEN the magnet and object matter?",
        "4. What might explain why some magnets are stronger than others?",
        "5. What questions do you have about magnetic forces?",
    ]

    y_pos = 1.2
    for q in questions:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.5), COLORS['light_gray_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.05), Inches(9.0), Inches(0.4),
                    q, font_size=11, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.55

    add_colored_shape(slide, Inches(0.2), Inches(4.0), Inches(9.6), Inches(0.7), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.03), Inches(9.2), Inches(0.65),
                "TIP: Magnetic force follows an INVERSE relationship with distance!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station1_intro_slide(prs):
    """Slide 5: Station 1 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['green_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 1: Magnetic Field Mapping",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | 6 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(4.65), Inches(1.8), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(1.8), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(4.3), Inches(0.3),
                "RESOURCES:", font_size=12, bold=True, color=COLORS['green_end'])
    resources = """• Bar magnets + iron filings
• Compass arrays
• PhET Magnets & Compasses
• Field line templates"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(4.3), Inches(1.2),
                resources, font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.05), Inches(1.15), Inches(4.7), Inches(1.8), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(5.25), Inches(1.25), Inches(4.3), Inches(0.3),
                "FOCUS:", font_size=12, bold=True, color=COLORS['blue_accent'])
    focus = """VISUALIZE magnetic fields:
• Field lines show direction
• Density shows strength
• Lines NEVER cross
• N→S outside magnet"""
    add_text_box(slide, Inches(5.25), Inches(1.6), Inches(4.3), Inches(1.2),
                focus, font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(3.1), Inches(9.6), Inches(0.65), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.13), Inches(9.2), Inches(0.6),
                "🔗 SPIRAL from C5: Like wave patterns, field patterns reveal invisible forces!",
                font_size=11, color=COLORS['dark_text'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(3.9), Inches(9.6), Inches(0.85), COLORS['light_pink_bg'])
    add_text_box(slide, Inches(0.4), Inches(4.0), Inches(9.2), Inches(0.25),
                "⚠️ MISCONCEPTION ALERT:", font_size=11, bold=True, color=COLORS['red_accent'])
    add_text_box(slide, Inches(0.4), Inches(4.3), Inches(9.2), Inches(0.4),
                "Magnets do NOT have unlimited strength! Force decreases rapidly with distance.",
                font_size=11, color=COLORS['dark_text'])


def add_station1_activity_slide(prs):
    """Slide 6: Station 1 Activity Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 1: Field Visualization",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    steps = [
        ("SPRINKLE", "Iron filings on paper over bar magnet - observe patterns"),
        ("TRACE", "Draw field lines following the filing pattern"),
        ("COMPASS", "Use compass array to confirm field direction"),
        ("COMPARE", "How do fields differ for strong vs weak magnets?"),
        ("MODEL", "Use PhET to see 3D field structure"),
    ]

    y_pos = 0.85
    colors = [COLORS['magnet_blue'], COLORS['field_purple'], COLORS['green_end'],
              COLORS['force_orange'], COLORS['teal']]
    for i, (step, desc) in enumerate(steps):
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.58), COLORS['light_gray_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.58), colors[i])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(1.8), Inches(0.48),
                    step, font_size=12, bold=True, color=colors[i],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(2.3), Inches(y_pos + 0.05), Inches(7.3), Inches(0.48),
                    desc, font_size=11, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.63

    add_colored_shape(slide, Inches(0.2), Inches(4.05), Inches(9.6), Inches(0.65), COLORS['magnet_blue'])
    add_text_box(slide, Inches(0.4), Inches(4.08), Inches(9.2), Inches(0.6),
                "KEY: Field lines are DENSE where force is STRONG, sparse where weak!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station2_intro_slide(prs):
    """Slide 7: Station 2 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['force_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 2: Force-Distance Investigation",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "20 Points | 5 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(1.0), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "RESOURCES:", font_size=12, bold=True, color=COLORS['force_orange'])
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(0.5),
                "Spring scales • Magnets at set distances • Data recording tables • Graphing tools",
                font_size=11, color=COLORS['dark_text'])

    add_text_box(slide, Inches(0.3), Inches(2.3), Inches(9.4), Inches(0.3),
                "THE INVERSE RELATIONSHIP:", font_size=13, bold=True, color=COLORS['force_orange'])

    add_colored_shape(slide, Inches(0.3), Inches(2.7), Inches(9.4), Inches(1.6), COLORS['light_magnet_bg'])
    inverse = """As DISTANCE doubles → Force becomes ~1/4 as strong (inverse square)

EXAMPLE:
• At 1 cm: Force = 100 units
• At 2 cm: Force = 25 units (1/4)
• At 3 cm: Force = 11 units (1/9)
• At 4 cm: Force = 6.25 units (1/16)

This is why magnets feel MUCH weaker just a little farther away!"""
    add_text_box(slide, Inches(0.5), Inches(2.8), Inches(9.0), Inches(1.4),
                inverse, font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(4.45), Inches(9.6), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.48), Inches(9.2), Inches(0.5),
                "FOCUS: Collect data, graph force vs distance, identify the inverse relationship",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station2_activity_slide(prs):
    """Slide 8: Station 2 Activity Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['force_orange'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 2: Data Collection & Graphing",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(0.3), Inches(0.85), Inches(9.4), Inches(0.3),
                "YOUR DATA TABLE:", font_size=13, bold=True, color=COLORS['force_orange'])

    # Data table example
    add_colored_shape(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(1.6), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(2.0), Inches(0.35),
                "Distance (cm)", font_size=11, bold=True, color=COLORS['dark_text'])
    add_text_box(slide, Inches(2.6), Inches(1.3), Inches(2.0), Inches(0.35),
                "Force (N)", font_size=11, bold=True, color=COLORS['dark_text'])
    add_text_box(slide, Inches(4.7), Inches(1.3), Inches(4.5), Inches(0.35),
                "Observations", font_size=11, bold=True, color=COLORS['dark_text'])

    data_rows = [
        ("1", "Measure", "Note: paperclip jumps to magnet"),
        ("2", "Measure", "Note: weaker pull"),
        ("4", "Measure", "Note: barely lifts paperclip"),
        ("8", "Measure", "Note: no visible effect"),
    ]

    y_pos = 1.7
    for dist, force, obs in data_rows:
        add_text_box(slide, Inches(0.5), Inches(y_pos), Inches(2.0), Inches(0.3),
                    dist, font_size=10, color=COLORS['dark_text'])
        add_text_box(slide, Inches(2.6), Inches(y_pos), Inches(2.0), Inches(0.3),
                    force, font_size=10, color=COLORS['dark_text'])
        add_text_box(slide, Inches(4.7), Inches(y_pos), Inches(4.5), Inches(0.3),
                    obs, font_size=10, color=COLORS['gray_text'])
        y_pos += 0.3

    add_colored_shape(slide, Inches(0.2), Inches(2.95), Inches(9.6), Inches(0.8), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.05), Inches(9.2), Inches(0.25),
                "GRAPH YOUR DATA:", font_size=12, bold=True, color=COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(3.35), Inches(9.2), Inches(0.35),
                "• X-axis: Distance (cm)  • Y-axis: Force (N)  • What shape is your curve?",
                font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(3.9), Inches(9.6), Inches(0.85), COLORS['magnet_blue'])
    add_text_box(slide, Inches(0.4), Inches(4.0), Inches(9.2), Inches(0.7),
                "KEY: Your graph should curve DOWN sharply - NOT a straight line!\nThis proves the INVERSE SQUARE relationship.",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station3_intro_slide(prs):
    """Slide 9: Station 3 Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['field_purple'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Station 3: Design a Magnetic Levitation Device",
                font_size=22, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "25 Points | 5 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(1.4), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(0.08), Inches(1.4), COLORS['field_purple'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "ENGINEERING CHALLENGE:", font_size=14, bold=True, color=COLORS['field_purple'])
    challenge = """Design a system that can levitate a small object using only magnets!
Applications: Maglev trains, frictionless bearings, floating displays.

Challenge: Magnetic levitation is UNSTABLE - objects tend to flip or slide sideways!"""
    add_text_box(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(0.9),
                challenge, font_size=12, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(2.7), Inches(4.7), Inches(1.1), COLORS['light_blue_bg'])
    add_text_box(slide, Inches(0.4), Inches(2.8), Inches(4.4), Inches(0.25),
                "CONSTRAINTS:", font_size=11, bold=True, color=COLORS['blue_accent'])
    add_text_box(slide, Inches(0.4), Inches(3.1), Inches(4.4), Inches(0.65),
                "• Only magnetic materials\n• Must float for 10+ seconds\n• No external power source",
                font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.05), Inches(2.7), Inches(4.7), Inches(1.1), COLORS['light_green_bg'])
    add_text_box(slide, Inches(5.25), Inches(2.8), Inches(4.3), Inches(0.25),
                "SUCCESS CRITERIA:", font_size=11, bold=True, color=COLORS['green_end'])
    add_text_box(slide, Inches(5.25), Inches(3.1), Inches(4.3), Inches(0.65),
                "• Stable levitation achieved\n• Can explain why it works\n• Identify force factors",
                font_size=10, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(3.95), Inches(9.6), Inches(0.7), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(3.98), Inches(9.2), Inches(0.65),
                "SEP-1: Asking Questions | What factors affect levitation stability?",
                font_size=11, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_station3_activity_slide(prs):
    """Slide 10: Station 3 Activity Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['field_purple'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Station 3: Levitation Design Process",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    add_text_box(slide, Inches(0.3), Inches(0.85), Inches(9.4), Inches(0.3),
                "APPROACHES TO TRY:", font_size=13, bold=True, color=COLORS['field_purple'])

    approaches = [
        ("Repulsion stacking", "Stack ring magnets on a pole - repelling forces create levitation"),
        ("Diamagnetic levitation", "Use bismuth/pyrolytic graphite to repel magnetic field"),
        ("Halbach array", "Arrange magnets to focus field on one side"),
        ("Spinning stabilization", "Spin the levitating magnet for gyroscopic stability"),
    ]

    y_pos = 1.2
    for approach, description in approaches:
        add_colored_shape(slide, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.48), COLORS['light_magnet_bg'])
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.02), Inches(2.5), Inches(0.44),
                    approach, font_size=11, bold=True, color=COLORS['field_purple'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.1), Inches(y_pos + 0.02), Inches(6.4), Inches(0.44),
                    description, font_size=10, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.53

    add_colored_shape(slide, Inches(0.2), Inches(3.4), Inches(9.6), Inches(1.35), COLORS['light_green_bg'])
    add_text_box(slide, Inches(0.4), Inches(3.5), Inches(9.2), Inches(0.25),
                "JUSTIFY YOUR DESIGN:", font_size=12, bold=True, color=COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(3.8), Inches(9.2), Inches(0.9),
                "• What force factors does your design manipulate?\n• Why is stability challenging?\n• How does distance affect your levitation height?",
                font_size=11, color=COLORS['dark_text'])


def add_concepts_slide(prs):
    """Slide 11: Key Concepts"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Key Concepts This Week",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    concepts = [
        ("Magnetic Fields", "Invisible force regions around magnets; visualized with iron filings"),
        ("Field Lines", "Show direction (N→S) and strength (dense = strong)"),
        ("Inverse Relationship", "Force decreases rapidly with distance (inverse square)"),
        ("Poles & Interactions", "Like poles repel, opposite poles attract"),
    ]

    y_pos = 0.9
    for title, detail in concepts:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.75), COLORS['light_magnet_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.75), COLORS['magnet_blue'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.08), Inches(2.8), Inches(0.6),
                    title, font_size=13, bold=True, color=COLORS['magnet_blue'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.3), Inches(y_pos + 0.08), Inches(6.3), Inches(0.6),
                    detail, font_size=11, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.82

    add_colored_shape(slide, Inches(0.2), Inches(4.2), Inches(9.6), Inches(0.55), COLORS['field_purple'])
    add_text_box(slide, Inches(0.4), Inches(4.23), Inches(9.2), Inches(0.5),
                "CCC-2: Cause and Effect | Distance CAUSES changes in force strength!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_vocabulary_slide(prs):
    """Slide 12: Vocabulary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['purple_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Essential Vocabulary",
                font_size=26, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    vocab = [
        ("Magnetic field", "Region around a magnet where magnetic force acts"),
        ("Field lines", "Visual representation of field direction and strength"),
        ("Magnetic pole", "Region where field is strongest (North or South)"),
        ("Inverse relationship", "As one variable increases, another decreases"),
        ("Ferromagnetic", "Materials that can be magnetized (iron, nickel, cobalt)"),
        ("Magnetize", "To make a material magnetic by aligning its domains"),
    ]

    y_pos = 0.9
    for term, definition in vocab:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.55), COLORS['light_purple_bg'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.5), Inches(0.45),
                    term, font_size=11, bold=True, color=COLORS['purple_end'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.0), Inches(y_pos + 0.05), Inches(6.6), Inches(0.45),
                    definition, font_size=10, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    add_colored_shape(slide, Inches(0.2), Inches(4.55), Inches(9.6), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(4.58), Inches(9.2), Inches(0.5),
                "📝 Add these terms to your NOTECARD for the exit ticket!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_exit_intro_slide(prs):
    """Slide 13: Exit Ticket Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['orange_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "Exit Ticket: Magnetic Force Integration",
                font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.63), Inches(9.3), Inches(0.35),
                "23 Points | 6 Questions",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(1.15), Inches(9.6), Inches(2.4), COLORS['light_orange_bg'])
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.3),
                "QUESTION BREAKDOWN:", font_size=14, bold=True, color=COLORS['orange_end'])

    breakdown = [
        ("2 NEW", "Magnetic fields, force-distance relationship, pole interactions"),
        ("2 SPIRAL", "Connect to C5 waves (invisible forces, patterns)"),
        ("1 INTEGRATION", "Apply field concepts to predict magnetic behavior"),
        ("1 SEP", "Asking questions about force factors"),
    ]

    y_pos = 1.65
    for qtype, desc in breakdown:
        add_text_box(slide, Inches(0.5), Inches(y_pos), Inches(1.8), Inches(0.45),
                    qtype, font_size=12, bold=True, color=COLORS['orange_end'])
        add_text_box(slide, Inches(2.4), Inches(y_pos), Inches(7.2), Inches(0.45),
                    desc, font_size=11, color=COLORS['dark_text'])
        y_pos += 0.5

    add_colored_shape(slide, Inches(0.2), Inches(3.7), Inches(9.6), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(3.73), Inches(9.2), Inches(0.5),
                "⏱️ TIME: ~15 minutes | Use your NOTECARD!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(4.4), Inches(9.6), Inches(0.55), COLORS['light_pink_bg'])
    add_text_box(slide, Inches(0.4), Inches(4.43), Inches(9.2), Inches(0.5),
                "⚠️ Remember: Magnets have LIMITED range - force drops rapidly with distance!",
                font_size=11, color=COLORS['dark_text'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_exit_tips_slide(prs):
    """Slide 14: Exit Ticket Tips"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['orange_end'])
    add_text_box(slide, Inches(0.35), Inches(0.18), Inches(9.3), Inches(0.5),
                "Exit Ticket Tips",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    tips = [
        "Draw field lines correctly: N to S, never crossing, dense near poles",
        "Remember inverse relationship: 2x distance = 1/4 force",
        "Connect to C5 waves: both are invisible forces with patterns",
        "Know pole interactions: like repel, opposite attract",
        "Check your notecard for key terms and relationships",
    ]

    y_pos = 0.85
    for i, tip in enumerate(tips, 1):
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.55), COLORS['light_gray_bg'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(0.5), Inches(0.45),
                    str(i), font_size=14, bold=True, color=COLORS['orange_end'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(1.0), Inches(y_pos + 0.05), Inches(8.6), Inches(0.45),
                    tip, font_size=12, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    add_colored_shape(slide, Inches(0.2), Inches(3.9), Inches(9.6), Inches(0.85), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(3.95), Inches(9.2), Inches(0.75),
                "You've got this! You mapped fields, measured forces, and designed levitators.\nNow show what you learned!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_notecard_slide(prs):
    """Slide 15: Notecard Reminder"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['magnet_blue'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "📝 Your Notecard Should Include:",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    items = [
        "How to draw field lines (N→S, dense near poles, never cross)",
        "Force-distance inverse relationship (2x distance = 1/4 force)",
        "Pole interactions (like repel, opposite attract)",
        "Key vocabulary: magnetic field, field lines, ferromagnetic",
        "Connection to C5: invisible forces that follow patterns",
    ]

    y_pos = 0.9
    for item in items:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.55), COLORS['light_magnet_bg'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(0.4), Inches(0.45),
                    "✓", font_size=14, bold=True, color=COLORS['green_end'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(0.9), Inches(y_pos + 0.05), Inches(8.7), Inches(0.45),
                    item, font_size=11, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.6

    add_colored_shape(slide, Inches(0.2), Inches(3.95), Inches(9.6), Inches(0.85), COLORS['field_purple'])
    add_text_box(slide, Inches(0.4), Inches(3.98), Inches(9.2), Inches(0.8),
                "Remember: Your notecard is the ONLY physical resource allowed!\nMake it count - organize it clearly!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def add_summary_slide(prs):
    """Slide 16: Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.6), COLORS['teal_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.5),
                "Week 1 Summary: What You Learned",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia",
                anchor=MSO_ANCHOR.MIDDLE)

    takeaways = [
        ("Magnetic Fields", "Invisible force regions visualized by iron filings"),
        ("Field Lines", "Show direction (N→S) and strength (density)"),
        ("Inverse Relationship", "Force drops rapidly with distance"),
        ("Engineering", "Levitation requires balancing forces precisely"),
    ]

    y_pos = 0.9
    for title, detail in takeaways:
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(9.6), Inches(0.65), COLORS['light_magnet_bg'])
        add_colored_shape(slide, Inches(0.2), Inches(y_pos), Inches(0.08), Inches(0.65), COLORS['magnet_blue'])
        add_text_box(slide, Inches(0.4), Inches(y_pos + 0.05), Inches(2.5), Inches(0.55),
                    title, font_size=13, bold=True, color=COLORS['magnet_blue'],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(3.0), Inches(y_pos + 0.05), Inches(6.6), Inches(0.55),
                    detail, font_size=12, color=COLORS['dark_text'],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_pos += 0.72

    add_colored_shape(slide, Inches(0.2), Inches(3.85), Inches(9.6), Inches(0.65), COLORS['force_orange'])
    add_text_box(slide, Inches(0.4), Inches(3.88), Inches(9.2), Inches(0.6),
                "NEXT WEEK: Electromagnetism - How spinning a magnet creates electricity!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_colored_shape(slide, Inches(0.2), Inches(4.65), Inches(9.6), Inches(0.55), COLORS['field_purple'])
    add_text_box(slide, Inches(0.4), Inches(4.68), Inches(9.2), Inches(0.5),
                "ANSWER: Stronger magnets have stronger fields that extend farther before weakening!",
                font_size=12, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


if __name__ == "__main__":
    prs = create_presentation()
    output_path = "/home/user/Kairos.Sci.Repo/content/grade8/cycle06/week1/G8_C6_W1_Slides_Final.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Created: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
