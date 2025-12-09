#!/usr/bin/env python3
"""
Create G7_C3_W3 Synthesis & Assessment presentation
Assessment week - different structure than regular weeks
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

COLORS = {
    'red_start': RGBColor(0xE5, 0x3E, 0x3E),
    'red_end': RGBColor(0xC5, 0x30, 0x30),
    'red_dark': RGBColor(0x9B, 0x2C, 0x2C),
    'purple_start': RGBColor(0x66, 0x7E, 0xEA),
    'purple_end': RGBColor(0x76, 0x4B, 0xA2),
    'green_start': RGBColor(0x48, 0xBB, 0x78),
    'green_end': RGBColor(0x27, 0x67, 0x49),
    'teal': RGBColor(0x38, 0xB2, 0xAC),
    'teal_dark': RGBColor(0x23, 0x4E, 0x52),
    'orange': RGBColor(0xD6, 0x9E, 0x2E),
    'dark_text': RGBColor(0x2D, 0x37, 0x48),
    'gray_text': RGBColor(0x4A, 0x55, 0x68),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'light_red_bg': RGBColor(0xFF, 0xF5, 0xF5),
    'light_purple_bg': RGBColor(0xFA, 0xF5, 0xFF),
    'light_green_bg': RGBColor(0xF0, 0xFF, 0xF4),
    'light_teal_bg': RGBColor(0xE6, 0xFF, 0xFA),
    'light_orange_bg': RGBColor(0xFF, 0xFB, 0xEB),
}


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    add_title_slide(prs)
    add_assessment_overview_slide(prs)
    add_review_guide_slide(prs)
    add_common_mistakes_slide(prs)
    add_part1_intro_slide(prs)
    add_part1_support_slide(prs)
    add_part2_intro_slide(prs)
    add_part2_support_slide(prs)
    add_part3_intro_slide(prs)
    add_part3_support_slide(prs)
    add_completion_slide(prs)

    return prs


def add_colored_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=None, align=PP_ALIGN.LEFT, font_name="Arial", anchor=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    if anchor:
        tf.anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    if color:
        p.font.color.rgb = color
    return txBox


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['red_end'])

    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(1.2),
                "📋 Week 3: Synthesis & Assessment 🎯",
                font_size=42, bold=True, color=COLORS['white'],
                align=PP_ALIGN.CENTER, font_name="Georgia")

    add_text_box(slide, Inches(0.5), Inches(2.9), Inches(9), Inches(0.5),
                "Grade 7 Science | Rosche | Kairos Academies",
                font_size=20, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(3.4), Inches(9), Inches(0.5),
                "Cycle 3 Cumulative Assessment | 100 Points Total | ~75 Minutes",
                font_size=16, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Teaser box (simplified - no border)
    teaser_bg = add_colored_shape(slide, Inches(2), Inches(4.2), Inches(6), Inches(0.8), COLORS['white'])
    add_text_box(slide, Inches(2.1), Inches(4.35), Inches(5.8), Inches(0.5),
                "⚠️ Assessment Week - Show What You Know!",
                font_size=16, color=COLORS['red_dark'], align=PP_ALIGN.CENTER)


def add_assessment_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.2), Inches(0.2), Inches(9.6), Inches(0.6), COLORS['red_end'])
    add_text_box(slide, Inches(0.4), Inches(0.28), Inches(9.2), Inches(0.5),
                "⚠️ Assessment Week - Important Information",
                font_size=24, bold=True, color=COLORS['white'], font_name="Georgia")

    add_colored_shape(slide, Inches(0.2), Inches(1.0), Inches(9.6), Inches(2.2), COLORS['light_red_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(1.0), Inches(0.08), Inches(2.2), COLORS['red_end'])

    info = """• This is your cumulative assessment for Cycle 3 Climate Change unit
• You may NOT go back once you submit each part
• Complete all 3 parts in order: Synthesis → Assessment → Misconception Check
• Take breaks between parts - your brain works better with rest!
• Ask Mr. Rosche if you need clarification on any question"""
    add_text_box(slide, Inches(0.5), Inches(1.15), Inches(9.0), Inches(2.0),
                info, font_size=15, color=COLORS['red_dark'])

    # Three parts overview
    add_colored_shape(slide, Inches(0.2), Inches(3.4), Inches(3.0), Inches(1.0), COLORS['purple_start'])
    add_text_box(slide, Inches(0.4), Inches(3.5), Inches(2.6), Inches(0.3),
                "🔗 Part 1: Synthesis", font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.4), Inches(3.85), Inches(2.6), Inches(0.45),
                "20 pts | ~15 min", font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(3.5), Inches(3.4), Inches(3.0), Inches(1.0), COLORS['red_end'])
    add_text_box(slide, Inches(3.7), Inches(3.5), Inches(2.6), Inches(0.3),
                "📝 Part 2: Assessment", font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.7), Inches(3.85), Inches(2.6), Inches(0.45),
                "60 pts | ~40 min", font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(6.8), Inches(3.4), Inches(3.0), Inches(1.0), COLORS['green_end'])
    add_text_box(slide, Inches(7.0), Inches(3.5), Inches(2.6), Inches(0.3),
                "🎯 Part 3: Misconceptions", font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.0), Inches(3.85), Inches(2.6), Inches(0.45),
                "20 pts | ~20 min", font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.2), Inches(4.6), Inches(9.6), Inches(0.6), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.68), Inches(9.2), Inches(0.4),
                "💡 Take 5-minute breaks between parts for best performance!",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_review_guide_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.2), Inches(0.15), Inches(9.6), Inches(0.5), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.4),
                "📚 What You Need to Know (Review Guide)",
                font_size=22, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Week 1
    add_colored_shape(slide, Inches(0.2), Inches(0.8), Inches(4.7), Inches(1.5), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(0.8), Inches(0.08), Inches(1.5), COLORS['purple_start'])
    add_text_box(slide, Inches(0.4), Inches(0.9), Inches(4.4), Inches(0.3),
                "Week 1: Greenhouse Effect", font_size=13, bold=True, color=COLORS['purple_start'])
    w1 = """• CO₂ absorbs IR and vibrates (doesn't break)
• 3+ atom molecules absorb IR
• Carbon cycle - atoms conserved
• 6CO₂ + 6H₂O → C₆H₁₂O₆ + 6O₂"""
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(4.4), Inches(0.95),
                w1, font_size=11, color=COLORS['dark_text'])

    # Week 2
    add_colored_shape(slide, Inches(5.1), Inches(0.8), Inches(4.7), Inches(1.5), COLORS['light_teal_bg'])
    add_colored_shape(slide, Inches(5.1), Inches(0.8), Inches(0.08), Inches(1.5), COLORS['teal'])
    add_text_box(slide, Inches(5.3), Inches(0.9), Inches(4.4), Inches(0.3),
                "Week 2: Feedback Loops", font_size=13, bold=True, color=COLORS['teal_dark'])
    w2 = """• Albedo = fraction of light reflected
• Positive feedback = change amplifies itself
• Ice-albedo: melting → dark ocean → more heat
• Carbon sink saturation: warm ocean absorbs less"""
    add_text_box(slide, Inches(5.3), Inches(1.25), Inches(4.4), Inches(0.95),
                w2, font_size=11, color=COLORS['dark_text'])

    # Cycle 2 Spiral
    add_colored_shape(slide, Inches(0.2), Inches(2.45), Inches(9.6), Inches(1.1), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(2.45), Inches(0.08), Inches(1.1), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(2.55), Inches(9.2), Inches(0.3),
                "Cycle 2 Spiral Review:", font_size=13, bold=True, color=COLORS['green_end'])
    c2 = """• Breaking bonds REQUIRES energy (not releases!)  •  Forming bonds RELEASES energy  •  Conservation of mass: atoms rearranged, never created/destroyed"""
    add_text_box(slide, Inches(0.4), Inches(2.9), Inches(9.2), Inches(0.55),
                c2, font_size=12, color=COLORS['dark_text'])

    # Key equations
    add_colored_shape(slide, Inches(0.2), Inches(3.7), Inches(9.6), Inches(1.0), COLORS['purple_start'])
    add_text_box(slide, Inches(0.4), Inches(3.8), Inches(9.2), Inches(0.25),
                "🔑 KEY EQUATIONS TO KNOW:", font_size=12, bold=True, color=COLORS['white'])
    add_text_box(slide, Inches(0.4), Inches(4.1), Inches(9.2), Inches(0.5),
                "Photosynthesis: 6CO₂ + 6H₂O + light → C₆H₁₂O₆ + 6O₂   |   Combustion: C₆H₁₂O₆ + 6O₂ → 6CO₂ + 6H₂O + energy",
                font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.2), Inches(4.85), Inches(9.6), Inches(0.55), COLORS['red_end'])
    add_text_box(slide, Inches(0.4), Inches(4.9), Inches(9.2), Inches(0.4),
                "📝 Use your notecard from W1 & W2 during the assessment!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_common_mistakes_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.2), Inches(0.15), Inches(9.6), Inches(0.5), COLORS['orange'])
    add_text_box(slide, Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.4),
                "⚠️ Common Mistakes to Avoid",
                font_size=22, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    mistakes = [
        ("\"Breaking bonds releases energy\"", "Breaking bonds REQUIRES energy; forming bonds releases energy"),
        ("\"Mass disappears in reactions\"", "Mass is conserved; atoms rearrange but are never destroyed"),
        ("\"Positive feedback is good\"", "Positive means AMPLIFYING (same direction), not beneficial"),
    ]

    y = 0.8
    for wrong, correct in mistakes:
        add_colored_shape(slide, Inches(0.2), Inches(y), Inches(4.7), Inches(1.0), COLORS['light_red_bg'])
        add_text_box(slide, Inches(0.4), Inches(y + 0.1), Inches(4.4), Inches(0.25),
                    "❌ WRONG:", font_size=10, bold=True, color=COLORS['red_end'])
        add_text_box(slide, Inches(0.4), Inches(y + 0.35), Inches(4.4), Inches(0.55),
                    wrong, font_size=11, color=COLORS['dark_text'])

        add_colored_shape(slide, Inches(5.1), Inches(y), Inches(4.7), Inches(1.0), COLORS['light_green_bg'])
        add_text_box(slide, Inches(5.3), Inches(y + 0.1), Inches(4.4), Inches(0.25),
                    "✅ CORRECT:", font_size=10, bold=True, color=COLORS['green_end'])
        add_text_box(slide, Inches(5.3), Inches(y + 0.35), Inches(4.4), Inches(0.55),
                    correct, font_size=11, color=COLORS['dark_text'])
        y += 1.15

    add_colored_shape(slide, Inches(0.2), Inches(4.3), Inches(9.6), Inches(0.9), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(4.4), Inches(9.2), Inches(0.7),
                "💡 REMEMBER: CO₂ molecules VIBRATE when they absorb infrared light - they don't break apart!\nThis vibration is what traps heat in the atmosphere.",
                font_size=13, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_part1_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['purple_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🔗 Part 1: Synthesis Review",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "20 Points | ~15 min | Connects Week 1 + Week 2", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "WHAT YOU'LL DO", font_size=14, bold=True, color=COLORS['purple_start'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.0), COLORS['light_purple_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(2.0), COLORS['purple_start'])

    steps = """1. Create a diagram connecting:
   • Greenhouse effect
   • Feedback loops
   • Human emissions
2. Identify the primary driver of climate change
3. Explain how the complete system works
4. Connect to Cycle 2 conservation concepts"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(1.8),
                steps, font_size=12, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(2.4), COLORS['purple_end'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "🎯 KEY CONNECTION", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.2), Inches(1.6), Inches(4.3), Inches(1.8),
                "How does the greenhouse effect (W1) connect to feedback loops (W2)?\n\nCO₂ traps heat →\nIce melts →\nAlbedo decreases →\nMore heat absorbed →\nMore ice melts...",
                font_size=13, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.15), Inches(3.7), Inches(9.7), Inches(0.5), COLORS['light_teal_bg'])
    add_text_box(slide, Inches(0.35), Inches(3.78), Inches(9.3), Inches(0.35),
                "💡 This is about CONNECTING concepts, not just remembering them!",
                font_size=13, color=COLORS['teal_dark'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.15), Inches(4.35), Inches(9.7), Inches(0.55), COLORS['purple_start'])
    add_text_box(slide, Inches(0.35), Inches(4.4), Inches(9.3), Inches(0.4),
                "📝 Complete Part 1 Form, then take a 5-minute break",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_part1_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['purple_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.4),
                "🔗 Part 1 – Synthesis Diagram Guide",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia")

    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(9.6), Inches(2.5), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(9.2), Inches(0.3),
                "🔄 The Complete Climate System (Use this for your diagram):", font_size=13, bold=True, color=COLORS['purple_start'])

    diagram = """Human activities (burning fossil fuels, deforestation)
                    ↓
           Release CO₂ into atmosphere
                    ↓
    CO₂ absorbs infrared radiation (VIBRATES, doesn't break)
                    ↓
           Heat is trapped → Global warming
                    ↓
    POSITIVE FEEDBACK #1: Ice melts → Dark ocean exposed → More absorption → More melting
                    ↓
    POSITIVE FEEDBACK #2: Warm oceans absorb LESS CO₂ → More stays in atmosphere"""
    add_text_box(slide, Inches(0.4), Inches(1.3), Inches(9.2), Inches(1.9),
                diagram, font_size=11, color=COLORS['dark_text'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.2), Inches(3.5), Inches(9.6), Inches(0.6), COLORS['green_end'])
    add_text_box(slide, Inches(0.4), Inches(3.58), Inches(9.2), Inches(0.4),
                "🔑 KEY: Atoms are CONSERVED. Carbon cycles through atmosphere, oceans, plants, and fossil fuels!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.15), Inches(4.25), Inches(9.7), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(4.3), Inches(9.3), Inches(0.4),
                "⏸️ After Part 1: Take a 5-minute break before Part 2!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_part2_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['red_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "📝 Part 2: Cumulative Assessment",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "60 Points | ~40 min | All Learning Targets", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Four sections
    sections = [
        ("Section A", "Greenhouse Effect (15 pts)", "Molecular mechanisms", COLORS['light_purple_bg'], COLORS['purple_start']),
        ("Section B", "Carbon Cycle (15 pts)", "Calculations + concepts", COLORS['light_teal_bg'], COLORS['teal']),
        ("Section C", "Albedo & Feedback (15 pts)", "System thinking", COLORS['light_green_bg'], COLORS['green_end']),
        ("Section D", "Integration (15 pts)", "Connecting everything", COLORS['light_red_bg'], COLORS['red_end']),
    ]

    x = 0.2
    for title, subtitle, desc, bg, accent in sections:
        add_colored_shape(slide, Inches(x), Inches(1.15), Inches(2.3), Inches(1.5), bg)
        add_colored_shape(slide, Inches(x), Inches(1.15), Inches(2.3), Inches(0.05), accent)
        add_text_box(slide, Inches(x + 0.1), Inches(1.25), Inches(2.1), Inches(0.25),
                    title, font_size=11, bold=True, color=accent)
        add_text_box(slide, Inches(x + 0.1), Inches(1.5), Inches(2.1), Inches(0.3),
                    subtitle, font_size=10, bold=True, color=COLORS['dark_text'])
        add_text_box(slide, Inches(x + 0.1), Inches(1.85), Inches(2.1), Inches(0.7),
                    desc, font_size=10, color=COLORS['gray_text'])
        x += 2.4

    add_colored_shape(slide, Inches(0.2), Inches(2.8), Inches(9.6), Inches(0.8), COLORS['light_red_bg'])
    add_colored_shape(slide, Inches(0.2), Inches(2.8), Inches(0.08), Inches(0.8), COLORS['red_end'])
    add_text_box(slide, Inches(0.4), Inches(2.9), Inches(9.2), Inches(0.6),
                "⚠️ WARNING: You cannot go back once you submit. Read each question carefully!\nMake sure to answer EVERY question before moving to the next section.",
                font_size=12, color=COLORS['red_dark'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.15), Inches(3.75), Inches(9.7), Inches(0.55), COLORS['red_end'])
    add_text_box(slide, Inches(0.35), Inches(3.8), Inches(9.3), Inches(0.4),
                "📝 Complete Part 2 Form - Take your time!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_part2_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['red_dark'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.4),
                "📝 Part 2 – Key Formulas & Quick Reference",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia")

    # Left column - Equations
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(4.7), Inches(2.0), COLORS['light_purple_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(4.4), Inches(0.25),
                "🧪 KEY EQUATIONS:", font_size=12, bold=True, color=COLORS['purple_start'])
    equations = """Photosynthesis:
6CO₂ + 6H₂O + light → C₆H₁₂O₆ + 6O₂

Combustion/Respiration:
C₆H₁₂O₆ + 6O₂ → 6CO₂ + 6H₂O + energy

Albedo:
0 = absorbs all | 1 = reflects all"""
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(4.4), Inches(1.5),
                equations, font_size=11, color=COLORS['dark_text'])

    # Right column - Concepts
    add_colored_shape(slide, Inches(5.1), Inches(0.85), Inches(4.7), Inches(2.0), COLORS['light_teal_bg'])
    add_text_box(slide, Inches(5.3), Inches(0.95), Inches(4.4), Inches(0.25),
                "💡 KEY CONCEPTS:", font_size=12, bold=True, color=COLORS['teal_dark'])
    concepts = """• CO₂ VIBRATES when absorbing IR (doesn't break)
• Breaking bonds = REQUIRES energy
• Forming bonds = RELEASES energy
• Positive feedback = AMPLIFIES change
• High albedo = reflects light = stays cool
• Low albedo = absorbs light = heats up"""
    add_text_box(slide, Inches(5.3), Inches(1.25), Inches(4.4), Inches(1.5),
                concepts, font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(3.0), Inches(9.6), Inches(0.6), COLORS['orange'])
    add_text_box(slide, Inches(0.4), Inches(3.08), Inches(9.2), Inches(0.4),
                "🔢 For calculations: Count atoms on BOTH sides - they must be EQUAL (conservation of mass)!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.15), Inches(3.75), Inches(9.7), Inches(0.55), COLORS['teal'])
    add_text_box(slide, Inches(0.35), Inches(3.8), Inches(9.3), Inches(0.4),
                "⏸️ After Part 2: Take a 5-minute break before Part 3!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_part3_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.85), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.45),
                "🎯 Part 3: Misconception Check",
                font_size=26, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.3),
                "20 Points | ~20 min | Final Chance!", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "WHAT THIS TESTS", font_size=14, bold=True, color=COLORS['green_end'])

    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(1.8), COLORS['light_green_bg'])
    add_colored_shape(slide, Inches(0.3), Inches(1.5), Inches(0.05), Inches(1.8), COLORS['green_end'])

    tests = """• Bond Energy: Breaking bonds requires energy
  (most common mistake!)

• Conservation of Mass: Atoms are never
  created or destroyed

• These questions catch common errors -
  think VERY carefully!"""
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(1.6),
                tests, font_size=12, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.0), Inches(1.1), Inches(4.7), Inches(2.2), COLORS['orange'])
    add_text_box(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.3),
                "⚠️ BIGGEST MISTAKE:", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.2), Inches(1.6), Inches(4.3), Inches(1.6),
                "\"Breaking bonds releases energy\"\n\n❌ WRONG!\n\n✅ Breaking bonds REQUIRES energy\n✅ Forming bonds RELEASES energy",
                font_size=13, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(0.15), Inches(3.5), Inches(9.7), Inches(0.55), COLORS['green_end'])
    add_text_box(slide, Inches(0.35), Inches(3.55), Inches(9.3), Inches(0.4),
                "📝 Complete Part 3 Form - Last section! You've got this!",
                font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_part3_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0.15), Inches(0.15), Inches(9.7), Inches(0.55), COLORS['green_start'])
    add_text_box(slide, Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.4),
                "🎯 Part 3 – Misconception Traps to Avoid",
                font_size=22, bold=True, color=COLORS['white'], font_name="Georgia")

    # Two-column comparison
    add_colored_shape(slide, Inches(0.2), Inches(0.85), Inches(4.7), Inches(1.5), COLORS['light_red_bg'])
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(4.4), Inches(0.25),
                "❌ TRAP ANSWERS (DON'T PICK):", font_size=12, bold=True, color=COLORS['red_end'])
    traps = """• "Energy is released when bonds break"
• "Mass is lost in combustion reactions"
• "Positive feedback means good outcomes"
• "CO₂ molecules break when heated" """
    add_text_box(slide, Inches(0.4), Inches(1.25), Inches(4.4), Inches(1.0),
                traps, font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(5.1), Inches(0.85), Inches(4.7), Inches(1.5), COLORS['light_green_bg'])
    add_text_box(slide, Inches(5.3), Inches(0.95), Inches(4.4), Inches(0.25),
                "✅ CORRECT ANSWERS (LOOK FOR):", font_size=12, bold=True, color=COLORS['green_end'])
    correct = """• "Energy is required to break bonds"
• "Mass is conserved in all reactions"
• "Positive feedback amplifies change"
• "CO₂ molecules vibrate when absorbing IR" """
    add_text_box(slide, Inches(5.3), Inches(1.25), Inches(4.4), Inches(1.0),
                correct, font_size=11, color=COLORS['dark_text'])

    add_colored_shape(slide, Inches(0.2), Inches(2.5), Inches(9.6), Inches(0.7), COLORS['teal'])
    add_text_box(slide, Inches(0.4), Inches(2.6), Inches(9.2), Inches(0.5),
                "💡 TIP: If an answer seems too simple or matches what you'd \"expect,\" be careful!\nThe correct answer in science is often counter-intuitive.",
                font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)


def add_completion_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_colored_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), COLORS['green_end'])

    add_text_box(slide, Inches(0.5), Inches(0.8), Inches(9), Inches(0.8),
                "🎉 Cycle 3 Complete!",
                font_size=40, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER, font_name="Georgia")

    add_text_box(slide, Inches(0.5), Inches(1.7), Inches(9), Inches(0.5),
                "Congratulations on finishing the Climate Change unit!",
                font_size=20, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_colored_shape(slide, Inches(2), Inches(2.4), Inches(6), Inches(2.2), COLORS['white'])
    add_text_box(slide, Inches(2.2), Inches(2.5), Inches(5.6), Inches(0.3),
                "Key Takeaways:", font_size=16, bold=True, color=COLORS['green_end'])
    takeaways = """• CO₂ absorbs IR → molecules vibrate → heat trapped
• Carbon is conserved through all Earth systems
• Positive feedback loops accelerate climate change
• Multiple feedbacks interact and amplify each other"""
    add_text_box(slide, Inches(2.2), Inches(2.85), Inches(5.6), Inches(1.6),
                takeaways, font_size=14, color=COLORS['dark_text'])

    add_text_box(slide, Inches(0.5), Inches(4.8), Inches(9), Inches(0.5),
                "Great work this cycle! 🌍",
                font_size=18, color=COLORS['white'], align=PP_ALIGN.CENTER)


if __name__ == "__main__":
    prs = create_presentation()
    output_path = "/home/user/Kairos.Sci.Repo/content/grade7/cycle03/week3/G7_C3_W3_Assessment_Slides_Final.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Created: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
